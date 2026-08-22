from __future__ import annotations

import importlib.util
from pathlib import Path

from PIL import Image
from pptx import Presentation


REPO = Path(__file__).resolve().parents[1]
SCRIPT = REPO / "scripts/presentations/build_phase19_genetic_support_deck.py"


def load_builder():
    spec = importlib.util.spec_from_file_location("phase19_deck_builder", SCRIPT)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def test_six_slide_native_figures_are_separate_packages():
    builder = load_builder()
    assert len(builder.FIG) == 6
    assert len({path.parent for path in builder.FIG.values()}) == 6
    for path in builder.FIG.values():
        assert path.exists()
        with Image.open(path) as image:
            assert image.size == (5580, 2115)


def test_builds_and_validates_external_facing_genetic_support_deck(tmp_path):
    builder = load_builder()
    output = tmp_path / "human_genetic_support_for_key_drivers.pptx"
    assert builder.build_deck(output) == output.resolve()
    builder.validate_deck(output)

    presentation = Presentation(output)
    assert builder.EXPECTED_SLIDE_COUNT == 28
    assert len(presentation.slides) == builder.EXPECTED_SLIDE_COUNT
    assert presentation.slide_width == builder.SLIDE_W
    assert presentation.slide_height == builder.SLIDE_H
    notes = [
        slide.notes_slide.notes_text_frame.text.strip()
        for slide in presentation.slides
    ]
    assert all(notes)
    assert all(
        len(note.split()) >= builder.MIN_SPEAKER_NOTE_WORDS
        for note in notes
    )
    assert all(
        all(heading in note for heading in builder.NOTE_SECTION_HEADINGS)
        for note in notes
    )

    text = "\n".join(slide_text(slide) for slide in presentation.slides)
    for title in [
        "The presentation moves from study design to evidence and future validation",
        "Study design and public data",
        "Evidence found across the gene list",
        "How the evidence guides future validation",
        "Supporting details",
        "Four gene regions stood out below 5×10⁻⁸; all 19 P values are shown",
        "APOE led the registered first-screen results across 47 settings",
        "RPS15 has three brain-QTL records for focused validation",
        "Data-package improvements that will strengthen reproducibility",
    ]:
        assert title in text

    for dataset_label in [
        "GCST90027158",
        "GCST90726396",
        "GCST90726397",
        "GCST90726398",
        "NG00184.v1",
        "eQTL Catalogue r7",
        "FunGen-xQTL",
        "GENCODE v44",
        "HGNC",
    ]:
        assert dataset_label in text

    assert "Result: evidence found, with focused validation next" in text
    assert "6 mitochondrial genes define a dedicated follow-up" in text
    assert "Four regions—near ANKRD11, APOE, COX7C, and RPS15—stood out" in text
    assert "The other 15 ranged from 2.929×10⁻⁶ to 2.931×10⁻⁴" in text
    assert "Nearby-AD P value reported" in text
    assert "Complete DNA-variant files can test whether both signals point to the same variant" in text
    assert "EVIDENCE TO CARRY FORWARD" in text
    assert "Next milestone: a fully documented, portable data package" in text

    slide_2_text = slide_text(presentation.slides[1])
    assert "NEXT VALIDATION" in slide_2_text
    assert "Same-variant test" in slide_2_text
    assert "RPS15 + ANKRD11" in slide_2_text
    slide_2_notes = presentation.slides[1].notes_slide.notes_text_frame.text
    assert "COX7C and SELENOW" in slide_2_notes
    assert "RPS15 and ANKRD11" in slide_2_notes
    assert "same-variant test" in slide_2_notes
    assert "nearby regional signal does not by itself identify" in slide_2_notes

    assert "Required gene-activity or reference data were missing" not in text
    assert "Complete DNA-variant files needed to compare the Alzheimer's and gene-activity signals were unavailable" not in text
    assert "What this study could not answer" not in text
    assert "WHAT THIS DOES NOT SHOW" not in text
    assert "No shared-signal probability was available" not in text
    assert "Status: the data package still needs repair" not in text
    assert "responsible genes remained unresolved" not in text
    assert "No strong signal found" not in text
    assert "No nearby AD signal" not in text
    assert "not enough for confirmation" not in text
    assert "negative checks" not in text
    assert "non-passing regions" not in text
    assert "did not pass" not in text
    assert "four regions passed" not in text
    assert "full_genetic_support_complete = FALSE" not in text
    assert "PP.H4 = 0" not in text
    assert "mtDNA genes were negative" not in text
    assert "Phase 18" not in text
    assert "Phase 19" not in text
