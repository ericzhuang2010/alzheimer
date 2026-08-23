from __future__ import annotations

import csv
import hashlib
import importlib.util
import sys
from pathlib import Path

from pptx import Presentation


REPO = Path(__file__).resolve().parents[2]
SCRIPT_PATH = REPO / (
    "scripts/presentations/append_mouse_dataset_validation_recommendation_slide.py"
)
SPEC = importlib.util.spec_from_file_location(
    "mouse_dataset_validation_recommendation_appender", SCRIPT_PATH
)
assert SPEC and SPEC.loader
appender = importlib.util.module_from_spec(SPEC)
sys.modules[SPEC.name] = appender
SPEC.loader.exec_module(appender)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle, delimiter="\t"))


def _visible_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def test_recommendation_source_and_content_contract() -> None:
    assert _sha256(appender.SOURCE_MD) == appender.SOURCE_MD_SHA256
    appender.base.validate_source_contract(appender.SOURCE_MD)

    source = appender.SOURCE_MD.read_text(encoding="utf-8")
    for phrase in (
        "These three datasets should be analyzed separately",
        "Their raw count matrices should **not** be merged",
        "Freeze the human PFC result first",
        "Analyze every study separately",
        "Triangulated external validation",
        "convergent support",
    ):
        assert phrase in source


def test_published_recommendation_slide_and_reports_are_current() -> None:
    deck = appender.TARGET_DECK
    assert deck.exists()
    appender.validate_final_deck(deck)

    prs = Presentation(deck)
    assert len(prs.slides) == 15
    slide = prs.slides[-1]
    visible = _visible_text(slide)
    assert appender.TITLE in visible
    for phrase in (
        "GSE185063",
        "GSE241553",
        "GSE163857",
        "GSE212606",
        "COMBINE CONCLUSIONS—NOT RAW COUNT MATRICES",
        "triangulated external validation",
        "convergent support",
    ):
        assert phrase in visible
    assert all(shape.shape_type != 13 for shape in slide.shapes)

    notes_frame = slide.notes_slide.notes_text_frame
    assert notes_frame is not None
    notes = notes_frame.text
    assert len(notes.split()) >= 80
    for heading in appender.NOTE_HEADINGS:
        assert heading in notes

    report_dir = appender.DEFAULT_REPORT_DIR
    stem = "mouse_dataset_comparison_08252026"
    status = _read_tsv(report_dir / f"{stem}_status.tsv")
    checks = _read_tsv(report_dir / f"{stem}_checks.tsv")
    manifest = {
        row["artifact_role"]: row
        for row in _read_tsv(report_dir / f"{stem}_input_manifest.tsv")
    }
    inventory = _read_tsv(report_dir / f"{stem}_slide_inventory.tsv")

    assert len(status) == 1
    assert status[0]["validation_status"] == "validated_complete"
    assert status[0]["visual_review_status"] == "complete"
    assert status[0]["slides"] == "15"
    assert status[0]["deck_path"] == (
        "docs/presentations/mouse_dataset_comparison_08252026.pptx"
    )
    assert status[0]["deck_sha256"] == _sha256(deck)
    assert len(status[0]["base_deck_sha256"]) == 64
    assert status[0]["base_deck_sha256"] != status[0]["deck_sha256"]

    assert len(inventory) == 15
    assert inventory[-1]["slide_number"] == "15"
    assert inventory[-1]["title"] == appender.TITLE
    assert inventory[-1]["picture_count"] == "0"

    assert all(row["passed"] == "True" for row in checks)
    assert manifest["source_markdown"]["sha256"] == appender.SOURCE_MD_SHA256
    assert manifest["append_script"]["sha256"] == _sha256(SCRIPT_PATH)
    assert manifest["deck"]["sha256"] == _sha256(deck)
