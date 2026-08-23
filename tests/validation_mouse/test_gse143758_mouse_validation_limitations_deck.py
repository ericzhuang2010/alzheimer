from __future__ import annotations

import csv
import hashlib
import importlib.util
from pathlib import Path

from pptx import Presentation


REPO = Path(__file__).resolve().parents[2]
SCRIPT = REPO / "scripts/presentations/build_gse143758_mouse_validation_limitations_deck.py"
SPEC = importlib.util.spec_from_file_location("gse143758_deck_builder", SCRIPT)
assert SPEC and SPEC.loader
builder = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(builder)


def read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle, delimiter="\t"))


def test_source_document_contract_and_hash() -> None:
    assert builder.SOURCE_DOC.exists()
    assert builder.sha256(builder.SOURCE_DOC) == builder.SOURCE_DOC_SHA256
    text = builder.validate_source_contract(builder.SOURCE_DOC)
    for phrase in (
        "37 GSM records",
        "8 independent mice",
        "one female WT and one female 5xFAD",
        "APOE x sex x disease interaction",
        "Main all-cell atlas and age course: hippocampus",
        "smaller astrocyte-only cortex/PFC subset",
        "cortex resources are primarily astrocyte-focused",
        "the same mice also contributed hippocampus",
        "direct PFC, APOE, or sex-interaction validation",
    ):
        assert phrase.lower() in text.lower()


def test_temporary_build_and_pending_report(tmp_path: Path) -> None:
    output = tmp_path / "review.pptx"
    report_dir = tmp_path / "report"
    built = builder.build_deck(output)
    builder.write_report(built, report_dir, visual_review_status="pending")
    builder.validate_deck(built)

    prs = Presentation(built)
    assert len(prs.slides) == 5
    assert [builder._slide_text(slide).splitlines()[1] for slide in prs.slides]
    assert all(
        sum(1 for shape in slide.shapes if shape.shape_type == 13) == 0
        for slide in prs.slides
    )
    for slide, title in zip(prs.slides, builder.TITLES):
        assert " ".join(title.split()) in " ".join(builder._slide_text(slide).split())
        note = slide.notes_slide.notes_text_frame.text
        assert all(heading in note for heading in builder.NOTE_HEADINGS)

    regional_text = builder._slide_text(prs.slides[3])
    for phrase in (
        "Cortex / PFC",
        "7-month hippocampus",
        "7- and 10-month astrocytes",
        "Some PFC samples reuse",
        "cross-region concordance",
        "not direct PFC replication",
    ):
        assert phrase.lower() in regional_text.lower()

    status = read_tsv(
        report_dir / "gse143758_mouse_validation_limitations_status.tsv"
    )[0]
    assert status["validation_status"] == "awaiting_visual_review"
    assert status["visual_review_status"] == "pending"
    assert status["deck_sha256"] == hashlib.sha256(built.read_bytes()).hexdigest()
    assert status["source_doc_sha256"] == builder.SOURCE_DOC_SHA256

    checks = read_tsv(
        report_dir / "gse143758_mouse_validation_limitations_checks.tsv"
    )
    visual = [row for row in checks if row["check_id"] == "visual_review"]
    assert len(visual) == 1 and visual[0]["passed"] == "False"
    assert all(
        row["passed"] == "True" for row in checks if row["check_id"] != "visual_review"
    )


def test_published_deck_and_report_are_complete() -> None:
    builder.validate_deck(builder.DEFAULT_OUT)
    report_dir = builder.DEFAULT_REPORT_DIR
    status = read_tsv(
        report_dir / "gse143758_mouse_validation_limitations_status.tsv"
    )[0]
    assert status["validation_status"] == "validated_complete"
    assert status["visual_review_status"] == "complete"
    assert status["slides"] == "5"
    assert status["deck_path"] == builder.display_path(builder.DEFAULT_OUT)
    assert status["deck_sha256"] == builder.sha256(builder.DEFAULT_OUT)
    assert status["source_doc_sha256"] == builder.SOURCE_DOC_SHA256

    checks = read_tsv(
        report_dir / "gse143758_mouse_validation_limitations_checks.tsv"
    )
    assert checks and all(row["passed"] == "True" for row in checks)

    manifest = read_tsv(
        report_dir / "gse143758_mouse_validation_limitations_input_manifest.tsv"
    )
    assert {row["artifact_role"] for row in manifest} == {
        "source_docx", "builder", "deck"
    }
    for row in manifest:
        path = REPO / row["path"] if not Path(row["path"]).is_absolute() else Path(row["path"])
        assert path.exists()
        assert int(row["bytes"]) == path.stat().st_size
        assert row["sha256"] == builder.sha256(path)
