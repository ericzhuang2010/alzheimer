from __future__ import annotations

import csv
import hashlib
import importlib.util
import sys
from pathlib import Path

from pptx import Presentation


REPO = Path(__file__).resolve().parents[2]
BUILDER_PATH = REPO / "scripts/presentations/build_mouse_dataset_comparison_deck.py"
SPEC = importlib.util.spec_from_file_location("mouse_dataset_comparison_builder", BUILDER_PATH)
assert SPEC and SPEC.loader
builder = importlib.util.module_from_spec(SPEC)
sys.modules[SPEC.name] = builder
SPEC.loader.exec_module(builder)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle, delimiter="\t"))


def test_source_and_criterion_contract() -> None:
    assert _sha256(builder.SOURCE_MD) == builder.SOURCE_MD_SHA256
    builder.validate_source_contract()
    assert len(builder.CRITERIA) == 8
    assert len(builder.DATASETS) == 12
    assert len(builder.TITLES) == 14
    assert len(builder.criterion_matrix_rows()) == 96
    assert all(len(profile.criteria) == 8 for profile in builder.DATASETS)

    by_accession = {profile.accessions: profile for profile in builder.DATASETS}
    assert by_accession["GSE163857"].criteria[0].state == "NOT MATCHED"
    assert by_accession["GSE163857"].criteria[1].state == "STRONG"
    assert by_accession["GSE163857"].criteria[2].state == "STRONG"
    assert by_accession["GSE163857"].criteria[3].state == "STRONG"
    gse143758 = by_accession["GSE143758"]
    assert gse143758.criteria[0].state == "PARTIAL"
    assert "broad atlas is hippocampal" in gse143758.criteria[0].fact
    assert "one WT and one 5xFAD" in gse143758.note_focus
    assert "not replicated all-cell PFC validation" in gse143758.safe_claim
    assert by_accession["GSE140399"].criteria[5].state == "PARTIAL"
    assert by_accession["GSE127884"].criteria[5].state == "NOT MATCHED"


def test_build_validate_and_reports(tmp_path: Path) -> None:
    deck = tmp_path / "comparison.pptx"
    report_dir = tmp_path / "report"
    built = builder.build_deck(deck)
    builder.write_report(built, report_dir, visual_review_status="pending")
    builder.validate_deck(built)

    prs = Presentation(built)
    assert len(prs.slides) == 14
    assert prs.slide_width == builder.SLIDE_W
    assert prs.slide_height == builder.SLIDE_H
    assert builder.SOURCE_MD_SHA256 in (prs.core_properties.comments or "")
    assert all(
        all(shape.shape_type != 13 for shape in slide.shapes)
        for slide in prs.slides
    )

    for slide in list(prs.slides)[2:]:
        visible = "\n".join(
            shape.text for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert all(criterion in visible for criterion in builder.CRITERIA)
        assert "BEST USE" in visible
        assert "DO NOT CLAIM" in visible

    all_text = "\n".join(
        "\n".join(
            shape.text for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        ) + "\n" + slide.notes_slide.notes_text_frame.text
        for slide in prs.slides
    )
    for forbidden in (
        "Cortex/PFC:", "cortex/PFC mice", "n=1/cell",
        "factorial cells", "control cells have n=1–3",
    ):
        assert forbidden not in all_text
    assert "all 8 APOE×sex×disease groups" in all_text
    assert "male APOE3 control group" in all_text
    assert "Cortical data exist, but use them only" in all_text

    matrix = _read_tsv(report_dir / "mouse_dataset_comparison_criterion_matrix.tsv")
    assert len(matrix) == 96
    assert {row["criterion"] for row in matrix} == set(builder.CRITERIA)
    assert {row["accession_group"] for row in matrix} == {
        profile.accessions for profile in builder.DATASETS
    }

    status = _read_tsv(report_dir / "mouse_dataset_comparison_status.tsv")
    assert len(status) == 1
    assert status[0]["validation_status"] == "awaiting_visual_review"
    assert status[0]["visual_review_status"] == "pending"
    assert status[0]["slides"] == "14"
    assert status[0]["criteria_slides"] == "2"
    assert status[0]["dataset_slides"] == "12"
    assert status[0]["criterion_rows"] == "96"
    assert status[0]["deck_sha256"] == _sha256(built)


def test_published_deck_and_report_are_current() -> None:
    assert builder.DEFAULT_OUT.exists()
    builder.validate_deck(builder.DEFAULT_OUT)
    status_path = builder.DEFAULT_REPORT_DIR / "mouse_dataset_comparison_status.tsv"
    matrix_path = builder.DEFAULT_REPORT_DIR / "mouse_dataset_comparison_criterion_matrix.tsv"
    manifest_path = builder.DEFAULT_REPORT_DIR / "mouse_dataset_comparison_input_manifest.tsv"
    assert status_path.exists() and matrix_path.exists() and manifest_path.exists()

    status = _read_tsv(status_path)
    assert len(status) == 1
    assert status[0]["validation_status"] == "validated_complete"
    assert status[0]["visual_review_status"] == "complete"
    assert status[0]["deck_path"] == "docs/presentations/mouse_dataset_comparison_08232026.pptx"
    assert status[0]["deck_sha256"] == _sha256(builder.DEFAULT_OUT)
    assert status[0]["source_markdown_sha256"] == builder.SOURCE_MD_SHA256

    manifest = {row["artifact_role"]: row for row in _read_tsv(manifest_path)}
    assert manifest["source_markdown"]["sha256"] == builder.SOURCE_MD_SHA256
    assert manifest["builder"]["sha256"] == _sha256(BUILDER_PATH)
    assert manifest["deck"]["sha256"] == _sha256(builder.DEFAULT_OUT)
    assert len(_read_tsv(matrix_path)) == 96
