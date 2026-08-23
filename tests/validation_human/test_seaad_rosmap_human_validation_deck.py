from __future__ import annotations

import csv
import hashlib
import importlib.util
import zipfile
from pathlib import Path

from pptx import Presentation


REPO = Path(__file__).resolve().parents[2]
SCRIPT = REPO / "scripts/presentations/build_seaad_rosmap_human_validation_deck.py"


def load_builder():
    spec = importlib.util.spec_from_file_location("seaad_human_validation_deck", SCRIPT)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle, delimiter="\t"))


def slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
    )


def test_derived_values_preserve_scientific_denominators():
    builder = load_builder()
    metrics = builder.derive_metrics()
    assert builder.DEFAULT_OUT == (
        REPO / "docs/presentations/seaad_rosmap_human_validation_08252026.pptx"
    )
    assert metrics["donors"] == 78
    assert metrics["supertypes"] == 129
    assert metrics["analysis_role"] == "post-hoc exploratory"
    assert metrics["query_rule_id"] == "fdr_only_query_sensitivity"
    assert metrics["donor_minimum"] == 3
    assert metrics["fdr_threshold"] == 0.05
    assert metrics["planned_contrasts"] == 774
    assert metrics["completed_contrasts"] == 381
    assert metrics["not_estimable_contrasts"] == 393
    assert dict(metrics["completed_groups"]) == {
        "F_e33": 105,
        "F_e4": 95,
        "M_e33": 104,
        "M_e4": 77,
    }
    assert metrics["planned_directions"] == 1548
    assert metrics["completed_directions"] == 762
    assert metrics["kda_calls"] == 42
    assert metrics["kda_return_positive_calls"] == 27
    assert metrics["kda_no_return_calls"] == 15
    assert metrics["significant_return_rows"] == 201
    assert metrics["candidate_units"] == 38_788
    assert metrics["selected_units"] == 11
    assert metrics["selected_genes"] == 9
    assert metrics["selected_mt_units"] == 8
    assert metrics["selected_non_mt_units"] == 3
    assert dict(metrics["list_state_counts"]) == {
        "ranked_candidates": 4,
        "no_passing_candidate": 6,
        "not_testable_no_included_runs": 4,
    }
    assert metrics["rosmap_units"] == 47
    assert metrics["rosmap_testable"] == 36
    assert metrics["strict_units"] == 6
    assert metrics["strict_genes"] == 4
    assert metrics["strict_gene_symbols"] == ["MT-CO2", "MT-CO3", "MT-CYB", "MT-ND5"]
    assert metrics["non_mt_rosmap_units"] == 21
    assert metrics["non_mt_testable"] == 17


def test_builds_validated_deck_with_notes_alt_text_and_reports(tmp_path):
    builder = load_builder()
    output = tmp_path / "seaad_rosmap_human_validation.pptx"
    report = tmp_path / "report"
    result = builder.build_deck(
        output, report, visual_review_status="complete"
    )
    assert result == output.resolve()
    builder.validate_deck(output)

    presentation = Presentation(output)
    assert len(presentation.slides) == 15
    assert presentation.slide_width == builder.SLIDE_W
    assert presentation.slide_height == builder.SLIDE_H
    assert len(builder.MAIN_TITLES) == 9
    assert len(builder.APPENDIX_TITLES) == 6

    for slide in presentation.slides:
        note = slide.notes_slide.notes_text_frame.text
        assert len(note.split()) >= builder.MIN_NOTE_WORDS
        assert all(heading in note for heading in builder.NOTE_HEADINGS)

    text = "\n".join(
        slide_text(slide) + "\n" + slide.notes_slide.notes_text_frame.text
        for slide in presentation.slides
    )
    for required in [
        "11 SEA-AD selected units",
        "Post-hoc exploratory",
        "same-network MT matches",
        "Only 42 directions produced mitochondrial gene sets large enough for KDA",
        "Six ROSMAP units reappear in the same neuronal network and driver class",
        "Zero non-MT overlap does not mean the biology is absent",
        "Four of six sex/APOE groups",
        "FDR < 0.05; no fold-change cutoff",
        "≥3 donors per disease arm",
        "27 of 42",
        "201",
        "Mitochondrial genes define the query",
    ]:
        assert required in text
    for forbidden in [
        "13 SEA-AD selected units",
        "520 directions",
        "Only 260 of 774",
        "29 with at least one",
        "13 with none",
        "208 return rows",
        "RPL30",
        "KANSL1L",
        "fewer than five independent donors",
        "completed primary analysis",
        "No optional sensitivity branch executed",
        "|log₂FC| > log₂(1.3)",
        "84 calls",
        "six unique strict genes",
        "failed replication",
        "1,548 KDA calls",
    ]:
        assert forbidden.lower() not in text.lower()

    pictures = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == 13
    ]
    assert len(pictures) == 8
    assert all(
        shape._element.xpath(".//p:cNvPr")[0].get("descr", "")
        for shape in pictures
    )
    for slide_number, figure_label in builder.FIGURE_SLIDES.items():
        slide_pictures = [
            shape
            for shape in presentation.slides[slide_number - 1].shapes
            if shape.shape_type == 13
        ]
        assert len(slide_pictures) == 1
        picture = slide_pictures[0]
        assert hashlib.sha256(picture.image.blob).hexdigest() == builder.sha256(
            builder.FIG[figure_label]
        )
        assert picture.crop_left == 0.0
        assert picture.crop_right == 0.0
        assert picture.crop_top == 0.0
        assert picture.crop_bottom == 0.0

    with zipfile.ZipFile(output) as archive:
        assert archive.testzip() is None
        media = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        embedded = {
            builder.hashlib.sha256(archive.read(name)).hexdigest()
            for name in media
        }
        assert {builder.sha256(path) for path in builder.FIG.values()} <= embedded

    status = read_tsv(report / "seaad_rosmap_human_validation_status.tsv")
    assert len(status) == 1
    assert status[0]["validation_status"] == "validated_complete"
    assert status[0]["visual_review_status"] == "complete"
    assert status[0]["schema_version"] == "seaad_rosmap_human_validation_deck_v2"
    assert status[0]["main_slides"] == "9"
    assert status[0]["appendix_slides"] == "6"
    assert status[0]["total_slides"] == "15"
    assert Path(status[0]["deck_path"]) == output.resolve()
    assert int(status[0]["deck_bytes"]) == output.stat().st_size
    assert status[0]["deck_sha256"] == builder.sha256(output)

    checks = read_tsv(report / "seaad_rosmap_human_validation_checks.tsv")
    assert all(row["passed"] == "True" for row in checks)
    expected_check_ids = {
        "pptx_integrity",
        "slide_count",
        "widescreen",
        "speaker_notes",
        "validated_figures",
        "embedded_media_identity",
        "slide_figure_identity",
        "planned_direction_arithmetic",
        "deg_completion_counts",
        "active_fdr_only_signal",
        "kda_call_arithmetic",
        "kda_return_counts",
        "selection_counts",
        "strict_overlap_counts",
        "testability_denominator",
        "active_tier_identity",
        "active_thresholds",
        "visual_review",
    }
    assert expected_check_ids <= {row["check_id"] for row in checks}
    inventory = read_tsv(report / "seaad_rosmap_human_validation_slide_inventory.tsv")
    assert len(inventory) == 15
    assert [row["section"] for row in inventory].count("main") == 9
    assert [row["section"] for row in inventory].count("appendix") == 6


def test_pending_build_is_not_reported_as_validated_complete(tmp_path):
    builder = load_builder()
    output = tmp_path / "seaad_rosmap_human_validation_08252026.pptx"
    report = tmp_path / "report"
    builder.build_deck(output, report, visual_review_status="pending")
    status = read_tsv(report / "seaad_rosmap_human_validation_status.tsv")
    assert status[0]["validation_status"] == "awaiting_visual_review"
    assert status[0]["visual_review_status"] == "pending"
    checks = read_tsv(report / "seaad_rosmap_human_validation_checks.tsv")
    visual = [row for row in checks if row["check_id"] == "visual_review"]
    assert visual == [{
        "check_id": "visual_review",
        "passed": "False",
        "detail": "Visual review remains pending",
    }]


def test_canonical_circles_are_the_only_circle_inputs():
    builder = load_builder()
    expected_parent = (
        REPO / "results/figures/validation_human/seaad_two_case_circular"
    )
    assert builder.FIG["mt_circle"].parent == expected_parent
    assert builder.FIG["non_mt_circle"].parent == expected_parent
    assert "circular_slide" not in str(builder.FIG["mt_circle"])
    assert "summary_slide" not in str(builder.FIG["non_mt_circle"])
