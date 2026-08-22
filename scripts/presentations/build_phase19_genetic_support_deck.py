#!/usr/bin/env python3
"""Build the human-genetic-support presentation.

The deck follows ``phase19_presentation_slide_design.md`` and embeds six
validated, slide-native genetic-support figures.  The main story is organized
into three sections with overview and divider slides; supporting slides retain
candidate lists, methods, dense matrices, and provenance caveats.
"""

from __future__ import annotations

import argparse
import hashlib
import os
import subprocess
import tempfile
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
FIG_ROOT = REPO / "results/figures/analysis/phase_19_genetic_support"
RESULTS = REPO / "results/minerva_production"
DEFAULT_OUT = REPO / "docs/presentations/human_genetic_support_for_key_drivers.pptx"

FIG = {
    "workflow": FIG_ROOT / "gated_workflow/genetic_support_gated_workflow.png",
    "tier1": FIG_ROOT / "tier1_summary/genetic_support_tier1_slide_summary.png",
    "csf": FIG_ROOT / "csf_outcome_summary/genetic_support_csf_outcome_summary.png",
    "non_apoe": FIG_ROOT / "non_apoe_evidence/genetic_support_non_apoe_evidence_cards.png",
    "tier2": FIG_ROOT / "tier2_route_attrition/genetic_support_tier2_route_attrition.png",
    "ad_pvalues": FIG_ROOT / "ad_nearby_pvalues/genetic_support_ad_nearby_pvalues.png",
}

AUX = {
    "tier1_matrix": RESULTS / "19_genetic_support_tier1/genetic_support_evidence_matrix.png",
    "tier1_loci": RESULTS / "19_genetic_support_tier1/genetic_support_locus_plots.pdf",
    "recovery_matrix": RESULTS / "19_genetic_support_tier2_recovery/recovery_evidence_matrix.png",
    "recovery_loci": RESULTS / "19_genetic_support_tier2_recovery/recovery_locus_plots.pdf",
    "csf_matrix": RESULTS / "19_genetic_support_endophenotype_gwas_qtl_extension/endophenotype_evidence_matrix.png",
}

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 35, 61)
NAVY_2 = RGBColor(30, 59, 91)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
AMBER = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(221, 229, 238)
MID = RGBColor(103, 116, 132)
DARK = RGBColor(47, 47, 47)
GRAY = RGBColor(77, 77, 77)
NO_SUPPORT = RGBColor(189, 189, 189)
PALE_BLUE = RGBColor(230, 242, 249)
PALE_GREEN = RGBColor(229, 244, 239)
PALE_AMBER = RGBColor(255, 246, 224)
PALE_RED = RGBColor(253, 235, 228)
PALE_GRAY = RGBColor(242, 244, 247)

FONT = "Arial"

MAIN_TITLES = [
    "Do genes highlighted by brain-cell networks also show inherited links to Alzheimer's disease?",
    "At a glance: APOE has multi-dataset support, with four focused validation targets",
    "The presentation moves from study design to evidence and future validation",
    "Study design and public data",
    "Network and genetic evidence provide complementary views of disease biology",
    "We combined five kinds of public data, each with a different job",
    "Gene-activity datasets add cell, tissue, RNA, and protein evidence",
    "A step-by-step design organizes evidence and future validation",
    "Evidence found across the gene list",
    "The first public-data screen highlighted APOE, COX7C, and SELENOW",
    "APOE was supported across Alzheimer's disease and all three spinal-fluid markers",
    "Four additional genes provide specific signals for focused validation",
    "The 54 planned comparisons map clear next steps for same-variant testing",
    "How the evidence guides future validation",
    "Exact P values and data coverage guide the next analyses",
    "The current study provides a focused foundation for broader validation",
    "Next: validate the strongest leads and broaden the genetic evidence",
    "The evidence highlights APOE and four focused paths for follow-up",
]

APPENDIX_TITLES = [
    "Supporting details",
    "The original list: 25 genes in 47 gene–network settings",
    "Dataset details: what each source was used for",
    "Pre-set rules made every evidence comparison consistent",
    "Four gene regions stood out below 5×10⁻⁸; all 19 P values are shown",
    "APOE led the registered first-screen results across 47 settings",
    "Four regions are focused priorities for gene-level validation",
    "APOE evidence across Alzheimer's disease and three spinal-fluid markers",
    "RPS15 has three brain-QTL records for focused validation",
    "Data-package improvements that will strengthen reproducibility",
]

MAIN_NOTES = [
    """Teaching goal: Introduce the question and distinguish cell-network importance from inherited genetic association. Earlier work used brain-cell networks to identify “key-driver” genes: genes that sit in influential positions within patterns of genes that turn on and off together. This presentation asks whether inherited DNA data provide an additional line of support for those genes.

Walk through the slide: Point to “25 genes from 47 cell-network results.” There are 25 unique genes, but some appeared in more than one cell type or network. Each gene in a particular network is counted as one gene–network setting, giving 47 settings in total. The list was chosen before this genetic follow-up, so the genetic results were not used to rewrite or rerank it. Then point to the evidence statement. APOE is the clearest lead because it connects to several public datasets, including Alzheimer’s disease itself and three biomarkers measured in cerebrospinal fluid. COX7C, SELENOW, RPS15, and ANKRD11 have specific signals that make them useful targets for focused validation.

Key idea to explain: Network evidence describes relationships among gene activities inside cells. Genetic evidence asks whether inherited DNA differences are statistically associated with disease or disease-related traits. Combining them gives two complementary views rather than one replacing the other.

Scientific boundary: The four follow-up genes do not all have the same kind or strength of evidence, so later slides describe each precisely. The node-and-line image is an abstract network symbol, not a biological pathway diagram. None of these results alone proves that a gene causes disease.

Transition: “I’ll begin with the full study in one picture, then unpack how each kind of evidence was obtained.”""",
    """Teaching goal: Give the audience a one-minute map of the study and show that the evidence types have different meanings.

Walk through the slide: Start at the upper left. The starting list contained 25 unique genes appearing in 47 gene-by-cell-network combinations. The same gene can appear in more than one combination. Move to “Nuclear analysis.” Nineteen genes are encoded on chromosomes in the cell nucleus and could be examined with the common-variant datasets used here. Six genes are encoded by mitochondrial DNA, a small separate genome inside mitochondria. Those six define a dedicated mitochondrial follow-up.

Point to APOE. It has the broadest support in this project: an established Alzheimer’s-disease link plus signals for amyloid-β42, total tau, and phosphorylated tau 181 in cerebrospinal fluid. These molecules are commonly studied as markers of Alzheimer-related biology. COX7C and SELENOW matched records in the registered public summary screen. RPS15 and ANKRD11 had very small P values for DNA regions near them.

Key idea to explain: A P value describes how surprising a result would be under a no-association model; smaller values draw attention. The “same-variant test” at lower right would ask whether an Alzheimer’s signal and a gene-activity signal are best explained by the same DNA variant.

Scientific boundary: A nearby regional signal does not by itself identify the responsible gene. A same-variant result would strengthen a biological connection, but would still need other evidence before making a causal claim.

Transition: “Now that we have the headline results, let’s see how the rest of the presentation is organized.”""",
    """Teaching goal: Help the audience follow the logic of the talk rather than treating the slides as a long list of datasets and genes.

Walk through the slide: Point to Section 1. Slides 4 through 8 explain the original gene list, the public datasets, and the order of the tests. This foundation matters because each database answers a different question. Point to Section 2. Slides 9 through 13 summarize all 47 gene–network settings, examine APOE across several datasets, and then present the particular signals involving COX7C, SELENOW, RPS15, and ANKRD11. Point to Section 3. Slides 14 through 18 show how exact P values and available data lead to concrete next steps, including same-variant comparisons, larger cell-specific studies, more ancestral populations, and additional genetic mechanisms.

Key idea to explain: Good scientific interpretation links each conclusion to the type of data that supports it. APOE has multi-dataset support, while the other four genes provide different and more focused starting points for validation.

Scientific boundary: The roadmap groups related topics for teaching purposes; it does not place every highlighted gene into one identical evidence category.

Transition: Point to the appendix ribbon and say, “Detailed cutoffs, all 19 nearby-Alzheimer’s P values, full evidence tables, and reproducibility checks are available at the end. First, let’s define the study design.”""",
    """Teaching goal: Prepare the audience for three recurring terms—GWAS, QTL, and CSF—and explain the scope of the study.

Walk through the slide: Begin with GWAS, or genome-wide association study. A GWAS compares DNA variants across many people to find variants associated with a trait such as Alzheimer’s disease. A DNA variant is a genome position at which people can carry different DNA letters. Next define QTL, or quantitative trait locus. A QTL connects a DNA variant to a measurable biological feature, such as the amount of RNA made from a gene, the way RNA is spliced, or the amount of a protein. Finally define CSF as cerebrospinal fluid, the fluid surrounding the brain and spinal cord. The study used public GWAS results for three CSF markers: amyloid-β42, total tau, and phosphorylated tau 181.

Key idea to explain: These datasets create a chain of questions. GWAS asks where inherited DNA relates to a trait. QTL data ask whether DNA in that region relates to gene activity. Comparing both can prioritize a possible gene-level explanation.

Scientific boundary: Association is a statistical link; it does not automatically identify the causal variant, gene, or mechanism. CSF markers reflect disease-related biology but are not the same outcome as a clinical diagnosis.

Transition: “With those definitions in place, we can compare what the original network analysis and the inherited-DNA analysis each contribute.”""",
    """Teaching goal: Show why network biology and inherited genetics can produce different but complementary information.

Walk through the slide: Point first to the green box. A cell-network analysis studies genes whose activity rises and falls together in a particular cell type or disease condition. A gene near the center may help organize or reflect an important cellular process. Here, that analysis produced 47 gene-network results representing 25 unique genes. Point next to the blue box. An inherited-DNA analysis asks whether DNA variants passed through families show a statistical relationship with Alzheimer’s disease or a related marker. Nineteen nuclear genes could be examined in the common-variant workflow. Six mitochondrial genes define a dedicated design because mitochondrial DNA is inherited and analyzed differently from nuclear chromosomes.

Key idea to explain: Emphasize “complementary questions.” A gene can be important in disease biology even if it does not change inherited Alzheimer’s risk. A gene can be central in diseased cells because it responds to inflammation, aging, treatment, or changes elsewhere in the system. Conversely, a genetic-risk region may act only in a cell state or time point that the network dataset did not capture. The two approaches therefore need not produce identical lists.

Scientific boundary: “Linked to inherited risk” means statistically associated, not proven to cause disease. A network-central gene may be a response to disease rather than an inherited cause, and a nearby DNA association may involve another gene in the region.

Transition: “To see how we investigated these complementary questions, let’s look at the five public-data layers and the job assigned to each.”""",
    """Teaching goal: Explain why five public-data layers were needed and what unique question each one answered.

Walk through the slide: Read the table by rows. The first row documents the 25 genes and 47 settings. GENCODE and HGNC standardize gene locations and official names, preventing duplicate labels and separating 19 nuclear genes from 6 mitochondrial genes. The second row is a quick public-summary screen. FunGen-xQTL collects previously computed results linking Alzheimer’s genetics to gene regulation. Its snapshot identifier is shown so the same release can be retrieved.

The third row is the Bellenguez Alzheimer’s GWAS, accession GCST90027158. GRCh38 is the genome coordinate system used to locate variants. For each nuclear gene, the screen examined one megabase—one million DNA letters—on either side. The fourth row contains brain QTL datasets, which ask whether variants affect RNA amount, RNA splicing, or protein amount. The final row covers amyloid-β42, total tau, and p-tau181 in cerebrospinal fluid. Nineteen genes times three markers produced 57 screens.

Key idea to explain: Dataset versions and accession codes provide an audit trail. “Summary statistics” are association results calculated across groups, not individual medical records.

Scientific boundary: A public-summary match is a lead that must be interpreted in its original context. A signal within the ±1-megabase window is regional evidence because several nearby genes can share the same associated region. The CSF studies listed here represent European-ancestry samples, which matters when considering how broadly results may generalize.

Transition: “The QTL layer deserves a closer look, because it is the bridge from a DNA variant to a possible change in gene activity.”""",
    """Teaching goal: Teach the three QTL types and show why cell type, tissue, and sample size affect interpretation.

Walk through the slide: Define an eQTL as a variant associated with RNA amount, an sQTL as a variant associated with RNA splicing, and a pQTL as a variant associated with protein amount. Move across the upper cards. Microglia are immune-like brain cells; Young 2019 includes 104 samples with expression and splicing QTLs. Aygun 2021 includes 73 neuron-like samples. Walker 2019 includes 211 mixed-brain samples. Mixed tissue gives broader coverage, although a cell-specific signal may be diluted or reflect several cell types. NG00184.v1 provides wider RNA, splicing, and protein coverage.

Point to the lower-left box. APOE follow-up used spinal-fluid pQTL data from 3,506 European-ancestry samples after earlier association signals crossed the pre-set screening references. At lower right, RPS15 had three distinct mixed-brain QTL tracks. These produced six positive setting rows because each track was evaluated in two network settings; they are not six independent studies.

Key idea to explain: A QTL is a statistical bridge from inherited DNA to a measurable molecular feature. Regulation can be cell-specific, so the biological source of the sample matters as much as the QTL label.

Scientific boundary: QTL associations do not by themselves prove a causal chain from variant to gene to disease. Neuron-like samples are not an exact match for every mature neuron subtype, and mixed-brain results do not identify the contributing cell type.

Transition: “Now we can combine the disease and gene-activity layers into one step-by-step workflow.”""",
    """Teaching goal: Explain the two analysis paths and why compatible variant-level inputs are required before a same-signal test.

Walk through the slide: Start at the dark box. The 47 gene–network pairs, representing 25 genes, were fixed before the genetic results were examined. Follow upper Lane A. It searches a fixed FunGen snapshot for previously computed Alzheimer’s fine-mapping and gene-regulation results. This lookup highlighted APOE, COX7C, and SELENOW across four gene–network pairs.

Trace Lane B from left to right. Step 1 asks whether a nearby region is associated with clinical Alzheimer’s disease or a CSF marker. Step 2 asks whether variants there also affect the candidate gene’s RNA, splicing, or protein. Step 3 aligns variant identifiers, genome coordinates, statistical models, ancestry, and linkage disequilibrium. Linkage disequilibrium, or LD, describes how nearby variants tend to be inherited together. Step 4 is a colocalization test. PP.H4 is the model-based probability that the disease and gene-activity associations share one underlying variant.

Key idea to explain: Each completed step adds a different layer. A regional disease signal plus a gene-activity signal is more informative than either alone; compatible full data enable the next same-variant comparison.

Scientific boundary: A public-data match depends on its original evidence type. PP.H4 requires complete, compatible variant-level inputs and is a probability under a statistical model, not proof of causation. The arrows show analysis order, not a biological causal pathway. The 6 mitochondrial genes across 20 settings require mitochondrial-specific data and methods.

Transition: “With the workflow established, we can now examine the evidence found across the complete gene list.”""",
    """Teaching goal: Preview the results section and set expectations for how different forms of evidence will be compared.

Walk through the slide: First, “All 47 settings” means 47 gene-in-network combinations representing 25 unique genes. The same gene may occur in more than one cell network, so a count of settings is not a count of independent genes or experiments. The next slide summarizes the registered public-data screen across this complete set.

Second, “APOE across datasets” signals the strongest integrated result here. APOE has an established inherited association with Alzheimer’s disease in the public data and met the study’s regional and gene-level screening rules for amyloid-β42, total tau, and p-tau181. Third, the four focused leads are COX7C, SELENOW, RPS15, and ANKRD11. COX7C and SELENOW appeared in the public-summary screen. RPS15 and ANKRD11 sit near notable Alzheimer’s regional signals, and RPS15 also has brain gene-activity records worth carrying forward.

Key idea to explain: The evidence is heterogeneous, meaning it comes from different methods and answers different questions. Accurate interpretation keeps those categories separate while using them to prioritize clear follow-up tests.

Scientific boundary: APOE’s broad support does not demonstrate every possible cell-specific regulatory mechanism. A regional signal near RPS15 or ANKRD11 may involve another nearby gene, so it is a reason to investigate rather than an assignment of responsibility.

Transition: “We will start with the full 47-setting scorecard, then zoom in from APOE to the four focused follow-up paths.”""",
    """Teaching goal: Show how the first screen was counted and explain why the three highlighted genes represent different strengths and types of genetic evidence.

Walk through the slide: Start with panel A. A “gene–network pair” means one gene considered in one brain-cell network. The same gene can therefore appear more than once. Across 47 pairs, the registered public-data search found one strong source-grade match and three weak source-grade matches; none was assigned the moderate source grade. The strong match was APOE in the astrocyte network. The three weak rows were COX7C in two networks and SELENOW in one network. Importantly, the two COX7C rows came from one bulk-brain RNA-splicing record, so they are two network applications of one observation, not two independent replications.

In panel B, APOE is linked directly to the Alzheimer’s-associated DNA variant rs429358. Its reported Alzheimer’s disease P value was about 1.88 × 10⁻¹⁵⁵, an exceptionally small value, and its source inclusion score was 1.0. COX7C was connected to rs2010322 in a bulk-brain splicing summary; that record’s Alzheimer’s P value was about 2.64 × 10⁻⁶. SELENOW appeared in a published list linking genetically predicted gene activity to Alzheimer’s disease, although that list did not provide the exact model score or neuron subtype.

Panel C shows constructive next routes. Sixteen nuclear genes can be checked in more datasets, while six mitochondrial-DNA genes require methods designed specifically for mitochondrial genetics.

Key idea to explain: A match is a useful lead, and the source grade describes only this registered search—not the gene’s complete biology.

Scientific boundary: A very small P value measures incompatibility with a statistical null model; it is not the probability that a gene causes disease and does not measure effect size. The source inclusion score is also a ranking metric, not a causal probability.

Transition: Next, we ask whether the strongest gene, APOE, also connects to biological markers measured in spinal fluid.""",
    """Teaching goal: Explain why testing biomarkers adds a second, biologically informative layer of support for APOE.

Walk through the slide: Cerebrospinal fluid, or CSF, surrounds the brain and spinal cord. The three biomarkers here are amyloid-β42, total tau, and phosphorylated tau at position 181, called p-tau181. Amyloid and tau are central features of Alzheimer’s pathology, so these quantitative measurements can connect inherited DNA differences to processes closer to disease biology than a case-versus-control diagnosis alone.

Each biomarker study included 18,948 participants of European ancestry. We tested 19 nuclear genes across three biomarkers, giving 19 × 3 = 57 gene–biomarker comparisons. APOE was the highlighted gene for all three biomarkers. In every case, it crossed two pre-set references: a nearby-variant screen using P < 5 × 10⁻⁸ and a corrected whole-gene test using P < 8.77 × 10⁻⁴. The regional minimum for amyloid-β42 was so small that the computer stored it as zero through numerical underflow; this means “smaller than the software could represent,” not a literal probability of zero. The regional minima were 5.4 × 10⁻¹⁶¹ for total tau and 3.27 × 10⁻¹⁷⁴ for p-tau181. The gene-body P value was 5 × 10⁻¹⁰ for each biomarker.

The gray blocks summarize the other 18 genes: 54 comparisons were at or above one or both screening references. They provide a transparent baseline for future datasets rather than a verdict on every possible genetic mechanism.

Key idea to explain: APOE appears consistently across diagnosis and three related biomarkers, which makes it a compelling multi-dataset finding.

Scientific boundary: These associations do not yet show how APOE changes a biomarker, whether one DNA variant explains both signals, or which brain cell carries the mechanism.

Transition: We now examine the distinct evidence found for four additional genes.""",
    """Teaching goal: Compare four non-APOE leads without treating all evidence types as equivalent.

Walk through the slide: Each card separates “what the data show” from “what to validate next.” For COX7C, one published bulk-brain record connects an Alzheimer’s signal with RNA splicing. The regional minimum was 8.579 × 10⁻¹⁴ within one megabase, led by rs62375397. Its Walker brain-expression P value, 2.583 × 10⁻³, was above that dataset’s pre-set reference of 5.167 × 10⁻⁶. Complete splicing and Alzheimer’s variant files can enable a shared-variant comparison.

SELENOW appears in a published predicted-expression gene list. Its nearby regional minimum was 6.410 × 10⁻⁵, above the conservative 5 × 10⁻⁸ reference. Recovering the model score, exact cell annotation, and variant-level data would make this lead more specific.

RPS15 combines two useful signals: a nearby Alzheimer’s minimum of 4.089 × 10⁻³⁰, led by rs12151021, and a bulk-neocortex expression P value of 2.120 × 10⁻⁶, below its 3.759 × 10⁻⁶ reference. Three additional bulk-brain QTL tracks reinforce RPS15 as a focused validation target. A QTL is a DNA region associated with a molecular feature such as RNA amount or splicing.

ANKRD11 has a nearby Alzheimer’s minimum of 1.283 × 10⁻¹¹, led by rs56407236. Its measured brain-expression P value was 1.819 × 10⁻⁴, above the 4.522 × 10⁻⁶ reference, so the regional result motivates more gene-specific testing.

For each of these four genes, zero of three spinal-fluid comparisons were below both pre-set references; those recorded results remain useful for comparison with future datasets.

Key idea to explain: These genes have concrete leads, but each lead calls for a different next experiment or dataset.

Scientific boundary: A small regional P value identifies an associated neighborhood, not necessarily the named nearby gene. P values from different tests are not effect sizes and should not be ranked as biological importance.

Transition: The next slide organizes every planned gene-activity comparison and shows exactly where shared-variant testing can advance.""",
    """Teaching goal: Explain the accounting of all 54 comparisons and introduce the logic of a shared-variant, or colocalization, test.

Walk through the slide: The 54 comparisons come from 27 nuclear gene–network pairs. Each pair had two questions: does nearby DNA variation relate to RNA amount, called an expression QTL or eQTL, and does it relate to how RNA pieces are joined, called a splicing QTL or sQTL? Thus, 27 expression tests plus 27 splicing tests equals 54.

The horizontal bar sorts each comparison into one status. These categories are mutually exclusive; they are not a sequence of 54 samples shrinking step by step. Forty-two comparisons had a nearby Alzheimer’s P value at or above the conservative 5 × 10⁻⁸ reference. Four had regional Alzheimer’s signals but gene-activity P values at or above their source-specific references. Six splicing comparisons have public-file coverage that can be clarified with fuller input files.

The most informative group is the two orange comparisons. APOE in the astrocyte question and RPS15 in the oligodendrocyte precursor cell, or OPC, question each had both a nearby Alzheimer’s signal and a bulk-brain gene-activity signal. These are priorities for a shared-variant test. Such a test asks whether the same underlying DNA variant could explain both association patterns, rather than merely observing two signals in the same broad region.

Panel B reports zero completed primary tests because the required fitted gene-activity model and matching variant-correlation reference were not jointly available. Therefore, the shared-variant probability is unavailable, not zero.

Key idea to explain: The study produced a complete map of where the evidence currently sits and which inputs unlock the next calculation.

Scientific boundary: Two signals in one region can arise from different nearby variants inherited together; only a matched shared-variant analysis can distinguish these possibilities.

Transition: We now turn from reporting results to interpreting what the P values and data coverage tell us to do next.""",
    """Teaching goal: Prepare the audience to interpret the next slides as a validation plan built from measured results, rather than as a simple supported-versus-unsupported scorecard.

Walk through the slide: This section has three questions. First, what do the P values show? A P value asks how surprising the observed data would be under a statistical model with no association. Smaller values indicate stronger statistical disagreement with that model, but they do not tell us the size of the biological effect or prove causation. We will keep every exact value visible and use the pre-set 5 × 10⁻⁸ reference as a conservative screening rule.

Second, what do the current data support? Here we distinguish several evidence levels: a disease-associated DNA region, a signal tied to a named gene’s activity, a signal repeated across datasets or traits, and a test showing that disease and gene activity may share one causal variant. These levels answer different questions, so they should not be collapsed into one label.

Third, what should be validated next? The data point toward concrete tasks: obtain complete variant-level files for APOE and RPS15, study the exact brain-cell types highlighted by the network analysis, repeat findings in independent and more diverse samples, and examine genetic mechanisms beyond common nuclear variants.

Key idea to explain: A carefully organized status map is itself useful: it converts each observation into a testable next step.

Scientific boundary: Values above a conservative screening reference are not statements that a gene has no biological role. Values below it locate statistical associations but do not by themselves name the causal gene or mechanism.

Transition: The next slide makes this logic concrete with all 19 exact regional P values, two priority shared-variant tests, and additional genetic routes.""",
    """Teaching goal: Teach the meaning of the conservative genome-wide reference and connect three evidence categories to specific follow-up analyses.

Walk through the slide: Begin with the first card. Every one of the 19 nuclear gene regions has a reported minimum Alzheimer’s P value. Four regions—near ANKRD11, APOE, COX7C, and RPS15—were below 5 × 10⁻⁸. The other 15 values ranged from 2.929 × 10⁻⁶ to 2.931 × 10⁻⁴ and are listed in the appendix. The very small cutoff is used because a genome-wide association study tests roughly a million independent genetic patterns. A simple multiple-testing idea is 0.05 divided by 1,000,000, which equals 5 × 10⁻⁸. This conservative rule greatly reduces chance findings when many tests are performed.

The second card highlights APOE and one RPS15–OPC comparison. Both had a nearby Alzheimer’s signal and a mixed-brain gene-activity signal. OPCs are oligodendrocyte precursor cells, cells that can develop into the myelin-producing oligodendrocytes of the brain. Complete variant-level gene-activity statistics, a fitted prediction model, and a matching reference for correlations among nearby variants can test whether one variant contributes to both patterns.

The third card reminds us that the current screen covers mainly common nuclear variants acting near a gene. Additional routes include mitochondrial DNA, rare and structural variants, distant gene regulation, gene-by-gene or gene-by-environment interactions, disease-stage effects, and regulation after RNA is produced.

Key idea to explain: The cutoff is a conservative screen, while the exact P values preserve information for future comparison and validation.

Scientific boundary: A P value is not the chance that the null hypothesis is true, and a regional minimum does not assign causality to the nearest gene.

Transition: These results define four practical opportunities to strengthen the study.""",
    """Teaching goal: Present the study’s current outputs as a strong foundation and explain how four kinds of validation add complementary information.

Walk through the slide: The first opportunity is to complete shared-variant tests for APOE and RPS15. This requires more than a list of significant variants. For each variant, we need the effect estimate, uncertainty, alleles, frequency, and sample size; a fitted gene-activity model; and a matching linkage-disequilibrium reference. Linkage disequilibrium means that nearby variants can be inherited together, creating correlated signals that the model must account for.

The second opportunity is larger, cell-specific studies. Astrocytes support neurons, microglia are immune-like brain cells, neurons carry electrical signals, and OPCs produce new oligodendrocytes. Bulk-brain tissue mixes these cells. Larger studies that measure each cell type separately can ask whether the activity signal truly occurs in the network context that first highlighted the gene.

The third opportunity is broader populations and independent samples. Patterns of nearby variant correlation differ among ancestral backgrounds. Studying multiple populations can improve resolution and show how broadly a result applies. Repeating the analysis in people who were not part of the discovery dataset provides a stronger test of reproducibility.

The fourth opportunity expands the mechanisms examined: rare variants, large structural changes, mitochondrial DNA, distant regulatory effects, interactions, disease timing, and laboratory perturbations. These can capture biology outside the common, nearby nuclear-variant screen.

Key idea to explain: Validation is a ladder: strengthen the statistical link, locate it in the correct cell, test its generality, and then examine mechanism.

Scientific boundary: Each new layer answers a different question; even a convincing genetic association is distinct from demonstrating the molecular pathway in a laboratory.

Transition: The next slide turns these opportunities into an ordered roadmap.""",
    """Teaching goal: Explain the order of the proposed work and why the roadmap moves from focused statistical validation to broad genetic and laboratory testing.

Walk through the slide: Step one is to deepen APOE and RPS15. These are the most efficient first targets because both already connect an Alzheimer’s-associated region with a gene-activity signal. The needed additions are complete gene-activity statistics, matched DNA-correlation reference data, and larger studies of the exact cell types—especially astrocytes for APOE and OPCs or inhibitory neurons for RPS15. With these inputs, a signal-aware shared-variant analysis can compare the full patterns across a region.

Step two broadens genetic testing. Protein-based and RNA-based gene tests can ask whether inherited variants predict protein abundance, RNA amount, or RNA splicing. Rare-variant tests combine uncommon changes within a gene, while interaction tests ask whether a variant’s effect depends on another factor. More traits—such as disease progression or brain pathology—and more populations can reveal signals that a single diagnosis dataset may not capture.

Step three adds independent and experimental validation. The six mitochondrial-DNA genes need a dedicated analysis of mitochondrial variants, copy number, and mixtures of mitochondrial DNA sequences within a person. Repeating the network analysis in a separate cohort tests reproducibility. Finally, changing a candidate gene’s activity in cells or model systems and measuring the predicted network response can directly test mechanism.

Key idea to explain: The order is intentional: resolve the most mature leads first, broaden the search second, and test biological causality with independent and laboratory approaches.

Scientific boundary: A shared-variant result would connect two association patterns, but it would still require cell-specific and experimental work to show the direction and molecular consequences.

Transition: The final slide summarizes what the current evidence contributes and the focused paths it opens.""",
    """Teaching goal: Leave the audience with a clear evidence hierarchy and a positive, scientifically careful summary.

Walk through the slide: The first take-home point is APOE. It has the most direct evidence: a fine-mapped Alzheimer’s variant was assigned to the gene, and APOE crossed both nearby-region and whole-gene screening references for amyloid-β42, total tau, and p-tau181. This agreement across diagnosis and three related biomarkers makes APOE the strongest multi-dataset result.

Second, COX7C and SELENOW were highlighted by public gene-activity summaries. COX7C has a bulk-brain RNA-splicing record and a notable regional Alzheimer’s P value. SELENOW appears in a predicted-expression gene list. These sources provide specific starting points for cell-matched, variant-level, and independent validation.

Third, the regions near RPS15 and ANKRD11 were below the conservative 5 × 10⁻⁸ reference. RPS15 also has mixed-brain RNA-amount and RNA-splicing signals, making it a priority for a shared-variant comparison. For ANKRD11, the strong regional association motivates gene-specific follow-up to determine which gene in that neighborhood is connected to the signal.

Fourth, all 19 exact regional P values are preserved. This matters because a transparent baseline lets later studies compare larger samples, additional traits, different ancestries, and new genetic mechanisms without changing the original rules after seeing the results.

Key idea to explain: Evidence was found at several levels—direct gene mapping, gene-activity summaries, and associated regions—and each level now has a focused validation route.

Scientific boundary: These levels should remain distinct. Regional proximity is not causal assignment, and association alone does not establish the exact brain cell or molecular mechanism.

Transition: Close by emphasizing that the study strengthens the APOE finding and turns four additional genes into concrete, testable next questions.""",
]

APPENDIX_NOTES = [
    """Teaching goal: This slide begins the appendix. The main presentation gave the overall story; the appendix shows the supporting information that lets an audience check how that story was built. Think of it as the methods and evidence folder behind a science-fair display. We will move through five kinds of detail: the original gene list, the public datasets, the decision rules, the complete results, and the reproducibility plan.

Walk through the slide: The phrase “original gene list” is important. The genes were chosen by the earlier brain-cell network analysis before these genetic results were examined. Keeping that list fixed helps prevent cherry-picking, which would mean changing the tested genes after seeing which ones produced attractive results. A “gene–network setting” means one gene considered in one brain-cell network. The same gene can therefore appear in more than one setting.

Key idea to explain: The dataset and rule slides answer two basic questions: Where did each number come from, and how was each comparison judged? The evidence slides then show exact P values and source records, including values on both sides of the conservative screening references. Finally, the reproducibility slide lists practical improvements to make the analysis package easier for another researcher to repeat.

Scientific boundary: Technical detail strengthens the positive findings. APOE has the clearest multi-dataset support, while COX7C, SELENOW, RPS15, and ANKRD11 each provide a specific lead for future validation. The appendix explains exactly what kind of lead each one provides and keeps regional, gene-level, and cell-specific claims separate.

Transition: We will begin with the fixed starting list, because every later count and comparison depends on knowing exactly which genes and cell networks were included.""",
    """Teaching goal: This table shows the starting list exactly as it came from the brain-cell network analysis. A biological network is a map of genes whose activity patterns are connected. A “driver” in this setting is a gene that is especially central to that network; it does not automatically mean that inherited DNA changes in the gene cause Alzheimer’s disease. Genetics provides a complementary test of that idea.

Walk through the slide: The rows are seven broad brain-cell networks. Astrocytes support neurons and help maintain the brain environment. Excitatory and inhibitory neurons send signals that increase or decrease activity. Microglia are the brain’s immune cells. OPCs, or oligodendrocyte precursor cells, can mature into oligodendrocytes, which make insulating myelin. Vasculature cells help form and regulate blood vessels.

Key idea to explain: There are 25 unique genes but 47 gene–network settings because several genes appear in multiple cell networks. For example, RPS15 appears in both inhibitory neurons and OPCs. Each appearance is a separate biological question, but repeated appearances are not independent genetic confirmations.

Scientific boundary: The middle column retains an original label, “mitochondria-related driver group.” That label describes the network pattern, not where every gene is encoded. COX7C and UQCR10, for example, are encoded in nuclear DNA. Only six genes here—MT-ATP6, MT-CO2, MT-CO3, MT-CYB, MT-ND4, and MT-ND5—are encoded by mitochondrial DNA. Those six need specialized mitochondrial methods. The remaining 19 nuclear genes were mapped to the GRCh38 human reference genome and screened in nearby regions extending one million DNA letters on either side of each gene.

Transition: Now that the starting units are clear, the next slide explains which public dataset answered each genetic question.""",
    """Teaching goal: This slide is a data map: each row names a source and the specific job it performed. The starting network results supplied the 47 gene–network settings. HGNC checked that each gene symbol was an approved human gene name, and GENCODE version 44 supplied stable gene identifiers and genomic positions on the GRCh38 reference genome.

Walk through the slide: The FunGen-xQTL snapshot was used for the first screen. A first screen is a quick, pre-planned search of already summarized public results; it can identify useful leads without claiming to be an exhaustive search of every genetic study. It combined information from Alzheimer’s genome-wide association studies, fine-mapping, gene-activity studies, and gene-based prediction studies.

The Bellenguez Alzheimer’s dataset, accession GCST90027158, supplied genome-wide summary statistics for clinical diagnosis. A genome-wide association study, or GWAS, compares DNA variants across many people to find variants associated with a trait. “Summary statistics” are variant-level results such as P values and effect estimates; they do not contain individual medical records. We used these data to inspect a window around each of the 19 nuclear genes.

Key idea to explain: NIAGADS NG00184 and eQTL Catalogue data asked whether nearby DNA variants were related to RNA amount, RNA splicing, or protein abundance in brain samples. Such relationships are called quantitative trait loci, or QTLs. Three additional GWAS Catalog datasets tested spinal-fluid amyloid-β42, total tau, and p-tau181, each in 18,948 participants. Focused APOE and RPS15 resources then provided follow-up evidence.

Scientific boundary: These sources are public summary data, so they support broad screening without exposing individual records. Different sources answer different questions, and a source useful for finding a nearby disease signal is not automatically sufficient for proving how a gene works in a specific cell type.

Transition: Recording accession numbers and versions makes clear which release produced each result; next, we will see the pre-set rules used to compare those results consistently.""",
    """Teaching goal: This slide explains the rules used to organize the evidence. They were set before reviewing the outcomes to treat every gene consistently. A P value measures how unusual a result at least this strong would be under a model with no association. It is not the probability that a gene is false, and a cutoff should not turn a continuous result into a statement about all biology. That is why exact P values are reported on both sides of each reference line.

Walk through the slide: For the nearby Alzheimer’s screen, the conservative genome-wide reference was P < 5×10⁻⁸. For a gene-activity dataset, the reference was 0.05 divided by the number of variants tested in that gene region. This division is a multiple-testing correction: testing many variants creates more chances for an apparently interesting result by coincidence. For the spinal-fluid gene tests, 0.05 was divided by 19 genes and three biomarkers, producing 8.77193×10⁻⁴.

Key idea to explain: An eQTL links a DNA variant to RNA amount, an sQTL links it to RNA splicing, and a pQTL links it to protein abundance. The three values 10⁻⁴, 10⁻⁴, and 5×10⁻⁶ are Bayesian starting assumptions, called priors, for a planned shared-signal analysis; they are not measured P values. PP.H4 is the model’s estimated probability that the disease and gene-activity signals share a causal variant, with 0.80 pre-set as the strong-signal reference.

Scientific boundary: PIP, VCP, and CL are scores or labels reported by source studies. They provide useful evidence, but they are not interchangeable with PP.H4. A small P value measures strength against a statistical null model; it does not by itself give effect size, direction, causal gene, or cell type.

Transition: The action labels show what comes next—broaden the search, report the measured value, complete missing inputs, harmonize datasets so variants line up, or use a mitochondrial-DNA-specific method. We can now apply these rules to all 19 nuclear-gene regions.""",
    """Teaching goal: This figure shows the smallest Alzheimer’s GWAS P value found in the region around each of the 19 nuclear genes. Each region includes the gene body plus one million DNA bases on either side. Showing all 19 values is useful because P values are continuous measurements: a value does not lose all information simply because it lies on one side of a screening reference.

Walk through the slide: Why is the reference as small as 5×10⁻⁸? A GWAS tests roughly a million largely independent parts of the genome. If one million tests each used 0.05, chance alone could create about 50,000 results below 0.05 when no true associations existed. A simple correction divides 0.05 by about one million, giving 5×10⁻⁸. This is therefore a conservative screening convention designed for very large numbers of comparisons.

Key idea to explain: Four regions stand out below that reference: the APOE region, the RPS15 region at P = 4.089×10⁻³⁰, the COX7C region at P = 8.579×10⁻¹⁴, and the ANKRD11 region at P = 1.283×10⁻¹¹. The stored APOE regional minimum is zero because of numerical underflow. Underflow means the true number was smaller than the file’s number format could represent, so it was rounded to zero; it does not mean the probability is literally zero.

Scientific boundary: These four regions are valuable validation priorities. However, a regional signal identifies a stretch of linked DNA, not automatically the labeled gene. Nearby genes and regulatory elements may share correlated variants. Gene-activity data, fine-mapping, and same-variant analysis are therefore needed to connect a regional signal to a particular gene and biological mechanism. The remaining 15 exact values provide a transparent baseline for broader future tests.

Transition: The next slide returns to the original 47 settings and shows what the compact first screen found before the deeper regional work.""",
    """Teaching goal: This slide gives the complete audit of the first screen. “Registered” means that the source files, matching rules, and evidence categories were chosen in advance. “First screen” means a compact lookup in public summary tables rather than a full reanalysis of every raw variant dataset. The matrix at left contains all 47 gene–network settings, so it is mainly a reference table. The larger count cards at right summarize it.

Walk through the slide: APOE received the one strong category because the public source directly mapped major Alzheimer’s evidence to APOE, including the well-known variant rs429358. Three setting rows received the limited category: COX7C in astrocytes, COX7C in inhibitory neurons, and SELENOW in excitatory neurons. The two COX7C rows arise from one bulk-brain splicing-QTL record projected onto two network questions. They should therefore be described as two relevant contexts supported by one source observation, not as two independent replications. SELENOW appeared in a public TWAS list. TWAS, or transcriptome-wide association study, tests whether genetically predicted gene activity is associated with a trait.

Key idea to explain: The 23 “broader-source follow-up” rows represent 16 nuclear genes without a direct match in this compact snapshot. That label is an invitation to examine additional sources and mechanisms, not a total-evidence verdict. The 20 mitochondrial-DNA follow-up rows represent six mitochondrial genes across multiple cell networks; the nuclear-DNA workflow was not designed for them.

Scientific boundary: The two COX7C setting rows share one source record, and TWAS-list membership is a lead rather than a complete mechanism. The later, targeted RPS15 analysis is deliberately separate from this fixed first-screen matrix. Keeping the original screen unchanged preserves a clear audit trail, while the newer RPS15 evidence can still guide focused validation.

Transition: We next move from the compact public-summary screen to detailed regional plots and the requirements for assigning a signal to a gene.""",
    """Teaching goal: This slide connects two views of the Alzheimer’s GWAS results. On the left, the matrix records the outcome of each planned gene-and-cell comparison. On the right, four locus plots zoom in on the regions near ANKRD11, APOE, COX7C, and RPS15—the four regional minima below the conservative 5×10⁻⁸ reference.

Walk through the slide: A locus plot usually places genomic position along the horizontal axis and −log10(P) on the vertical axis. This transformation makes tiny P values easier to see: for example, P = 10⁻¹⁰ becomes a height of 10. Each dot represents a DNA variant. Nearby variants often rise together because they are inherited in correlated blocks, a pattern called linkage disequilibrium, or LD. Therefore, the tallest dot identifies an associated region but may not identify the functional variant or responsible gene.

Key idea to explain: The matrix tracks a step-by-step route. First, the region needs an Alzheimer’s signal. Next, the candidate gene needs a measurable gene-activity signal, such as an eQTL or sQTL. Finally, both datasets must describe matching variants and have compatible statistical models and LD reference data. When those ingredients are available, a shared-variant analysis—often called colocalization—can estimate whether the disease signal and gene-activity signal are likely driven by the same underlying DNA change.

Scientific boundary: The four regions are positive, testable leads, but they represent different evidence levels. APOE already has direct gene-level support. RPS15 also has brain QTL leads. COX7C has one public splicing-QTL record, while ANKRD11 currently has regional proximity. A nearby peak alone does not assign the association to the named candidate.

Transition: The next slide shows why APOE stands apart: its evidence connects clinical Alzheimer’s disease with three disease-related spinal-fluid traits.""",
    """Teaching goal: This slide brings several kinds of APOE evidence together. The locus plot at left shows direct Alzheimer’s mapping around APOE, including rs429358, a coding variant that changes the APOE protein and is a well-established Alzheimer’s risk marker. In the source fine-mapping results, rs429358 was included in a 95% credible set. A credible set is a group of variants that a statistical model considers sufficient to contain the causal variant with a stated probability, assuming the model is appropriate.

Walk through the slide: The right-hand matrix shows 19 nuclear genes tested against three spinal-fluid biomarkers, for 57 gene–biomarker combinations. Cerebrospinal fluid, or CSF, surrounds the brain and spinal cord. Amyloid-β42 relates to amyloid plaque biology; total tau reflects the overall amount of tau protein; and p-tau181 measures tau phosphorylated at position 181, a change associated with Alzheimer’s pathology. APOE was the gene with qualifying regional and gene-based results for all three biomarkers.

Key idea to explain: The small table separates two questions. “Strongest nearby P” is the most associated single variant in the region. “Gene-level P” comes from MAGMA, which combines information across variants assigned to a gene; the final column also includes nearby DNA within 10,000 bases. All three APOE gene-level values are far below the pre-set 8.77193×10⁻⁴ reference.

Scientific boundary: For amyloid-β42, the regional P value was stored as zero because of numerical underflow: the true value was smaller than the software format could hold, not literally zero. This multi-trait pattern strengthens APOE as the clearest genetic lead, but a strong gene-level association does not by itself specify which molecular pathway or brain-cell type carries the effect.

Transition: Future same-variant and cell-specific analyses can test whether a particular regulatory signal acts through APOE expression, splicing, or protein abundance in astrocytes. We next apply the same careful logic to the most interesting non-APOE follow-up, RPS15.""",
    """Teaching goal: This slide explains why RPS15 is a useful follow-up candidate. A quantitative trait locus, or QTL, is a DNA region where genetic variation is associated with a measurable molecular trait. An eQTL relates to RNA amount, while an sQTL relates to how RNA pieces are spliced together. The audit planned 37 source-and-context comparisons, measured 31 of them, and found six positive setting rows. Those six rows reduce to three unique bulk-brain source results because each result was relevant to both the OPC and inhibitory-neuron network questions.

Walk through the slide: The three unique tracks are shown in the center. The MSBB BA36 eQTL has P = 2.41403×10⁻⁷. BA36 is a region of temporal cortex. The ROSMAP dorsolateral prefrontal cortex, or DLPFC, sQTL has P = 3.86842×10⁻³⁰, and the ROSMAP posterior-cingulate sQTL has P = 3.30886×10⁻⁷. “Adjusted” values control for the many QTL comparisons. PIP means posterior inclusion probability: within a fine-mapping model, it estimates how strongly a variant belongs to a causal signal. The two ROSMAP tracks have maximum PIP values of 1.0 and 0.910283.

Key idea to explain: These are meaningful gene-activity leads, and the RPS15 neighborhood also has a strong Alzheimer’s regional signal. The zero in the final count card means that a formal shared-variant test has not yet been completed; it is not a probability of zero.

Scientific boundary: The public QTL results came from mixed brain tissue rather than purified OPCs or inhibitory neurons, and complete fitted models plus source-matched LD were unavailable. Six positive setting rows therefore represent three unique source tracks, not six independent discoveries.

Transition: The focused next step is to obtain matched cell-specific data and test whether the Alzheimer’s and RPS15 QTL patterns point to the same DNA variant. The final slide explains how a stronger data package will make that follow-up easier to reproduce.""",
    """Teaching goal: Reproducibility means that another researcher can obtain the same inputs, follow the documented steps, and recover the same results. This final slide turns the file audit into a practical improvement plan. These actions do not change the reported gene-level conclusions; they make the evidence trail easier to verify, reuse, and extend.

Walk through the slide: First, two result files were intended to contain zero rows, but their compressed gzip structure is incomplete. They should be rebuilt as valid empty tables so standard software can open and test them. Second, four APOE protein-QTL files described in the execution report should be added to the formal inventory with their download locations, sizes, and checksums. A checksum is a digital fingerprint used to verify that a file has not changed.

Third, large inputs that were streamed or stored outside the repository need clear retrieval instructions. Fourth, absolute paths beginning with an original computer’s /home/... directory should become project-relative paths, which work after the project folder is moved to another computer. Fifth, two recorded case-and-control counts for the Bellenguez Alzheimer’s dataset should be reconciled so future statistical models use one confirmed sample description.

Key idea to explain: Finally, nine Walker brain eQTL rows are labeled as splicing QTLs in one sensitivity file. The dataset registry identifies them as expression QTLs, so correcting the labels will align the metadata with the actual analysis. Metadata are data about data—for example, the trait, tissue, version, and analysis type.

Scientific boundary: These are documentation, file-format, and metadata improvements. They do not change which genes were highlighted by the current analyses, but they should be completed before a future researcher treats the package as fully portable.

Transition: Together, these six repairs produce valid files, complete inventories, consistent labels, and traceable sources—a strong foundation for independent validation and future same-variant analyses. This closes the presentation by turning the findings into a concrete next-work checklist.""",
]

EXPECTED_TITLES = MAIN_TITLES + APPENDIX_TITLES
EXPECTED_SLIDE_COUNT = len(EXPECTED_TITLES)
NOTE_SECTION_HEADINGS = (
    "Teaching goal:",
    "Walk through the slide:",
    "Key idea to explain:",
    "Scientific boundary:",
    "Transition:",
)
MIN_SPEAKER_NOTE_WORDS = 160


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def stroke(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_run(run, *, size: float, color: RGBColor, bold: bool = False,
            italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             color: RGBColor = WHITE, outline: RGBColor | None = LIGHT,
             radius: bool = True, transparency: int = 0,
             line_width: float = 1.0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color, transparency)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline, line_width)
    return shape


def add_circle(slide, x: float, y: float, d: float, color: RGBColor,
               outline: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline)
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *,
             size: float = 16, color: RGBColor = DARK, bold: bool = False,
             italic: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.03):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic)
    return box


def add_rich_text(slide, spans: list[tuple[str, dict]], x: float, y: float,
                  w: float, h: float, *, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    for text_value, style in spans:
        run = p.add_run()
        run.text = text_value
        set_run(
            run,
            size=style.get("size", 16),
            color=style.get("color", DARK),
            bold=style.get("bold", False),
            italic=style.get("italic", False),
        )
    return box


def add_bullets(slide, items: list[str], x: float, y: float, w: float, *,
                size: float = 15, accent: RGBColor = BLUE,
                color: RGBColor = DARK, line_h: float = 0.56) -> None:
    for index, item in enumerate(items):
        cy = y + index * line_h
        add_circle(slide, x, cy + 0.13, 0.09, accent)
        add_text(slide, item, x + 0.20, cy, w - 0.20, line_h,
                 size=size, color=color, valign=MSO_ANCHOR.MIDDLE)


def add_connector(slide, x1: float, y1: float, x2: float, y2: float,
                  color: RGBColor = BLUE, width: float = 2.0):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    stroke(shape, color, width)
    shape.line.end_arrowhead = True
    return shape


def set_alt_text(shape, title: str, description: str) -> None:
    props = shape._element.xpath(".//p:cNvPr")
    if props:
        props[0].set("name", title)
        props[0].set("descr", description)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float,
                        *, alt: str):
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    picture_w, picture_h = image_w * scale, image_h * scale
    picture_x = x + (w - picture_w) / 2
    picture_y = y + (h - picture_h) / 2
    picture = slide.shapes.add_picture(
        str(path), Inches(picture_x), Inches(picture_y),
        Inches(picture_w), Inches(picture_h)
    )
    set_alt_text(picture, alt, alt)
    return picture


def new_slide(prs: Presentation, *, bg: RGBColor = OFF_WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_header(slide, kicker: str, title: str, page_no: int, *,
               accent: RGBColor = BLUE, subtitle: str | None = None) -> None:
    add_text(slide, kicker.upper(), 0.55, 0.18, 5.6, 0.24,
             size=9.5, color=accent, bold=True)
    title_size = 24.5 if len(title) < 70 else 22.0
    add_text(slide, title, 0.55, 0.44, 11.95, 0.56,
             size=title_size, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, 0.57, 1.00, 11.65, 0.27,
                 size=10.7, color=MID)
    add_text(slide, f"{page_no:02d}", 12.42, 0.21, 0.36, 0.20,
             size=9, color=MID, bold=True, align=PP_ALIGN.RIGHT)


def add_source(slide, text_value: str) -> None:
    add_text(slide, text_value, 0.55, 7.25, 12.15, 0.14,
             size=6.8, color=MID)


def add_note(slide, note: str) -> None:
    slide.notes_slide.notes_text_frame.text = note


def add_ribbon(slide, text_value: str, *, y: float = 6.15,
               accent: RGBColor = BLUE, fill_color: RGBColor = NAVY) -> None:
    add_rect(slide, 0.55, y, 12.23, 0.55, color=fill_color,
             outline=None, radius=False)
    add_rect(slide, 0.55, y, 0.10, 0.55, color=accent,
             outline=None, radius=False)
    add_text(slide, text_value, 0.80, y + 0.10, 11.72, 0.31,
             size=12.0, color=WHITE, bold=True,
             valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def add_figure_slide(prs: Presentation, *, page_no: int, title: str,
                     figure: Path, alt: str, source: str, note: str,
                     ribbon: str, accent: RGBColor = BLUE,
                     eyebrow: str = "Main result") -> None:
    slide = new_slide(prs, bg=WHITE)
    add_header(slide, eyebrow, title, page_no, accent=accent)
    add_picture_contain(slide, figure, 0.55, 1.16, 12.23, 4.64, alt=alt)
    add_ribbon(slide, ribbon, y=6.12, accent=accent)
    add_source(slide, source)
    add_note(slide, note)


def render_pdf_page(pdf_path: Path, page: int, output_path: Path) -> None:
    command = [
        "gs", "-dSAFER", "-dBATCH", "-dNOPAUSE", "-sDEVICE=pngalpha",
        "-r300", f"-dFirstPage={page}", f"-dLastPage={page}",
        f"-sOutputFile={output_path}", str(pdf_path),
    ]
    subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if not output_path.exists() or output_path.stat().st_size == 0:
        raise RuntimeError(f"Ghostscript did not create {output_path}")


def add_chip(slide, label: str, x: float, y: float, w: float, *,
             accent: RGBColor, bg: RGBColor) -> None:
    add_rect(slide, x, y, w, 0.64, color=bg, outline=accent,
             line_width=1.5)
    add_circle(slide, x + 0.18, y + 0.22, 0.18, accent)
    add_text(slide, label, x + 0.50, y + 0.16, w - 0.67, 0.30,
             size=13.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_three_column_cards(slide, cards: list[tuple[str, str, str, RGBColor, RGBColor]],
                           *, y: float = 1.47, h: float = 4.83) -> None:
    card_w = 3.83
    for index, (number, title, body, accent, bg) in enumerate(cards):
        x = 0.70 + index * 4.10
        add_rect(slide, x, y, card_w, h, color=WHITE, outline=LIGHT)
        add_rect(slide, x, y, card_w, 0.80, color=bg, outline=None,
                 radius=False)
        add_text(slide, number, x + 0.22, y + 0.16, 0.45, 0.38,
                 size=16.0, color=accent, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, x + 0.74, y + 0.13, card_w - 0.96, 0.46,
                 size=17.0, color=NAVY, bold=True,
                 valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x + 0.27, y + 1.12, card_w - 0.54, h - 1.40,
                 size=14.0, color=DARK)


def add_section_divider(prs: Presentation, *, marker: str, eyebrow: str,
                        title: str, subtitle: str, topics: list[str],
                        page_no: int, accent: RGBColor, note: str) -> None:
    """Add a dark narrative divider between presentation sections."""
    slide = new_slide(prs, bg=NAVY)
    add_rect(slide, 0, 0, 13.333, 7.5, color=NAVY, outline=None,
             radius=False)
    add_text(slide, eyebrow.upper(), 0.76, 0.63, 3.60, 0.28,
             size=11.0, color=accent, bold=True)
    add_rect(slide, 0.76, 1.28, 0.11, 2.26, color=accent,
             outline=None, radius=False)
    add_text(slide, title, 1.16, 1.36, 8.35, 1.50,
             size=32.0, color=WHITE, bold=True,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, subtitle, 1.18, 3.18, 7.90, 0.92,
             size=15.3, color=RGBColor(204, 219, 234))

    add_text(slide, marker, 9.55, 0.88, 2.52, 1.92,
             size=104.0 if len(marker) <= 2 else 82.0,
             color=NAVY_2, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    add_circle(slide, 11.63, 1.06, 0.56, accent)
    add_circle(slide, 10.98, 2.93, 0.28, SKY)
    add_rect(slide, 10.90, 2.12, 1.20, 0.035,
             color=RGBColor(105, 137, 169), outline=None, radius=False)
    add_rect(slide, 11.07, 2.14, 0.035, 0.91,
             color=RGBColor(105, 137, 169), outline=None, radius=False)

    add_text(slide, "IN THIS SECTION", 0.78, 4.67, 2.10, 0.25,
             size=9.8, color=SKY, bold=True)
    topic_x = [0.78, 4.20, 7.62]
    for index, topic in enumerate(topics[:3]):
        x = topic_x[index]
        add_rect(slide, x, 5.12, 3.10, 0.76,
                 color=NAVY_2, outline=accent)
        add_text(slide, topic, x + 0.18, 5.30, 2.74, 0.32,
                 size=11.4, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_text(slide, f"{page_no:02d}", 12.38, 0.34, 0.42, 0.22,
             size=9.0, color=RGBColor(145, 171, 197), bold=True,
             align=PP_ALIGN.RIGHT)
    add_note(slide, note)


def build_deck(output_path: Path = DEFAULT_OUT) -> Path:
    for path in [*FIG.values(), *AUX.values()]:
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Inherited Genetic Evidence for Genes Highlighted by Cell Networks"
    prs.core_properties.subject = "Public-data study of 25 genes from 47 cell-network results"
    prs.core_properties.author = "Alzheimer project analysis team"
    prs.core_properties.keywords = "Alzheimer, genetics, QTL, APOE, key drivers"
    prs.core_properties.comments = (
        "Generated from validated genetic-support figures and minerva_production result bundles."
    )

    with tempfile.TemporaryDirectory(prefix="phase19_deck_assets_") as temp_dir:
        temp = Path(temp_dir)
        apoe_locus = temp / "apoe_tier1_locus_page1.png"
        recovery_loci = temp / "recovery_loci_page1.png"
        render_pdf_page(AUX["tier1_loci"], 1, apoe_locus)
        render_pdf_page(AUX["recovery_loci"], 1, recovery_loci)

        # 1 — title
        slide = new_slide(prs, bg=NAVY)
        add_rect(slide, 0, 0, 13.333, 7.5, color=NAVY, outline=None,
                 radius=False)
        add_text(slide, "GENES, NETWORKS, AND INHERITED RISK", 0.76, 0.67,
                 7.2, 0.28, size=11.5, color=SKY, bold=True)
        add_text(slide, "Do genes highlighted by brain-cell networks also show\ninherited links to Alzheimer's disease?",
                 0.76, 1.25, 8.25, 1.78, size=31.0, color=WHITE,
                 bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide,
                 "A public-data study of 25 genes from 47 cell-network results",
                 0.78, 3.38, 7.75, 0.62, size=17.0,
                 color=RGBColor(210, 224, 239))
        add_rect(slide, 0.78, 4.30, 7.68, 0.06, color=BLUE,
                 outline=None, radius=False)
        add_text(slide,
                 "APOE shows multi-dataset support; four additional genes have specific evidence for focused validation.",
                 0.78, 4.65, 7.80, 0.75, size=14.5,
                 color=RGBColor(191, 210, 229))

        # Abstract network motif at right, deliberately non-biological.
        nodes = [
            (9.25, 1.05, 0.92, BLUE), (10.75, 1.54, 0.62, AMBER),
            (11.77, 0.90, 0.42, SKY), (9.55, 2.82, 0.50, PURPLE),
            (11.20, 3.06, 0.82, BLUE), (10.28, 4.56, 0.42, VERMILION),
            (11.77, 4.85, 0.62, TEAL),
        ]
        for x, y, diameter, color in nodes:
            add_circle(slide, x, y, diameter, color)
        for x1, y1, x2, y2 in [
            (9.72, 1.48, 11.04, 1.83), (11.06, 1.95, 11.56, 3.28),
            (9.80, 3.04, 11.43, 3.46), (10.02, 3.30, 10.47, 4.70),
            (11.53, 3.70, 12.06, 5.14), (10.66, 4.75, 11.98, 5.12),
        ]:
            add_connector(slide, x1, y1, x2, y2,
                          RGBColor(133, 164, 194), width=1.5)

        add_text(slide, "STUDY DESIGN  •  DATASETS  •  EVIDENCE  •  INTERPRETATION",
                 0.78, 6.35, 7.80, 0.25, size=10.0,
                 color=RGBColor(145, 171, 197), bold=True)
        add_text(slide, "Public-data update • 21 August 2026", 0.78, 7.03,
                 4.2, 0.18, size=8.5, color=RGBColor(145, 171, 197))
        add_note(slide, MAIN_NOTES[0])

        # 2 — executive overview
        slide = new_slide(prs)
        add_header(slide, "Executive overview", MAIN_TITLES[1], 2, accent=BLUE)
        overview_cards = [
            ("STARTING LIST", "25 genes", "47 gene-by-cell-network combinations", TEAL, PALE_GREEN),
            ("NUCLEAR ANALYSIS", "19 genes", "6 mitochondrial genes define a dedicated follow-up", PURPLE, PALE_BLUE),
            ("MULTI-DATASET RESULT", "APOE", "Linked to Alzheimer's disease and 3 spinal-fluid markers", BLUE, PALE_BLUE),
            ("FIRST-SCREEN MATCHES", "COX7C + SELENOW", "Matched public AD and gene-regulation summaries", AMBER, PALE_AMBER),
            ("REGIONAL PRIORITIES", "RPS15 + ANKRD11", "Regional minima: 4.089×10⁻³⁰ and 1.283×10⁻¹¹", VERMILION, PALE_RED),
            ("NEXT VALIDATION", "Same-variant test", "Complete DNA-variant files can test whether both signals point to the same variant", TEAL, PALE_GREEN),
        ]
        for index, (label, value, detail, accent, bg) in enumerate(overview_cards):
            row, col = divmod(index, 3)
            x = 0.70 + col * 4.10
            y = 1.38 + row * 2.23
            add_rect(slide, x, y, 3.83, 1.82, color=WHITE, outline=LIGHT)
            add_rect(slide, x, y, 0.11, 1.82, color=accent,
                     outline=None, radius=False)
            add_text(slide, label, x + 0.30, y + 0.20, 3.20, 0.22,
                     size=9.0, color=accent, bold=True)
            add_text(slide, value, x + 0.30, y + 0.58, 3.20, 0.46,
                     size=20.0 if len(value) < 15 else 16.5,
                     color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, detail, x + 0.30, y + 1.18, 3.20, 0.42,
                     size=10.4, color=GRAY, valign=MSO_ANCHOR.MIDDLE)
        add_ribbon(slide,
                   "The genetic study adds an independent evidence layer and identifies focused validation targets.",
                   y=6.34, accent=BLUE)
        add_source(slide, "Source: consolidated genetic-support summary; formal Tier 1, recovery, and CSF status tables")
        add_note(slide, MAIN_NOTES[1])

        # 3 — presentation roadmap
        slide = new_slide(prs)
        add_header(slide, "Presentation roadmap", MAIN_TITLES[2], 3,
                   accent=TEAL)
        agenda_rows = [
            (
                "01", "Study design and public data", "Slides 04–08",
                "What the two analyses ask, which public datasets were used, and how each check worked.",
                TEAL, PALE_GREEN,
            ),
            (
                "02", "Evidence found across the gene list", "Slides 09–13",
                "APOE's multi-dataset support plus specific signals for four additional genes.",
                BLUE, PALE_BLUE,
            ),
            (
                "03", "How the evidence guides validation", "Slides 14–18",
                "How exact P values and dataset coverage define the most useful next steps.",
                VERMILION, PALE_RED,
            ),
        ]
        for index, (number, title, slide_range, body, accent, bg) in enumerate(agenda_rows):
            y = 1.42 + index * 1.64
            add_rect(slide, 0.74, y, 11.86, 1.37, color=WHITE, outline=LIGHT)
            add_rect(slide, 0.74, y, 0.10, 1.37, color=accent,
                     outline=None, radius=False)
            add_circle(slide, 1.10, y + 0.34, 0.68, accent)
            add_text(slide, number, 1.10, y + 0.53, 0.68, 0.23,
                     size=11.5, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, title, 2.04, y + 0.21, 6.85, 0.36,
                     size=18.0, color=NAVY, bold=True)
            add_text(slide, slide_range, 9.70, y + 0.23, 2.38, 0.27,
                     size=10.2, color=accent, bold=True,
                     align=PP_ALIGN.RIGHT)
            add_text(slide, body, 2.05, y + 0.72, 9.90, 0.42,
                     size=11.4, color=GRAY)
        add_ribbon(slide,
                   "Exact cutoffs, all 19 nearby-AD P values, full evidence tables, and reproducibility checks are in the appendix.",
                   y=6.50, accent=TEAL)
        add_note(slide, MAIN_NOTES[2])

        # 4 — section divider: study design and datasets
        add_section_divider(
            prs, marker="01", eyebrow="Section 01", title=MAIN_TITLES[3],
            subtitle="What was tested, which public resources were used, and how the evidence was checked step by step.",
            topics=["GWAS: DNA–trait links", "QTL: DNA–gene activity links", "CSF: spinal fluid"],
            page_no=4, accent=TEAL, note=MAIN_NOTES[3],
        )

        # 5 — complementary questions
        slide = new_slide(prs)
        add_header(slide, "Framing", MAIN_TITLES[4], 5, accent=TEAL)
        add_rect(slide, 0.73, 1.42, 4.95, 2.04, color=PALE_GREEN,
                 outline=TEAL, line_width=1.5)
        add_text(slide, "CELL-NETWORK ANALYSIS", 1.03, 1.72, 2.65, 0.26,
                 size=11.0, color=TEAL, bold=True)
        add_text(slide, "Which genes are central in diseased cells?", 1.03, 2.08,
                 4.28, 0.58, size=17.0, color=NAVY, bold=True)
        add_text(slide, "47 cell-network results  •  25 genes", 1.03, 2.78,
                 4.25, 0.32, size=15.5, color=GRAY, bold=True)

        add_connector(slide, 5.85, 2.44, 7.30, 2.44, color=BLUE, width=2.5)
        add_text(slide, "complementary\nquestions", 5.76, 2.69, 1.61, 0.62,
                 size=10.0, color=MID, bold=True, align=PP_ALIGN.CENTER)

        add_rect(slide, 7.45, 1.42, 5.13, 2.04, color=PALE_BLUE,
                 outline=BLUE, line_width=1.5)
        add_text(slide, "INHERITED-DNA ANALYSIS", 7.77, 1.72, 2.70, 0.26,
                 size=11.0, color=BLUE, bold=True)
        add_text(slide, "Which genes are linked to inherited AD risk?", 7.77, 2.08,
                 4.45, 0.58, size=17.0, color=NAVY, bold=True)
        add_text(slide, "19 nuclear genes analyzed  •  6 mitochondrial genes define a dedicated next study",
                 7.77, 2.78, 4.45, 0.48, size=12.5, color=GRAY, bold=True)

        add_bullets(slide, [
            "The network analysis finds genes at the center of disease-related activity in cells.",
            "The genetic analysis asks whether inherited DNA differences near a gene are linked to Alzheimer's disease or disease markers.",
            "Together, the two analyses connect disease-related cell activity with inherited-risk evidence.",
        ], 1.02, 4.15, 11.25, size=16.0, accent=BLUE, line_h=0.69)
        add_ribbon(slide,
                   "Network biology and genetics are complementary; each provides useful evidence for prioritizing follow-up.",
                   y=6.42, accent=TEAL)
        add_source(slide, "Source: genetic-support consolidated summary §1.1; call_key_driver_returns.tsv")
        add_note(slide, MAIN_NOTES[4])

        # 6 — dataset portfolio
        slide = new_slide(prs, bg=WHITE)
        add_header(slide, "Public datasets", MAIN_TITLES[5], 6,
                   accent=BLUE)
        data_layers = [
            (
                "Starting gene list",
                "Network-analysis list • GENCODE v44 • HGNC 2026-06-05",
                "Keep the same 47 settings / 25 genes; map 19 nuclear genes and flag 6 mitochondrial genes.",
                TEAL,
            ),
            (
                "Published summary lists",
                "FunGen-xQTL snapshot f6f63fc… • six public files (~8.74 MiB)",
                "Screen earlier Alzheimer's and gene-regulation results and identify records for deeper source-level validation.",
                AMBER,
            ),
            (
                "Alzheimer's disease study",
                "Bellenguez 2022 • GCST90027158 • complete GRCh38 summary statistics",
                "Look for disease-linked DNA variants within ±1 Mb (one million DNA letters) of each nuclear gene.",
                BLUE,
            ),
            (
                "Brain gene-activity data",
                "NIAGADS NG00184.v1 • eQTL Catalogue r7",
                "Find DNA variants linked to RNA amount, RNA splicing, or protein amount in brain samples.",
                VERMILION,
            ),
            (
                "Spinal-fluid markers",
                "GCST90726396 / GCST90726397 / GCST90726398 • European N=18,948 each",
                "Test amyloid-β42, total tau, and p-tau181 across 19 genes × 3 markers = 57 checks.",
                PURPLE,
            ),
        ]
        col_x = [0.55, 2.55, 7.25]
        col_w = [2.00, 4.70, 5.53]
        for col, header in enumerate(["Data type", "Dataset / version", "What it tells us"]):
            add_rect(slide, col_x[col], 1.26, col_w[col], 0.50,
                     color=NAVY, outline=WHITE, radius=False)
            add_text(slide, header, col_x[col] + 0.13, 1.38,
                     col_w[col] - 0.26, 0.24, size=10.7,
                     color=WHITE, bold=True)
        for row_index, (layer, resource, role, accent) in enumerate(data_layers):
            y = 1.76 + row_index * 0.91
            bg = WHITE if row_index % 2 == 0 else PALE_GRAY
            for col, value in enumerate([layer, resource, role]):
                add_rect(slide, col_x[col], y, col_w[col], 0.91,
                         color=bg, outline=LIGHT, radius=False)
                add_text(slide, value, col_x[col] + 0.14, y + 0.10,
                         col_w[col] - 0.28, 0.67,
                         size=10.3 if col else 10.8,
                         color=NAVY if col == 0 else DARK,
                         bold=(col == 0), valign=MSO_ANCHOR.MIDDLE)
            add_rect(slide, 0.55, y, 0.07, 0.91, color=accent,
                     outline=None, radius=False)
        add_ribbon(slide,
                   "Summary data = group-level results, not individual records • CSF = fluid around the brain and spinal cord.",
                   y=6.45, accent=BLUE)
        add_source(slide, "Source: candidate, Tier 1, Tier 2, recovery, and CSF dataset registries and input inventories")
        add_note(slide, MAIN_NOTES[5])

        # 7 — QTL resource detail
        slide = new_slide(prs)
        add_header(slide, "Gene-activity data", MAIN_TITLES[6], 7,
                   accent=VERMILION)
        qtl_cards = [
            (
                "MICROGLIA", "Young 2019 • 104 samples",
                "eQTL (RNA amount): QTD000559\nsQTL (RNA splicing): QTD000563\n\nSame cell type",
                TEAL, PALE_GREEN,
            ),
            (
                "NEURON-LIKE CELLS", "Aygun 2021 • 73 samples",
                "eQTL (RNA amount): QTD000569\nsQTL (RNA splicing): QTD000573\n\nRelated neuron samples",
                BLUE, PALE_BLUE,
            ),
            (
                "MIXED BRAIN TISSUE", "Walker 2019 • 211 samples",
                "eQTL (RNA amount): QTD000579\nsQTL (RNA splicing): QTD000583\n\nBroad tissue view; exact-cell follow-up can add detail",
                AMBER, PALE_AMBER,
            ),
            (
                "BROADER BRAIN DATA", "NG00184.v1",
                "eQTL: RNA amount\nsQTL: RNA splicing\npQTL: protein amount\n\nWider RNA and protein coverage for future comparisons",
                VERMILION, PALE_RED,
            ),
        ]
        for index, (label, title, body, accent, bg) in enumerate(qtl_cards):
            x = 0.55 + index * 3.07
            add_rect(slide, x, 1.35, 2.82, 3.30, color=WHITE, outline=LIGHT)
            add_rect(slide, x, 1.35, 2.82, 0.62, color=bg,
                     outline=None, radius=False)
            add_text(slide, label, x + 0.16, 1.54, 2.50, 0.20,
                     size=8.4, color=accent, bold=True,
                     align=PP_ALIGN.CENTER)
            add_text(slide, title, x + 0.18, 2.16, 2.46, 0.43,
                     size=15.0, color=NAVY, bold=True,
                     align=PP_ALIGN.CENTER)
            add_text(slide, body, x + 0.22, 2.85, 2.38, 1.52,
                     size=10.4, color=DARK, align=PP_ALIGN.CENTER)

        add_rect(slide, 0.55, 4.91, 6.00, 1.04, color=WHITE, outline=BLUE)
        add_text(slide, "APOE PROTEIN FOLLOW-UP (pQTL)", 0.78, 5.10, 2.80, 0.20,
                 size=9.0, color=BLUE, bold=True)
        add_text(slide,
                 "NG00130.v2 • 3,506 European samples • spinal-fluid protein data used to follow up APOE after its earlier P values crossed their screening references.",
                 0.78, 5.42, 5.45, 0.37, size=10.3, color=DARK)
        add_rect(slide, 6.78, 4.91, 6.00, 1.04, color=WHITE,
                 outline=VERMILION)
        add_text(slide, "RPS15 FOLLOW-UP", 7.01, 5.10, 2.30, 0.20,
                 size=9.0, color=VERMILION, bold=True)
        add_text(slide,
                 "Three distinct mixed-brain gene-activity tracks produced six positive setting rows because each track was counted in two network settings.",
                 7.01, 5.42, 5.45, 0.37, size=10.3, color=DARK)
        add_ribbon(slide,
                   "These resources identified gene-activity leads for APOE and RPS15; matched variant-level data can support the next comparison.",
                   y=6.32, accent=VERMILION)
        add_source(slide, "Source: recovery_dataset_registry.tsv; endophenotype_dataset_registry.tsv; targeted RPS15 audit")
        add_note(slide, MAIN_NOTES[6])

        # 8 — workflow figure
        add_figure_slide(
            prs, page_no=8, title=MAIN_TITLES[7], figure=FIG["workflow"],
            alt="Step-by-step workflow for checking public genetic evidence and comparing disease signals with gene-activity signals",
            source="Source: generated workflow package; analysis contracts, registries, and route manifests",
            note=MAIN_NOTES[7],
            ribbon="Each step adds a layer of evidence; matched variant-level data enable the final same-variant comparison.",
            accent=BLUE,
        )

        # 9 — section divider: genetic evidence
        add_section_divider(
            prs, marker="02", eyebrow="Section 02", title=MAIN_TITLES[8],
            subtitle="We summarize all 47 gene–network settings, highlight APOE's multi-dataset support, and identify four additional follow-up leads.",
            topics=["All 47 settings", "APOE across datasets", "Four focused follow-up leads"],
            page_no=9, accent=BLUE, note=MAIN_NOTES[8],
        )

        # 10–13 — four result figures, full-width and untrimmed.
        add_figure_slide(
            prs, page_no=10, title=MAIN_TITLES[9], figure=FIG["tier1"],
            alt="First-screen scorecard for all 47 gene–network settings",
            source="Source: Tier 1 genetic_support_evidence_summary.tsv and genetic_support_status.tsv",
            note=MAIN_NOTES[9],
            ribbon="This first screen highlighted APOE plus source matches for COX7C and SELENOW; a separate follow-up added RPS15 evidence.",
            accent=BLUE,
        )
        add_figure_slide(
            prs, page_no=11, title=MAIN_TITLES[10], figure=FIG["csf"],
            alt="Spinal-fluid marker results showing APOE linked across all three markers",
            source="Source: CSF endophenotype_gate_decisions.tsv and MAGMA candidate-gene results",
            note=MAIN_NOTES[10],
            ribbon="APOE crossed both the nearby-region and gene-based screening references for amyloid-β42, total tau, and p-tau181.",
            accent=BLUE,
        )
        add_figure_slide(
            prs, page_no=12, title=MAIN_TITLES[11], figure=FIG["non_apoe"],
            alt="Evidence summaries for COX7C, SELENOW, RPS15, and ANKRD11",
            source="Source: generated non-APOE plot data; Tier 1, recovery, RPS15 audit, and CSF result bundles",
            note=MAIN_NOTES[11],
            ribbon="COX7C and SELENOW had public gene-level records; COX7C, RPS15, and ANKRD11 also had notable nearby-region P values.",
            accent=AMBER,
        )
        add_figure_slide(
            prs, page_no=13, title=MAIN_TITLES[12], figure=FIG["tier2"],
            alt="Status summary for 54 planned gene-activity comparisons, including two routes with both disease-region and gene-activity signals",
            source="Source: Tier 2 recovery_route_decisions.tsv and header-only recovery_colocalization.tsv.gz",
            note=MAIN_NOTES[12],
            ribbon="APOE and OPC RPS15 had both nearby-AD and mixed-brain gene-activity signals, making them priorities for same-variant testing.",
            accent=VERMILION,
        )

        # 14 — section divider: interpretation and next steps
        add_section_divider(
            prs, marker="03", eyebrow="Section 03", title=MAIN_TITLES[13],
            subtitle="Exact P values and clear test-status labels turn the current evidence into a focused validation plan.",
            topics=["What the P values show", "What current data support", "Next validation steps"],
            page_no=14, accent=VERMILION, note=MAIN_NOTES[13],
        )

        # 15 — distinguish result categories
        slide = new_slide(prs)
        add_header(slide, "Interpretation", MAIN_TITLES[14], 15,
                   accent=VERMILION)
        cards = [
            (
                "01", "Exact P values for all 19 regions",
                "All 19 regions have a reported minimum P value. ANKRD11, APOE, COX7C, and RPS15 were below the conservative 5×10⁻⁸ reference.\n\nThe other 15 ranged from 2.929×10⁻⁶ to 2.931×10⁻⁴; every value is in the appendix.\n\n5×10⁻⁸ ≈ 0.05 ÷ 1,000,000 tests.",
                BLUE, PALE_BLUE,
            ),
            (
                "02", "Two priority same-variant tests",
                "APOE and one OPC RPS15 comparison had both a nearby Alzheimer's signal and a mixed-brain gene-activity signal.\n\nComplete variant-level gene-activity files and matched DNA-reference data can test whether the same variant contributes to both.",
                TEAL, PALE_GREEN,
            ),
            (
                "03", "More genetic routes to explore",
                "Mitochondrial DNA, rare or large DNA changes, distant gene control, gene interactions, disease timing, and changes after RNA is made offer complementary directions for future validation.",
                PURPLE, PALE_GRAY,
            ),
        ]
        add_three_column_cards(slide, cards, y=1.42, h=4.75)
        add_ribbon(slide, "The exact P values provide a transparent baseline; APOE and RPS15 offer the clearest next same-variant tests.",
                   y=6.40, accent=TEAL)
        add_source(slide, "Source: recovery_regional_gwas_summary.tsv; CSF gate decisions; Tier 1 mtDNA evidence states")
        add_note(slide, MAIN_NOTES[14])

        # 16 — explicit limitations
        slide = new_slide(prs)
        add_header(slide, "Future validation", MAIN_TITLES[15], 16,
                   accent=TEAL)
        limitations = [
            (
                "01", "Complete same-variant tests",
                "APOE and RPS15 are priority targets for complete gene-activity and matched DNA-reference datasets.",
                TEAL, PALE_GREEN,
            ),
            (
                "02", "Larger cell-specific studies",
                "Larger studies of astrocytes, OPCs, neurons, and microglia can test the gene-activity leads in the cell settings identified by the network analysis.",
                AMBER, PALE_AMBER,
            ),
            (
                "03", "Broader populations",
                "Studies across more ancestral backgrounds—and fully independent samples—can show how widely the findings apply.",
                BLUE, PALE_BLUE,
            ),
            (
                "04", "More genetic mechanisms",
                "Rare, large, mitochondrial, and distant DNA effects, gene interactions, disease timing, and laboratory experiments can extend the current common-variant screen.",
                PURPLE, PALE_GRAY,
            ),
        ]
        for index, (number, title, body, accent, bg) in enumerate(limitations):
            row, col = divmod(index, 2)
            x = 0.68 + col * 6.12
            y = 1.40 + row * 2.33
            add_rect(slide, x, y, 5.83, 2.02, color=WHITE, outline=LIGHT)
            add_rect(slide, x, y, 0.10, 2.02, color=accent,
                     outline=None, radius=False)
            add_text(slide, number, x + 0.27, y + 0.22, 0.55, 0.30,
                     size=14.0, color=accent, bold=True,
                     align=PP_ALIGN.CENTER)
            add_text(slide, title, x + 0.96, y + 0.20, 4.48, 0.36,
                     size=16.5, color=NAVY, bold=True)
            add_text(slide, body, x + 0.96, y + 0.78, 4.45, 0.94,
                     size=11.6, color=DARK)
        add_ribbon(slide,
                   "The evidence strongly supports APOE and prioritizes COX7C, SELENOW, RPS15, and ANKRD11 for different kinds of follow-up.",
                   y=6.34, accent=TEAL)
        add_source(slide, "Source: consolidated genetic-support summary §§5–6; route, QTL, ancestry, and provenance audits")
        add_note(slide, MAIN_NOTES[15])

        # 17 — roadmap
        slide = new_slide(prs)
        add_header(slide, "Next steps", MAIN_TITLES[16], 17, accent=TEAL)
        roadmap = [
            (
                "01", "Deepen APOE and RPS15",
                "Complete gene-activity datasets\n+ matching DNA-reference data\n+ larger studies of exact cell types",
                VERMILION, PALE_RED,
            ),
            (
                "02", "Broaden genetic testing",
                "Protein- and RNA-based gene tests\n+ rare DNA changes and interactions\n+ more traits and populations",
                BLUE, PALE_BLUE,
            ),
            (
                "03", "Add independent and lab validation",
                "Dedicated mitochondrial-DNA study\n+ repeat the network study in new data\n+ change gene activity in lab experiments",
                TEAL, PALE_GREEN,
            ),
        ]
        add_three_column_cards(slide, roadmap, y=1.50, h=4.55)
        for index in range(2):
            start_x = 4.58 + index * 4.10
            add_connector(slide, start_x, 3.83, start_x + 0.45, 3.83,
                          color=BLUE, width=2.0)
        add_ribbon(slide,
                   "Matched variant-level gene-activity and DNA-reference data can unlock same-variant analyses for APOE and RPS15.",
                   y=6.37, accent=TEAL)
        add_source(slide, "Source: genetic-support consolidated summary §6.2; bundle-repair actions are detailed in the appendix")
        add_note(slide, MAIN_NOTES[16])

        # 18 — close
        slide = new_slide(prs, bg=NAVY)
        add_rect(slide, 0, 0, 13.333, 7.5, color=NAVY, outline=None,
                 radius=False)
        add_text(slide, "TAKE-HOME MESSAGE", 0.72, 0.50, 4.5, 0.26,
                 size=10.5, color=SKY, bold=True)
        add_text(slide, MAIN_TITLES[17], 0.72, 0.92, 11.8, 0.75,
                 size=29.0, color=WHITE, bold=True)
        takeaways = [
            ("01", "APOE", "Direct AD mapping plus matching results for all three spinal-fluid markers.", BLUE),
            ("02", "COX7C + SELENOW", "Public gene-activity records identified both; COX7C also had a notable nearby-region P value.", AMBER),
            ("03", "RPS15 + ANKRD11", "Both regions had P < 5×10⁻⁸; RPS15 also had mixed-brain gene-activity signals.", VERMILION),
            ("04", "Broader gene list", "All exact regional P values now form a transparent baseline for future studies.", TEAL),
        ]
        for index, (number, title, body, accent) in enumerate(takeaways):
            row = index // 2
            col = index % 2
            x = 0.72 + col * 6.10
            y = 2.05 + row * 1.62
            add_rect(slide, x, y, 5.72, 1.28, color=NAVY_2,
                     outline=accent, line_width=1.4)
            add_text(slide, number, x + 0.24, y + 0.19, 0.52, 0.36,
                     size=16.0, color=accent, bold=True,
                     align=PP_ALIGN.CENTER)
            add_text(slide, title, x + 0.92, y + 0.17, 4.50, 0.33,
                     size=17.0, color=WHITE, bold=True)
            add_text(slide, body, x + 0.92, y + 0.63, 4.48, 0.42,
                     size=12.0, color=RGBColor(207, 222, 237))
        add_rect(slide, 3.52, 5.76, 6.29, 0.73, color=WHITE,
                 outline=TEAL, line_width=1.6)
        add_text(slide, "Result: evidence found, with focused validation next", 3.76, 5.96,
                 5.81, 0.31, size=15.5, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide,
                 "Genetics strengthens the APOE finding and guides targeted follow-up across the broader network list.",
                 2.15, 6.78, 9.03, 0.30, size=12.5,
                 color=RGBColor(173, 197, 220), align=PP_ALIGN.CENTER)
        add_text(slide, "18", 12.42, 0.32, 0.36, 0.20,
                 size=9, color=RGBColor(145, 171, 197), bold=True,
                 align=PP_ALIGN.RIGHT)
        add_note(slide, MAIN_NOTES[17])

        # 19 — appendix divider
        add_section_divider(
            prs, marker="A", eyebrow="Appendix", title=APPENDIX_TITLES[0],
            subtitle="The exact gene list, dataset versions, pre-set rules, all 19 nearby-AD P values, full evidence tables, and a clear plan for future validation.",
            topics=["Original gene list", "Data + rules", "Full evidence + reproducibility"],
            page_no=19, accent=SKY,
            note=APPENDIX_NOTES[0],
        )

        # 20 — appendix candidate list
        slide = new_slide(prs, bg=WHITE)
        add_header(slide, "Appendix • original gene list", APPENDIX_TITLES[1], 20,
                   accent=TEAL)
        headers = ["Cell network", "Mitochondria-related group (original order)", "Other driver group (original order)"]
        rows = [
            ("Astrocytes", "MT-CO2, MT-CO3, MT-ATP6, COX7C, COX4I1", "RPL11, RPLP1, RPL15, APOE, LAPTM4A"),
            ("Excitatory neurons", "MT-CO2, UQCR10, COX4I1, COX6B1, MT-CYB", "RPL11, RPS13, SELENOW, LAMTOR5, DYNLT1"),
            ("Inhibitory neurons", "MT-CO2, MT-CO3, MT-CYB, MT-ND5, COX7C", "RPS15, LAMTOR5, RPLP1, ATP6V1F, RPL38"),
            ("Microglia", "MT-CO2, MT-ND4", "RPL11"),
            ("OPCs", "MT-CO3, MT-CO2, MT-ND4", "RPS15, FTL, ANKRD11, NCOA1"),
            ("Oligodendrocytes", "MT-CO2, MT-ND4", "RPL11"),
            ("Vasculature cells", "MT-CO3, MT-CO2, MT-ATP6, MT-ND4", "None"),
        ]
        col_x = [0.55, 2.45, 7.55]
        col_w = [1.90, 5.10, 5.23]
        table_y = 1.30
        for col, header in enumerate(headers):
            add_rect(slide, col_x[col], table_y, col_w[col], 0.55,
                     color=NAVY, outline=WHITE, radius=False)
            add_text(slide, header, col_x[col] + 0.12, table_y + 0.13,
                     col_w[col] - 0.24, 0.26, size=11.0, color=WHITE,
                     bold=True, valign=MSO_ANCHOR.MIDDLE)
        for row_index, row in enumerate(rows):
            y = table_y + 0.55 + row_index * 0.69
            bg = WHITE if row_index % 2 == 0 else PALE_GRAY
            for col, value in enumerate(row):
                add_rect(slide, col_x[col], y, col_w[col], 0.69,
                         color=bg, outline=LIGHT, radius=False)
                add_text(slide, value, col_x[col] + 0.12, y + 0.10,
                         col_w[col] - 0.24, 0.47,
                         size=10.8 if col else 11.3, color=NAVY if col == 0 else DARK,
                         bold=(col == 0), valign=MSO_ANCHOR.MIDDLE)
        add_ribbon(slide,
                   "‘MT driver’ was a label in the original network analysis; COX7C and UQCR10 are nuclear genes. Six genes are encoded by mitochondrial DNA.",
                   y=6.83, accent=TEAL)
        add_note(slide, APPENDIX_NOTES[1])

        # 21 — dataset inventory
        slide = new_slide(prs)
        add_header(slide, "Appendix • data inventory", APPENDIX_TITLES[2], 21,
                   accent=BLUE)
        inventory = [
            ("Starting gene list", "Original network study • GENCODE v44 • HGNC 2026-06-05", "Keep 47 settings; map gene names and DNA locations"),
            ("Published summaries", "FunGen-xQTL snapshot f6f63fc…", "Quick screen of earlier AD and gene-regulation results"),
            ("Alzheimer's disease", "Bellenguez • GCST90027158", "Check DNA within ±1 Mb of the 19 nuclear genes"),
            ("Brain gene activity", "NG00184.v1 • eQTL Catalogue r7", "Check DNA links to RNA amount, RNA splicing, and protein"),
            ("Spinal-fluid markers", "GCST90726396 / 397 / 398 • N=18,948 each", "Test amyloid-β42, total tau, and p-tau181"),
            ("Focused follow-up", "NG00130.v2 APOE protein QTL • local NG00184 RPS15", "Study the open APOE and RPS15 questions"),
        ]
        x_values = [0.55, 2.72, 7.44]
        widths = [2.17, 4.72, 5.34]
        for col, header in enumerate(["Data type", "Dataset / version", "What it tells us"]):
            add_rect(slide, x_values[col], 1.30, widths[col], 0.55,
                     color=NAVY, outline=WHITE, radius=False)
            add_text(slide, header, x_values[col] + 0.13, 1.43,
                     widths[col] - 0.26, 0.26, size=11.2,
                     color=WHITE, bold=True)
        for row_index, row in enumerate(inventory):
            y = 1.85 + row_index * 0.78
            bg = WHITE if row_index % 2 == 0 else PALE_GRAY
            for col, value in enumerate(row):
                add_rect(slide, x_values[col], y, widths[col], 0.78,
                         color=bg, outline=LIGHT, radius=False)
                add_text(slide, value, x_values[col] + 0.13, y + 0.10,
                         widths[col] - 0.26, 0.54,
                         size=11.1 if col else 11.5,
                         color=NAVY if col == 0 else DARK,
                         bold=(col == 0), valign=MSO_ANCHOR.MIDDLE)
        add_ribbon(slide,
                   "Public summary statistics support gene-level screening without individual records; Bellenguez sample counts will be added after two bundle fields are reconciled.",
                   y=6.78, accent=BLUE)
        add_note(slide, APPENDIX_NOTES[2])

        # 22 — thresholds and states
        slide = new_slide(prs, bg=WHITE)
        add_header(slide, "Appendix • decision rules", APPENDIX_TITLES[3], 22,
                   accent=BLUE,
                   subtitle="Cutoffs were pre-set screening references; exact P values are reported on both sides of each cutoff.")
        thresholds = [
            ("Nearby-AD regional screen", "P < 5×10⁻⁸"),
            ("Gene-activity link", "P < 0.05 / N tested"),
            ("Spinal-fluid gene test", "P < 8.77193×10⁻⁴"),
            ("Model assumptions", "10⁻⁴ / 10⁻⁴ / 5×10⁻⁶"),
            ("Strong shared signal", "Estimated probability ≥ 0.80"),
        ]
        for index, (label, value) in enumerate(thresholds):
            x = 0.55 + index * 2.47
            width = 2.26
            add_rect(slide, x, 1.27, width, 1.04, color=PALE_BLUE,
                     outline=BLUE)
            add_text(slide, label.upper(), x + 0.13, 1.43, width - 0.26,
                     0.20, size=8.8, color=BLUE, bold=True,
                     align=PP_ALIGN.CENTER)
            add_text(slide, value, x + 0.13, 1.79, width - 0.26, 0.28,
                     size=12.2, color=NAVY, bold=True,
                     align=PP_ALIGN.CENTER)
        states = [
            ("Broaden public-data search", "The registered quick screen had no direct match; additional sources can extend the search."),
            ("Nearby-AD P value reported", "The exact regional minimum is shown relative to the conservative 5×10⁻⁸ reference."),
            ("Gene-activity P value reported", "The measured value is shown relative to its pre-set source-specific reference."),
            ("Validate with complete inputs", "A complete measurement, file, model, or dataset description can enable the next test."),
            ("Harmonize both datasets", "Both signals existed; matching variants, models, and reference data can enable comparison."),
            ("Use an mtDNA-specific method", "Mitochondrial genes have a dedicated future testing route."),
        ]
        for index, (state, meaning) in enumerate(states):
            col = index % 2
            row = index // 2
            x = 0.62 + col * 6.18
            y = 2.70 + row * 1.15
            accent = TEAL if state in {"Validate with complete inputs", "Harmonize both datasets", "Use an mtDNA-specific method"} else BLUE
            add_rect(slide, x, y, 5.86, 0.96, color=WHITE,
                     outline=LIGHT)
            add_rect(slide, x, y, 0.09, 0.96, color=accent,
                     outline=None, radius=False)
            add_text(slide, state, x + 0.25, y + 0.15, 2.50, 0.27,
                     size=11.3, color=NAVY, bold=True)
            add_text(slide, meaning, x + 2.72, y + 0.12, 2.83, 0.54,
                     size=9.8, color=DARK)
        add_ribbon(slide,
                   "PIP, VCP, and CL are scores from the source studies. PP.H4 estimates whether the disease and gene-activity signals share a DNA variant.",
                   y=6.63, accent=BLUE)
        add_note(slide, APPENDIX_NOTES[3])

        # 23 — exact nearby-AD P values and threshold explanation
        add_figure_slide(
            prs, page_no=23, title=APPENDIX_TITLES[4],
            figure=FIG["ad_pvalues"],
            alt="Exact minimum nearby Alzheimer’s-association P values for all 19 nuclear-gene regions, shown relative to the conservative 5×10⁻⁸ screening cutoff",
            source="Source: recovery_regional_gwas_summary.tsv; Bellenguez GWAS Catalog GCST90027158",
            note=APPENDIX_NOTES[4],
            ribbon="Four regions—near ANKRD11, APOE, COX7C, and RPS15—stood out below 5×10⁻⁸ and now prioritize gene-level validation.",
            accent=AMBER,
            eyebrow="Appendix • Alzheimer's P values",
        )

        # 24 — Tier 1 matrix
        slide = new_slide(prs)
        add_header(slide, "Appendix • full Tier 1 audit", APPENDIX_TITLES[5], 24,
                   accent=BLUE)
        add_picture_contain(slide, AUX["tier1_matrix"], 0.55, 1.18, 4.22, 5.88,
                            alt="Full first-screen evidence matrix for 47 gene–network settings")
        # The published matrix carries an internal analysis-number title. Mask
        # only that title band while preserving the complete 47-row matrix.
        add_rect(slide, 0.55, 1.18, 4.22, 0.17, color=WHITE,
                 outline=None, radius=False)
        add_text(slide, "First-screen matrix: 47 settings", 0.73, 1.205,
                 3.86, 0.12, size=7.8, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        add_rect(slide, 5.04, 1.34, 7.72, 4.86, color=WHITE, outline=LIGHT)
        add_text(slide, "COUNTS ACROSS 47 SETTINGS", 5.37, 1.67, 3.95, 0.24,
                 size=10.0, color=BLUE, bold=True)
        count_cards = [
            ("1", "Strong", "APOE", BLUE, PALE_BLUE),
            ("0", "Moderate", "None", TEAL, PALE_GREEN),
            ("3", "Limited", "COX7C ×2; SELENOW ×1", AMBER, PALE_AMBER),
            ("23", "Broader-source\nfollow-up", "16 nuclear genes", NO_SUPPORT, PALE_GRAY),
            ("20", "mtDNA-specific\nfollow-up", "6 mitochondrial genes", TEAL, PALE_GREEN),
        ]
        for index, (value, label, detail, accent, bg) in enumerate(count_cards):
            row = index // 3
            col = index % 3
            x = 5.34 + col * 2.35
            y = 2.08 + row * 1.38
            width = 2.08
            add_rect(slide, x, y, width, 1.12, color=bg, outline=accent)
            add_text(slide, value, x + 0.14, y + 0.15, 0.62, 0.48,
                     size=21.0, color=accent, bold=True,
                     valign=MSO_ANCHOR.MIDDLE)
            label_height = 0.46 if row == 1 else 0.28
            detail_y = y + 0.72 if row == 1 else y + 0.55
            detail_height = 0.22 if row == 1 else 0.34
            add_text(slide, label, x + 0.79, y + 0.12, 1.12, label_height,
                     size=9.4 if row == 1 else 10.5, color=NAVY, bold=True)
            add_text(slide, detail, x + 0.79, detail_y, 1.12, detail_height,
                     size=7.8 if row == 1 else 8.7, color=GRAY)
        add_text(slide,
                 "COX7C's two limited settings come from one public result, not two independent confirmations. The extra RPS15 study is not included here.",
                 5.38, 5.06, 6.92, 0.70, size=12.0, color=DARK)
        add_ribbon(slide,
                   "Evidence found: APOE led the screen, with public-data matches for COX7C and SELENOW; other routes define broader or specialized validation.",
                   y=6.52, accent=BLUE)
        add_source(slide, "Source: Tier 1 genetic_support_evidence_matrix.png and genetic_support_evidence_summary.tsv")
        add_note(slide, APPENDIX_NOTES[5])

        # 25 — recovery matrix and loci
        slide = new_slide(prs, bg=WHITE)
        add_header(slide, "Appendix • Tier 2 recovery detail", APPENDIX_TITLES[6], 25,
                   accent=BLUE)
        add_picture_contain(slide, AUX["recovery_matrix"], 0.45, 1.16, 4.20, 5.72,
                            alt="Detailed comparison outcomes across nuclear gene–network settings")
        add_rect(slide, 0.45, 1.43, 4.20, 0.30, color=WHITE,
                 outline=None, radius=False)
        add_text(slide, "Detailed comparison outcomes", 0.70, 1.515,
                 3.70, 0.13, size=7.8, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        add_picture_contain(slide, recovery_loci, 4.86, 1.32, 7.94, 4.76,
                            alt="Four regional AD locus plots for ANKRD11, APOE, COX7C, and RPS15")
        add_rect(slide, 5.13, 6.12, 7.42, 0.68, color=PALE_BLUE,
                 outline=BLUE)
        add_text(slide,
                 "Four regional minima were below 5×10⁻⁸: ANKRD11, APOE, COX7C, and RPS15. These regions are priorities for shared-variant validation.",
                 5.37, 6.30, 6.94, 0.28, size=11.5, color=NAVY,
                 bold=True, align=PP_ALIGN.CENTER)
        add_source(slide, "Source: recovery_evidence_matrix.png; recovery_locus_plots.pdf; recovery_route_decisions.tsv")
        add_note(slide, APPENDIX_NOTES[6])

        # 26 — APOE locus + CSF matrix and MAGMA table
        slide = new_slide(prs)
        add_header(slide, "Appendix • APOE detail", APPENDIX_TITLES[7], 26,
                   accent=BLUE)
        add_picture_contain(slide, apoe_locus, 0.55, 1.23, 7.55, 3.60,
                            alt="Tier 1 APOE locus plot showing direct rs429358 entries")
        add_picture_contain(slide, AUX["csf_matrix"], 8.42, 1.18, 4.30, 5.76,
                            alt="Spinal-fluid marker evidence for 19 nuclear genes across three biomarkers")
        add_rect(slide, 8.42, 1.55, 4.30, 0.32, color=WHITE,
                 outline=None, radius=False)
        add_text(slide, "Spinal-fluid marker evidence", 8.66, 1.64,
                 3.82, 0.14, size=7.6, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        # The published matrix includes an internal workflow label in its raster
        # footer. Keep the scientific panel unchanged while masking that label
        # in the external-facing deck.
        add_rect(slide, 8.42, 6.40, 4.30, 0.22, color=WHITE,
                 outline=None, radius=False)
        add_rect(slide, 0.67, 5.00, 7.18, 1.60, color=WHITE, outline=LIGHT)
        add_text(slide, "MARKER", 0.91, 5.18, 1.47, 0.21,
                 size=9.2, color=BLUE, bold=True)
        add_text(slide, "STRONGEST NEARBY P", 2.54, 5.18, 1.55, 0.21,
                 size=9.2, color=BLUE, bold=True)
        add_text(slide, "GENE-LEVEL P", 4.26, 5.18, 1.68, 0.21,
                 size=9.2, color=BLUE, bold=True)
        add_text(slide, "GENE + NEARBY DNA", 6.10, 5.18, 1.45, 0.21,
                 size=9.2, color=BLUE, bold=True)
        magma_rows = [
            ("Amyloid-β42", "stored as 0*", "5.0×10⁻¹⁰", "2.3037×10⁻¹⁴"),
            ("Total tau", "5.4×10⁻¹⁶¹", "5.0×10⁻¹⁰", "1.2218×10⁻¹³"),
            ("p-tau181", "3.27×10⁻¹⁷⁴", "5.0×10⁻¹⁰", "5.0×10⁻¹⁰"),
        ]
        for index, row in enumerate(magma_rows):
            y = 5.52 + index * 0.31
            for col, (x, width) in enumerate([(0.91, 1.47), (2.54, 1.55), (4.26, 1.68), (6.10, 1.45)]):
                add_text(slide, row[col], x, y, width, 0.22,
                         size=9.2, color=NAVY if col == 0 else DARK,
                         bold=(col == 0), valign=MSO_ANCHOR.MIDDLE)
        add_text(slide,
                 "* Amyloid-β42 P was smaller than the file's number format could store; it is reported as underflow rather than literal zero, so its numerical region plot is not shown.",
                 0.91, 6.52, 6.52, 0.20, size=7.7, color=VERMILION,
                 italic=True)
        add_source(slide, "Source: Tier 1 common_variant_evidence.tsv.gz and locus plot page 1; CSF gate and MAGMA result tables")
        add_note(slide, APPENDIX_NOTES[7])

        # 27 — RPS15 audit
        slide = new_slide(prs, bg=WHITE)
        add_header(slide, "Appendix • targeted public-data audit", APPENDIX_TITLES[8], 27,
                   accent=AMBER)
        metrics = [
            ("37", "planned comparisons", BLUE),
            ("31", "comparisons measured", BLUE),
            ("6", "positive setting rows", AMBER),
            ("3", "unique mixed-brain results", AMBER),
            ("0", "shared-variant tests to date", TEAL),
        ]
        for index, (value, label, accent) in enumerate(metrics):
            x = 0.57 + index * 2.52
            add_rect(slide, x, 1.29, 2.27, 1.06, color=PALE_GRAY,
                     outline=accent)
            add_text(slide, value, x + 0.14, 1.46, 0.67, 0.44,
                     size=23.0, color=accent, bold=True)
            add_text(slide, label, x + 0.87, 1.48, 1.20, 0.42,
                     size=10.2, color=NAVY, bold=True)
        tracks = [
            ("MSBB BA36 RNA-amount QTL", "P = 2.41403×10⁻⁷", "Adjusted for many tests: 0.00209909"),
            ("ROSMAP DLPFC RNA-splicing QTL", "P = 3.86842×10⁻³⁰", "Adjusted: 1.26188×10⁻²⁶ • variant probability: 1.0"),
            ("ROSMAP posterior-cingulate RNA-splicing QTL", "P = 3.30886×10⁻⁷", "Adjusted: 0.000858045 • variant probability: 0.910283"),
        ]
        for index, (track, p_value, detail) in enumerate(tracks):
            y = 2.78 + index * 1.02
            add_rect(slide, 0.72, y, 8.14, 0.78, color=WHITE, outline=LIGHT)
            add_rect(slide, 0.72, y, 0.09, 0.78, color=AMBER,
                     outline=None, radius=False)
            add_text(slide, track, 1.01, y + 0.13, 2.84, 0.30,
                     size=12.5, color=NAVY, bold=True)
            add_text(slide, p_value, 3.95, y + 0.13, 2.03, 0.30,
                     size=11.6, color=DARK, bold=True)
            add_text(slide, detail, 6.05, y + 0.10, 2.53, 0.52,
                     size=8.5, color=GRAY)
        add_rect(slide, 9.22, 2.78, 3.40, 2.82, color=PALE_GREEN,
                 outline=TEAL, line_width=1.5)
        add_text(slide, "EVIDENCE TO CARRY FORWARD", 9.52, 3.08, 2.79, 0.24,
                 size=10.0, color=TEAL, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide,
                 "Three unique mixed-brain RPS15 results: one RNA-amount and two RNA-splicing.\n\nThese sources form six rows across OPC and inhibitory-neuron questions.\n\nNext: test matched cell-specific data for a shared AD–gene-activity variant.",
                 9.54, 3.38, 2.76, 1.92, size=10.4, color=NAVY,
                 bold=True, align=PP_ALIGN.CENTER)
        add_ribbon(slide,
                   "RPS15 now has three public-data leads; matched cell-specific data and a shared-variant test can add the next validation layer.",
                   y=6.35, accent=AMBER)
        add_source(slide, "Source: opc_rps15_evidence_summary.tsv and opc_rps15_qtl_audit.tsv")
        add_note(slide, APPENDIX_NOTES[8])

        # 28 — provenance and reproducibility
        slide = new_slide(prs)
        add_header(slide, "Appendix • reproducibility plan", APPENDIX_TITLES[9], 28,
                   accent=TEAL)
        caveats = [
            ("Rebuild two empty-result files", "Regenerate the two declared zero-row outputs so they are valid compressed files."),
            ("Add four APOE protein files", "Extend the published inventory with the four NG00130.v2 files described in the execution report."),
            ("Document large-input locations", "Record where streamed or externally stored large inputs can be retrieved and verified."),
            ("Make saved paths portable", "Replace the original /home/... paths in the RPS15 report with project-relative locations."),
            ("Reconcile sample counts", "Resolve the two Bellenguez case/control counts recorded in different bundle fields."),
            ("Correct nine QTL labels", "Relabel nine expression-QTL rows currently marked as splicing QTL; the registry already has the correct type."),
        ]
        for index, (title, body) in enumerate(caveats):
            col = index % 2
            row = index // 2
            x = 0.64 + col * 6.18
            y = 1.28 + row * 1.50
            add_rect(slide, x, y, 5.84, 1.24, color=WHITE, outline=LIGHT)
            add_rect(slide, x, y, 0.10, 1.24, color=TEAL,
                     outline=None, radius=False)
            add_text(slide, title, x + 0.30, y + 0.18, 2.12, 0.30,
                     size=12.4, color=NAVY, bold=True)
            add_text(slide, body, x + 2.52, y + 0.15, 3.02, 0.75,
                     size=9.8, color=DARK)
        add_rect(slide, 1.76, 5.95, 9.82, 0.70, color=NAVY,
                 outline=TEAL, line_width=1.4)
        add_text(slide,
                 "Reproducibility plan: rebuild files • complete inventories • document sources • fix labels and paths • reconcile metadata",
                 2.02, 6.15, 9.30, 0.28, size=12.0, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Next milestone: a fully documented, portable data package", 3.67, 6.82,
                 6.00, 0.23, size=12.0, color=TEAL,
                 bold=True, align=PP_ALIGN.CENTER)
        add_source(slide, "Source: genetic-support consolidated summary §5.3 and bundle-integrity audit")
        add_note(slide, APPENDIX_NOTES[9])

        output_path = output_path.resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        temporary_output = output_path.with_name(f".{output_path.name}.tmp")
        prs.save(temporary_output)
        os.replace(temporary_output, output_path)

    validate_deck(output_path)
    return output_path


def _all_slide_text(slide) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            values.append(shape.text)
    return "\n".join(values)


def validate_deck(path: Path) -> None:
    if not path.exists() or path.stat().st_size < 100_000:
        raise AssertionError(f"Deck is missing or unexpectedly small: {path}")

    prs = Presentation(path)
    if len(prs.slides) != EXPECTED_SLIDE_COUNT:
        raise AssertionError(
            f"Expected {EXPECTED_SLIDE_COUNT} slides, found {len(prs.slides)}"
        )
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        raise AssertionError("Deck is not 13.333333 × 7.5 inch widescreen")

    all_text: list[str] = []
    picture_alt: list[str] = []
    for index, (slide, expected_title) in enumerate(zip(prs.slides, EXPECTED_TITLES), start=1):
        slide_text = _all_slide_text(slide)
        all_text.append(slide_text)
        normalized_slide_text = " ".join(slide_text.split())
        normalized_title = " ".join(expected_title.split())
        if normalized_title not in normalized_slide_text:
            raise AssertionError(f"Slide {index} is missing expected title: {expected_title}")
        note_text = slide.notes_slide.notes_text_frame.text.strip()
        if not note_text:
            raise AssertionError(f"Slide {index} has no speaker note")
        if len(note_text.split()) < MIN_SPEAKER_NOTE_WORDS:
            raise AssertionError(
                f"Slide {index} speaker note is too short for tutorial use: "
                f"{len(note_text.split())} words"
            )
        for heading in NOTE_SECTION_HEADINGS:
            if heading not in note_text:
                raise AssertionError(
                    f"Slide {index} speaker note is missing tutorial section: {heading}"
                )
        all_text.append(note_text)
        for shape in slide.shapes:
            tolerance = Inches(0.02)
            if shape.left < -tolerance or shape.top < -tolerance:
                raise AssertionError(f"Slide {index} has a shape outside the top/left bound")
            if shape.left + shape.width > SLIDE_W + tolerance:
                raise AssertionError(f"Slide {index} has a shape beyond the right bound")
            if shape.top + shape.height > SLIDE_H + tolerance:
                raise AssertionError(f"Slide {index} has a shape beyond the bottom bound")
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                props = shape._element.xpath(".//p:cNvPr")
                description = props[0].get("descr", "") if props else ""
                if not description:
                    raise AssertionError(f"Slide {index} contains an image without alt text")
                picture_alt.append(description)

    joined_text = "\n".join(all_text)
    forbidden = [
        "PP.H4 = 0", "PP.H4=0", "mtDNA genes were negative",
        "mtDNA genes tested negative", "validated mechanism", "proved causal",
        "not enough evidence", "not enough for confirmation", "true negative screen",
        "negative checks", "non-passing regions", "did not pass",
        "four regions passed", "no strong nearby AD variant signal",
        "what this study could not answer", "what this does not show",
        "responsible genes remained unresolved",
        "status: the data package still needs repair",
        "no shared-signal probability was available",
        "Phase 18", "Phase 19", "phase18", "phase19",
    ]
    for phrase in forbidden:
        if phrase.lower() in joined_text.lower():
            raise AssertionError(f"Forbidden presentation wording found: {phrase}")

    expected_alt_fragments = [
        "Step-by-step workflow",
        "First-screen scorecard",
        "Spinal-fluid marker results",
        "Evidence summaries for COX7C",
        "Status summary for 54 planned gene-activity comparisons",
        "Exact minimum nearby",
    ]
    for fragment in expected_alt_fragments:
        if not any(fragment in alt for alt in picture_alt):
            raise AssertionError(f"Missing required figure alt text: {fragment}")

    with zipfile.ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise AssertionError("PPTX ZIP integrity check failed")
        media_members = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        media_hashes = {
            hashlib.sha256(archive.read(name)).hexdigest() for name in media_members
        }
        for label, source in FIG.items():
            if sha256(source) not in media_hashes:
                raise AssertionError(f"Required source figure is not embedded byte-for-byte: {label}")
        slide_members = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        if len(slide_members) != EXPECTED_SLIDE_COUNT:
            raise AssertionError(
                "PPTX package does not contain exactly "
                f"{EXPECTED_SLIDE_COUNT} slide XML parts"
            )
        presentation_xml = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith(".xml")
        ).lower()
        for internal_label in ("phase 18", "phase 19", "phase18", "phase19"):
            if internal_label in presentation_xml:
                raise AssertionError(
                    f"Internal phase label remains in presentation XML: {internal_label}"
                )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT,
                        help=f"Output PPTX path (default: {DEFAULT_OUT})")
    parser.add_argument("--validate-only", action="store_true",
                        help="Validate an existing --output without rebuilding it")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if args.validate_only:
        validate_deck(args.output.resolve())
        print(f"Validated: {args.output.resolve()}")
    else:
        path = build_deck(args.output)
        print(f"Built and validated: {path}")


if __name__ == "__main__":
    main()
