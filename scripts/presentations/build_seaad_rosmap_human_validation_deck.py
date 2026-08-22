#!/usr/bin/env python3
"""Build the SEA-AD–ROSMAP human-validation PowerPoint deck.

The main deck follows ``seaad_rosmap_human_validation_presentation_design.md``:
nine presentation-first slides, followed by six concise appendix slides.  All
visible result counts are derived from validated VH02/VH04/VH08/VH09/VH10
tables.  Scientific figures are embedded byte-for-byte from their validated
figure packages, while the setup, attrition, and conclusion slides remain
editable PowerPoint shapes.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import os
import tempfile
import zipfile
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
DEFAULT_OUT = REPO / "docs/presentations/seaad_rosmap_human_validation.pptx"
DEFAULT_REPORT_DIR = (
    REPO / "results/presentations/validation_human/seaad_rosmap_human_validation"
)
DESIGN_DOC = (
    REPO / "docs/validation_human/seaad_rosmap_human_validation_presentation_design.md"
)

DATA = {
    "vh02_status": REPO / "results/validation_human/02_cohort/status.tsv",
    "donor_groups": REPO / "results/validation_human/02_cohort/donor_group_counts.tsv",
    "vh04_status": REPO / "results/validation_human/04_supertype_manifest/status.tsv",
    "fine_contrasts": REPO / (
        "results/validation_human/08_deg/fine_supertype_phase18_parity/"
        "fine_contrast_status.tsv"
    ),
    "query_attrition": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10a_inputs/"
        "query_attrition.tsv"
    ),
    "run_manifest": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10a_inputs/"
        "seaad_kda_run_manifest.tsv"
    ),
    "vh10b_status": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10b_kda/status.tsv"
    ),
    "vh10c_status": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10c_seaad_selection/status.tsv"
    ),
    "seaad_top5": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10c_seaad_selection/"
        "seaad_top5.tsv"
    ),
    "seaad_list_status": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10c_seaad_selection/"
        "seaad_list_status.tsv"
    ),
    "vh09_status": REPO / "results/validation_human/09_rosmap_kda_candidates/status.tsv",
    "vh10d_status": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10d_overlap/status.tsv"
    ),
    "candidate_overlap": REPO / (
        "results/validation_human/10_seaad_kda_rediscovery/10d_overlap/"
        "rosmap_seaad_candidate_overlap.tsv"
    ),
    "seaad_config": REPO / "scripts/validation_human/seaad_deg_config.yml",
    "kda_config": REPO / "scripts/validation_human/seaad_phase18_validation_config.yml",
}

FIG = {
    "setup": REPO / (
        "results/figures/validation_human/seaad_rosmap_validation_setup/"
        "seaad_rosmap_validation_setup.png"
    ),
    "mt_circle": REPO / (
        "results/figures/validation_human/seaad_two_case_circular/"
        "seaad_mt_driver_circular.png"
    ),
    "non_mt_circle": REPO / (
        "results/figures/validation_human/seaad_two_case_circular/"
        "seaad_non_mt_driver_circular.png"
    ),
    "strict_overlap": REPO / (
        "results/figures/validation_human/seaad_rosmap_strict_overlap_ranks/"
        "seaad_rosmap_strict_overlap_ranks.png"
    ),
    "gene_overlap": REPO / (
        "results/figures/validation_human/seaad_rosmap_top_driver_gene_overlap_slide/"
        "seaad_rosmap_top_driver_gene_overlap_slide.png"
    ),
    "non_mt_diagnostic": REPO / (
        "results/figures/validation_human/seaad_rosmap_non_mt_diagnostic/"
        "seaad_rosmap_non_mt_diagnostic.png"
    ),
    "deg_landscape": REPO / (
        "results/figures/validation_human/seaad_fine_deg_landscape/"
        "seaad_fine_deg_landscape.png"
    ),
    "kda_outcomes": REPO / (
        "results/figures/validation_human/seaad_kda_call_outcomes/"
        "seaad_kda_call_outcomes.png"
    ),
}

FIG_STATUS = {
    "setup": FIG["setup"].with_name("seaad_rosmap_validation_setup_status.tsv"),
    "circles": FIG["mt_circle"].with_name("seaad_two_case_circular_status.tsv"),
    "strict_overlap": FIG["strict_overlap"].with_name(
        "seaad_rosmap_strict_overlap_ranks_status.tsv"
    ),
    "gene_overlap": FIG["gene_overlap"].with_name(
        "seaad_rosmap_top_driver_gene_overlap_slide_status.tsv"
    ),
    "non_mt_diagnostic": FIG["non_mt_diagnostic"].with_name(
        "seaad_rosmap_non_mt_diagnostic_status.tsv"
    ),
    "deg_landscape": FIG["deg_landscape"].with_name(
        "seaad_fine_deg_landscape_status.tsv"
    ),
    "kda_outcomes": FIG["kda_outcomes"].with_name(
        "seaad_kda_call_outcomes_status.tsv"
    ),
}

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)
MAIN_SLIDES = 9
APPENDIX_SLIDES = 6
EXPECTED_SLIDE_COUNT = MAIN_SLIDES + APPENDIX_SLIDES

NAVY = RGBColor(15, 35, 61)
NAVY_2 = RGBColor(30, 59, 91)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
ORANGE = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(221, 229, 238)
MID = RGBColor(103, 116, 132)
DARK = RGBColor(36, 43, 51)
GRAY = RGBColor(79, 89, 101)
PALE_BLUE = RGBColor(225, 239, 248)
PALE_GREEN = RGBColor(225, 244, 238)
PALE_ORANGE = RGBColor(255, 244, 218)
PALE_RED = RGBColor(253, 235, 228)
PALE_GRAY = RGBColor(240, 243, 246)
FONT = "Arial"

NOTE_HEADINGS = ("What to point at:", "Main takeaway:", "Boundary / transition:")
MIN_NOTE_WORDS = 55

MAIN_TITLES = [
    "SEA-AD independently recovers a focused neuronal mitochondrial signal from ROSMAP",
    "SEA-AD evidence was independent; the network and KDA rules were shared",
    "Only 42 directions produced mitochondrial gene sets large enough for KDA",
    "SEA-AD MT drivers concentrate in excitatory and inhibitory neurons",
    "SEA-AD selected five non-MT drivers across three networks",
    "Six ROSMAP units reappear in the same neuronal network and driver class",
    "Ignoring network, all six SEA-AD MT genes occur in ROSMAP; non-MT overlap is zero",
    "Zero non-MT overlap does not mean the biology is absent",
    "SEA-AD supports a focused neuronal MT signal; broader validation remains incomplete",
]

APPENDIX_TITLES = [
    "Detailed validation setup",
    "Fine-supertype DEG landscape",
    "KDA call outcomes",
    "Query and selection rules",
    "Selected drivers and ROSMAP testability",
    "Provenance and interpretation limits",
]

ALL_TITLES = MAIN_TITLES + APPENDIX_TITLES


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def display_path(path: Path) -> str:
    """Use a repository-relative path when possible, otherwise an absolute path."""
    try:
        return str(path.resolve().relative_to(REPO))
    except ValueError:
        return str(path.resolve())


def read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle, delimiter="\t"))


def as_bool(value: str) -> bool:
    normalized = str(value).strip().lower()
    if normalized in {"true", "t", "1", "yes"}:
        return True
    if normalized in {"false", "f", "0", "no", "", "na"}:
        return False
    raise ValueError(f"Cannot interpret boolean value: {value!r}")


def one_row(path: Path) -> dict[str, str]:
    rows = read_tsv(path)
    if len(rows) != 1:
        raise AssertionError(f"Expected one row in {path}, found {len(rows)}")
    return rows[0]


def assert_validated_status(path: Path, *, figure: bool = False) -> dict[str, str]:
    row = one_row(path)
    if row.get("validation_status") != "validated_complete":
        raise AssertionError(f"Input is not validated_complete: {path}")
    if row.get("failed_checks", "") not in {"", "0"}:
        raise AssertionError(f"Input has failed checks: {path}")
    if figure:
        if row.get("visual_review_status") != "complete":
            raise AssertionError(f"Figure lacks completed visual review: {path}")
        if row.get("failed_blocking_checks", "0") != "0":
            raise AssertionError(f"Figure has blocking failures: {path}")
        if row.get("pending_nonblocking_checks", "0") != "0":
            raise AssertionError(f"Figure has pending checks: {path}")
    return row


def derive_metrics() -> dict[str, object]:
    """Read validated tables and return the values used by the deck."""
    for path in [*DATA.values(), *FIG.values(), *FIG_STATUS.values(), DESIGN_DOC]:
        if not path.exists() or path.stat().st_size == 0:
            raise FileNotFoundError(path)

    vh02 = assert_validated_status(DATA["vh02_status"])
    vh04 = assert_validated_status(DATA["vh04_status"])
    vh10b = assert_validated_status(DATA["vh10b_status"])
    vh10c = assert_validated_status(DATA["vh10c_status"])
    vh09 = assert_validated_status(DATA["vh09_status"])
    vh10d = assert_validated_status(DATA["vh10d_status"])
    for status in FIG_STATUS.values():
        assert_validated_status(status, figure=True)

    donors = read_tsv(DATA["donor_groups"])
    donor_counts: dict[str, dict[str, int]] = defaultdict(dict)
    for row in donors:
        donor_counts[row["signature_group"]][row["diagnosis"]] = int(row["donors"])
    expected_groups = ["F_e2", "F_e33", "F_e4", "M_e2", "M_e33", "M_e4"]
    if sorted(donor_counts) != sorted(expected_groups):
        raise AssertionError("Unexpected SEA-AD sex/APOE group set")
    dementia = sum(values["Dementia"] for values in donor_counts.values())
    no_dementia = sum(values["No dementia"] for values in donor_counts.values())
    if (dementia, no_dementia) != (37, 41):
        raise AssertionError("SEA-AD disease-arm counts changed")

    fine = read_tsv(DATA["fine_contrasts"])
    fine_status = Counter(row["terminal_status"] for row in fine)
    completed_groups = Counter(
        row["signature_group"] for row in fine if row["terminal_status"] == "completed"
    )
    if len(fine) != 774 or fine_status != Counter({"not_estimable": 514, "completed": 260}):
        raise AssertionError("Fine-contrast structural/completion contract changed")
    if completed_groups != Counter({"F_e33": 100, "M_e33": 92, "F_e4": 68}):
        raise AssertionError("Completed fine-contrast group contract changed")

    attrition_rows = read_tsv(DATA["query_attrition"])
    attrition = {row["terminal_status"]: int(row["direction_slots"]) for row in attrition_rows}
    expected_attrition = {
        "source_contrast_not_estimable": 1028,
        "query_empty": 462,
        "query_below_minimum": 16,
        "eligible_small_query": 21,
        "eligible_phase18_sized": 21,
    }
    if attrition != expected_attrition:
        raise AssertionError(f"Query attrition changed: {attrition}")
    planned_directions = sum(attrition.values())
    completed_directions = planned_directions - attrition["source_contrast_not_estimable"]
    kda_calls = attrition["eligible_small_query"] + attrition["eligible_phase18_sized"]
    if (planned_directions, completed_directions, kda_calls) != (1548, 520, 42):
        raise AssertionError("Direction/call arithmetic changed")

    top5_rows = read_tsv(DATA["seaad_top5"])
    selected = [
        row for row in top5_rows
        if row["list_status"] == "ranked_candidates"
        and row["current_symbol"] not in {"", "NA"}
    ]
    selected_units = len(selected)
    selected_genes = len({row["current_symbol"] for row in selected})
    selected_classes = Counter(row["case_id"] for row in selected)
    if (selected_units, selected_genes) != (13, 11):
        raise AssertionError("SEA-AD selected unit/gene count changed")
    if selected_classes != Counter({"mt_driver": 8, "non_mt_driver": 5}):
        raise AssertionError("SEA-AD MT/non-MT selected split changed")
    if int(vh10c["selected_top5_units"]) != selected_units:
        raise AssertionError("VH10C status does not match selected display rows")

    list_status = read_tsv(DATA["seaad_list_status"])
    list_state_counts = Counter(row["list_status"] for row in list_status)
    if list_state_counts != Counter({
        "ranked_candidates": 5,
        "no_passing_candidate": 5,
        "not_testable_no_included_runs": 4,
    }):
        raise AssertionError("SEA-AD list-state contract changed")

    overlap = read_tsv(DATA["candidate_overlap"])
    rosmap_selected = [row for row in overlap if as_bool(row["rosmap_top5"])]
    testable = [row for row in rosmap_selected if as_bool(row["in_common_assessable_universe"])]
    strict = [row for row in rosmap_selected if as_bool(row["seaad_top5"])]
    strict_genes = {row["gene"] for row in strict}
    strict_classes = Counter(row["case_id"] for row in strict)
    replication = Counter(row["replication_status"] for row in rosmap_selected)
    if (len(rosmap_selected), len(testable), len(strict), len(strict_genes)) != (47, 36, 6, 4):
        raise AssertionError("Strict overlap/testability contract changed")
    if strict_classes != Counter({"mt_driver": 6}):
        raise AssertionError("Strict overlap is no longer entirely MT")
    if replication != Counter({
        "tested_not_selected": 30,
        "not_testable": 11,
        "rediscovered_top5": 6,
    }):
        raise AssertionError("ROSMAP selected-unit fate contract changed")

    non_mt_rosmap = [row for row in rosmap_selected if row["case_id"] == "non_mt_driver"]
    non_mt_testable = [row for row in non_mt_rosmap if as_bool(row["in_common_assessable_universe"])]
    if (len(non_mt_rosmap), len(non_mt_testable)) != (21, 17):
        raise AssertionError("Non-MT testability contract changed")

    metrics: dict[str, object] = {
        "donors": int(vh02["analysis_donors"]),
        "dementia_donors": dementia,
        "no_dementia_donors": no_dementia,
        "donor_counts": donor_counts,
        "selected_nuclei": int(vh04["selected_nuclei"]),
        "supertypes": int(vh04["included_supertypes"]),
        "planned_contrasts": len(fine),
        "completed_contrasts": fine_status["completed"],
        "not_estimable_contrasts": fine_status["not_estimable"],
        "completed_groups": completed_groups,
        "planned_directions": planned_directions,
        "completed_directions": completed_directions,
        "kda_calls": kda_calls,
        "selected_units": selected_units,
        "selected_genes": selected_genes,
        "selected_mt_units": selected_classes["mt_driver"],
        "selected_non_mt_units": selected_classes["non_mt_driver"],
        "rosmap_units": len(rosmap_selected),
        "rosmap_genes": int(vh09["selected_unique_genes"]),
        "rosmap_testable": len(testable),
        "rosmap_not_testable": len(rosmap_selected) - len(testable),
        "strict_units": len(strict),
        "strict_genes": len(strict_genes),
        "strict_gene_symbols": sorted(strict_genes),
        "non_mt_rosmap_units": len(non_mt_rosmap),
        "non_mt_testable": len(non_mt_testable),
        "non_mt_not_testable": len(non_mt_rosmap) - len(non_mt_testable),
        "non_mt_same_network_support": 4,
        "non_mt_no_same_network_support": 13,
        "selected_rows": selected,
        "list_state_counts": list_state_counts,
        "kda_return_positive_calls": int(vh10b["completed_significant_calls"]),
        "kda_no_return_calls": int(vh10b["completed_no_significant_calls"]),
        "candidate_units": int(vh10c["candidate_units"]),
    }
    if int(vh10d["rosmap_selected_units"]) != metrics["rosmap_units"]:
        raise AssertionError("VH10D ROSMAP selected count mismatch")
    if int(vh10d["rosmap_testable_selected_units"]) != metrics["rosmap_testable"]:
        raise AssertionError("VH10D testable count mismatch")
    if int(vh10d["strict_shared_top5_units"]) != metrics["strict_units"]:
        raise AssertionError("VH10D strict shared count mismatch")
    return metrics


def fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def stroke(shape, color: RGBColor, width: float = 1.0, dash: str | None = None) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = dash


def set_run(run, *, size: float, color: RGBColor, bold: bool = False,
            italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             color: RGBColor = WHITE, outline: RGBColor | None = LIGHT,
             radius: bool = True, transparency: int = 0,
             line_width: float = 1.0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color, transparency)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline, line_width)
    return shape


def add_circle(slide, x: float, y: float, d: float, *,
               color: RGBColor, outline: RGBColor | None = None,
               line_width: float = 1.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline, line_width)
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *,
             size: float = 18, color: RGBColor = DARK, bold: bool = False,
             italic: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.03):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic)
    return box


def add_rich_text(slide, spans: list[tuple[str, dict]], x: float, y: float,
                  w: float, h: float, *, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    for value, style in spans:
        run = paragraph.add_run()
        run.text = value
        set_run(
            run,
            size=style.get("size", 18),
            color=style.get("color", DARK),
            bold=style.get("bold", False),
            italic=style.get("italic", False),
        )
    return box


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, *,
                  color: RGBColor = BLUE, width: float = 2.5,
                  arrow: bool = True):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    stroke(shape, color, width)
    if arrow:
        shape.line.end_arrowhead = True
    return shape


def set_alt_text(shape, title: str, description: str) -> None:
    props = shape._element.xpath(".//p:cNvPr")
    if props:
        props[0].set("name", title)
        props[0].set("descr", description)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float,
                        *, alt: str):
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    picture_w, picture_h = image_w * scale, image_h * scale
    picture_x = x + (w - picture_w) / 2
    picture_y = y + (h - picture_h) / 2
    picture = slide.shapes.add_picture(
        str(path), Inches(picture_x), Inches(picture_y),
        Inches(picture_w), Inches(picture_h)
    )
    set_alt_text(picture, alt, alt)
    return picture


def new_slide(prs: Presentation, *, bg: RGBColor = WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_header(slide, kicker: str, title: str, page_no: int, *,
               accent: RGBColor = TEAL, subtitle: str | None = None) -> None:
    add_text(slide, kicker.upper(), 0.55, 0.17, 5.6, 0.25,
             size=10.5, color=accent, bold=True)
    if len(title) <= 62:
        title_size = 29.0
    elif len(title) <= 82:
        title_size = 26.0
    else:
        title_size = 23.5
    title_h = 0.73 if len(title) <= 82 else 0.86
    add_text(slide, title, 0.55, 0.43, 12.02, title_h,
             size=title_size, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, 0.57, 1.13, 11.8, 0.34,
                 size=12.5, color=MID)
    add_text(slide, f"{page_no:02d}", 12.49, 0.19, 0.29, 0.19,
             size=9.0, color=MID, bold=True, align=PP_ALIGN.RIGHT)


def add_source(slide, value: str) -> None:
    add_text(slide, value, 0.55, 7.25, 12.16, 0.15,
             size=7.0, color=MID)


def add_note(slide, note: str) -> None:
    slide.notes_slide.notes_text_frame.text = note


def add_metric_card(slide, value: str, label: str, x: float, y: float, w: float, *,
                    accent: RGBColor, background: RGBColor) -> None:
    add_rect(slide, x, y, w, 1.45, color=background, outline=accent,
             line_width=1.7)
    add_text(slide, value, x + 0.22, y + 0.19, w - 0.44, 0.58,
             size=31.0, color=accent, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, label, x + 0.25, y + 0.82, w - 0.50, 0.40,
             size=15.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def add_number_step(slide, number: str, label: str, x: float, y: float, w: float,
                    *, accent: RGBColor, background: RGBColor) -> None:
    add_rect(slide, x, y, w, 2.02, color=background, outline=accent,
             line_width=1.8)
    add_text(slide, number, x + 0.18, y + 0.28, w - 0.36, 0.70,
             size=35.0, color=accent, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, label, x + 0.26, y + 1.08, w - 0.52, 0.66,
             size=16.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def add_native_card(slide, title: str, body: str, x: float, y: float, w: float, h: float,
                    *, accent: RGBColor, background: RGBColor, number: str | None = None) -> None:
    add_rect(slide, x, y, w, h, color=WHITE, outline=LIGHT, line_width=1.2)
    add_rect(slide, x, y, w, 0.20, color=accent, outline=None, radius=False)
    if number:
        add_circle(slide, x + 0.22, y + 0.36, 0.48, color=background,
                   outline=accent, line_width=1.5)
        add_text(slide, number, x + 0.25, y + 0.43, 0.42, 0.26,
                 size=15.0, color=accent, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        title_x = x + 0.82
        title_w = w - 1.05
    else:
        title_x = x + 0.27
        title_w = w - 0.54
    add_text(slide, title, title_x, y + 0.34, title_w, 0.68,
             size=18.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, body, x + 0.29, y + 1.19, w - 0.58, h - 1.43,
             size=15.0, color=DARK, valign=MSO_ANCHOR.TOP)


def figure_slide(prs: Presentation, *, page_no: int, kicker: str, title: str,
                 figure: Path, alt: str, source: str, note: str,
                 accent: RGBColor = TEAL) -> None:
    slide = new_slide(prs)
    add_header(slide, kicker, title, page_no, accent=accent)
    add_picture_contain(slide, figure, 0.48, 1.20, 12.37, 5.84, alt=alt)
    add_source(slide, source)
    add_note(slide, note)


def full_canvas_circle_slide(prs: Presentation, *, page_no: int, figure: Path,
                             alt: str, source: str, note: str) -> None:
    """Preserve the canonical circular asset without creating a second version."""
    slide = new_slide(prs)
    add_picture_contain(slide, figure, 0.61, 0.04, 12.10, 7.20, alt=alt)
    add_text(slide, f"{page_no:02d}", 12.78, 0.14, 0.24, 0.18,
             size=8.0, color=MID, bold=True, align=PP_ALIGN.RIGHT)
    add_source(slide, source)
    add_note(slide, note)


def _notes() -> list[str]:
    return [
        """What to point at: Start with the two large numbers. SEA-AD produced 13 selected network–gene–class units, and six of those match a frozen ROSMAP unit in the same broad network and driver class.

Main takeaway: Independent SEA-AD expression evidence recovers a focused neuronal mitochondrial signal rather than the entire ROSMAP list. MT here means the frozen core-MitoCarta driver class.

Boundary / transition: A unit is a network plus gene plus driver class, so one gene may count in more than one network. Next I will show how the independent evidence and shared analysis scaffold were separated.""",
        """What to point at: Follow the four boxes from SEA-AD donors to donor-level expression, signed mitochondrial DEG queries, KDA on the matching frozen network, and the post-freeze ROSMAP comparison. The SEA-AD candidate list was fixed before candidate identities were used for comparison.

Main takeaway: Donors, expression values, differential-expression results, and queries came independently from SEA-AD. The seven broad networks and the KDA and selection rules were deliberately held fixed for comparability.

Boundary / transition: This is not a new SEA-AD network reconstruction. It is independent expression evidence evaluated on a shared frozen scaffold. Next I will explain why only a small fraction of planned comparisons became KDA calls.""",
        """What to point at: Read the three boxes left to right: 1,548 planned fine-cell-type by group by direction combinations, 520 directions from completed DEG contrasts, and 42 mitochondrial gene sets large enough for KDA.

Main takeaway: A planned direction is not an executed KDA call. Donor support first limited which contrasts could be fitted; then most completed directions had too few effective mitochondrial query genes.

Boundary / transition: Independent donors determine biological replication. More nuclei from the same donors improve measurement but do not add independent samples. The next two slides show the drivers selected from the 42 executed calls.""",
        """What to point at: Focus on the Excitatory and Inhibitory arcs. Eight selected MT units represent six genes. MT-CO2 and MT-CYB appear in both neuronal networks, which is why two recurrence curves cross the center.

Main takeaway: The SEA-AD MT signal is compact and neuronal rather than broadly distributed across all seven networks.

Boundary / transition: The curves indicate the same selected gene appearing in two network lists; they are not biological network edges. Gray and hatched slots distinguish testable lists with no passing driver from networks with no included run. Next we show the SEA-AD non-MT selections.""",
        """What to point at: The selected non-MT genes are HGSNAT in Excitatory neurons, BEX3, RPS27A, and RPL30 in Inhibitory neurons, and KANSL1L in Oligodendrocytes. The other list positions are deliberately not filled.

Main takeaway: SEA-AD did select non-MT drivers; the later zero overlap means these five genes differ from the final ROSMAP non-MT list, not that SEA-AD had no non-MT result.

Boundary / transition: Non-MT means outside the frozen core-MitoCarta class, not unrelated to mitochondria. OPC and Vasculature lists lacked included KDA runs and should not be treated as negative results. Next we move to the prespecified strict comparison.""",
        """What to point at: The top bands separate MT and non-MT. Then follow the rank lines in Excitatory and Inhibitory neurons. Six same-network, same-gene, same-class units recur; because MT-CO2 and MT-CYB recur in both networks, these are four unique gene symbols.

Main takeaway: The primary cross-cohort endpoint is six strict neuronal MT matches and zero strict non-MT matches. Thirty-six of the 47 frozen ROSMAP units had a matching testable SEA-AD universe.

Boundary / transition: The printed p-values are nominal per-list overlap tests. Eleven ROSMAP units had no eligible SEA-AD run and are not negative replications. The next slide deliberately collapses network identity for a simpler gene-level view.""",
        """What to point at: In the MT panel, the SEA-AD circle is fully contained inside the ROSMAP circle: all six SEA-AD MT genes occur somewhere in ROSMAP. In the non-MT panel, the two gene sets are disjoint.

Main takeaway: The descriptive gene-level view shows six shared MT genes and no shared non-MT genes after each gene is counted once, regardless of network.

Boundary / transition: This is secondary to the strict endpoint. MT-ATP6 and MT-ND4 overlap only after networks are ignored, so the six common genes here are not the same as six strict units. No gene-level overlap p-value is claimed. Next we unpack the non-MT mismatch.""",
        """What to point at: Panel A follows the 21 frozen ROSMAP non-MT units. Four OPC units were not testable, 17 had a matching assessable network, four had one SEA-AD supporting run, and none passed final selection across runs. Panel B looks in the reverse direction at the five SEA-AD non-MT units.

Main takeaway: Zero final overlap reflects sparse, differently distributed cross-run evidence; it is not evidence that the genes or biology are absent.

Boundary / transition: Several donor strata could not be estimated because one disease arm had fewer than five independent donors. That reduced matching evidence but is not claimed as the sole cause. The zero occurred before the top-five display cap. We finish with the restrained conclusion.""",
        """What to point at: Read the three cards as supported, not established, and next step. The supported result is the same-network neuronal MT recurrence. Broad non-MT replication and untestable groups remain unresolved.

Main takeaway: SEA-AD provides focused independent support for a neuronal mitochondrial signal, while the wider ROSMAP driver list is not broadly reproduced under this design.

Boundary / transition: Shared networks and selection machinery limit how independent the full analysis is, and KDA prioritization is not causal proof. The practical next step is better donor-balanced external cohorts and, where possible, independently reconstructed networks. The appendix contains the detailed setup and audit evidence.""",
        """What to point at: This is the detailed version of the setup. The top lane shows independent SEA-AD evidence from donors through DEG queries and KDA selection. The lower ROSMAP lane enters only after the SEA-AD list is frozen for comparison.

Main takeaway: The compact main-deck workflow is backed by a checksum-validated, auditable implementation with explicit separation between cohort-specific evidence and shared technical assets.

Boundary / transition: The 1,548 count is a structural grid, not a run count. The original ROSMAP grid spans 54 fine types and six groups, while the SEA-AD design uses its own 129 supertypes. The next appendix slide shows where the DEG signal was concentrated.""",
        """What to point at: The heatmap shows completed fine-supertype contrasts and the distribution of Phase-18-parity DEG incidences. The bottom summary makes the dominant male APOE ε3/ε3 contribution visible.

Main takeaway: Only 260 of 774 fine contrasts completed, and the usable DEG signal was strongly concentrated in a subset of neuronal contrasts, especially male APOE ε3/ε3.

Boundary / transition: Counts here are feature-by-contrast incidences, not unique genes, and they are upstream of symbol deduplication, MitoCarta filtering, and network intersection. The next slide shows what happened after runnable queries entered KDA.""",
        """What to point at: Panel A separates calls with at least one significant return from calls with none. Panel B shows that 42 calls produced run-level evidence, while final across-run selection reduced the result to 13 selected units.

Main takeaway: Within-run evidence was common in the neuronal networks, but final selection across runs was intentionally much stricter.

Boundary / transition: A run-level significant return is not a final selected driver. The 208 return rows are not unique genes, and within-run q-values differ from aggregate ACAT/BH q-values. The next appendix slide summarizes the frozen rules without the implementation detail.""",
        """What to point at: Read the four cards in order: define a signed core-MitoCarta DEG query, run KDA only when the effective query is large enough, aggregate evidence across runs, and display no more than five genes without backfilling.

Main takeaway: SEA-AD changed only the runnable query floor to three genes; ROSMAP used ten. Coverage, conservative support, aggregate correction, driver classes, and ranking remain aligned with the frozen Phase 18 selection logic.

Boundary / transition: The query consists of mitochondrial DEG genes, but candidate drivers can be any assessable network gene and are later classified as MT or non-MT. These rules prioritize evidence; they do not establish causality. The next appendix slide lists the selected SEA-AD units and testability summary.""",
        """What to point at: The left panel lists every SEA-AD selected unit by network and class. The right panel shows the ROSMAP comparison denominator and the three SEA-AD list states.

Main takeaway: All 13 SEA-AD passing units were displayed because no list exceeded the five-gene cap. Of 47 frozen ROSMAP units, 36 were testable and six were strict matches; the 11 untestable units came from networks without an eligible SEA-AD run.

Boundary / transition: Tested but not selected and not testable are different states. The table summarizes the prespecified network–gene–class endpoint, not gene-level overlap. The final appendix slide records provenance and interpretation limits.""",
        """What to point at: The left column lists the evidence controls: validated source tables, validated figure packages, byte-identical embedded images, and speaker notes. The right column lists the main boundaries.

Main takeaway: The deck is a reproducible view of the completed primary analysis. It does not rely on VH05/VH06 QC figures or an unexecuted sensitivity branch.

Boundary / transition: The shared network scaffold is not an independently reconstructed network. Uneven donor coverage limits testability, the compact transfer does not support new gene-level plots for every intermediate, and KDA is a prioritization method rather than causal proof. These points define the appropriate scope of the conclusions.""",
    ]


def build_deck(output_path: Path = DEFAULT_OUT, report_dir: Path = DEFAULT_REPORT_DIR,
               *, visual_review_status: str = "pending") -> Path:
    if visual_review_status not in {"pending", "complete"}:
        raise ValueError("visual_review_status must be 'pending' or 'complete'")
    metrics = derive_metrics()
    notes = _notes()
    if len(notes) != EXPECTED_SLIDE_COUNT:
        raise AssertionError("Speaker-note inventory does not match slide inventory")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Cross-cohort rediscovery of ROSMAP key drivers in SEA-AD"
    prs.core_properties.subject = (
        "Independent SEA-AD expression evidence on a shared frozen network/KDA scaffold"
    )
    prs.core_properties.author = "Alzheimer project analysis team"
    prs.core_properties.keywords = "SEA-AD, ROSMAP, human validation, KDA, MitoCarta"
    prs.core_properties.comments = (
        "Generated from validated VH02/VH04/VH08/VH09/VH10 tables and validated figure packages."
    )

    # 1 — result-first opening.
    slide = new_slide(prs, bg=OFF_WHITE)
    add_text(slide, "HUMAN VALIDATION", 0.68, 0.50, 4.2, 0.28,
             size=12.0, color=TEAL, bold=True)
    add_text(slide, MAIN_TITLES[0], 0.68, 1.03, 11.65, 1.30,
             size=32.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(
        slide,
        "Independent SEA-AD expression evidence analyzed on a shared, frozen network/KDA scaffold",
        0.71, 2.49, 11.65, 0.45, size=17.0, color=GRAY,
    )
    add_metric_card(slide, str(metrics["selected_units"]),
                    "SEA-AD selected units", 1.35, 3.40, 4.58,
                    accent=TEAL, background=PALE_GREEN)
    add_metric_card(slide, str(metrics["strict_units"]),
                    "same-network MT matches", 7.37, 3.40, 4.58,
                    accent=NAVY, background=PALE_BLUE)
    add_rect(slide, 1.35, 5.34, 10.60, 0.78, color=WHITE, outline=LIGHT)
    add_text(slide, "Focused neuronal MT rediscovery—not broad reproduction of every ROSMAP driver",
             1.63, 5.57, 10.04, 0.32, size=17.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_source(slide, "Source: validated VH10C SEA-AD selection and VH10D strict overlap outputs")
    add_note(slide, notes[0])

    # 2 — editable workflow.
    slide = new_slide(prs)
    add_header(slide, "Study setup", MAIN_TITLES[1], 2, accent=TEAL)
    xs = [0.48, 3.66, 6.84, 10.02]
    workflow = [
        ("SEA-AD donors", f"{metrics['donors']} donors\nDementia vs No dementia", TEAL, PALE_GREEN),
        ("Donor-level expression", f"{metrics['supertypes']} supertypes × 6 groups\nSigned mitochondrial DEG query", BLUE, PALE_BLUE),
        ("Frozen KDA scaffold", "Matching broad network\nSame KDA and selection rules", NAVY, PALE_GRAY),
        ("Freeze, then compare", "SEA-AD list fixed first\nROSMAP opened for comparison", ORANGE, PALE_ORANGE),
    ]
    for index, (title, body, accent, background) in enumerate(workflow, start=1):
        add_native_card(slide, title, body, xs[index - 1], 1.57, 2.74, 3.31,
                        accent=accent, background=background, number=str(index))
        if index < 4:
            add_connector(slide, xs[index - 1] + 2.78, 3.20, xs[index] - 0.06, 3.20,
                          color=NAVY, width=2.2)
    add_rect(slide, 3.66, 5.18, 5.92, 0.72, color=PALE_GRAY, outline=NAVY)
    add_text(slide, "Shared: 7 broad networks • MitoCarta annotation • fKDA • selection math",
             3.89, 5.39, 5.46, 0.30, size=14.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Independent: donors • expression • DEG results • signed queries • frozen SEA-AD list",
             1.20, 6.30, 10.93, 0.36, size=16.0, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)
    add_source(slide, "Source: validated VH02/VH04/VH08/VH10A–C outputs; frozen network and selector authorities")
    add_note(slide, notes[1])

    # 3 — editable attrition.
    slide = new_slide(prs)
    add_header(slide, "From design to executed calls", MAIN_TITLES[2], 3, accent=BLUE)
    add_number_step(slide, f"{metrics['planned_directions']:,}",
                    "planned fine cell type × group × direction combinations",
                    0.55, 1.72, 3.45, accent=BLUE, background=PALE_BLUE)
    add_number_step(slide, f"{metrics['completed_directions']:,}",
                    "directions from completed DEG contrasts",
                    4.94, 1.72, 3.45, accent=PURPLE, background=PALE_GRAY)
    add_number_step(slide, f"{metrics['kda_calls']}",
                    "mitochondrial gene sets large enough for KDA",
                    9.33, 1.72, 3.45, accent=TEAL, background=PALE_GREEN)
    add_connector(slide, 4.04, 2.73, 4.87, 2.73, color=NAVY, width=2.5)
    add_connector(slide, 8.43, 2.73, 9.26, 2.73, color=NAVY, width=2.5)
    add_rect(slide, 1.08, 4.38, 11.17, 1.23, color=OFF_WHITE, outline=LIGHT)
    add_text(slide, "Only three of six sex/APOE groups contributed completed fine-supertype contrasts",
             1.37, 4.61, 10.59, 0.34, size=17.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Female ε3/ε3  •  Female ε4  •  Male ε3/ε3",
             1.37, 5.07, 10.59, 0.28, size=16.0, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)
    add_rect(slide, 1.62, 6.02, 10.10, 0.55, color=NAVY, outline=None, radius=False)
    add_text(slide, "Independent donors determine contrast estimability; more nuclei do not add biological replicates",
             1.88, 6.15, 9.58, 0.26, size=14.0, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_source(slide, "Source: VH08 fine_contrast_status.tsv; VH10A query_attrition.tsv")
    add_note(slide, notes[2])

    # 4–5 — retain the sole canonical circular figures, with no alternate rendering.
    full_canvas_circle_slide(
        prs, page_no=4, figure=FIG["mt_circle"],
        alt=("Canonical SEA-AD MT driver circular figure showing eight selected network–gene units, "
             "six genes, and recurrence of MT-CO2 and MT-CYB across neuronal networks"),
        source="Source: validated VH10C selection; canonical seaad_two_case_circular figure package",
        note=notes[3],
    )
    full_canvas_circle_slide(
        prs, page_no=5, figure=FIG["non_mt_circle"],
        alt=("Canonical SEA-AD non-MT driver circular figure showing HGSNAT, BEX3, RPS27A, "
             "RPL30, and KANSL1L across three broad networks"),
        source="Source: validated VH10C selection; canonical seaad_two_case_circular figure package",
        note=notes[4],
    )

    # 6–8 — validated slide-native figures.
    figure_slide(
        prs, page_no=6, kicker="Primary cross-cohort endpoint", title=MAIN_TITLES[5],
        figure=FIG["strict_overlap"],
        alt=("Strict ROSMAP versus SEA-AD overlap rank figure: six shared network–gene–class units, "
             "four unique genes, all in neuronal MT lists; 36 of 47 ROSMAP units testable"),
        source="Source: validated VH09 frozen ROSMAP units and VH10D strict-overlap tables",
        note=notes[5], accent=NAVY,
    )
    figure_slide(
        prs, page_no=7, kicker="Secondary descriptive view", title=MAIN_TITLES[6],
        figure=FIG["gene_overlap"],
        alt=("Gene-level overlap figure with networks collapsed: six common MT genes, four ROSMAP-only "
             "MT genes, no SEA-AD-only MT genes, and disjoint non-MT gene sets"),
        source="Source: validated VH09/VH10C selected lists and VH10D gene-level overlap",
        note=notes[6], accent=PURPLE,
    )
    figure_slide(
        prs, page_no=8, kicker="Interpreting the non-MT result", title=MAIN_TITLES[7],
        figure=FIG["non_mt_diagnostic"],
        alt=("Non-MT diagnostic tracing 21 frozen ROSMAP units to four untestable, 17 testable, "
             "four with one SEA-AD support, and zero passing final selection, with reverse lookup of "
             "five SEA-AD non-MT units in ROSMAP"),
        source="Source: validated VH09, VH10A/B/D outputs and frozen Phase 18 call-return authority",
        note=notes[7], accent=ORANGE,
    )

    # 9 — restrained close.
    slide = new_slide(prs, bg=OFF_WHITE)
    add_header(slide, "Take-home message", MAIN_TITLES[8], 9, accent=TEAL)
    cards = [
        ("Supported", "Same-network neuronal MT rediscovery", TEAL, PALE_GREEN),
        ("Not established", "Replication of the non-MT list or untestable groups", ORANGE, PALE_ORANGE),
        ("Next step", "Better donor-balanced cohorts and independently reconstructed networks", BLUE, PALE_BLUE),
    ]
    for index, (title, body, accent, bg) in enumerate(cards):
        add_native_card(slide, title, body, 0.69 + index * 4.22, 1.68, 3.84, 3.62,
                        accent=accent, background=bg)
    add_rect(slide, 1.03, 5.83, 11.27, 0.72, color=NAVY, outline=None, radius=False)
    add_text(slide, "Focused cross-cohort support—not causal proof or broad validation of every driver",
             1.30, 6.04, 10.73, 0.30, size=16.0, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_source(slide, "Source: synthesis of validated VH08–VH10 results")
    add_note(slide, notes[8])

    # Appendix A1 — detailed setup.
    figure_slide(
        prs, page_no=10, kicker="Appendix A1", title=APPENDIX_TITLES[0],
        figure=FIG["setup"],
        alt=("Detailed SEA-AD validation setup showing independent donor-level DEG queries, frozen shared "
             "networks and KDA rules, SEA-AD selection freeze, and post-freeze ROSMAP comparison"),
        source="Source: validated VH02/VH04/VH08/VH09/VH10 inputs; setup figure package",
        note=notes[9], accent=GRAY,
    )
    figure_slide(
        prs, page_no=11, kicker="Appendix A2", title=APPENDIX_TITLES[1],
        figure=FIG["deg_landscape"],
        alt=("Fine-supertype differential-expression landscape across six sex/APOE groups, with completed "
             "contrasts, Phase-18-parity feature incidences, and strong male APOE epsilon3/epsilon3 concentration"),
        source="Source: validated VH08 fine-supertype DEG summaries",
        note=notes[10], accent=PURPLE,
    )
    figure_slide(
        prs, page_no=12, kicker="Appendix A3", title=APPENDIX_TITLES[2],
        figure=FIG["kda_outcomes"],
        alt=("KDA outcome figure showing 42 calls, 29 with at least one significant return, 13 with none, "
             "and the distinct across-run aggregation to 13 selected driver units"),
        source="Source: validated VH10A–C KDA and selection outputs",
        note=notes[11], accent=BLUE,
    )

    # Appendix A4 — method summary as editable cards.
    slide = new_slide(prs)
    add_header(slide, "Appendix A4", APPENDIX_TITLES[3], 13, accent=GRAY)
    rules = [
        ("Define query", "Signed core-MitoCarta DEG set\nFDR < 0.05 and |log₂FC| > log₂(1.3)", TEAL, PALE_GREEN),
        ("Run KDA", "Intersect with induced-network background\nSEA-AD ≥3 genes; ROSMAP ≥10", BLUE, PALE_BLUE),
        ("Select across runs", "Coverage ≥0.80 • ≥1 support run\naggregate ACAT/BH q ≤0.05", NAVY, PALE_GRAY),
        ("Display", "Rank by q, p, then symbol\nmaximum five; no backfill", ORANGE, PALE_ORANGE),
    ]
    for index, (title, body, accent, bg) in enumerate(rules, start=1):
        x = 0.50 + (index - 1) * 3.19
        add_native_card(slide, title, body, x, 1.55, 2.78, 3.35,
                        accent=accent, background=bg, number=str(index))
        if index < 4:
            add_connector(slide, x + 2.83, 3.18, x + 3.11, 3.18,
                          color=NAVY, width=1.8)
    add_rect(slide, 1.22, 5.36, 10.89, 0.95, color=OFF_WHITE, outline=LIGHT)
    add_text(slide,
             "Mitochondrial genes define the query; every assessable network gene can be a candidate driver",
             1.52, 5.64, 10.29, 0.38, size=16.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_source(slide, "Source: frozen SEA-AD validation config and Phase 18 KDA/selection authorities")
    add_note(slide, notes[12])

    # Appendix A5 — compact selected/testability table.
    slide = new_slide(prs)
    add_header(slide, "Appendix A5", APPENDIX_TITLES[4], 14, accent=GRAY)
    add_rect(slide, 0.55, 1.42, 7.55, 5.34, color=WHITE, outline=LIGHT)
    add_text(slide, "All 13 SEA-AD selected units", 0.82, 1.66, 6.99, 0.39,
             size=20.0, color=NAVY, bold=True)
    selected_by_key: dict[tuple[str, str], list[dict[str, str]]] = defaultdict(list)
    for row in metrics["selected_rows"]:  # type: ignore[index]
        selected_by_key[(row["broad_network"], row["case_id"])].append(row)
    display_rows = [
        ("Excitatory", "MT", ["MT-CO2", "MT-CYB", "MT-ND4", "MT-ATP6"]),
        ("Excitatory", "non-MT", ["HGSNAT"]),
        ("Inhibitory", "MT", ["MT-CO2", "MT-ND5", "MT-CO3", "MT-CYB"]),
        ("Inhibitory", "non-MT", ["BEX3", "RPS27A", "RPL30"]),
        ("Oligodendrocytes", "non-MT", ["KANSL1L"]),
    ]
    y = 2.28
    for index, (network, driver_class, genes) in enumerate(display_rows):
        if index % 2 == 0:
            add_rect(slide, 0.76, y - 0.05, 7.12, 0.72,
                     color=PALE_GRAY, outline=None, radius=False)
        add_text(slide, network, 0.90, y + 0.09, 1.85, 0.28,
                 size=14.0, color=NAVY, bold=True)
        class_color = TEAL if driver_class == "MT" else PURPLE
        add_text(slide, driver_class, 2.82, y + 0.09, 0.79, 0.28,
                 size=14.0, color=class_color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, " • ".join(genes), 3.73, y + 0.07, 3.88, 0.32,
                 size=13.0, color=DARK)
        y += 0.80
    add_text(slide, "Ranked lists are capped at five and are not backfilled.",
             0.84, 6.38, 6.98, 0.24, size=10.5, color=MID, italic=True)

    add_rect(slide, 8.38, 1.42, 4.40, 5.34, color=OFF_WHITE, outline=LIGHT)
    add_text(slide, "ROSMAP comparison", 8.67, 1.66, 3.82, 0.39,
             size=20.0, color=NAVY, bold=True)
    metrics_right = [
        (str(metrics["rosmap_units"]), "frozen selected units", ORANGE),
        (str(metrics["rosmap_testable"]), "testable in SEA-AD", TEAL),
        (str(metrics["strict_units"]), "strict shared units", NAVY),
        (str(metrics["rosmap_not_testable"]), "no eligible SEA-AD run", GRAY),
    ]
    for index, (value, label, accent) in enumerate(metrics_right):
        yy = 2.25 + index * 0.86
        add_text(slide, value, 8.72, yy, 0.78, 0.47,
                 size=24.0, color=accent, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, label, 9.66, yy + 0.08, 2.63, 0.31,
                 size=14.0, color=DARK, bold=True)
    add_text(slide, "SEA-AD list states", 8.70, 5.79, 3.68, 0.29,
             size=14.0, color=NAVY, bold=True)
    add_text(slide, "5 ranked  •  5 tested/no passing  •  4 no runs",
             8.70, 6.15, 3.68, 0.30, size=11.6, color=GRAY)
    add_source(slide, "Source: VH10C seaad_top5/list_status; VH10D candidate-overlap status")
    add_note(slide, notes[13])

    # Appendix A6 — provenance and limits.
    slide = new_slide(prs, bg=OFF_WHITE)
    add_header(slide, "Appendix A6", APPENDIX_TITLES[5], 15, accent=GRAY)
    add_native_card(
        slide, "Reproducibility controls",
        "• Validated source status and checksum gates\n"
        "• Validated color/grayscale figure packages\n"
        "• Embedded PNGs verified byte-for-byte\n"
        "• Source lines, alt text, and speaker notes\n"
        "• Deck inventory, checks, and SHA-256 status",
        0.73, 1.50, 5.89, 4.76, accent=TEAL, background=PALE_GREEN,
    )
    add_native_card(
        slide, "Interpretation limits",
        "• Shared broad networks and selector\n"
        "• Uneven independent-donor coverage\n"
        "• SEA-AD query floor 3 vs ROSMAP 10\n"
        "• No optional sensitivity branch executed\n"
        "• KDA prioritization is not causal proof",
        6.71, 1.50, 5.89, 4.76, accent=ORANGE, background=PALE_ORANGE,
    )
    add_rect(slide, 1.29, 6.51, 10.75, 0.45, color=NAVY, outline=None, radius=False)
    add_text(slide, "VH05/VH06 QC figures are not required for the scientific claims in this deck",
             1.56, 6.60, 10.21, 0.24, size=13.0, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_source(slide, "Source: validated input/figure manifests and deck build report")
    add_note(slide, notes[14])

    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output_path.stem}.", suffix=".pptx", dir=output_path.parent,
        delete=False,
    ) as handle:
        temporary = Path(handle.name)
    try:
        prs.save(temporary)
        os.replace(temporary, output_path)
        output_path.chmod(0o644)
    finally:
        temporary.unlink(missing_ok=True)

    validate_deck(output_path)
    write_reports(
        output_path, report_dir.resolve(), metrics,
        visual_review_status=visual_review_status,
    )
    return output_path


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
    )


def _picture_alt(shape) -> str:
    props = shape._element.xpath(".//p:cNvPr")
    return props[0].get("descr", "") if props else ""


def validate_deck(path: Path) -> None:
    if not path.exists() or path.stat().st_size < 500_000:
        raise AssertionError(f"Deck is missing or unexpectedly small: {path}")
    metrics = derive_metrics()
    prs = Presentation(path)
    if len(prs.slides) != EXPECTED_SLIDE_COUNT:
        raise AssertionError(
            f"Expected {EXPECTED_SLIDE_COUNT} slides, found {len(prs.slides)}"
        )
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        raise AssertionError("Deck is not 13.333333 × 7.5 inch widescreen")

    all_text: list[str] = []
    picture_alt: list[str] = []
    picture_count = 0
    for index, slide in enumerate(prs.slides, start=1):
        slide_text = _slide_text(slide)
        note_text = slide.notes_slide.notes_text_frame.text.strip()
        if not note_text:
            raise AssertionError(f"Slide {index} has no speaker notes")
        if len(note_text.split()) < MIN_NOTE_WORDS:
            raise AssertionError(
                f"Slide {index} notes are too short: {len(note_text.split())} words"
            )
        for heading in NOTE_HEADINGS:
            if heading not in note_text:
                raise AssertionError(f"Slide {index} notes missing heading: {heading}")
        all_text.extend([slide_text, note_text])
        tolerance = Inches(0.02)
        for shape in slide.shapes:
            if shape.left < -tolerance or shape.top < -tolerance:
                raise AssertionError(f"Slide {index} has shape outside top/left bounds")
            if shape.left + shape.width > SLIDE_W + tolerance:
                raise AssertionError(f"Slide {index} has shape beyond right bound")
            if shape.top + shape.height > SLIDE_H + tolerance:
                raise AssertionError(f"Slide {index} has shape beyond bottom bound")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                alt = _picture_alt(shape)
                if not alt:
                    raise AssertionError(f"Slide {index} picture lacks alt text")
                picture_alt.append(alt)

    if picture_count != 8:
        raise AssertionError(f"Expected 8 embedded pictures, found {picture_count}")

    slide8_pictures = [
        shape
        for shape in prs.slides[7].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if len(slide8_pictures) != 1:
        raise AssertionError(
            f"Expected one diagnostic picture on Slide 8, found {len(slide8_pictures)}"
        )
    slide8_picture = slide8_pictures[0]
    slide8_hash = hashlib.sha256(slide8_picture.image.blob).hexdigest()
    expected_slide8_hash = sha256(FIG["non_mt_diagnostic"])
    if slide8_hash != expected_slide8_hash:
        raise AssertionError("Slide 8 does not embed the canonical non-MT diagnostic")
    if "Non-MT diagnostic" not in _picture_alt(slide8_picture):
        raise AssertionError("Slide 8 diagnostic alt text changed")

    joined = "\n".join(all_text)
    required_text = [
        "13 SEA-AD selected units",
        "same-network MT matches",
        "1,548",
        "520",
        "42",
        "Mitochondrial genes define the query",
        "47",
        "36",
        "strict shared units",
        "KDA prioritization is not causal proof",
    ]
    for value in required_text:
        if value not in joined:
            raise AssertionError(f"Required deck text missing: {value}")
    forbidden = [
        "84 calls",
        "six unique strict genes",
        "failed replication",
        "proved causal",
        "biologically absent",
        "1,548 KDA calls",
        "FDR-only sensitivity executed",
    ]
    for value in forbidden:
        if value.lower() in joined.lower():
            raise AssertionError(f"Forbidden deck wording found: {value}")

    expected_alt_fragments = [
        "Canonical SEA-AD MT driver circular",
        "Canonical SEA-AD non-MT driver circular",
        "Strict ROSMAP versus SEA-AD overlap",
        "Gene-level overlap figure",
        "Non-MT diagnostic",
        "Detailed SEA-AD validation setup",
        "Fine-supertype differential-expression",
        "KDA outcome figure",
    ]
    for fragment in expected_alt_fragments:
        if not any(fragment in alt for alt in picture_alt):
            raise AssertionError(f"Expected figure alt text missing: {fragment}")

    # Native-slide titles are editable.  The two circular slides intentionally retain
    # their canonical baked title to avoid creating a second figure version.
    editable_title_indices = [0, 1, 2, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14]
    for index in editable_title_indices:
        if ALL_TITLES[index] not in _slide_text(prs.slides[index]):
            raise AssertionError(f"Slide {index + 1} missing expected editable title")

    with zipfile.ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise AssertionError("PPTX ZIP integrity failed")
        slide_xml = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        if len(slide_xml) != EXPECTED_SLIDE_COUNT:
            raise AssertionError("PPTX slide XML count mismatch")
        media_members = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        media_hashes = {hashlib.sha256(archive.read(name)).hexdigest() for name in media_members}
        for label, figure in FIG.items():
            if sha256(figure) not in media_hashes:
                raise AssertionError(f"Figure not embedded byte-for-byte: {label}")

    # Reassert the easily confused scientific denominators at validation time.
    if metrics["strict_units"] != 6 or metrics["strict_genes"] != 4:
        raise AssertionError("Strict unit/gene distinction changed")
    if metrics["rosmap_testable"] != 36 or metrics["rosmap_units"] != 47:
        raise AssertionError("ROSMAP testability denominator changed")


def _write_tsv(path: Path, rows: Iterable[dict[str, object]], fields: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(f".{path.name}.tmp")
    with temporary.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, delimiter="\t", fieldnames=fields, lineterminator="\n")
        writer.writeheader()
        for row in rows:
            writer.writerow({field: row.get(field, "") for field in fields})
    os.replace(temporary, path)


def write_reports(deck_path: Path, report_dir: Path, metrics: dict[str, object], *,
                  visual_review_status: str) -> None:
    report_dir.mkdir(parents=True, exist_ok=True)
    input_rows: list[dict[str, object]] = []
    inputs = {
        "design": DESIGN_DOC,
        "builder": Path(__file__).resolve(),
        **{f"data:{label}": path for label, path in DATA.items()},
        **{f"figure:{label}": path for label, path in FIG.items()},
        **{f"figure_status:{label}": path for label, path in FIG_STATUS.items()},
    }
    for role, path in inputs.items():
        input_rows.append({
            "role": role,
            "path": display_path(path),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        })
    _write_tsv(
        report_dir / "seaad_rosmap_human_validation_input_manifest.tsv",
        input_rows, ["role", "path", "bytes", "sha256"],
    )

    prs = Presentation(deck_path)
    slide8_pictures = [
        shape
        for shape in prs.slides[7].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    slide8_identity = (
        len(slide8_pictures) == 1
        and hashlib.sha256(slide8_pictures[0].image.blob).hexdigest()
        == sha256(FIG["non_mt_diagnostic"])
    )
    inventory_rows: list[dict[str, object]] = []
    for index, slide in enumerate(prs.slides, start=1):
        pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        inventory_rows.append({
            "slide_number": index,
            "section": "main" if index <= MAIN_SLIDES else "appendix",
            "title": ALL_TITLES[index - 1],
            "picture_count": len(pictures),
            "speaker_note_words": len(slide.notes_slide.notes_text_frame.text.split()),
        })
    _write_tsv(
        report_dir / "seaad_rosmap_human_validation_slide_inventory.tsv",
        inventory_rows,
        ["slide_number", "section", "title", "picture_count", "speaker_note_words"],
    )

    checks = [
        ("pptx_integrity", True, "PPTX ZIP and python-pptx validation passed"),
        ("slide_count", len(prs.slides) == EXPECTED_SLIDE_COUNT,
         f"{MAIN_SLIDES} main + {APPENDIX_SLIDES} appendix"),
        ("widescreen", prs.slide_width == SLIDE_W and prs.slide_height == SLIDE_H,
         "13.333333 × 7.5 inches"),
        ("speaker_notes", all(row["speaker_note_words"] >= MIN_NOTE_WORDS for row in inventory_rows),
         "Every slide has three-section notes"),
        ("validated_figures", True, "All embedded figure-package statuses validated_complete/complete"),
        ("embedded_media_identity", True, "All eight distinct PNG inputs embedded byte-for-byte"),
        ("slide8_non_mt_diagnostic_identity", slide8_identity,
         "Slide 8 embeds the canonical non-MT diagnostic byte-for-byte"),
        ("planned_direction_arithmetic", metrics["planned_directions"] == 1548,
         "1548 = 1028 not estimable + 520 completed-source directions"),
        ("kda_call_arithmetic", metrics["kda_calls"] == 42,
         "42 = 21 size 3–9 + 21 size ≥10"),
        ("selection_counts", metrics["selected_units"] == 13 and metrics["selected_genes"] == 11,
         "13 SEA-AD units / 11 genes"),
        ("strict_overlap_counts", metrics["strict_units"] == 6 and metrics["strict_genes"] == 4,
         "6 network–gene–class units / 4 unique genes"),
        ("testability_denominator", metrics["rosmap_testable"] == 36 and metrics["rosmap_units"] == 47,
         "36 of 47 frozen ROSMAP units testable"),
        ("no_unexecuted_sensitivity", True, "Deck contains no 84-call or executed FDR-only claim"),
        ("visual_review", visual_review_status == "complete",
         "PowerPoint PDF reviewed slide-by-slide in color and grayscale"
         if visual_review_status == "complete" else "Visual review remains pending"),
    ]
    check_rows = [
        {"check_id": check_id, "passed": passed, "detail": detail}
        for check_id, passed, detail in checks
    ]
    _write_tsv(
        report_dir / "seaad_rosmap_human_validation_checks.tsv",
        check_rows, ["check_id", "passed", "detail"],
    )
    blocking_checks = [
        row for row in check_rows
        if row["check_id"] != "visual_review" or visual_review_status == "complete"
    ]
    if not all(row["passed"] for row in blocking_checks):
        raise AssertionError("Deck report contains failed checks")

    status_rows = [{
        "schema_version": "seaad_rosmap_human_validation_deck_v1",
        "deck_id": "seaad_rosmap_human_validation",
        "validation_status": "validated_complete",
        "visual_review_status": visual_review_status,
        "main_slides": MAIN_SLIDES,
        "appendix_slides": APPENDIX_SLIDES,
        "total_slides": EXPECTED_SLIDE_COUNT,
        "input_files": len(input_rows),
        "checks": len(check_rows),
        "deck_path": display_path(deck_path),
        "deck_bytes": deck_path.stat().st_size,
        "deck_sha256": sha256(deck_path),
        "completed_utc": datetime.now(timezone.utc).isoformat(),
    }]
    _write_tsv(
        report_dir / "seaad_rosmap_human_validation_status.tsv",
        status_rows,
        [
            "schema_version", "deck_id", "validation_status", "visual_review_status",
            "main_slides", "appendix_slides", "total_slides", "input_files", "checks",
            "deck_path", "deck_bytes", "deck_sha256", "completed_utc",
        ],
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--report-dir", type=Path, default=DEFAULT_REPORT_DIR)
    parser.add_argument(
        "--visual-review-status", choices=("pending", "complete"), default="pending",
        help="Set to complete only after reviewing a PowerPoint/Keynote render in color and grayscale",
    )
    parser.add_argument("--validate-only", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = args.output.resolve()
    if args.validate_only:
        validate_deck(output)
        print(f"Validated: {output}")
        return
    built = build_deck(
        output, args.report_dir.resolve(),
        visual_review_status=args.visual_review_status,
    )
    print(f"Built and validated: {built}")


if __name__ == "__main__":
    main()
