#!/usr/bin/env python3
"""Rename the audience-facing Modality criterion to Data type in the deck."""

from __future__ import annotations

import argparse
import os
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation


REPO = Path(__file__).resolve().parents[2]
DEFAULT_DECK = REPO / "docs/presentations/mouse_dataset_comparison_08252026.pptx"
OLD_LABEL = "Modality"
NEW_LABEL = "Data type"
EXPECTED_REPLACEMENTS = 13


def _shape_signature(shape) -> tuple:
    return (
        shape.name,
        str(shape.shape_type),
        int(shape.left),
        int(shape.top),
        int(shape.width),
        int(shape.height),
        shape.text if getattr(shape, "has_text_frame", False) else "",
    )


def _deck_signature(prs: Presentation) -> tuple:
    return tuple(
        (
            tuple(_shape_signature(shape) for shape in slide.shapes),
            slide.notes_slide.notes_text_frame.text
            if slide.notes_slide.notes_text_frame is not None else "",
        )
        for slide in prs.slides
    )


def _replace_exact_label(shape) -> bool:
    if not getattr(shape, "has_text_frame", False) or shape.text != OLD_LABEL:
        return False
    runs = [
        run
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text
    ]
    if len(runs) != 1 or runs[0].text != OLD_LABEL:
        raise AssertionError(
            f"Unexpected run structure for {shape.name!r}: {shape.text!r}"
        )
    runs[0].text = NEW_LABEL
    return True


def validate_deck(path: Path, *, before_signature: tuple | None = None) -> None:
    if not path.exists() or path.stat().st_size < 50_000:
        raise AssertionError(f"Deck missing or unexpectedly small: {path}")
    prs = Presentation(path)
    if len(prs.slides) != 15:
        raise AssertionError(f"Expected 15 slides, found {len(prs.slides)}")

    occurrences = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if OLD_LABEL.lower() in shape.text.lower():
                raise AssertionError(
                    f"Old label remains on slide {slide_index}: {shape.text!r}"
                )
            if shape.text == NEW_LABEL:
                occurrences.append((slide_index, shape_index))
    if len(occurrences) != EXPECTED_REPLACEMENTS:
        raise AssertionError(
            f"Expected {EXPECTED_REPLACEMENTS} Data type labels, found "
            f"{len(occurrences)}"
        )
    if [slide for slide, _ in occurrences] != list(range(2, 15)):
        raise AssertionError(
            f"Unexpected slides containing Data type: {occurrences}"
        )

    if before_signature is not None:
        after = _deck_signature(prs)
        if len(after) != len(before_signature):
            raise AssertionError("Slide count changed")
        changes = []
        for slide_index, (before_slide, after_slide) in enumerate(
            zip(before_signature, after), start=1
        ):
            before_shapes, before_notes = before_slide
            after_shapes, after_notes = after_slide
            if before_notes != after_notes:
                changes.append((slide_index, "notes"))
            if len(before_shapes) != len(after_shapes):
                changes.append((slide_index, "shape_count"))
                continue
            for shape_index, (before_shape, after_shape) in enumerate(
                zip(before_shapes, after_shapes)
            ):
                if before_shape == after_shape:
                    continue
                before_text = before_shape[-1]
                after_text = after_shape[-1]
                geometry_same = before_shape[:-1] == after_shape[:-1]
                if not (
                    geometry_same
                    and before_text == OLD_LABEL
                    and after_text == NEW_LABEL
                ):
                    changes.append((slide_index, f"shape_{shape_index}"))
        if changes:
            raise AssertionError(f"Unexpected deck changes: {changes}")

    with zipfile.ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise AssertionError("PPTX ZIP integrity check failed")


def rename_labels(path: Path) -> Path:
    path = path.resolve()
    if not path.exists():
        raise FileNotFoundError(path)
    prs = Presentation(path)
    before = _deck_signature(prs)
    count = sum(
        _replace_exact_label(shape)
        for slide in prs.slides
        for shape in slide.shapes
    )
    if count == 0:
        validate_deck(path)
        return path
    if count != EXPECTED_REPLACEMENTS:
        raise AssertionError(
            f"Expected {EXPECTED_REPLACEMENTS} replacements, made {count}"
        )

    with tempfile.NamedTemporaryFile(
        prefix=f".{path.stem}.", suffix=".tmp.pptx",
        dir=path.parent, delete=False,
    ) as handle:
        temp_path = Path(handle.name)
    try:
        prs.save(temp_path)
        validate_deck(temp_path, before_signature=before)
        os.replace(temp_path, path)
    finally:
        if temp_path.exists():
            temp_path.unlink()
    validate_deck(path)
    return path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--validate-only", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    path = args.deck.resolve()
    if args.validate_only:
        validate_deck(path)
        print(f"Validated: {path}")
        return
    updated = rename_labels(path)
    print(f"Updated and validated: {updated}")


if __name__ == "__main__":
    main()
