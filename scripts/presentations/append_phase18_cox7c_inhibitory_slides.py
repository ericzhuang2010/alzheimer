#!/usr/bin/env python3
"""Append the COX7C inhibitory-neuron pathway and STRING slides."""

from __future__ import annotations

import argparse
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN


REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO))

from scripts.presentations.append_phase18_cox7c_astrocyte_slides import (  # noqa: E402
    replace_text_preserve_style,
)
from scripts.presentations.build_phase18_key_driver_selection_deck import (  # noqa: E402
    BLUE,
    GOLD,
    GRAY,
    LIGHT,
    NAVY,
    PALE_BLUE,
    TEAL,
    VERMILION,
    WHITE,
    add_bullets,
    add_header,
    add_metric,
    add_panel_title,
    add_picture_contain,
    add_rect,
    add_source,
    add_text,
    new_slide,
    trim_white,
)


DEFAULT_DECK = REPO / "docs/presentations/key_driver_selection_analysis 08182026.pptx"
FIG_ROOT = REPO / "results/figures/analysis/phase_18_key_driver_selection/COX7C/inhibitory"
PATHWAY_FIG = FIG_ROOT / "phase18_cox7c_inhibitory_consensus_network_pathways.png"
STRING_FIG = FIG_ROOT / "stringdb_full_medium_conf.png"


def append_slides(input_path: Path, output_path: Path) -> Path:
    for path in (input_path, PATHWAY_FIG, STRING_FIG):
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation(input_path)
    if len(prs.slides) != 27:
        raise RuntimeError(f"Expected a 27-slide source deck, found {len(prs.slides)}")

    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    if "COX7C • INHIBITORY NEURONS" in all_text:
        raise RuntimeError("COX7C inhibitory-neuron slides already appear to be present")

    replace_text_preserve_style(prs.slides[1], "Slides 16–27", "Slides 16–29")
    replace_text_preserve_style(
        prs.slides[8],
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–27 show APOE, RPL11, and COX7C network/protein examples",
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–29 show APOE, RPL11, and COX7C network/protein examples",
    )

    with tempfile.TemporaryDirectory(prefix="cox7c_inhibitory_deck_assets_") as temp_dir:
        assets = Path(temp_dir)
        pathway = trim_white(PATHWAY_FIG, assets)
        string = trim_white(STRING_FIG, assets)

        # 28 — COX7C inhibitory-neuron directed pathway figure
        slide = new_slide(prs)
        add_header(
            slide,
            "COX7C • inhibitory neurons",
            "The inhibitory-neuron signal is stronger and fully stable",
            28,
            accent=GOLD,
            subtitle="Six conservative-support runs yield a broader consensus than astrocytes, with complete leave-one-fine-cell-type candidate retention.",
        )
        add_picture_contain(
            slide,
            pathway,
            0.10,
            1.22,
            8.92,
            5.92,
            alt="COX7C-centered inhibitory-neuron consensus network with directed Bayesian-network edges and pathway outlines",
        )
        add_rect(slide, 9.18, 1.42, 3.62, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.45, 1.70, 3.07, accent=GOLD)
        add_metric(slide, "6", "conservative-support inhibitory runs", 9.45, 2.18, 1.39, accent=GOLD)
        add_metric(slide, "6.24×10⁻⁶", "aggregate ACAT q", 10.99, 2.18, 1.54, accent=BLUE)
        add_panel_title(slide, "Network readout", 9.45, 3.48, 3.07, accent=BLUE)
        add_bullets(
            slide,
            [
                "33 nodes and 34 directed edges retain 15 mitochondrial query hits; candidate-retention fraction is 1.00.",
                "Upstream chain: RPLP1 → RPS15 → COX7C—distinct from the astrocyte RPL11 route.",
                "ETC/OXPHOS: 9 genes, FDR 6.11×10⁻⁸; complex-I biogenesis: 4 genes, FDR 0.0149.",
                "ROS detoxification is contextual (FDR 0.192); 1/6 preserves all hits, while 2/6 retains eight.",
            ],
            9.45,
            3.96,
            2.94,
            size=10.5,
            line_h=0.58,
            accent=BLUE,
        )
        add_text(
            slide,
            "Interpretation: a stable respiratory core with exploratory single-run branches; five AD-down and one AD-up contexts are descriptive, not interactions.",
            9.48,
            6.43,
            2.88,
            0.32,
            size=9.2,
            color=NAVY,
            bold=True,
        )
        add_source(
            slide,
            "Source: phase18_cox7c_inhibitory_consensus_network_pathways.png, caption, methods, ORA table, and gene-by-gene analysis",
        )

        # 29 — COX7C inhibitory-neuron STRING figure and validation interpretation
        slide = new_slide(prs)
        add_header(
            slide,
            "COX7C • inhibitory neurons",
            "STRING reinforces the respiratory core—not edge direction",
            29,
            accent=GOLD,
            subtitle="Medium-confidence STRING associations are undirected and not inhibitory-neuron- or AD-specific.",
        )
        add_rect(slide, 0.42, 1.39, 6.37, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(
            slide,
            string,
            0.64,
            1.62,
            5.93,
            4.98,
            alt="STRING medium-confidence functional association network for COX7C and inhibitory-neuron consensus-neighborhood proteins",
        )
        add_rect(slide, 7.04, 1.39, 5.76, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 7.34, 1.72, 5.15, accent=GOLD)
        add_bullets(
            slide,
            [
                "COX7C joins a dense OXPHOS core containing complex I, III, IV, and V proteins.",
                "SOD1 and TXN connect to the core, but neither STRING connectivity nor the contextual ROS outline proves pathway activation.",
                "A translation/mitochondrial-ribosome branch includes RPS15, RPLP1, MRPS14, MRPL50, SRP14, and TPT1.",
                "Weakly connected or isolated proteins show that STRING does not validate every Bayesian-neighborhood edge.",
            ],
            7.34,
            2.20,
            5.05,
            size=10.8,
            line_h=0.56,
            accent=GOLD,
        )
        add_panel_title(slide, "Cross-network validation", 7.34, 4.64, 5.15, accent=TEAL)
        add_text(
            slide,
            "The stable inhibitory result and the weaker astrocyte result independently nominate COX7C within two different network models. This strengthens prioritization, but a shared bulk-sQTL and expected respiratory-complex connectivity do not establish causality.",
            7.35,
            5.10,
            5.02,
            0.70,
            size=10.5,
            color=GRAY,
        )
        add_rect(slide, 7.35, 5.91, 4.98, 0.59, color=PALE_BLUE, outline=BLUE)
        add_text(
            slide,
            "Use mild perturbation in inhibitory subtypes; compare COX7C with COX4I1 and assay the exact target module, complex-IV assembly/activity, mitonuclear stoichiometry, respiration, ROS, and survival.",
            7.57,
            6.00,
            4.55,
            0.40,
            size=9.3,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "A specific module response at a dose that avoids generic respiratory collapse is the decisive result.",
            7.43,
            6.61,
            4.82,
            0.17,
            size=8.7,
            color=VERMILION,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_source(
            slide,
            "Sources: stringdb_full_medium_conf.png • phase18_key_driver_gene_by_gene_initial_analysis.md • cross-validation and human-genetic-support reviews",
        )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    validate_output(output_path)
    return output_path


def validate_output(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 29:
        raise RuntimeError(f"Expected 29 slides, found {len(prs.slides)}")
    image_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    if image_shapes != 18:
        raise RuntimeError(f"Expected 18 figure images, found {image_shapes}")

    endings = []
    for slide_index in range(len(prs.slides) - 2, len(prs.slides)):
        slide = prs.slides[slide_index]
        endings.append(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
        )
    if not all("COX7C • INHIBITORY NEURONS" in text for text in endings):
        raise RuntimeError("The final two slides are not the expected COX7C inhibitory-neuron slides")

    with zipfile.ZipFile(path) as archive:
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels"))
        )
    if "filter_attrition" in package_text:
        raise RuntimeError("Deprecated filter_attrition content found in the PPTX package")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    args = parser.parse_args()
    result = append_slides(args.input, args.output)
    print(result)
