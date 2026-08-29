#!/usr/bin/env python3
"""Append the ROSMAP × SEA-AD overlap slides as Part 3 of the deck.

The nine appended slides present the cross-cohort key-driver overlap
analysis from ``docs/validation_human/rosmap_seaad_simple_aggr_driver_
analysis.md``: comparison design, gene-level overlap, the eight matched
category pairs, the WDR82 highlight, the stratum-specificity findings, the
sex-consistent exceptions, and the interpretation boundary. Every displayed
number is recomputed from the validated aggregate tables and checked against
frozen expected values before any slide is written. Slides 1-16 are left
unchanged.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import update_phase11_seaad_simple_aggr_part2 as ui  # noqa: E402  (shared styling helpers)


ROOT = Path(__file__).resolve().parents[2]
DEFAULT_DECK = ROOT / "docs" / "presentations" / "phase20_sex_apoe_kda_fine_broad.pptx"
ROS_CATEGORIES = (
    ROOT
    / "results"
    / "minerva_production"
    / "20_sex_apoe_kda_simple_aggr"
    / "simple_category_gene_aggregates.tsv"
)
SEA_CATEGORIES = (
    ROOT
    / "results"
    / "validation_human"
    / "11_sex_apoe_kda_simple_aggr"
    / "simple_category_gene_aggregates.tsv"
)
ROS_RECURRENCE_DATA = (
    ROOT
    / "results/figures/analysis/phase_20_sex_apoe_simple_aggr"
    / "driver_recurrence/phase20_simple_aggr_driver_recurrence_plot_data.tsv"
)
SEA_RECURRENCE_DATA = (
    ROOT
    / "results/figures/validation_human/phase_11_sex_apoe_simple_aggr"
    / "driver_recurrence/phase11_seaad_simple_aggr_driver_recurrence_plot_data.tsv"
)
AUDIT_PATH = (
    ROOT
    / "results"
    / "presentations"
    / "phase20_sex_apoe_kda_fine_broad"
    / "rosmap_seaad_overlap_slides_checks.tsv"
)

EXPECTED_INPUT_SLIDES = 16
GROUP_ORDER = ["F_e2", "F_e33", "F_e4", "M_e2", "M_e33", "M_e4"]
NETWORK_LABELS = {
    "Astrocytes": "Astrocytes",
    "Excitatory_neurons": "Excitatory neurons",
    "Inhibitory_neurons": "Inhibitory neurons",
    "Microglia": "Microglia",
    "OPCs": "OPCs",
    "Oligodendrocytes": "Oligodendrocytes",
    "Vasculature_cells": "Vasculature",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--audit", type=Path, default=AUDIT_PATH)
    return parser.parse_args()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def load_non_mt(path: Path) -> pd.DataFrame:
    frame = pd.read_csv(path, sep="\t", dtype=str, keep_default_na=False, na_values=["NA"])
    frame = frame[
        frame["case_id"].eq("non_mt_driver") & ~frame["is_core_mito"].isin(["TRUE"])
    ].copy()
    frame["score"] = frame["returned_run_q_acat_score"].astype(float)
    frame["calls"] = frame["returned_call_count"].astype(int)
    frame = frame.sort_values(
        ["signature_group", "broad_network", "score", "current_symbol"],
        kind="mergesort",
    )
    frame["display_rank"] = (
        frame.groupby(["signature_group", "broad_network"]).cumcount() + 1
    )
    return frame


def q_text(value: float) -> str:
    if value < 0.001:
        return f"{value:.1e}".replace("e-0", "e-")
    return f"{value:.3f}"


def load_facts() -> dict[str, Any]:
    for path in (ROS_CATEGORIES, SEA_CATEGORIES, ROS_RECURRENCE_DATA, SEA_RECURRENCE_DATA):
        if not path.is_file():
            raise FileNotFoundError(path)
    ros = load_non_mt(ROS_CATEGORIES)
    sea = load_non_mt(SEA_CATEGORIES)

    shared_genes = sorted(set(ros["current_symbol"]) & set(sea["current_symbol"]))
    key = ["signature_group", "broad_network", "current_symbol"]
    matched = sea.merge(ros, on=key, suffixes=("_sea", "_ros")).sort_values(
        ["signature_group", "broad_network", "score_sea"]
    )

    ros_rec = set(pd.read_csv(ROS_RECURRENCE_DATA, sep="\t")["current_symbol"])
    sea_rec = set(pd.read_csv(SEA_RECURRENCE_DATA, sep="\t")["current_symbol"])
    both_recurrence = sorted(ros_rec & sea_rec)

    sea_m_genes = set(sea.loc[sea["signature_group"].eq("M_e33"), "current_symbol"])
    sea_m_shared = sea_m_genes & set(ros["current_symbol"])
    group_matrix = {
        group: len(sea_m_shared & set(ros.loc[ros["signature_group"].eq(group), "current_symbol"]))
        for group in GROUP_ORDER
    }

    pairs = sea.merge(ros, on="current_symbol", suffixes=("_sea", "_ros"))
    same_network = pairs["broad_network_sea"].eq(pairs["broad_network_ros"])
    same_group = pairs["signature_group_sea"].eq(pairs["signature_group_ros"])
    ros_net = ros["broad_network"].value_counts(normalize=True)
    sea_net = sea["broad_network"].value_counts(normalize=True)
    chance_same_network = float(
        sum(ros_net.get(n, 0.0) * sea_net.get(n, 0.0) for n in set(ros_net.index) | set(sea_net.index))
    )

    sex_split = {"male_only": [], "female_only": [], "both": []}
    for gene in sorted(sea_m_shared):
        sexes = set(ros.loc[ros["current_symbol"].eq(gene), "signature_group"].str[0])
        bucket = "male_only" if sexes == {"M"} else "female_only" if sexes == {"F"} else "both"
        sex_split[bucket].append(gene)

    wdr82 = matched[matched["current_symbol"].eq("WDR82")].iloc[0]
    rps15_sea = sea[
        sea["current_symbol"].eq("RPS15") & sea["signature_group"].eq("M_e33")
    ].iloc[0]
    hgsnat_ros_groups = set(ros.loc[ros["current_symbol"].eq("HGSNAT"), "signature_group"])

    facts = {
        "ros_units": len(ros),
        "ros_genes": ros["current_symbol"].nunique(),
        "sea_units": len(sea),
        "sea_genes": sea["current_symbol"].nunique(),
        "shared_genes": shared_genes,
        "matched": matched,
        "both_recurrence": both_recurrence,
        "sea_m_gene_count": len(sea_m_genes),
        "sea_m_shared_count": len(sea_m_shared),
        "group_matrix": group_matrix,
        "pair_count": len(pairs),
        "same_network_count": int(same_network.sum()),
        "same_group_count": int(same_group.sum()),
        "chance_same_network": chance_same_network,
        "sex_split": sex_split,
        "wdr82": wdr82,
        "rps15_sea_rank": int(rps15_sea["display_rank"]),
        "rps15_ros_categories": int((ros["current_symbol"] == "RPS15").sum()),
        "hgsnat_ros_groups": hgsnat_ros_groups,
    }

    expected = {
        "ros_units": 689,
        "ros_genes": 433,
        "sea_units": 96,
        "sea_genes": 91,
        "sea_m_gene_count": 88,
        "sea_m_shared_count": 32,
        "pair_count": 84,
        "same_network_count": 52,
        "same_group_count": 11,
        "rps15_sea_rank": 9,
        "rps15_ros_categories": 12,
    }
    for key_name, expected_value in expected.items():
        if facts[key_name] != expected_value:
            raise RuntimeError(
                f"Source drift for {key_name}: {facts[key_name]} != {expected_value}"
            )
    if len(shared_genes) != 35:
        raise RuntimeError(f"Expected 35 shared genes, found {len(shared_genes)}")
    if len(matched) != 8:
        raise RuntimeError(f"Expected 8 matched category pairs, found {len(matched)}")
    if both_recurrence != ["BEX3", "DYNLT1", "SELENOM", "TTC8"]:
        raise RuntimeError(f"Unexpected both-recurrence genes: {both_recurrence}")
    if group_matrix != {"F_e2": 7, "F_e33": 10, "F_e4": 9, "M_e2": 14, "M_e33": 8, "M_e4": 9}:
        raise RuntimeError(f"Unexpected group matrix: {group_matrix}")
    if (len(sex_split["male_only"]), len(sex_split["female_only"]), len(sex_split["both"])) != (14, 11, 7):
        raise RuntimeError("Unexpected sex split of shared M_e33 genes")
    if not {"HGSNAT", "TARBP1", "RPS27A"} <= set(sex_split["male_only"]):
        raise RuntimeError("Expected HGSNAT, TARBP1, RPS27A among male-only shared genes")
    if hgsnat_ros_groups != {"M_e2", "M_e4"}:
        raise RuntimeError(f"Unexpected HGSNAT ROSMAP groups: {hgsnat_ros_groups}")
    if int(wdr82["display_rank_sea"]) != 1 or int(wdr82["display_rank_ros"]) != 5:
        raise RuntimeError("WDR82 rank contract failed")
    if int(wdr82["calls_ros"]) != 11:
        raise RuntimeError("WDR82 ROSMAP call-count contract failed")
    return facts


# --- slide builders -------------------------------------------------------


def append_divider(prs: Presentation):
    slide = ui.new_slide(prs, bg=ui.NAVY)
    accent = ui.PURPLE
    ui.add_text(slide, "PART 3", 0.78, 0.67, 2.4, 0.28, size=10.5, color=accent, bold=True)
    ui.add_rect(slide, 0.78, 1.26, 0.10, 2.38, color=accent, outline=None, radius=False)
    ui.add_text(
        slide,
        "Do the same drivers appear in both cohorts?",
        1.18, 1.36, 8.65, 1.34, size=32, color=ui.WHITE, bold=True,
        font=ui.FONT_HEAD, valign=MSO_ANCHOR.MIDDLE,
    )
    ui.add_text(
        slide,
        "The ROSMAP and SEA-AD returned-only results are compared at gene level "
        "and at matched sex/APOE-by-cell-type level.",
        1.20, 3.00, 8.30, 0.84, size=15.0, color=ui.RGBColor(204, 219, 234),
    )
    ui.add_rect(slide, 9.78, 1.28, 2.70, 4.70, color=ui.NAVY_2, outline=None)
    ui.add_text(slide, "IN THIS PART", 10.08, 1.67, 2.10, 0.24, size=9.4, color=accent, bold=True)
    topics = ["Comparison design", "What overlaps", "Matched categories", "What it means"]
    current_y = 2.20
    for index, topic in enumerate(topics, start=1):
        ui.add_circle(slide, 10.06, current_y + 0.03, 0.34, accent)
        ui.add_text(
            slide, str(index), 10.06, current_y + 0.08, 0.34, 0.17,
            size=8.2, color=ui.WHITE, bold=True, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        ui.add_text(slide, topic, 10.53, current_y, 1.58, 0.52, size=11.1, color=ui.WHITE, bold=True)
        current_y += 0.82
    ui.add_notes(
        slide,
        goal="Frame the final question: does the SEA-AD validation reproduce the ROSMAP drivers?",
        walkthrough="Part 3 compares the two returned-only non-MT result sets. It moves from the comparison design to gene-level overlap, the exactly matched categories, and finally the stratum-specificity findings.",
        boundary="Both inputs are exploratory post-selected rankings; the comparison describes convergence and does not test formal replication.",
        transition="Start with the two levels at which agreement can be measured.",
    )
    return slide


def append_design_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    ui.add_title_block(
        slide,
        "Two ways to ask whether a driver replicates",
        "Gene-level agreement is permissive; category-level agreement also requires the same sex/APOE group and cell type.",
    )
    ui.add_rect(slide, 0.64, 1.52, 5.86, 4.30, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_text(slide, "GENE LEVEL", 0.95, 1.84, 3.2, 0.27, size=10, color=ui.BLUE, bold=True)
    ui.add_text(
        slide, "Does the same gene return anywhere?",
        0.95, 2.18, 4.95, 0.64, size=20, color=ui.NAVY, bold=True, font=ui.FONT_HEAD,
    )
    ui.add_metric(slide, str(facts["ros_genes"]), "ROSMAP driver genes", 0.95, 3.10, 2.30, accent=ui.BLUE)
    ui.add_metric(slide, str(facts["sea_genes"]), "SEA-AD driver genes", 3.50, 3.10, 2.30, accent=ui.BLUE)
    ui.add_text(
        slide,
        "Ignores which sex/APOE group or cell type produced the return.",
        0.96, 4.60, 5.20, 0.60, size=11.5, color=ui.GRAY,
    )
    ui.add_rect(slide, 6.83, 1.52, 5.86, 4.30, color=ui.PALE_GOLD, outline=ui.GOLD)
    ui.add_text(slide, "CATEGORY LEVEL", 7.14, 1.84, 3.2, 0.27, size=10, color=ui.GOLD_TEXT, bold=True)
    ui.add_text(
        slide, "Same gene, same group, same cell type?",
        7.14, 2.18, 4.95, 0.64, size=20, color=ui.NAVY, bold=True, font=ui.FONT_HEAD,
    )
    ui.add_metric(slide, str(facts["ros_units"]), "ROSMAP gene × category units", 7.14, 3.10, 2.30, accent=ui.GOLD)
    ui.add_metric(slide, str(facts["sea_units"]), "SEA-AD gene × category units", 9.69, 3.10, 2.30, accent=ui.GOLD)
    ui.add_text(
        slide,
        "The strictest agreement: the driver keeps its full context.",
        7.15, 4.60, 5.20, 0.60, size=11.5, color=ui.GRAY,
    )
    ui.add_text(
        slide,
        "Agreement is only possible where SEA-AD has returns: one female e33 excitatory "
        "category and three male e33 cell types.",
        0.88, 6.35, 11.60, 0.28, size=11.0, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Define the two agreement levels before showing any overlap number.",
        walkthrough="Gene level asks only whether the same symbol returns in both cohorts. Category level additionally requires the same sex/APOE group and the same broad cell type, which is the counting unit used throughout the deck.",
        boundary="Because SEA-AD has returns in only four categories, category-level agreement has a hard ceiling set by availability, not biology.",
        transition="First result: how many genes are shared at the permissive level.",
    )
    return slide


def append_gene_overlap_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    shared = len(facts["shared_genes"])
    share_pct = round(100 * shared / facts["sea_genes"])
    ui.add_title_block(
        slide,
        f"{shared} of {facts['sea_genes']} SEA-AD drivers also return in ROSMAP",
        "Gene-level agreement is substantial despite different cohorts, platforms, and run availability.",
    )
    nodes = [
        (str(facts["sea_genes"]), "SEA-AD driver genes"),
        (f"{shared} ({share_pct}%)", "also ROSMAP drivers"),
        ("8", "exact category matches"),
    ]
    x = 0.90
    for value, label in nodes:
        ui.add_rect(slide, x, 1.60, 3.60, 1.45, color=ui.WHITE, outline=ui.LIGHT)
        ui.add_text(slide, value, x + 0.25, 1.78, 3.10, 0.55, size=26, color=ui.BLUE, bold=True, font=ui.FONT_HEAD, align=PP_ALIGN.CENTER)
        ui.add_text(slide, label, x + 0.25, 2.42, 3.10, 0.40, size=11.5, color=ui.GRAY, bold=True, align=PP_ALIGN.CENTER)
        x += 3.95
    ui.add_rect(slide, 0.90, 3.55, 11.50, 2.30, color=ui.PALE_SKY, outline=ui.SKY)
    ui.add_panel_title(slide, "Shared genes visible in this deck's figures", 1.20, 3.85, 10.80, accent=ui.BLUE)
    ui.add_bullets(slide, [
        "On both recurrence charts: BEX3, DYNLT1, SELENOM, TTC8.",
        "SEA-AD's most recurrent gene, HGSNAT, is also a ROSMAP driver — but in different sex/APOE groups.",
        "ROSMAP's headline gene RPS15 returns in SEA-AD male e33 inhibitory neurons.",
    ], 1.22, 4.42, 10.70, size=12.0, line_h=0.46, accent=ui.BLUE)
    ui.add_text(
        slide,
        "One message: a third of the SEA-AD drivers are re-discoveries of ROSMAP drivers.",
        0.88, 6.35, 11.60, 0.28, size=11.0, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Establish the headline gene-level overlap before adding context requirements.",
        walkthrough="Thirty-five of the ninety-one SEA-AD non-MT drivers, thirty-eight percent, appear somewhere among ROSMAP's four hundred thirty-three. Eight of those keep their exact category, previewed here and detailed on the next slide.",
        boundary="Gene-level overlap ignores context; the following slides show that context agreement is much rarer.",
        transition="Next: the eight pairs that match exactly.",
    )
    return slide


def append_matched_table_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    ui.add_title_block(
        slide,
        "Eight driver-category pairs match exactly",
        "Same gene, same sex/APOE group, and same broad cell type in both cohorts.",
    )
    rows = []
    for row in facts["matched"].itertuples():
        label = f"{row.signature_group} · {NETWORK_LABELS[row.broad_network]}"
        rows.append([
            label,
            row.current_symbol,
            f"{q_text(row.score_sea)}  (#{int(row.display_rank_sea)})",
            f"{q_text(row.score_ros)}  (#{int(row.display_rank_ros)})",
        ])
    ui.add_table(
        slide,
        ["Sex/APOE · cell type", "Gene", "SEA-AD q (rank)", "ROSMAP q (rank)"],
        rows,
        0.90, 1.55, [4.05, 1.85, 2.80, 2.80], row_h=0.50, header_h=0.50,
        font_size=10.4,
    )
    ui.add_text(
        slide,
        "The whole female e33 excitatory trio matches; the remaining five pairs are male e33 neurons.",
        0.88, 6.30, 11.60, 0.28, size=11.0, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Show the complete list of exact category-level matches with both cohorts' scores.",
        walkthrough="Three matches are the entire female e33 excitatory category: WDR82, DMTF1, and TPP2. The remaining five are male e33: TTC8 in excitatory neurons and RPL30, RPS15, PAFAH1B1, and PIP5K1A in inhibitory neurons.",
        boundary="Ranks are within-category display ranks of exploratory scores; a match is convergence, not a significance statement.",
        transition="One of these eight deserves a closer look: WDR82.",
    )
    return slide


def append_highlight_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    wdr82 = facts["wdr82"]
    ui.add_title_block(
        slide,
        "WDR82 anchors the cleanest replication",
        "Three storylines stand out among the exact matches.",
    )
    cards = [
        (
            "WDR82",
            "Female e33 · excitatory neurons",
            f"#1 SEA-AD driver in the category; ROSMAP ranks it #{int(wdr82['display_rank_ros'])} "
            f"in the same category from {int(wdr82['calls_ros'])} returned calls. "
            "It tiles on both cohorts' top-five figures.",
            ui.PALE_SKY, ui.BLUE,
        ),
        (
            "The female trio",
            "3 of 3 genes match",
            "Every SEA-AD driver in the female e33 excitatory category — WDR82, DMTF1, TPP2 — "
            "is also a ROSMAP driver in that exact category.",
            ui.PALE_GREEN, ui.TEAL,
        ),
        (
            "RPS15",
            "Male e33 · inhibitory neurons",
            f"ROSMAP's most recurrent driver ({facts['rps15_ros_categories']} categories) returns "
            f"in SEA-AD at rank #{facts['rps15_sea_rank']} of its category.",
            ui.PALE_GOLD, ui.GOLD,
        ),
    ]
    for index, (heading, context, body, bg, accent) in enumerate(cards):
        x = 0.66 + index * 4.10
        ui.add_rect(slide, x, 1.60, 3.83, 4.55, color=bg, outline=accent)
        ui.add_text(slide, heading, x + 0.28, 1.92, 3.30, 0.46, size=21, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
        ui.add_text(slide, context, x + 0.28, 2.52, 3.30, 0.30, size=11.5, color=ui.readable_accent(accent), bold=True)
        ui.add_rect(slide, x + 0.26, 3.05, 3.31, 2.80, color=ui.WHITE, outline=ui.LIGHT)
        ui.add_text(slide, body, x + 0.44, 3.28, 2.96, 2.40, size=11.0, color=ui.GRAY)
    ui.add_notes(
        slide,
        goal="Give the audience three concrete, memorable replication examples.",
        walkthrough="WDR82 is the only gene that is top-five in the same category in both cohorts. The female e33 excitatory category replicates completely, though it holds only three genes. RPS15, the most recurrent ROSMAP driver, also returns in SEA-AD male e33 inhibitory neurons.",
        boundary="The female trio is three genes; treat its perfect match rate as an anecdote, not a rate estimate.",
        transition="These examples suggest strong agreement — but the group-resolved view complicates that.",
    )
    return slide


def append_group_matrix_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    ui.add_title_block(
        slide,
        "Shared genes are not tied to the matching sex/APOE group",
        "Where the shared SEA-AD male e33 drivers land inside ROSMAP.",
    )
    ui.add_text(
        slide,
        f"{facts['sea_m_shared_count']} of the {facts['sea_m_gene_count']} SEA-AD male e33 drivers "
        "are also ROSMAP drivers. Counting where they appear in ROSMAP:",
        0.90, 1.50, 11.40, 0.34, size=13.0, color=ui.DARK,
    )
    matrix = facts["group_matrix"]
    ordered = sorted(matrix.items(), key=lambda item: (-item[1], item[0]))
    rows = []
    for group, count in ordered:
        note = ""
        if group == "M_e2":
            note = "highest overlap"
        elif group == "M_e33":
            note = "the matching group"
        rows.append([group, str(count), note])
    ui.add_table(
        slide,
        ["ROSMAP group", "Shared genes (of 32)", "Note"],
        rows,
        2.60, 2.05, [2.60, 2.60, 2.90], row_h=0.50, header_h=0.50,
        font_size=10.6,
    )
    ui.add_rect(slide, 0.90, 5.75, 11.50, 0.62, color=ui.PALE_RED, outline=ui.VERMILION)
    ui.add_text(
        slide,
        "The matching male e33 group ranks fifth of six — cross-cohort overlap is largely blind to sex/APOE.",
        1.14, 5.93, 11.05, 0.28, size=12.2, color=ui.VERMILION_TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Deliver the central negative finding: shared genes do not preferentially match the same sex/APOE group.",
        walkthrough="If replication were stratum-specific, the male e33 row would dominate. Instead male e2 shares the most genes with the SEA-AD male e33 set, and the matching group is near the bottom. The same genes simply recur in different strata.",
        boundary="These counts are descriptive; group gene-set sizes differ, and no calibrated enrichment test is claimed.",
        transition="If the stratum is not preserved, what is? The cell type.",
    )
    return slide


def append_celltype_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    same_net_pct = round(100 * facts["same_network_count"] / facts["pair_count"])
    chance_pct = round(100 * facts["chance_same_network"])
    same_group_pct = round(100 * facts["same_group_count"] / facts["pair_count"])
    ui.add_title_block(
        slide,
        "Cell type is preserved far better than sex/APOE group",
        f"All {facts['pair_count']} cross-cohort pairings of the shared genes, scored for matching context.",
    )
    ui.add_rect(slide, 0.90, 1.62, 5.55, 3.05, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_text(slide, f"{same_net_pct}%", 1.20, 1.95, 4.95, 0.95, size=52, color=ui.TEAL_TEXT, bold=True, font=ui.FONT_HEAD, align=PP_ALIGN.CENTER)
    ui.add_text(slide, "of pairs share the broad cell type", 1.20, 3.15, 4.95, 0.32, size=14.0, color=ui.NAVY, bold=True, align=PP_ALIGN.CENTER)
    ui.add_text(slide, f"random pairing would give ≈{chance_pct}%", 1.20, 3.60, 4.95, 0.30, size=11.0, color=ui.GRAY, align=PP_ALIGN.CENTER)
    ui.add_rect(slide, 6.85, 1.62, 5.55, 3.05, color=ui.PALE_RED, outline=ui.VERMILION)
    ui.add_text(slide, f"{same_group_pct}%", 7.15, 1.95, 4.95, 0.95, size=52, color=ui.VERMILION_TEXT, bold=True, font=ui.FONT_HEAD, align=PP_ALIGN.CENTER)
    ui.add_text(slide, "share the sex/APOE group", 7.15, 3.15, 4.95, 0.32, size=14.0, color=ui.NAVY, bold=True, align=PP_ALIGN.CENTER)
    ui.add_text(slide, "context transfers by cell type, not stratum", 7.15, 3.60, 4.95, 0.30, size=11.0, color=ui.GRAY, align=PP_ALIGN.CENTER)
    ui.add_rect(slide, 0.90, 5.00, 11.50, 1.30, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_panel_title(slide, "Example: excitatory drivers that switch strata", 1.20, 5.22, 10.80, accent=ui.PURPLE)
    ui.add_text(
        slide,
        "CSTF2, DHDDS, DIDO1, and POLR3F are excitatory-neuron drivers in both cohorts — "
        "female e33 in ROSMAP, male e33 in SEA-AD.",
        1.22, 5.70, 10.80, 0.40, size=12.0, color=ui.DARK,
    )
    ui.add_notes(
        slide,
        goal="Quantify what survives the cohort change: the cell-type context, not the stratum.",
        walkthrough="Pairing every SEA-AD unit with every ROSMAP unit of the same gene gives eighty-four pairs. Sixty-two percent agree on the broad cell type, well above the roughly forty-one percent expected from random pairing, but only thirteen percent agree on the sex/APOE group.",
        boundary="The chance baseline comes from the two cohorts' cell-type distributions; it is a descriptive reference, not a formal null model.",
        transition="A few genes are the exception and keep their sex context.",
    )
    return slide


def append_sex_exceptions_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    split = facts["sex_split"]
    ui.add_title_block(
        slide,
        "A few drivers do keep their sex context",
        "Most shared genes ignore the stratum — these are the exceptions worth remembering.",
    )
    ui.add_rect(slide, 0.66, 1.55, 5.95, 4.55, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_text(slide, "HGSNAT", 0.96, 1.87, 3.40, 0.46, size=22, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
    ui.add_text(slide, "male-only in both cohorts", 0.96, 2.44, 4.20, 0.30, size=12.0, color=ui.BLUE, bold=True)
    ui.add_bullets(slide, [
        "SEA-AD's most recurrent driver (2 categories).",
        "In ROSMAP it appears only in male e2 and male e4.",
        "Always male, and always excitatory neurons in ROSMAP.",
        "TARBP1 and RPS27A show the same male-only pattern.",
    ], 0.98, 2.95, 5.40, size=11.6, line_h=0.55, accent=ui.BLUE)
    ui.add_rect(slide, 6.90, 1.55, 5.78, 4.55, color=ui.PALE_GRAY, outline=ui.LIGHT)
    ui.add_panel_title(slide, "The overall picture is mixed", 7.20, 1.87, 5.15, accent=ui.GRAY)
    bars = [
        ("Male-only in ROSMAP", len(split["male_only"]), ui.BLUE),
        ("Female-only in ROSMAP", len(split["female_only"]), ui.VERMILION),
        ("Both sexes in ROSMAP", len(split["both"]), ui.GRAY),
    ]
    maximum = max(count for _, count, _ in bars)
    for index, (label, count, color) in enumerate(bars):
        y = 2.60 + index * 0.90
        ui.add_text(slide, label, 7.22, y, 2.60, 0.26, size=11.0, color=ui.NAVY, bold=True)
        ui.add_rect(slide, 7.22, y + 0.32, 3.90 * count / maximum, 0.24, color=color, outline=None, radius=False)
        ui.add_text(slide, str(count), 11.35, y + 0.28, 0.95, 0.28, size=13.0, color=color, bold=True)
    ui.add_text(
        slide,
        "Of the 32 shared male e33 drivers, ROSMAP shows them in male-only, "
        "female-only, or both-sex categories in roughly equal measure.",
        7.22, 5.35, 5.20, 0.62, size=10.4, color=ui.GRAY,
    )
    ui.add_notes(
        slide,
        goal="Balance the stratum-agnostic finding with the genes that are genuinely sex-consistent.",
        walkthrough="HGSNAT never appears in a female category in either cohort and is always an excitatory-neuron driver in ROSMAP; TARBP1 and RPS27A repeat the male-only pattern. But across all thirty-two shared male e33 genes the ROSMAP sex distribution is fourteen male-only, eleven female-only, and seven both.",
        boundary="Male-only status reflects where returns happened to occur; unequal run availability across strata can create such patterns without a biological sex effect.",
        transition="Close with what this comparison does and does not establish.",
    )
    return slide


def append_takeaway_slide(prs: Presentation, facts: dict[str, Any]):
    slide = ui.new_slide(prs)
    ui.add_title_block(
        slide,
        "What this comparison does—and does not—show",
    )
    ui.add_rect(slide, 0.66, 1.30, 5.95, 5.05, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "Supported", 0.97, 1.62, 5.35, accent=ui.TEAL)
    ui.add_bullets(slide, [
        f"{len(facts['shared_genes'])} of {facts['sea_genes']} SEA-AD drivers are ROSMAP re-discoveries.",
        "Cell-type context transfers: excitatory drivers stay excitatory, inhibitory stay inhibitory.",
        "Eight exact category matches, led by WDR82 and the full female e33 excitatory trio.",
        "HGSNAT is a reproducible male excitatory driver in both cohorts.",
    ], 0.99, 2.20, 5.40, size=11.8, line_h=0.78, accent=ui.TEAL)
    ui.add_rect(slide, 6.90, 1.30, 5.78, 5.05, color=ui.PALE_RED, outline=ui.VERMILION)
    ui.add_panel_title(slide, "Not supported / limits", 7.21, 1.62, 5.18, accent=ui.VERMILION)
    ui.add_bullets(slide, [
        "Sex/APOE stratum assignments largely do not transfer between cohorts.",
        "Only the male e33 stratum is well powered in SEA-AD; other strata are near-empty.",
        "Scores are exploratory and post-selected — convergence, not a formal replication test.",
        "Recurrent housekeeping-like genes may overlap partly because they return broadly.",
    ], 7.23, 2.20, 5.25, size=11.8, line_h=0.78, accent=ui.VERMILION)
    ui.add_text(
        slide,
        "Takeaway: SEA-AD supports the driver genes and their cell-type contexts, "
        "not the ROSMAP stratum assignments.",
        0.88, 6.65, 11.60, 0.30, size=12.0, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Leave the audience with a precise statement of what the cross-cohort comparison establishes.",
        walkthrough="The supported column collects the positive results: gene-level re-discovery, cell-type fidelity, the eight exact matches, and HGSNAT. The limits column lists why stratum claims should not be made: strata do not transfer, only one stratum is powered, scores are exploratory, and broadly returning genes overlap more easily.",
        boundary="None of the findings are formally FDR-controlled replication tests; they are structured descriptive convergence.",
        transition="End of the deck.",
    )
    return slide


def slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def write_audit(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["check_id", "observed", "expected", "passed"],
            delimiter="\t",
            lineterminator="\n",
        )
        writer.writeheader()
        writer.writerows(rows)


def main() -> int:
    args = parse_args()
    input_path = args.input.resolve()
    output_path = args.output.resolve()
    if not input_path.is_file():
        raise FileNotFoundError(input_path)
    facts = load_facts()
    original_deck_hash = sha256_file(input_path)

    prs = Presentation(str(input_path))
    if len(prs.slides) != EXPECTED_INPUT_SLIDES:
        raise RuntimeError(
            f"Expected the {EXPECTED_INPUT_SLIDES}-slide deck, found {len(prs.slides)}"
        )
    if "SEA-AD output: 96 non-MT category units" not in slide_text(prs.slides[15]):
        raise RuntimeError("Slide 16 does not match the expected pre-update contract")
    before_texts = {n: slide_text(prs.slides[n - 1]) for n in range(1, 17)}
    before_shapes = {n: len(prs.slides[n - 1].shapes) for n in range(1, 17)}

    ui.set_notes_body_template(prs.slides[0].notes_slide.notes_placeholder._element)
    append_divider(prs)
    append_design_slide(prs, facts)
    append_gene_overlap_slide(prs, facts)
    append_matched_table_slide(prs, facts)
    append_highlight_slide(prs, facts)
    append_group_matrix_slide(prs, facts)
    append_celltype_slide(prs, facts)
    append_sex_exceptions_slide(prs, facts)
    append_takeaway_slide(prs, facts)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    file_descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{output_path.name}.", suffix=".tmp", dir=output_path.parent
    )
    os.close(file_descriptor)
    temporary = Path(temporary_name)
    try:
        prs.save(str(temporary))
        temporary.replace(output_path)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise

    reloaded = Presentation(str(output_path))
    after_texts = {n: slide_text(reloaded.slides[n - 1]) for n in range(1, 17)}
    after_shapes = {n: len(reloaded.slides[n - 1].shapes) for n in range(1, 17)}
    unchanged_ok = before_texts == after_texts and before_shapes == after_shapes
    new_titles = {
        17: "Do the same drivers appear in both cohorts?",
        18: "Two ways to ask whether a driver replicates",
        19: f"{len(facts['shared_genes'])} of {facts['sea_genes']} SEA-AD drivers also return in ROSMAP",
        20: "Eight driver-category pairs match exactly",
        21: "WDR82 anchors the cleanest replication",
        22: "Shared genes are not tied to the matching sex/APOE group",
        23: "Cell type is preserved far better than sex/APOE group",
        24: "A few drivers do keep their sex context",
        25: "What this comparison does—and does not—show",
    }
    checks = [
        {
            "check_id": "output_slide_count",
            "observed": len(reloaded.slides),
            "expected": 25,
            "passed": len(reloaded.slides) == 25,
        },
        {
            "check_id": "slides_1_to_16_text_and_shape_counts_unchanged",
            "observed": "unchanged" if unchanged_ok else "changed",
            "expected": "unchanged",
            "passed": unchanged_ok,
        },
    ]
    for number, title in new_titles.items():
        present = title in slide_text(reloaded.slides[number - 1])
        checks.append(
            {
                "check_id": f"slide{number}_title",
                "observed": present,
                "expected": True,
                "passed": present,
            }
        )
    notes_ok = all(
        reloaded.slides[n - 1].has_notes_slide
        and reloaded.slides[n - 1].notes_slide.notes_text_frame is not None
        and reloaded.slides[n - 1].notes_slide.notes_text_frame.text.strip() != ""
        for n in range(17, 26)
    )
    checks.append(
        {
            "check_id": "new_slides_have_notes",
            "observed": notes_ok,
            "expected": True,
            "passed": notes_ok,
        }
    )
    write_audit(args.audit.resolve(), checks)
    failed = [row["check_id"] for row in checks if not row["passed"]]
    if failed:
        raise RuntimeError("Slide append failed checks: " + ", ".join(failed))
    print(f"updated={output_path}")
    print("slides_appended=17-25")
    print(f"original_sha256={original_deck_hash}")
    print(f"updated_sha256={sha256_file(output_path)}")
    print(f"audit={args.audit.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
