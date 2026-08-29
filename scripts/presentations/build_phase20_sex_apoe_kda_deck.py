#!/usr/bin/env python3
"""Build a teaching deck for fine-cell and broad-cell Phase 20 KDA.

The deck is generated from validated project-local production tables and the
four canonical Phase 20 result figures.  It explains each change of counting
unit and each filter before presenting recurrence and top-five results.
"""

from __future__ import annotations

import argparse
import csv
import gzip
import hashlib
import tempfile
import zipfile
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
DEFAULT_OUT = REPO / "docs" / "presentations" / "phase20_sex_apoe_kda_fine_broad.pptx"
AUDIT_DIR = REPO / "results" / "presentations" / "phase20_sex_apoe_kda_fine_broad"

FINE = REPO / "results" / "minerva_production" / "20_sex_apoe_kda"
BROAD = REPO / "results" / "minerva_production" / "20_sex_apoe_kda_broad"
DOCS = REPO / "docs" / "phase_20_sex_apoe_kda"

FIG = {
    "fine_recurrence": REPO
    / "results/figures/analysis/phase_20_sex_apoe_kda/driver_recurrence/phase20_driver_recurrence.png",
    "fine_top5": REPO
    / "results/figures/analysis/phase_20_sex_apoe_kda/top5_candidates/phase20_top5_candidates.png",
    "broad_recurrence": REPO
    / "results/figures/analysis/phase_20_sex_apoe_kda_broad/driver_recurrence/phase20_broad_driver_recurrence.png",
    "broad_top5": REPO
    / "results/figures/analysis/phase_20_sex_apoe_kda_broad/top5_candidates/phase20_broad_top5_candidates.png",
}

INPUTS = {
    "fine_status": FINE / "phase20_status.tsv",
    "fine_funnel": FINE / "phase20_filter_funnel.tsv",
    "fine_manifest": FINE / "phase20_category_manifest.tsv",
    "fine_run_manifest": FINE / "00_inputs" / "phase20_source_run_manifest.tsv",
    "fine_candidate_tests": FINE / "00_inputs" / "phase20_source_candidate_tests.tsv.gz",
    "fine_aggregates": FINE / "phase20_driver_aggregates.tsv.gz",
    "fine_candidates": FINE / "phase20_relaxed_candidates.tsv",
    "fine_strict_candidates": FINE / "phase20_strict_non_mt_reference_candidates.tsv",
    "fine_exploratory": FINE / "phase20_exploratory_leads.tsv",
    "fine_top5_data": FINE / "phase20_top5_summary.tsv",
    "broad_status": BROAD / "phase20_broad_status.tsv",
    "broad_funnel": BROAD / "phase20_broad_filter_funnel.tsv",
    "broad_directions": BROAD / "phase20_broad_direction_manifest.tsv",
    "broad_candidates": BROAD / "phase20_broad_non_mt_candidates.tsv",
    "broad_all_tests": BROAD / "phase20_broad_all_candidate_tests.tsv.gz",
    "broad_top5_data": BROAD / "phase20_broad_top5_summary.tsv",
    "broad_funnel_doc": DOCS / "phase20_broad_funnel_explained.md",
    "fine_recurrence": FIG["fine_recurrence"],
    "fine_top5": FIG["fine_top5"],
    "broad_recurrence": FIG["broad_recurrence"],
    "broad_top5": FIG["broad_top5"],
}

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 35, 61)
NAVY_2 = RGBColor(30, 59, 91)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
GOLD = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
TEAL_TEXT = RGBColor(0, 104, 77)
GOLD_TEXT = RGBColor(137, 86, 0)
VERMILION_TEXT = RGBColor(158, 58, 0)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(224, 231, 239)
MID = RGBColor(108, 121, 136)
DARK = RGBColor(31, 39, 48)
GRAY = RGBColor(78, 90, 104)
PALE_BLUE = RGBColor(228, 240, 248)
PALE_SKY = RGBColor(235, 247, 252)
PALE_GREEN = RGBColor(228, 244, 238)
PALE_GOLD = RGBColor(253, 242, 217)
PALE_RED = RGBColor(253, 235, 228)
PALE_PURPLE = RGBColor(241, 235, 247)
PALE_GRAY = RGBColor(241, 244, 247)

FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def readable_accent(color: RGBColor) -> RGBColor:
    if color == TEAL:
        return TEAL_TEXT
    if color == GOLD:
        return GOLD_TEXT
    if color == VERMILION:
        return VERMILION_TEXT
    return color


def iter_tsv(path: Path):
    opener = gzip.open if path.suffix == ".gz" else open
    with opener(path, "rt", encoding="utf-8", newline="") as handle:
        yield from csv.DictReader(handle, delimiter="\t")


def read_tsv(path: Path) -> list[dict[str, str]]:
    return list(iter_tsv(path))


def write_tsv(rows: Sequence[dict[str, Any]], path: Path, fields: Sequence[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=fields,
            delimiter="\t",
            lineterminator="\n",
            extrasaction="ignore",
        )
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def truth(value: Any) -> bool:
    return str(value).strip().upper() in {"TRUE", "T", "1", "YES"}


def count_text(value: int) -> str:
    return f"{value:,}"


def optional_float(value: Any) -> float | None:
    text = str(value).strip()
    if text in {"", "NA", "None"}:
        return None
    return float(text)


def fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def line(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_run(
    run,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    font: str = FONT_BODY,
) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15,
    color: RGBColor = DARK,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.03,
    font: str = FONT_BODY,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic, font=font)
    return box


def add_rich_text(
    slide,
    spans: Sequence[tuple[str, dict[str, Any]]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    frame.margin_top = frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    for text, style in spans:
        run = paragraph.add_run()
        run.text = text
        set_run(
            run,
            size=style.get("size", 15),
            color=style.get("color", DARK),
            bold=style.get("bold", False),
            italic=style.get("italic", False),
            font=style.get("font", FONT_BODY),
        )
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    color: RGBColor = WHITE,
    outline: RGBColor | None = LIGHT,
    radius: bool = True,
    transparency: int = 0,
    width: float = 1.0,
):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color, transparency)
    if outline is None:
        shape.line.fill.background()
    else:
        line(shape, outline, width)
    return shape


def add_circle(
    slide,
    x: float,
    y: float,
    diameter: float,
    color: RGBColor,
    *,
    outline: RGBColor | None = None,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(diameter),
        Inches(diameter),
    )
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        line(shape, outline)
    return shape


def add_bullets(
    slide,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    *,
    size: float = 13.2,
    color: RGBColor = DARK,
    accent: RGBColor = BLUE,
    line_h: float = 0.55,
) -> None:
    current_y = y
    for item in items:
        add_circle(slide, x, current_y + 0.14, 0.085, accent)
        add_text(
            slide,
            item,
            x + 0.19,
            current_y,
            w - 0.19,
            line_h,
            size=size,
            color=color,
        )
        current_y += line_h


def add_metric(
    slide,
    value: str,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: RGBColor = BLUE,
    bg: RGBColor = WHITE,
    note: str | None = None,
) -> None:
    height = 1.12 if note is None else 1.40
    add_rect(slide, x, y, w, height, color=bg, outline=LIGHT)
    value_size = 25 if len(value) <= 9 else 18
    add_text(
        slide,
        value,
        x + 0.16,
        y + 0.10,
        w - 0.32,
        0.47,
        size=value_size,
        color=readable_accent(accent),
        bold=True,
        font=FONT_HEAD,
    )
    add_text(
        slide,
        label,
        x + 0.16,
        y + 0.60,
        w - 0.32,
        0.43,
        size=8.8 if len(label) > 27 else 9.6,
        color=GRAY,
        bold=True,
    )
    if note:
        add_text(
            slide,
            note,
            x + 0.16,
            y + 1.08,
            w - 0.32,
            0.20,
            size=7.8,
            color=MID,
        )


def add_panel_title(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: RGBColor,
) -> None:
    add_rect(slide, x, y + 0.03, 0.07, 0.31, color=accent, outline=None, radius=False)
    add_text(
        slide,
        text,
        x + 0.16,
        y,
        w - 0.16,
        0.36,
        size=15.2,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
    )


def add_takeaway(
    slide,
    text: str,
    *,
    accent: RGBColor,
    y: float = 6.73,
) -> None:
    add_rect(slide, 0.55, y, 12.22, 0.43, color=NAVY, outline=None, radius=False)
    add_rect(slide, 0.55, y, 0.10, 0.43, color=accent, outline=None, radius=False)
    add_text(
        slide,
        text,
        0.79,
        y + 0.07,
        11.75,
        0.28,
        size=10.7,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    widths: Sequence[float],
    *,
    row_h: float = 0.48,
    header_h: float = 0.48,
    accent: RGBColor = BLUE,
    font_size: float = 9.6,
    highlight_last: bool = False,
) -> None:
    current_x = x
    for header, width_value in zip(headers, widths, strict=True):
        add_rect(
            slide,
            current_x,
            y,
            width_value,
            header_h,
            color=NAVY,
            outline=WHITE,
            radius=False,
            width=0.6,
        )
        add_text(
            slide,
            header,
            current_x + 0.06,
            y + 0.06,
            width_value - 0.12,
            header_h - 0.10,
            size=font_size - 0.4,
            color=WHITE,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        current_x += width_value
    for row_index, row in enumerate(rows):
        current_x = x
        bg = PALE_GREEN if highlight_last and row_index == len(rows) - 1 else (
            WHITE if row_index % 2 == 0 else PALE_GRAY
        )
        for value, width_value in zip(row, widths, strict=True):
            add_rect(
                slide,
                current_x,
                y + header_h + row_index * row_h,
                width_value,
                row_h,
                color=bg,
                outline=LIGHT,
                radius=False,
                width=0.6,
            )
            add_text(
                slide,
                value,
                current_x + 0.06,
                y + header_h + row_index * row_h + 0.05,
                width_value - 0.12,
                row_h - 0.08,
                size=font_size,
                color=DARK,
                bold=highlight_last and row_index == len(rows) - 1,
                valign=MSO_ANCHOR.MIDDLE,
            )
            current_x += width_value


def add_flow_row(
    slide,
    nodes: Sequence[tuple[str, str]],
    y: float,
    *,
    accent: RGBColor,
    start_x: float = 0.65,
    total_w: float = 12.0,
    node_h: float = 1.16,
) -> None:
    gap = 0.28
    node_w = (total_w - gap * (len(nodes) - 1)) / len(nodes)
    for index, (value, label) in enumerate(nodes):
        x = start_x + index * (node_w + gap)
        add_rect(slide, x, y, node_w, node_h, color=WHITE, outline=LIGHT)
        add_text(
            slide,
            value,
            x + 0.10,
            y + 0.13,
            node_w - 0.20,
            0.40,
            size=20 if len(value) < 8 else 16.5,
            color=readable_accent(accent),
            bold=True,
            align=PP_ALIGN.CENTER,
            font=FONT_HEAD,
        )
        add_text(
            slide,
            label,
            x + 0.10,
            y + 0.58,
            node_w - 0.20,
            0.42,
            size=9.4,
            color=GRAY,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if index < len(nodes) - 1:
            arrow_x = x + node_w + 0.055
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(arrow_x),
                Inches(y + 0.43),
                Inches(0.17),
                Inches(0.30),
            )
            fill(arrow, accent)
            arrow.line.fill.background()


def set_alt_text(shape, title: str, description: str) -> None:
    c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("name", title)
    c_nv_pr.set("descr", description)


def add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    alt: str,
):
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    picture_w, picture_h = image_w * scale, image_h * scale
    picture_x = x + (w - picture_w) / 2
    picture_y = y + (h - picture_h) / 2
    picture = slide.shapes.add_picture(
        str(path),
        Inches(picture_x),
        Inches(picture_y),
        Inches(picture_w),
        Inches(picture_h),
    )
    set_alt_text(picture, alt, alt)
    return picture


def add_header(
    slide,
    kicker: str,
    title: str,
    slide_no: int,
    *,
    accent: RGBColor,
    subtitle: str | None = None,
) -> None:
    title_size = 21.0 if len(title) >= 64 else 24.5
    add_text(slide, kicker.upper(), 0.58, 0.20, 4.4, 0.24, size=9.2, color=readable_accent(accent), bold=True)
    add_text(
        slide,
        title,
        0.58,
        0.45,
        11.65,
        0.52,
        size=title_size,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
    )
    if subtitle:
        add_text(slide, subtitle, 0.60, 0.97, 11.55, 0.30, size=11.2, color=GRAY)


def add_source(slide, source: str, slide_no: int) -> None:
    add_text(slide, f"Source: {source}", 0.57, 7.24, 11.78, 0.16, size=7.0, color=GRAY)


def add_notes(
    slide,
    *,
    goal: str,
    walkthrough: str,
    boundary: str,
    transition: str,
) -> None:
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        raise RuntimeError("Notes placeholder is unavailable")
    frame.text = (
        f"Teaching goal: {goal}\n\n"
        f"Walk through: {walkthrough}\n\n"
        f"Scientific boundary: {boundary}\n\n"
        f"Transition: {transition}"
    )


def new_slide(prs: Presentation, *, bg: RGBColor = OFF_WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def content_slide(
    prs: Presentation,
    *,
    title: str,
    kicker: str,
    slide_no: int,
    accent: RGBColor,
    subtitle: str | None = None,
):
    slide = new_slide(prs)
    add_header(slide, kicker, title, slide_no, accent=accent, subtitle=subtitle)
    return slide


def finalize_slide(
    slide,
    *,
    slide_no: int,
    source: str,
    goal: str,
    walkthrough: str,
    boundary: str,
    transition: str,
) -> None:
    add_source(slide, source, slide_no)
    add_notes(
        slide,
        goal=goal,
        walkthrough=walkthrough,
        boundary=boundary,
        transition=transition,
    )


def section_slide(
    prs: Presentation,
    *,
    slide_no: int,
    section_no: int,
    title: str,
    subtitle: str,
    topics: Sequence[str],
    accent: RGBColor,
    source: str,
    notes: dict[str, str],
):
    slide = new_slide(prs, bg=NAVY)
    add_text(
        slide,
        f"PART {section_no}",
        0.78,
        0.67,
        2.4,
        0.28,
        size=10.5,
        color=accent,
        bold=True,
    )
    add_rect(slide, 0.78, 1.26, 0.10, 2.38, color=accent, outline=None, radius=False)
    add_text(
        slide,
        title,
        1.18,
        1.36,
        8.65,
        1.34,
        size=32,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, subtitle, 1.20, 3.00, 8.30, 0.84, size=15.0, color=RGBColor(204, 219, 234))
    add_rect(slide, 9.78, 1.28, 2.70, 4.70, color=NAVY_2, outline=None)
    add_text(slide, "IN THIS PART", 10.08, 1.67, 2.10, 0.24, size=9.4, color=accent, bold=True)
    current_y = 2.20
    for index, topic in enumerate(topics, start=1):
        add_circle(slide, 10.06, current_y + 0.03, 0.34, accent)
        add_text(
            slide,
            str(index),
            10.06,
            current_y + 0.08,
            0.34,
            0.17,
            size=8.2,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, topic, 10.53, current_y, 1.58, 0.52, size=11.1, color=WHITE, bold=True)
        current_y += 0.82
    add_text(slide, f"Source: {source}", 0.78, 7.19, 11.2, 0.14, size=6.5, color=RGBColor(157, 176, 195))
    add_notes(slide, **notes)
    return slide


def load_facts() -> dict[str, Any]:
    missing = [str(path) for path in INPUTS.values() if not path.is_file()]
    if missing:
        raise FileNotFoundError("Missing presentation input(s): " + ", ".join(missing))

    fine_status = read_tsv(INPUTS["fine_status"])[0]
    fine_funnel = read_tsv(INPUTS["fine_funnel"])
    fine_manifest = read_tsv(INPUTS["fine_manifest"])
    fine_run_manifest = read_tsv(INPUTS["fine_run_manifest"])
    fine_candidates = read_tsv(INPUTS["fine_candidates"])
    fine_strict_candidates = read_tsv(INPUTS["fine_strict_candidates"])
    fine_exploratory = read_tsv(INPUTS["fine_exploratory"])
    fine_top5 = read_tsv(INPUTS["fine_top5_data"])
    broad_status = read_tsv(INPUTS["broad_status"])[0]
    broad_funnel = read_tsv(INPUTS["broad_funnel"])
    broad_directions = read_tsv(INPUTS["broad_directions"])
    broad_candidates = read_tsv(INPUTS["broad_candidates"])
    broad_all_tests = read_tsv(INPUTS["broad_all_tests"])
    broad_top5 = read_tsv(INPUTS["broad_top5_data"])

    fine_overall = {
        row["stage_id"]: int(row["retained_count"])
        for row in fine_funnel
        if row["scope"] == "overall"
    }
    broad_primary = {
        row["stage"]: int(row["passing_units"])
        for row in broad_funnel
        if row["funnel_scope"] == "primary_candidate_selection"
        and row["query_tier"] == "relaxed"
    }
    completed_broad = [
        row for row in broad_directions if row["terminal_status"].startswith("completed")
    ]
    ranked_fine_top5 = [row for row in fine_top5 if row["list_status"] == "ranked_candidates"]

    included_fine_runs = [row for row in fine_run_manifest if truth(row["phase20_included"])]
    fine_type_count = len({row["fine_cell_type"] for row in fine_run_manifest})
    fine_group_count = len({row["signature_group"] for row in fine_run_manifest})
    fine_direction_count = len({row["signature_direction"] for row in fine_run_manifest})
    fine_network_count = len({row["broad_network"] for row in fine_run_manifest})
    included_fine_types = {row["fine_cell_type"] for row in included_fine_runs}
    effective_query_sizes = [int(row["effective_query_genes"]) for row in fine_run_manifest]
    fine_run_counts = {
        "planned_slots": len(fine_run_manifest),
        "comparisons": len(
            {(row["fine_cell_type"], row["signature_group"]) for row in fine_run_manifest}
        ),
        "source_not_validated": sum(
            row["eligibility_status"] == "source_contrast_not_validated"
            for row in fine_run_manifest
        ),
        "effective_query_zero": sum(
            int(row["effective_query_genes"]) == 0
            and row["eligibility_status"] != "source_contrast_not_validated"
            for row in fine_run_manifest
        ),
        "effective_query_one_two": sum(
            1 <= int(row["effective_query_genes"]) <= 2 for row in fine_run_manifest
        ),
        "effective_query_three_nine": sum(
            3 <= int(row["effective_query_genes"]) <= 9 for row in fine_run_manifest
        ),
        "effective_query_ten_plus": sum(
            int(row["effective_query_genes"]) >= 10 for row in fine_run_manifest
        ),
        "included_runs": len(included_fine_runs),
        "included_fine_types": len(included_fine_types),
        "completed_significant": sum(
            row["terminal_status"] == "completed_significant" for row in included_fine_runs
        ),
        "completed_empty": sum(
            row["terminal_status"] == "completed_no_significant"
            for row in included_fine_runs
        ),
    }

    fine_evidence: Counter[str] = Counter()
    fine_source_cases: set[str] = set()
    relaxed_support_unit_keys: set[tuple[str, str, str]] = set()
    strict_support_unit_keys: set[tuple[str, str, str]] = set()
    for row in iter_tsv(INPUTS["fine_candidate_tests"]):
        fine_evidence["opportunities"] += 1
        fine_source_cases.add(row["case_id"])
        is_non_mt = row["case_id"] == "non_mt_driver"
        if is_non_mt == truth(row["is_core_mito"]):
            fine_evidence["case_annotation_mismatches"] += 1
        fine_evidence["non_mt_rows" if is_non_mt else "mt_rows"] += 1
        if truth(row["usable_test"]):
            fine_evidence["usable"] += 1
            if truth(row["explicit_family_member"]):
                fine_evidence["explicit"] += 1
            else:
                fine_evidence["implicit"] += 1
        else:
            fine_evidence["absent"] += 1
        original_q = optional_float(row["original_run_q"])
        if original_q is not None and original_q <= 0.05:
            fine_evidence["stock_q05"] += 1
        if not is_non_mt or not truth(row["explicit_family_member"]):
            continue
        fine_evidence["non_mt_explicit"] += 1
        overlap = optional_float(row["other_query_overlap"])
        fold = optional_float(row["final_fold_enrichment"])
        final_q = optional_float(row["final_run_q"])
        if overlap is not None and overlap > 0:
            fine_evidence["positive_overlap"] += 1
        if overlap is not None and overlap >= 2:
            fine_evidence["overlap_ge2"] += 1
        if overlap is not None and overlap >= 2 and fold is not None and fold > 1:
            fine_evidence["overlap_ge2_fold_gt1"] += 1
            key = (row["signature_group"], row["broad_network"], row["current_symbol"])
            if final_q is not None and final_q <= 0.10:
                fine_evidence["relaxed_support_events"] += 1
                relaxed_support_unit_keys.add(key)
            if final_q is not None and final_q <= 0.05:
                fine_evidence["strict_support_events"] += 1
                strict_support_unit_keys.add(key)

    coverage_symbols: set[str] = set()
    support_symbols: set[str] = set()
    support_categories: set[tuple[str, str]] = set()
    aggregate_symbols: set[str] = set()
    support_units = 0
    aggregate_rows = 0
    zero_usable_units = 0
    partial_below_relaxed = 0
    strict_coverage_units = 0
    strict_coverage_support_units = 0
    for row in iter_tsv(INPUTS["fine_aggregates"]):
        aggregate_rows += 1
        aggregate_symbols.add(row["current_symbol"])
        usable_runs = int(row["usable_run_count"])
        if usable_runs == 0:
            zero_usable_units += 1
        elif not truth(row["relaxed_coverage_pass"]):
            partial_below_relaxed += 1
        if truth(row["relaxed_coverage_pass"]):
            coverage_symbols.add(row["current_symbol"])
        if truth(row["relaxed_coverage_pass"]) and truth(row["relaxed_support_pass"]):
            support_units += 1
            support_symbols.add(row["current_symbol"])
            support_categories.add((row["signature_group"], row["broad_network"]))
        if truth(row["strict_coverage_pass"]):
            strict_coverage_units += 1
            if truth(row["conservative_support_pass"]):
                strict_coverage_support_units += 1

    fine_candidate_categories = {
        (row["signature_group"], row["broad_network"]) for row in fine_candidates
    }
    fine_candidate_groups = Counter(row["signature_group"] for row in fine_candidates)
    fine_candidate_networks = Counter(row["broad_network"] for row in fine_candidates)
    recurrence_counts = Counter(row["current_symbol"] for row in fine_candidates)
    recurrence_details = {
        gene: {
            "categories": recurrence_counts[gene],
            "groups": len(
                {row["signature_group"] for row in fine_candidates if row["current_symbol"] == gene}
            ),
            "networks": len(
                {row["broad_network"] for row in fine_candidates if row["current_symbol"] == gene}
            ),
            "strict": sum(
                truth(row["strict_non_mt_reference"])
                for row in fine_candidates
                if row["current_symbol"] == gene
            ),
        }
        for gene in recurrence_counts
    }
    fine_analyzable_categories = sum(
        int(row["included_run_count"]) > 0 for row in fine_manifest
    )
    fine_empty_categories = len(fine_manifest) - fine_analyzable_categories
    fine_analyzable_no_candidate = sum(
        int(row["included_run_count"]) > 0 and int(row["relaxed_candidate_count"]) == 0
        for row in fine_manifest
    )

    fine_numbers: dict[str, Any] = {
        **fine_run_counts,
        "fine_type_count": fine_type_count,
        "group_count": fine_group_count,
        "direction_count": fine_direction_count,
        "network_count": fine_network_count,
        "minimum_query": min(int(row["effective_query_genes"]) for row in included_fine_runs),
        "opportunities": fine_evidence["opportunities"],
        "explicit_tests": fine_evidence["explicit"],
        "implicit_rows": fine_evidence["implicit"],
        "absent_rows": fine_evidence["absent"],
        "stock_q05_rows": fine_evidence["stock_q05"],
        "mt_rows": fine_evidence["mt_rows"],
        "non_mt_rows": fine_evidence["non_mt_rows"],
        "non_mt_explicit": fine_evidence["non_mt_explicit"],
        "positive_overlap": fine_evidence["positive_overlap"],
        "overlap_ge2": fine_evidence["overlap_ge2"],
        "overlap_ge2_fold_gt1": fine_evidence["overlap_ge2_fold_gt1"],
        "relaxed_support_events": fine_evidence["relaxed_support_events"],
        "strict_support_events": fine_evidence["strict_support_events"],
        "support_precoverage_units": len(relaxed_support_unit_keys),
        "strict_support_precoverage_units": len(strict_support_unit_keys),
        "aggregate_rows": aggregate_rows,
        "aggregate_genes": len(aggregate_symbols),
        "zero_usable_units": zero_usable_units,
        "partial_below_relaxed": partial_below_relaxed,
        "relaxed_coverage_units": fine_overall["coverage_at_least_0_50"],
        "relaxed_coverage_genes": len(coverage_symbols),
        "relaxed_support_units": support_units,
        "relaxed_support_genes": len(support_symbols),
        "relaxed_support_categories": len(support_categories),
        "strict_coverage_units": strict_coverage_units,
        "strict_coverage_support_units": strict_coverage_support_units,
        "relaxed_candidates": len(fine_candidates),
        "relaxed_genes": len({row["current_symbol"] for row in fine_candidates}),
        "relaxed_categories": len(fine_candidate_categories),
        "strict_candidates": len(fine_strict_candidates),
        "strict_genes": len({row["current_symbol"] for row in fine_strict_candidates}),
        "strict_categories": len(
            {(row["signature_group"], row["broad_network"]) for row in fine_strict_candidates}
        ),
        "exploratory_candidates": len(fine_exploratory),
        "exploratory_genes": len({row["current_symbol"] for row in fine_exploratory}),
        "exploratory_categories": len(
            {(row["signature_group"], row["broad_network"]) for row in fine_exploratory}
        ),
        "top5_rows": len(ranked_fine_top5),
        "top5_strict_rows": sum(
            truth(row["strict_non_mt_reference"]) for row in ranked_fine_top5
        ),
        "structural_categories": len(fine_manifest),
        "analyzable_categories": fine_analyzable_categories,
        "empty_categories": fine_empty_categories,
        "analyzable_no_candidate": fine_analyzable_no_candidate,
        "candidate_groups": fine_candidate_groups,
        "candidate_networks": fine_candidate_networks,
        "recurrence": recurrence_counts,
        "recurrence_details": recurrence_details,
    }

    broad_stock_rows = [row for row in broad_all_tests if truth(row["stock_fkda_q05_return"])]

    assertions = [
        (fine_status["validation_status"] == "validated_complete", "fine status"),
        (int(fine_status["failed_checks"]) == 0, "fine checks"),
        (fine_status["analysis_id"] == "phase20_sex_apoe_kda_v2", "fine analysis id"),
        (
            fine_status["task_mode"] == "validated_phase12_min3_reaggregation",
            "fine min3 task mode",
        ),
        (fine_status["source_validation_status"] == "validated_complete", "fine source status"),
        (int(fine_status["included_runs"]) == 295, "fine included runs"),
        (int(fine_status["aggregate_rows"]) == 259548, "fine aggregate rows"),
        (int(fine_status["strict_candidates"]) == 58, "fine status strict candidates"),
        (int(fine_status["relaxed_candidates"]) == 74, "fine relaxed candidates"),
        (int(fine_status["exploratory_leads"]) == 15, "fine status exploratory leads"),
        (int(fine_status["relaxed_categories"]) == 16, "fine status relaxed categories"),
        (fine_run_counts["planned_slots"] == 648, "fine planned slots"),
        (fine_run_counts["included_runs"] == 295, "fine manifest included runs"),
        (fine_run_counts["source_not_validated"] == 6, "fine invalid source slots"),
        (fine_run_counts["effective_query_zero"] == 246, "fine effective query zero slots"),
        (fine_run_counts["effective_query_one_two"] == 101, "fine query 1-2 slots"),
        (fine_run_counts["effective_query_three_nine"] == 134, "fine query 3-9 runs"),
        (fine_run_counts["effective_query_ten_plus"] == 161, "fine query 10+ runs"),
        (fine_run_counts["completed_significant"] == 221, "fine significant-return runs"),
        (fine_run_counts["completed_empty"] == 74, "fine completed-empty runs"),
        (min(effective_query_sizes) == 0, "fine structural query floor"),
        (fine_numbers["minimum_query"] == 3, "fine included query floor"),
        (
            all(
                truth(row["phase20_included"])
                == (
                    row["eligibility_status"] == "eligible"
                    and row["terminal_status"].startswith("completed")
                    and int(row["effective_query_genes"]) >= 3
                )
                for row in fine_run_manifest
            ),
            "fine inclusion rule",
        ),
        (fine_evidence["opportunities"] == 2623910, "fine opportunity rows"),
        (fine_evidence["explicit"] == 108537, "fine explicit rows"),
        (fine_evidence["implicit"] == 2202083, "fine implicit rows"),
        (fine_evidence["absent"] == 313290, "fine absent rows"),
        (fine_evidence["non_mt_rows"] == 2411256, "fine non-MT rows"),
        (fine_evidence["mt_rows"] == 212654, "fine MT rows"),
        (fine_source_cases == {"mt_driver", "non_mt_driver"}, "fine source case taxonomy"),
        (fine_evidence["case_annotation_mismatches"] == 0, "fine source case annotation"),
        (fine_evidence["stock_q05"] == 2494, "fine stock return rows"),
        (fine_evidence["relaxed_support_events"] == 864, "fine relaxed support events"),
        (fine_evidence["strict_support_events"] == 593, "fine strict support events"),
        (fine_overall == {
            "input_non_mt_units": 259548,
            "coverage_at_least_0_50": 233368,
            "relaxed_run_support": 500,
            "category_q_at_most_0_10": 74,
        }, "fine overall funnel"),
        (len(fine_candidates) == 74, "fine candidate rows"),
        (len({row["current_symbol"] for row in fine_candidates}) == 37, "fine genes"),
        (len(fine_candidate_categories) == 16, "fine candidate categories"),
        (len(fine_strict_candidates) == 58, "fine strict rows"),
        (len(fine_exploratory) == 15, "fine exploratory-only rows"),
        (len(ranked_fine_top5) == 48, "fine displayed rows"),
        (sum(truth(row["strict_non_mt_reference"]) for row in ranked_fine_top5) == 41, "fine displayed strict"),
        (len(fine_manifest) == 42, "fine structural categories"),
        (fine_analyzable_categories == 38, "fine analyzable categories"),
        (fine_empty_categories == 4, "fine empty categories"),
        (len(included_fine_types) == 41, "fine included cell types"),
        (len(coverage_symbols) == 11232, "fine coverage-pass genes"),
        (support_units == 500, "fine coverage and support units"),
        (len(support_symbols) == 265, "fine support genes"),
        (len(support_categories) == 30, "fine support categories"),
        (aggregate_rows == 259548, "fine aggregate stream rows"),
        (
            all(
                row["case_id"] == "non_mt_driver" and not truth(row["is_core_mito"])
                for row in fine_candidates
                + fine_strict_candidates
                + fine_exploratory
                + ranked_fine_top5
            ),
            "fine non-MT output taxonomy",
        ),
        (broad_status["validation_status"] == "validated_complete", "broad status"),
        (int(broad_status["failed_checks"]) == 0, "broad checks"),
        (broad_status["aggregation_method"] == "none", "broad no aggregation"),
        (int(broad_status["primary_completed_runs"]) == 3, "broad completed runs"),
        (broad_primary == {
            "all_explicit_candidates": 377,
            "non_core_mt_candidates": 362,
            "query_overlap_ge2": 15,
            "fold_enrichment_gt1": 15,
            "non_mt_run_q_le0_10": 12,
        }, "broad primary funnel"),
        (len(completed_broad) == 3, "broad completed run rows"),
        (len(broad_candidates) == 12, "broad candidate rows"),
        (len({row["current_symbol"] for row in broad_candidates}) == 12, "broad genes"),
        (sum(truth(row["strict_direct_reference"]) for row in broad_candidates) == 9, "broad strict rows"),
        (len(broad_top5) == 7, "broad displayed rows"),
        (len(broad_stock_rows) == 11, "broad stock rows"),
        ({row["current_symbol"] for row in broad_stock_rows if truth(row["is_core_mito"])} == {"GLDC", "ME3"}, "broad stock core-MT rows"),
    ]
    failures = [label for passed, label in assertions if not passed]
    if failures:
        raise RuntimeError("Validated source facts drifted: " + ", ".join(failures))

    broad_run_map = {row["kda_run_id"]: row for row in completed_broad}
    return {
        "fine": fine_numbers,
        "fine_status": fine_status,
        "fine_manifest": fine_manifest,
        "fine_candidates": fine_candidates,
        "fine_top5": ranked_fine_top5,
        "fine_recurrence": recurrence_counts,
        "broad_status": broad_status,
        "broad_directions": broad_directions,
        "completed_broad": completed_broad,
        "broad_run_map": broad_run_map,
        "broad_candidates": broad_candidates,
        "broad_top5": broad_top5,
    }


def build_deck(output_path: Path) -> tuple[Path, list[dict[str, Any]]]:
    facts = load_facts()
    fine = facts["fine"]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    properties = prs.core_properties
    properties.title = "Fine-cell and broad-cell sex/APOE key-driver analysis"
    properties.subject = "Phase 20 KDA setup, filtering funnels, supported candidates, recurrence, and top-five results"
    properties.author = "Alzheimer project analysis team"
    properties.last_modified_by = "OpenAI Codex"
    properties.keywords = "Alzheimer disease, sex, APOE, key-driver analysis, KDA, ACAT, mitochondrial genes"
    properties.language = "en-US"
    properties.category = "Scientific presentation"
    properties.comments = "Generated from validated project-local Phase 20 production releases."
    generated_at = datetime.now(timezone.utc).replace(tzinfo=None)
    properties.created = generated_at
    properties.modified = generated_at
    metadata: list[dict[str, Any]] = []

    def register(slide, title: str, section: str) -> None:
        metadata.append({"slide": slide, "title": title, "section": section})

    # 1 — title
    title = "Fine-cell and broad-cell sex/APOE key-driver analysis"
    slide = new_slide(prs, bg=NAVY)
    add_text(slide, "PHASE 20", 0.78, 0.67, 2.3, 0.28, size=11, color=SKY, bold=True)
    add_text(
        slide,
        title,
        0.78,
        1.20,
        11.25,
        1.38,
        size=34,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "How mitochondrial DEG queries become non-core-MT network-driver candidates",
        0.82,
        2.72,
        10.8,
        0.48,
        size=17,
        color=RGBColor(204, 219, 234),
    )
    add_rect(slide, 0.82, 3.68, 5.72, 2.05, color=NAVY_2, outline=BLUE)
    add_text(slide, "PART I", 1.10, 3.98, 1.1, 0.23, size=9.5, color=SKY, bold=True)
    add_text(slide, "Fine-cell reaggregation", 1.10, 4.31, 4.75, 0.39, size=22, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(
        slide,
        f"{count_text(fine['included_runs'])} included runs → coverage + support + ACAT → {count_text(fine['relaxed_candidates'])} candidate units",
        1.10, 4.87, 4.83, 0.47, size=12.2, color=RGBColor(210, 224, 238),
    )
    add_rect(slide, 6.78, 3.68, 5.72, 2.05, color=NAVY_2, outline=TEAL)
    add_text(slide, "PART II", 7.06, 3.98, 1.1, 0.23, size=9.5, color=TEAL, bold=True)
    add_text(slide, "Direct broad-cell KDA", 7.06, 4.31, 4.80, 0.39, size=22, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(slide, "3 completed runs → within-run BH → 12 direct candidate rows", 7.06, 4.87, 4.83, 0.47, size=12.2, color=RGBColor(210, 224, 238))
    add_text(slide, "Validated local-production releases • 28 August 2026", 0.82, 6.55, 7.0, 0.24, size=10, color=RGBColor(157, 176, 195))
    add_text(slide, "Source: Phase 20 fine and broad validated production releases", 0.80, 7.19, 11.2, 0.14, size=6.5, color=RGBColor(157, 176, 195))
    add_notes(
        slide,
        goal="Introduce the two Phase 20 branches and signal that they use different evidence structures.",
        walkthrough="The fine-cell branch starts with many fine-cell and direction-specific KDA runs and reaggregates them within sex/APOE by broad-network categories. The broad branch starts from donor-level broad-cell DEGs and runs direct direction-specific KDA without ACAT.",
        boundary="The two candidate counts are not interchangeable because one is a gene-by-category aggregate unit and the other is a direct gene-by-run row.",
        transition="First compare the branches at a glance, then unpack every reduction step separately.",
    )
    register(slide, title, "Overview")

    # 2 — branch comparison and roadmap
    title = "Two branches answer related—but not identical—questions"
    slide = content_slide(prs, title=title, kicker="Roadmap", slide_no=2, accent=GOLD,
                          subtitle="Both start with mitochondrial DEG queries; only the fine-cell branch combines evidence across runs.")
    add_rect(slide, 0.64, 1.46, 5.86, 4.73, color=PALE_BLUE, outline=BLUE)
    add_text(slide, "FINE-CELL REAGGREGATION", 0.95, 1.78, 3.2, 0.27, size=10, color=BLUE, bold=True)
    add_text(slide, "Where does a gene recur across fine-cell evidence?", 0.95, 2.12, 4.95, 0.64, size=20, color=NAVY, bold=True, font=FONT_HEAD)
    add_metric(slide, count_text(fine["planned_slots"]), "planned directional slots", 0.95, 2.95, 1.55, accent=BLUE, bg=WHITE)
    add_metric(slide, count_text(fine["included_runs"]), "included KDA runs", 2.68, 2.95, 1.55, accent=BLUE, bg=WHITE)
    add_metric(slide, count_text(fine["relaxed_candidates"]), "gene × category units", 4.41, 2.95, 1.55, accent=BLUE, bg=WHITE)
    add_bullets(slide, [
        "Groups fine-cell runs within sex/APOE × broad network.",
        "Uses coverage, individual-run support, ACAT, then category BH.",
        f"Final {count_text(fine['relaxed_candidates'])} units represent {count_text(fine['relaxed_genes'])} distinct genes in {count_text(fine['relaxed_categories'])} categories.",
    ], 0.96, 4.38, 4.96, size=11.5, line_h=0.48, accent=BLUE)
    add_rect(slide, 6.83, 1.46, 5.86, 4.73, color=PALE_GREEN, outline=TEAL)
    add_text(slide, "DIRECT BROAD-CELL KDA", 7.14, 1.78, 3.2, 0.27, size=10, color=TEAL, bold=True)
    add_text(slide, "Which genes pass within one broad-cell direction?", 7.14, 2.12, 4.95, 0.64, size=20, color=NAVY, bold=True, font=FONT_HEAD)
    add_metric(slide, "84", "planned directional slots", 7.14, 2.95, 1.55, accent=TEAL, bg=WHITE)
    add_metric(slide, "3", "completed KDA runs", 8.87, 2.95, 1.55, accent=TEAL, bg=WHITE)
    add_metric(slide, "12", "direct candidate rows", 10.60, 2.95, 1.55, accent=TEAL, bg=WHITE)
    add_bullets(slide, [
        "Keeps broad-cell × group × direction results separate.",
        "Uses complete explicit tests and within-run non-core-MT BH.",
        "No coverage, ACAT, recurrence gate, or combined up/down q.",
    ], 7.15, 4.38, 4.96, size=11.5, line_h=0.48, accent=TEAL)
    add_takeaway(slide, "Always read a number together with its counting unit and branch-specific q value.", accent=GOLD)
    finalize_slide(
        slide, slide_no=2,
        source="phase20_status.tsv; phase20_filter_funnel.tsv; corresponding broad production tables",
        goal="Give the audience a mental map before the detailed funnels.",
        walkthrough=f"The left card describes cross-run reaggregation from {count_text(fine['included_runs'])} fine-cell runs. The right card describes three direct broad-cell runs. The starting queries are related, but the final units, multiplicity corrections, and interpretation differ.",
        boundary="A missing broad candidate is not equivalent to absence from an ACAT aggregate, and a fine candidate unit is not one independent gene-level experiment.",
        transition="Begin Part I by defining the counting units used in the fine-cell branch.",
    )
    register(slide, title, "Overview")

    # 3 — fine section divider
    title = "Fine-cell sex/APOE KDA reaggregation"
    slide = section_slide(
        prs, slide_no=3, section_no=1, title=title,
        subtitle=f"{count_text(fine['fine_type_count'])} fine cell types define the planned universe; {count_text(fine['included_runs'])} included runs from {count_text(fine['included_fine_types'])} fine types are aggregated within six groups and seven broad networks.",
        topics=["Run universe", "Filter funnel", "ACAT candidates", "Result figures"],
        accent=BLUE,
        source="phase20_status.tsv; phase20_source_run_manifest.tsv",
        notes={
            "goal": "Introduce the fine-cell branch and its sequence from setup to results.",
            "walkthrough": "The section first defines comparisons, directional slots, KDA calls, included runs, and final candidate units. It then explains DEG construction, run eligibility, complete evidence, coverage, support, ACAT, category BH, and display-only ranking.",
            "boundary": "The branch reconstructs complete evidence from the validated Phase 12 calls at effective query n ≥ 3; it does not rerun KDA or pool sex/APOE groups.",
            "transition": "Start with the four counting units that appear in the funnel.",
        },
    )
    register(slide, title, "Fine-cell")

    _append_fine_slides_4_18(prs, facts, register)

    # 19 — broad section divider
    title = "Direct broad-cell sex/APOE KDA"
    slide = section_slide(
        prs, slide_no=19, section_no=2, title=title,
        subtitle="Seven broad cell types are tested directly within six sex/APOE groups; each eligible direction remains its own KDA result.",
        topics=["Run universe", "Direct funnel", "Supported genes", "Result figures"],
        accent=TEAL,
        source="docs/phase_20_sex_apoe_kda/phase20_broad_funnel_explained.md",
        notes={
            "goal": "Introduce the broad-cell branch as a direct, within-run analysis rather than a reaggregation.",
            "walkthrough": "This section repeats the setup-to-results sequence for the broad-cell production release. It explains the 42 donor-level DEG contrasts, their 84 directional slots, query formation, the three completed KDA runs, within-run multiplicity correction, and the twelve direct candidate rows.",
            "boundary": "The broad branch does not pool runs with ACAT, calculate cross-run coverage, or require recurrence across broad-cell categories.",
            "transition": "Begin by defining the direct branch's counting units and why its final rows differ from the fine-cell units.",
        },
    )
    register(slide, title, "Broad-cell")

    # 20 — broad counting units
    title = "Broad-cell candidates stay attached to one cell type, group, and direction"
    slide = content_slide(prs, title=title, kicker="Broad setup", slide_no=20, accent=TEAL,
                          subtitle="There is no cross-run aggregation step in the direct branch.")
    cards = [
        ("1", "DEG contrast", "broad type × sex/APOE", "One AD-vs-NCI contrast from a donor-level broad-cell model", PALE_GREEN),
        ("2", "Directional slot", "+ AD-up or AD-down", "A possible core-MT query", PALE_SKY),
        ("3", "KDA run", "eligible slot with query n ≥ 3", "One network and one direction", PALE_BLUE),
        ("4", "Candidate row", "gene × broad × group × direction", "One recomputed non-MT run q", PALE_GOLD),
    ]
    for index, (number, term, key, detail, bg) in enumerate(cards):
        x = 0.66 + index * 3.10
        add_rect(slide, x, 1.50, 2.83, 4.63, color=bg, outline=TEAL if index < 3 else GOLD)
        add_circle(slide, x + 0.22, 1.78, 0.48, TEAL if index < 3 else GOLD)
        add_text(slide, number, x + 0.22, 1.88, 0.48, 0.18, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, term, x + 0.22, 2.48, 2.36, 0.60, size=19, color=NAVY, bold=True, font=FONT_HEAD)
        add_text(slide, key, x + 0.22, 3.30, 2.36, 0.65, size=12.2, color=TEAL if index < 3 else GOLD, bold=True)
        add_rect(slide, x + 0.20, 4.20, 2.40, 1.20, color=WHITE, outline=LIGHT)
        add_text(slide, detail, x + 0.34, 4.38, 2.12, 0.86, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_takeaway(slide, "A broad candidate is a gene selected from one direction-specific KDA run—not an ACAT-combined gene.", accent=TEAL)
    finalize_slide(
        slide, slide_no=20,
        source="phase20_broad_funnel_explained.md, analysis-unit definitions",
        goal="Define the four broad-branch units before presenting the numerical funnel.",
        walkthrough="A donor-level broad-cell model supplies an AD-versus-NCI stratum contrast, whose sign creates two potential core-mitochondrial queries. A slot is called only when at least three query genes remain. Every selected row stays linked to one broad cell type, group, and direction.",
        boundary="No candidate row combines an up and down direction or borrows evidence from another broad cell type or group.",
        transition="Use those definitions to partition all 84 planned directional slots.",
    )
    register(slide, title, "Broad-cell")

    # 21 — broad run universe
    title = "7 broad cell types create 84 slots—but only 3 direct KDA runs"
    slide = content_slide(prs, title=title, kicker="Broad runs", slide_no=21, accent=TEAL,
                          subtitle="The 84 slots are planned opportunities; they are not 84 executed KDA jobs.")
    add_rect(slide, 0.70, 1.47, 3.70, 1.31, color=PALE_GREEN, outline=TEAL)
    add_text(slide, "7 broad cell types", 0.95, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "× 6 groups = 42 DEG contrasts", 0.95, 2.18, 3.20, 0.30, size=12.3, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.82, 1.47, 3.70, 1.31, color=PALE_SKY, outline=SKY)
    add_text(slide, "42 planned contrasts", 5.07, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "× up/down = 84 predeclared slots", 5.07, 2.18, 3.20, 0.30, size=11.4, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.94, 1.47, 3.70, 1.31, color=PALE_GOLD, outline=GOLD)
    add_text(slide, "3 completed runs", 9.19, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "all three use AD-down queries", 9.19, 2.18, 3.20, 0.30, size=12.3, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_table(
        slide,
        ["Mutually exclusive relaxed-tier outcome", "Slots", "Interpretation"],
        [
            ["Source DEG contrast not estimable", "4", "2 contrasts × 2 directions"],
            ["No provisional directional core-MT query", "61", "No query genes pass DEG rules"],
            ["Effective query empty after network mapping", "4", "All provisional genes are lost"],
            ["Effective query size 1–2", "12", "Below direct KDA minimum"],
            ["Effective query size ≥3", "3", "KDA completed"],
            ["TOTAL", "84", "3 completed runs"],
        ],
        0.86, 3.17, [6.35, 1.22, 3.95], row_h=0.48, header_h=0.47,
        accent=TEAL, font_size=9.6, highlight_last=True,
    )
    add_takeaway(slide, "Most slots disappear before KDA because the relaxed broad DEG query is absent or too small.", accent=TEAL)
    finalize_slide(
        slide, slide_no=21,
        source="phase20_broad_direction_manifest.tsv; phase20_broad_funnel_explained.md",
        goal="Show exactly how 84 planned broad directional slots become three completed KDA runs.",
        walkthrough="The five outcomes are mutually exclusive. Four slots inherit two non-estimable source contrasts. Sixty-one have no provisional directional core-MT query, four lose every query gene during network mapping, and twelve retain only one or two genes. The remaining three meet the direct minimum of three genes.",
        boundary="A skipped slot has no KDA result. It must not be counted as a completed null run or assigned an artificial p value.",
        transition="Zoom out to the two-stage broad funnel: first construct runs, then select genes within each completed run.",
    )
    register(slide, title, "Broad-cell")

    # 22 — broad funnel overview
    title = "Broad-cell funnel at a glance: 84 slots to 12 direct candidate rows"
    slide = content_slide(prs, title=title, kicker="Broad funnel overview", slide_no=22, accent=TEAL,
                          subtitle="The upper row counts slots; the lower row counts explicit gene tests or candidate rows.")
    add_flow_row(slide, [
        ("84", "directional slots"), ("80", "estimable-source slots"),
        ("19", "with provisional query"), ("15", "non-empty after mapping"),
        ("3", "completed KDA runs"),
    ], 1.63, accent=TEAL)
    add_text(slide, "RUN CONSTRUCTION", 0.69, 1.36, 2.1, 0.18, size=8.7, color=TEAL, bold=True)
    add_flow_row(slide, [
        ("377", "explicit tested rows"), ("362", "non-core-MT family"),
        ("15", "query overlap ≥2"), ("15", "fold enrichment >1"),
        ("12", "within-run q ≤0.10"),
    ], 4.13, accent=GOLD)
    add_text(slide, "DIRECT CANDIDATE SELECTION", 0.69, 3.87, 2.9, 0.18, size=8.7, color=GOLD, bold=True)
    add_rect(slide, 0.82, 3.14, 11.70, 0.58, color=PALE_PURPLE, outline=PURPLE)
    add_text(slide, "Counting unit changes once: slots/runs → explicit gene × run rows", 1.04, 3.30, 11.25, 0.25, size=12.2, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Within-run BH is calculated before overlap and enrichment gates; there is no ACAT stage.", accent=GOLD)
    finalize_slide(
        slide, slide_no=22,
        source="phase20_broad_filter_funnel.tsv; phase20_broad_direction_manifest.tsv",
        goal="Give a complete overview of broad run construction and direct candidate selection.",
        walkthrough="The top row shows how source estimability, DEG query availability, network mapping, and the minimum effective query size leave three runs. The lower row begins from all explicit KDA-tested rows, removes core mitochondrial proteins, and then applies support, enrichment, and within-run q thresholds.",
        boundary="Percentages should be calculated only within a row and with the correct denominator. The transition from three runs to 377 rows changes the counting unit.",
        transition="Return to the source model to explain how donor-level DEGs become provisional broad queries.",
    )
    register(slide, title, "Broad-cell")

    # 23 — broad DEG gate
    title = "Donor-level DEG rules leave 65 provisional core-MT memberships"
    slide = content_slide(prs, title=title, kicker="Broad filter 1 • DEG query", slide_no=23, accent=TEAL,
                          subtitle="Each donor is a replicate; donor × broad-cell samples require enough nuclei and both diagnosis groups.")
    add_metric(slide, "786,242", "tested gene × contrast rows", 0.68, 1.49, 2.18, accent=TEAL)
    add_metric(slide, "2,336", "relaxed DEG memberships", 3.03, 1.49, 2.18, accent=TEAL)
    add_metric(slide, "65", "provisional core-MT memberships", 5.38, 1.49, 2.18, accent=TEAL)
    add_metric(slide, "35", "effective query memberships", 7.73, 1.49, 2.18, accent=GOLD)
    add_metric(slide, "3", "eligible directional queries", 10.08, 1.49, 2.18, accent=GOLD)
    add_rect(slide, 0.72, 3.08, 5.86, 3.05, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Source contrast must be estimable", 1.03, 3.40, 5.25, accent=TEAL)
    add_bullets(slide, [
        "Aggregate counts by donor × broad cell type.",
        "Keep a donor-cell sample only when it has ≥20 nuclei.",
        "Fit 7 shared broad-cell models; test up to 6 stratum contrasts each.",
        "Require ≥5 AD + ≥5 NCI donors: 40/42 contrasts are estimable.",
    ], 1.03, 3.93, 5.05, size=10.6, line_h=0.46, accent=TEAL)
    add_rect(slide, 6.82, 3.08, 5.80, 3.05, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "Provisional query rules", 7.13, 3.40, 5.19, accent=GOLD)
    add_bullets(slide, [
        "edgeR filterByExpr() defines the tested-gene universe.",
        "Relaxed DEG: BH q≤0.10 + |log2FC|≥log2(1.2).",
        "Signed 2,336: 983 AD-up + 1,353 AD-down memberships.",
        "Core-MT intersection: 20 up + 45 down = 65.",
        "Network mapping removes 30; 35 memberships remain.",
    ], 7.13, 3.88, 5.00, size=9.9, line_h=0.39, accent=GOLD)
    add_text(slide, "Broad-cell DEGs can reflect within-fine-type expression and shifts in fine-type composition.", 1.10, 6.30, 11.10, 0.24, size=9.6, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "A contrast may be not estimable before testing; this is not a zero-DEG result or shared-model failure.", accent=VERMILION)
    finalize_slide(
        slide, slide_no=23,
        source="phase20_broad_funnel_explained.md, source-DEG and query-construction sections",
        goal="Explain both source estimability and gene-level DEG filtering in the broad branch.",
        walkthrough="One joint edgeR model is fit per broad cell type, with donor-level pseudobulk samples as replicates. filterByExpr defines the tested universe before up to six stratum contrasts are evaluated. Relaxed thresholds yield 2,336 signed memberships; the core-MT intersection yields 65, and network mapping removes 30.",
        boundary="Two vascular contrasts fail donor-count eligibility even though the shared vascular model fits; they create four unavailable directions. Broad-cell changes can also mix expression shifts within fine types with fine-type composition shifts.",
        transition="Inspect the three surviving directional queries and their exact KDA accounting.",
    )
    register(slide, title, "Broad-cell")

    # 24 — exact broad runs
    title = "Only three relaxed queries reach the direct KDA minimum"
    slide = content_slide(prs, title=title, kicker="Broad filter 2 • query eligibility", slide_no=24, accent=TEAL,
                          subtitle="All three completed runs are AD-down mitochondrial queries; one completes with no candidates.")
    add_table(
        slide,
        ["Broad × group × direction", "Provisional", "Effective", "Background", "Explicit tests", "Relaxed rows"],
        [
            ["Astrocytes × F_e4 × AD-down", "18", "13", "7,828", "160", "2"],
            ["Astrocytes × M_e33 × AD-down", "5", "3", "7,828", "80", "0"],
            ["OPCs × F_e4 × AD-down", "6", "4", "7,817", "137", "10"],
        ],
        0.72, 1.53, [4.10, 1.18, 1.10, 1.26, 1.35, 1.24], row_h=0.62, header_h=0.54,
        accent=TEAL, font_size=9.7,
    )
    add_rect(slide, 0.72, 4.34, 3.68, 1.61, color=PALE_RED, outline=VERMILION)
    add_panel_title(slide, "Mapping loss", 1.01, 4.62, 3.08, accent=VERMILION)
    add_text(slide, "These runs start with 29 provisional memberships; 9 are lost in network mapping, leaving 20 effective query genes.", 1.03, 5.10, 3.05, 0.60, size=10.5, color=GRAY)
    add_rect(slide, 4.82, 4.34, 3.68, 1.61, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Minimum size", 5.11, 4.62, 3.08, accent=TEAL)
    add_text(slide, "The plan permits n≥3. The n=3 and n=4 runs are flagged small_query_3_9; only Astrocytes F_e4 has n≥10.", 5.13, 5.10, 3.05, 0.60, size=10.4, color=GRAY)
    add_rect(slide, 8.92, 4.34, 3.68, 1.61, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "Valid empty result", 9.21, 4.62, 3.08, accent=GOLD)
    add_text(slide, "Astrocytes × M_e33 completes successfully; its best non-core-MT row has q=0.1195, so zero rows pass q ≤0.10.", 9.23, 5.10, 3.05, 0.60, size=10.7, color=GRAY)
    add_takeaway(slide, "Completed-empty is evidence from a valid run; skipped is the absence of a runnable query.", accent=TEAL)
    finalize_slide(
        slide, slide_no=24,
        source="phase20_broad_direction_manifest.tsv; phase20_broad_all_candidate_tests.tsv.gz",
        goal="Make the three completed broad runs and their query sizes concrete.",
        walkthrough="The table traces provisional core-MT memberships through network mapping, then shows each matching background, explicit-test family, and relaxed-candidate yield. Both Phase 20 branches use the validated execution floor of three effective query genes; the direct n=3 and n=4 runs retain their small-query caution labels.",
        boundary="The n=3 and n=4 runs remain labeled small_query_3_9 and warrant extra caution. Astrocytes M_e33 is not a failure: it completed, produced explicit tests, and simply has no row at q≤0.10.",
        transition="Explain what the 377 explicit tests represent and why broad multiplicity correction stays within run.",
    )
    register(slide, title, "Broad-cell")

    # 25 — explicit KDA evidence
    title = "Broad KDA evaluates 377 explicit network-neighborhood tests"
    slide = content_slide(prs, title=title, kicker="Broad filter 3 • KDA evidence", slide_no=25, accent=TEAL,
                          subtitle="Only reachable/testable candidate genes receive explicit enrichment tests; the background is not the BH family.")
    add_metric(slide, "23,473", "background gene × run opportunities", 0.73, 1.53, 2.52, accent=TEAL,
               note="3 run-specific backgrounds")
    add_metric(slide, "377", "explicit KDA-tested rows", 3.53, 1.53, 2.52, accent=TEAL,
               note="160 + 80 + 137")
    add_metric(slide, "362", "non-core-MT BH rows", 6.33, 1.53, 2.52, accent=GOLD,
               note="151 + 78 + 133")
    add_metric(slide, "11", "stock significant returns", 9.13, 1.53, 2.52, accent=PURPLE,
               note="audit field, not final gate")
    add_rect(slide, 0.73, 3.38, 3.56, 2.45, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "One explicit row", 1.02, 3.70, 2.96, accent=TEAL)
    add_text(slide, "A gene reachable within 3 undirected hops with a valid directed layer 1–3 neighborhood, explicitly tested against one run's query.", 1.03, 4.25, 2.95, 1.00, size=10.6, color=GRAY)
    add_rect(slide, 4.61, 3.38, 3.56, 2.45, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "BH family", 4.90, 3.70, 2.96, accent=BLUE)
    add_text(slide, "Remove core mitochondrial proteins, then recompute BH separately across every remaining explicit row in that run—before selection gates.", 4.91, 4.25, 2.95, 1.00, size=11.3, color=GRAY)
    add_rect(slide, 8.49, 3.38, 3.56, 2.45, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "What is absent", 8.78, 3.70, 2.96, accent=GOLD)
    add_text(slide, "23,096 background opportunities lie outside the explicit family. They are neither failed candidates nor implicit P=1 rows.", 8.79, 4.25, 2.95, 1.00, size=10.6, color=GRAY)
    add_takeaway(slide, "11 stock q≤0.05 rows = 9 strict non-core-MT candidates + core-MT GLDC and ME3; final selection uses recomputed q.", accent=PURPLE)
    finalize_slide(
        slide, slide_no=25,
        source="phase20_broad_all_candidate_tests.tsv.gz; phase20_broad_direction_manifest.tsv",
        goal="Explain the broad KDA test universe and the within-run multiple-testing family.",
        walkthrough="Three run-specific backgrounds contribute 23,473 gene-by-run opportunities, but only 377 reachable genes with valid directed neighborhoods receive explicit tests. The other 23,096 are outside the explicit family. Removing fifteen core-MT rows leaves run-specific BH families of 151, 78, and 133.",
        boundary="The eleven stock returns include core-MT GLDC and ME3. Excluding them reconciles the stock count to nine strict non-core-MT candidates; stock return status is still only an audit field.",
        transition="Apply the biological-support and q gates to the complete non-core-MT families.",
    )
    register(slide, title, "Broad-cell")

    # 26 — broad candidate gates
    title = "Support and within-run q reduce 362 non-core-MT rows to 12"
    slide = content_slide(prs, title=title, kicker="Broad filters 4–6 • candidate selection", slide_no=26, accent=GOLD,
                          subtitle="BH q values are computed on full non-core-MT run families before post-BH gates are applied.")
    add_flow_row(slide, [
        ("362", "non-core-MT family"), ("15", "query overlap ≥2"),
        ("15", "fold enrichment >1"), ("12", "relaxed q ≤0.10"),
        ("9", "strict q ≤0.05"),
    ], 1.64, accent=GOLD)
    explanations = [
        ("Non-core-MT", "Exclude only core mitochondrial proteins; extended-tier NCOA1 remains eligible.", TEAL, PALE_GREEN),
        ("Overlap ≥2", "Require at least two effective query genes inside the candidate's tested neighborhood.", BLUE, PALE_BLUE),
        ("Enrichment >1", "Require more overlap than expected under the run-specific network background.", PURPLE, PALE_PURPLE),
        ("Run q", "Relaxed main list uses q ≤0.10; q ≤0.05 marks the strict reference subset.", GOLD, PALE_GOLD),
    ]
    for index, (heading, body, accent, bg) in enumerate(explanations):
        x = 0.67 + index * 3.11
        add_rect(slide, x, 3.40, 2.84, 2.48, color=bg, outline=accent)
        add_panel_title(slide, heading, x + 0.25, 3.73, 2.35, accent=accent)
        add_text(slide, body, x + 0.27, 4.30, 2.30, 1.10, size=10.7, color=GRAY)
    add_takeaway(slide, "After defining the non-core-MT family, compute BH before filtering on overlap, enrichment, or q.", accent=VERMILION)
    finalize_slide(
        slide, slide_no=26,
        source="phase20_broad_filter_funnel.tsv; phase20_broad_non_mt_candidates.tsv",
        goal="Explain every direct broad candidate gate and its numerical effect.",
        walkthrough="Core-MT genes are first removed to define the inferential family. Fifteen of 362 non-core-MT rows overlap at least two query genes; all fifteen have enrichment above one. Twelve pass q≤0.10 and nine pass q≤0.05. BH was computed across each full non-core-MT run family before these gates.",
        boundary="NCOA1 is labeled mitochondrial-extended but is not a core mitochondrial protein, so it legitimately remains in the non-core-MT candidate universe.",
        transition="Name the supported genes and separate strict rows, relaxed-only rows, and the completed-empty run.",
    )
    register(slide, title, "Broad-cell")

    # 27 — broad supported genes
    title = "Twelve broad candidates occur in two runs; nine pass the strict reference"
    slide = content_slide(prs, title=title, kicker="Broad results • supported genes", slide_no=27, accent=TEAL,
                          subtitle="AD-down labels the mitochondrial query direction; candidate drivers need not be downregulated or DEGs.")
    add_rect(slide, 0.68, 1.46, 3.58, 4.81, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Astrocytes • F_e4 • AD-down", 0.98, 1.78, 2.97, accent=TEAL)
    add_metric(slide, "2", "relaxed • both strict", 0.98, 2.30, 1.45, accent=TEAL, bg=WHITE)
    add_text(slide, "ELL2\nSLC44A3", 1.02, 3.75, 2.84, 1.10, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Both supported by AHCYL1 and PNPLA8 overlap.", 1.01, 5.26, 2.90, 0.55, size=10.3, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.54, 1.46, 3.58, 4.81, color=PALE_GRAY, outline=MID)
    add_panel_title(slide, "Astrocytes • M_e33 • AD-down", 4.84, 1.78, 2.97, accent=MID)
    add_metric(slide, "0", "completed-empty", 4.84, 2.30, 1.45, accent=MID, bg=WHITE)
    add_text(slide, "No q ≤0.10 rows", 4.89, 3.84, 2.84, 0.46, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Best row: NEK6, q=0.1195. The run completed successfully.", 4.87, 4.76, 2.90, 0.84, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.40, 1.46, 4.25, 4.81, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "OPCs • F_e4 • AD-down", 8.70, 1.78, 3.64, accent=GOLD)
    add_metric(slide, "10", "relaxed • 7 strict", 8.70, 2.30, 1.45, accent=GOLD, bg=WHITE)
    add_text(slide, "STRICT", 8.72, 3.72, 0.80, 0.20, size=8.8, color=TEAL, bold=True)
    add_text(slide, "CAMK2D • RAPGEF4 • RAB3IP • FOXN3\nAC092691.1 • FAM13A • NCOA1", 8.72, 4.00, 3.45, 0.76, size=11.0, color=NAVY, bold=True)
    add_text(slide, "RELAXED-ONLY", 8.72, 5.00, 1.18, 0.20, size=8.8, color=VERMILION, bold=True)
    add_text(slide, "FGF14 • GRID1 • DENND1A", 8.72, 5.28, 3.43, 0.40, size=11.0, color=NAVY, bold=True)
    add_takeaway(slide, "The broad result is 12 unique genes because none recurs across the two candidate-bearing runs.", accent=TEAL)
    finalize_slide(
        slide, slide_no=27,
        source="phase20_broad_non_mt_candidates.tsv; phase20_broad_direction_manifest.tsv",
        goal="List every supported broad candidate and make the tier distinction explicit.",
        walkthrough="The Astrocytes F_e4 run contributes two strict candidates. Astrocytes M_e33 is a valid completed-empty run. OPCs F_e4 contributes ten relaxed candidates: seven strict and three relaxed-only. Because the two candidate-bearing runs share no gene, twelve rows also equal twelve unique symbols.",
        boundary="AD-down identifies the mitochondrial query direction. It does not assert that the candidate driver itself is downregulated or is a DEG.",
        transition="Use the requested recurrence-style figure to show that q strength, rather than cross-run recurrence, differentiates these genes.",
    )
    register(slide, title, "Broad-cell")

    # 28 — broad recurrence figure
    title = "Broad candidates do not recur across runs; bars show within-run q strength"
    slide = content_slide(prs, title=title, kicker="Broad results • driver recurrence", slide_no=28, accent=TEAL,
                          subtitle="This recurrence analogue keeps all 12 unique genes and encodes −log10(non-core-MT run q).")
    add_rect(slide, 0.47, 1.27, 7.20, 5.40, color=WHITE, outline=LIGHT)
    add_picture_contain(slide, FIG["broad_recurrence"], 0.58, 1.37, 6.98, 5.18,
                        alt="Broad-cell Phase 20 horizontal bar chart of twelve direct non-core-mitochondrial candidates, with blue strict rows, hatched orange relaxed-only rows, and bar length equal to minus log10 within-run q")
    add_rect(slide, 7.93, 1.42, 4.73, 4.98, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "How to read it", 8.23, 1.74, 4.10, accent=TEAL)
    add_bullets(slide, [
        "Each gene appears in exactly one completed broad run.",
        "Bar length = −log10(non-core-MT run q), not recurrence count.",
        "Run labels/panels identify Astrocytes F_e4 versus OPCs F_e4.",
        "Blue = strict q≤0.05; hatched orange = relaxed-only.",
        "Nine genes are strict q ≤0.05; three OPC genes are relaxed-only.",
        "Astrocytes M_e33 is absent because its candidate list is empty.",
    ], 8.23, 2.25, 4.02, size=10.2, line_h=0.50, accent=TEAL)
    add_rect(slide, 8.22, 5.56, 4.10, 0.56, color=WHITE, outline=LIGHT)
    add_text(slide, "No broad recurrence gate was applied.", 8.41, 5.70, 3.72, 0.27, size=10.3, color=VERMILION, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "The figure answers ‘how strong within this run?’—not ‘in how many runs did this gene recur?’", accent=TEAL)
    finalize_slide(
        slide, slide_no=28,
        source="results/figures/analysis/phase_20_sex_apoe_kda_broad/driver_recurrence figure bundle",
        goal="Interpret the requested broad recurrence-style figure without implying cross-run recurrence.",
        walkthrough="All twelve candidate symbols are unique to one candidate-bearing run. Bar length shows minus log ten of recomputed within-run q; run headings separate Astrocytes F_e4 and OPCs F_e4. Blue bars pass strict q≤0.05, while hatched orange bars are relaxed-only.",
        boundary="This is not evidence that candidates reproduce across groups, directions, or networks, and no recurrence threshold enters the broad candidate definition.",
        transition="The next requested figure applies the display-only top-five limit within each completed broad-cell × group × direction run.",
    )
    register(slide, title, "Broad-cell")

    # 29 — broad top-five figure
    title = "The broad top-five view displays 7 strict candidates across 3 run rows"
    slide = content_slide(prs, title=title, kicker="Broad results • top-five candidates", slide_no=29, accent=TEAL,
                          subtitle="Two Astrocyte candidates, five OPC candidates, and a deliberately blank completed run.")
    add_rect(slide, 0.47, 1.31, 8.22, 5.31, color=WHITE, outline=LIGHT)
    add_picture_contain(slide, FIG["broad_top5"], 0.58, 1.41, 8.00, 5.10,
                        alt="Broad-cell Phase 20 three-row tile chart showing two Astrocytes F_e4 candidates, a blank Astrocytes M_e33 row, and five OPCs F_e4 candidates")
    add_rect(slide, 8.94, 1.45, 3.72, 4.95, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "Display contract", 9.23, 1.76, 3.13, accent=GOLD)
    add_metric(slide, "7", "displayed candidates", 9.23, 2.26, 1.33, accent=GOLD, bg=WHITE)
    add_metric(slide, "3", "completed run rows", 10.77, 2.26, 1.33, accent=GOLD, bg=WHITE)
    add_bullets(slide, [
        "Order: q, raw P, larger overlap, larger enrichment, then symbol.",
        "All seven displayed genes pass strict q ≤0.05.",
        "Up to five are shown per run; no failing row is backfilled.",
        "The full candidate table still contains all 12 relaxed rows.",
    ], 9.23, 3.68, 2.97, size=10.6, line_h=0.54, accent=GOLD)
    add_takeaway(slide, "The blank M_e33 row makes a valid zero visible; the top-five cap changes presentation only.", accent=GOLD)
    finalize_slide(
        slide, slide_no=29,
        source="results/figures/analysis/phase_20_sex_apoe_kda_broad/top5_candidates figure bundle; phase20_broad_top5_summary.tsv",
        goal="Interpret the requested broad top-five tile chart and its blank completed run.",
        walkthrough="The figure keeps one row for each completed broad-cell by group by direction run. Astrocytes F_e4 contributes two tiles, Astrocytes M_e33 stays blank, and OPCs F_e4 contributes its first five ranks. Ranking uses q, raw P, larger overlap, larger fold enrichment, then symbol.",
        boundary="Only five of the ten relaxed OPC candidates are visible. The remaining five are not removed from the production candidate table, and the blank row is not filled with a q-failing gene.",
        transition="Close by comparing the two branches filter by filter and stating what may be concluded from each.",
    )
    register(slide, title, "Broad-cell")

    # 30 — direct comparison
    title = "The same words can refer to different evidence in the two branches"
    slide = content_slide(prs, title=title, kicker="Synthesis", slide_no=30, accent=PURPLE,
                          subtitle="Use branch, counting unit, and multiplicity method together when reporting a result.")
    add_table(
        slide,
        ["Question", "Fine-cell reaggregation", "Direct broad-cell KDA"],
        [
            ["Starting design", f"{fine['fine_type_count']} fine types × {fine['group_count']} groups", "7 broad types × 6 groups"],
            ["Planned directions", f"{count_text(fine['planned_slots'])} slots", "84 slots"],
            ["KDA execution rule", "Called at effective query ≥3", "Called at effective query ≥3"],
            ["Analysis inclusion", f"All completed query ≥3: {count_text(fine['included_runs'])} runs", "All completed query ≥3: 3 runs"],
            ["Evidence completion", "Explicit + implicit P=1 + absent", "Explicit KDA rows only"],
            ["Multiple testing", "ACAT across runs, then category BH", "BH separately within each run"],
            ["Support requirement", "≥1 supporting run (q≤.10 main)", "Overlap ≥2 in that run"],
            ["Final main result", f"{fine['relaxed_candidates']} gene × category units / {fine['relaxed_genes']} genes", "12 gene × run rows / 12 genes"],
            ["Direction meaning", "Both directions may feed a category", "One retained direction per row"],
        ],
        0.68, 1.41, [2.50, 4.95, 4.95], row_h=0.51, header_h=0.51,
        accent=PURPLE, font_size=9.55,
    )
    add_takeaway(slide, "Fine asks whether evidence combines across eligible fine runs; broad asks whether one direct run supports the gene.", accent=PURPLE)
    finalize_slide(
        slide, slide_no=30,
        source="validated fine and broad production status, run-manifest, and funnel tables",
        goal="Put every potentially confusing term side by side across the two branches.",
        walkthrough=f"The comparison highlights the shared directional query construction and minimum effective query size of three, then shows where the branches diverge. Fine analysis retains all {count_text(fine['included_runs'])} completed directional calls and creates complete cross-run evidence for ACAT: explicit tests, implicit P=1 rows, and absent-from-background states. Broad analysis retains its three direct runs and controls multiplicity within each explicit family.",
        boundary="The final counts should never be pooled, ranked against each other, or interpreted as equivalent evidence units.",
        transition="End with the core interpretation rules and the exact production artifacts that support the deck.",
    )
    register(slide, title, "Synthesis")

    # 31 — take-home and source map
    title = "Take-home: follow the unit, then the filter, then the q value"
    slide = content_slide(prs, title=title, kicker="Take-home", slide_no=31, accent=GOLD,
                          subtitle="The deck is a teaching layer over validated, auditable production tables and current canonical figures.")
    takeaways = [
        ("1", "A run is one cell type × sex/APOE × direction KDA analysis.", "A planned slot becomes a run only when the effective query meets that branch's execution rule."),
        ("2", "A candidate is a network gene—not necessarily a DEG or query gene.", "Its neighborhood is enriched for the mitochondrial query; AD-up/down labels the query direction."),
        ("3", "Fine and broad q values answer different questions.", "Fine combines run p values with ACAT then controls category BH; broad controls BH within one run."),
        ("4", "Skipped, completed-empty, and filtered-out are distinct outcomes.", "Unavailable models and unrunnable queries carry no KDA evidence; completed rows can validly yield zero candidates."),
    ]
    for index, (number, heading, body) in enumerate(takeaways):
        y = 1.40 + index * 1.18
        add_circle(slide, 0.74, y + 0.06, 0.48, GOLD)
        add_text(slide, number, 0.74, y + 0.16, 0.48, 0.18, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, heading, 1.42, y, 6.78, 0.33, size=14.2, color=NAVY, bold=True)
        add_text(slide, body, 1.42, y + 0.42, 6.72, 0.51, size=10.6, color=GRAY)
    add_rect(slide, 8.56, 1.40, 4.02, 4.74, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "Audit trail", 8.88, 1.73, 3.39, accent=BLUE)
    add_text(slide, "Fine branch", 8.89, 2.29, 1.12, 0.22, size=9.4, color=BLUE, bold=True)
    add_text(slide, "phase20_status.tsv\nphase20_filter_funnel.tsv\nphase20_relaxed_candidates.tsv", 8.89, 2.61, 3.23, 0.86, size=10.4, color=NAVY, bold=True)
    add_text(slide, "Broad branch", 8.89, 3.75, 1.12, 0.22, size=9.4, color=TEAL, bold=True)
    add_text(slide, "phase20_broad_status.tsv\nphase20_broad_filter_funnel.tsv\nphase20_broad_non_mt_candidates.tsv", 8.89, 4.07, 3.23, 0.86, size=10.4, color=NAVY, bold=True)
    add_text(slide, "Speaker notes on every slide preserve teaching goals, interpretation boundaries, and transitions.", 8.89, 5.28, 3.20, 0.56, size=9.8, color=GRAY, italic=True)
    add_takeaway(slide, f"Report the final unit explicitly: {fine['relaxed_candidates']} fine gene × category units versus 12 broad gene × run rows.", accent=GOLD)
    finalize_slide(
        slide, slide_no=31,
        source="validated Phase 20 fine and broad production releases; canonical figure bundles",
        goal="Leave the audience with four interpretation rules and a compact audit-source map.",
        walkthrough="Reinforce that run, query, and candidate are different objects. Then emphasize that q values are branch-specific, and that unavailable, skipped, completed-empty, and statistically filtered outcomes must stay distinct. The source card identifies the central production tables behind all stated counts.",
        boundary="The deck summarizes association and network-enrichment evidence. It does not establish causal regulation or formal sex-by-APOE interaction effects.",
        transition="Use the linked production tables for row-level follow-up and the canonical PNG bundles for figure-specific provenance.",
    )
    register(slide, title, "Synthesis")

    if len(metadata) != 31:
        raise RuntimeError(f"Expected 31 slides, built {len(metadata)}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path, metadata


def slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            chunks.append(shape.text)
    return "\n".join(chunks)


def validate_deck(output_path: Path, metadata: Sequence[dict[str, Any]]) -> list[dict[str, Any]]:
    checks: list[dict[str, Any]] = []

    def check(check_id: str, observed: Any, expected: Any, passed: bool) -> None:
        checks.append({
            "schema_version": "phase20_deck_checks_v1",
            "check_id": check_id,
            "observed": str(observed),
            "expected": str(expected),
            "passed": str(bool(passed)).upper(),
        })

    check("pptx_exists", output_path.is_file(), True, output_path.is_file())
    prs = Presentation(output_path)
    check("slide_count", len(prs.slides), 31, len(prs.slides) == 31)
    check("widescreen_width", prs.slide_width, SLIDE_W, prs.slide_width == SLIDE_W)
    check("widescreen_height", prs.slide_height, SLIDE_H, prs.slide_height == SLIDE_H)

    source_slides = 0
    noted_slides = 0
    bounded_slides = 0
    picture_count = 0
    alt_text_count = 0
    upper_right_number_count = 0
    for slide in prs.slides:
        text_value = slide_text(slide)
        if "Source:" in text_value:
            source_slides += 1
        notes_frame = slide.notes_slide.notes_text_frame
        notes_text = notes_frame.text if notes_frame is not None else ""
        if all(marker in notes_text for marker in (
            "Teaching goal:", "Walk through:", "Scientific boundary:", "Transition:"
        )) and len(notes_text.split()) >= 45:
            noted_slides += 1
        in_bounds = True
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width + 10 or shape.top + shape.height > prs.slide_height + 10:
                in_bounds = False
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
                if c_nv_pr.get("descr", "").strip():
                    alt_text_count += 1
            if (
                shape.left >= Inches(12.20)
                and shape.top <= Inches(0.50)
                and hasattr(shape, "text")
                and shape.text.strip().isdigit()
            ):
                upper_right_number_count += 1
        if in_bounds:
            bounded_slides += 1

    check("source_line_on_every_slide", source_slides, 31, source_slides == 31)
    check("structured_notes_on_every_slide", noted_slides, 31, noted_slides == 31)
    check("all_shapes_within_slide", bounded_slides, 31, bounded_slides == 31)
    check("canonical_figure_picture_count", picture_count, 4, picture_count == 4)
    check("figure_alt_text_count", alt_text_count, 4, alt_text_count == 4)
    check("upper_right_slide_numbers", upper_right_number_count, 0, upper_right_number_count == 0)

    with zipfile.ZipFile(output_path) as archive:
        bad_member = archive.testzip()
        slide_xml = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        media_hashes = {
            hashlib.sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith("ppt/media/")
        }
    check("zip_integrity", bad_member, None, bad_member is None)
    check("slide_xml_count", len(slide_xml), 31, len(slide_xml) == 31)
    canonical_hashes = {sha256(path) for path in FIG.values()}
    embedded_count = len(canonical_hashes & media_hashes)
    check("canonical_pngs_embedded_byte_exact", embedded_count, 4, embedded_count == 4)

    expected_titles = [entry["title"] for entry in metadata]
    title_hits = sum(expected in slide_text(slide) for expected, slide in zip(expected_titles, prs.slides, strict=True))
    check("registered_title_sequence", title_hits, 31, title_hits == 31)
    check("core_title_metadata", prs.core_properties.title, "Fine-cell and broad-cell sex/APOE key-driver analysis", prs.core_properties.title == "Fine-cell and broad-cell sex/APOE key-driver analysis")
    check("core_author_metadata", prs.core_properties.author, "Alzheimer project analysis team", prs.core_properties.author == "Alzheimer project analysis team")
    return checks


def write_audits(
    output_path: Path,
    metadata: Sequence[dict[str, Any]],
    checks: Sequence[dict[str, Any]],
) -> None:
    AUDIT_DIR.mkdir(parents=True, exist_ok=True)
    manifest_rows = []
    for role, path in [*INPUTS.items(), ("builder", Path(__file__).resolve()), ("deck", output_path)]:
        manifest_rows.append({
            "schema_version": "phase20_deck_input_manifest_v1",
            "role": role,
            "path": str(path),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        })
    write_tsv(
        manifest_rows,
        AUDIT_DIR / "phase20_deck_input_manifest.tsv",
        ["schema_version", "role", "path", "bytes", "sha256"],
    )

    prs = Presentation(output_path)
    inventory = []
    for index, (slide, entry) in enumerate(zip(prs.slides, metadata, strict=True), start=1):
        notes_frame = slide.notes_slide.notes_text_frame
        notes_text = notes_frame.text if notes_frame is not None else ""
        inventory.append({
            "schema_version": "phase20_deck_slide_inventory_v1",
            "slide_number": index,
            "section": entry["section"],
            "title": entry["title"],
            "shape_count": len(slide.shapes),
            "picture_count": sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes),
            "speaker_note_words": len(notes_text.split()),
        })
    write_tsv(
        inventory,
        AUDIT_DIR / "phase20_deck_slide_inventory.tsv",
        ["schema_version", "slide_number", "section", "title", "shape_count", "picture_count", "speaker_note_words"],
    )
    write_tsv(
        checks,
        AUDIT_DIR / "phase20_deck_checks.tsv",
        ["schema_version", "check_id", "observed", "expected", "passed"],
    )
    failed = sum(row["passed"] != "TRUE" for row in checks)
    status_rows = [{
        "schema_version": "phase20_deck_status_v1",
        "output_path": str(output_path),
        "slide_count": len(prs.slides),
        "failed_checks": failed,
        "validation_status": "validated_complete" if failed == 0 else "validation_failed",
        "pptx_bytes": output_path.stat().st_size,
        "pptx_sha256": sha256(output_path),
    }]
    write_tsv(
        status_rows,
        AUDIT_DIR / "phase20_deck_status.tsv",
        ["schema_version", "output_path", "slide_count", "failed_checks", "validation_status", "pptx_bytes", "pptx_sha256"],
    )
    if failed:
        failed_ids = [row["check_id"] for row in checks if row["passed"] != "TRUE"]
        raise RuntimeError("Deck validation failed: " + ", ".join(failed_ids))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT, help="Destination PPTX path")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output_path, metadata = build_deck(args.output.resolve())
    checks = validate_deck(output_path, metadata)
    write_audits(output_path, metadata, checks)
    print(f"Built {output_path}")
    print(f"Slides: {len(metadata)}")
    print(f"Validation checks: {len(checks)} passed")
    print(f"Audit directory: {AUDIT_DIR}")


def _append_fine_slides_4_18(prs: Presentation, facts: dict[str, Any], register) -> None:
    fine = facts["fine"]

    # 4 — fine counting units
    title = "Four counting units prevent the funnel from being misread"
    slide = content_slide(prs, title=title, kicker="Fine setup", slide_no=4, accent=BLUE,
                          subtitle="The funnel is not one continuously shrinking number of genes.")
    cards = [
        ("1", "DEG comparison", "fine cell type × sex/APOE", "One signed AD-versus-NCI result", PALE_BLUE),
        ("2", "Directional slot", "+ AD-up or AD-down", "A planned query opportunity", PALE_SKY),
        ("3", "KDA run", "validated completed call", "Effective query n ≥ 3 enters Phase 20", PALE_GREEN),
        ("4", "Candidate unit", "gene × group × broad network", "One or more fine runs; either or both directions", PALE_GOLD),
    ]
    for index, (number, term, key, detail, bg) in enumerate(cards):
        x = 0.66 + index * 3.10
        add_rect(slide, x, 1.50, 2.83, 4.63, color=bg, outline=BLUE if index < 3 else GOLD)
        add_circle(slide, x + 0.22, 1.78, 0.48, BLUE if index < 3 else GOLD)
        add_text(slide, number, x + 0.22, 1.88, 0.48, 0.18, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, term, x + 0.22, 2.48, 2.36, 0.60, size=19, color=NAVY, bold=True, font=FONT_HEAD)
        add_text(slide, key, x + 0.22, 3.30, 2.36, 0.62, size=13, color=BLUE if index < 3 else GOLD, bold=True)
        add_rect(slide, x + 0.20, 4.20, 2.40, 1.20, color=WHITE, outline=LIGHT)
        add_text(slide, detail, x + 0.36, 4.47, 2.08, 0.67, size=11.3, color=GRAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_takeaway(slide, "A query gene seeds KDA; a candidate driver is a network gene whose neighborhood is tested.", accent=BLUE)
    finalize_slide(
        slide, slide_no=4,
        source="phase20_source_run_manifest.tsv; phase20_methods.md, analysis-unit definitions",
        goal="Separate the comparison, slot, run, and candidate-unit concepts before showing counts.",
        walkthrough="Move from left to right. A comparison is one model result. Its sign creates two planned slots. A slot becomes a KDA call when at least three effective query genes remain, and every validated completed call at that execution floor enters the cross-run aggregate. The final unit adds gene identity to group and broad network.",
        boundary="The word candidate does not mean mitochondrial DEG. Candidate drivers may be non-DEG, non-mitochondrial network genes.",
        transition=f"Apply those definitions to the complete {count_text(fine['planned_slots'])}-slot run universe.",
    )
    register(slide, title, "Fine-cell")

    # 5 — fine run universe
    title = f"{fine['fine_type_count']} fine cell types create {count_text(fine['planned_slots'])} planned directional slots"
    slide = content_slide(prs, title=title, kicker="Fine runs", slide_no=5, accent=BLUE,
                          subtitle="Up and down are two query directions from one signed DEG comparison—not two model fits.")
    add_rect(slide, 0.68, 1.47, 3.66, 1.32, color=PALE_BLUE, outline=BLUE)
    add_text(slide, f"{fine['fine_type_count']} fine cell types", 0.92, 1.74, 3.18, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"× {fine['group_count']} sex/APOE groups = {count_text(fine['comparisons'])} comparisons", 0.92, 2.18, 3.18, 0.30, size=12.4, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.72, 1.47, 3.86, 1.32, color=PALE_SKY, outline=SKY)
    add_text(slide, f"{count_text(fine['comparisons'])} signed results", 4.96, 1.74, 3.38, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"× up/down = {count_text(fine['planned_slots'])} directional slots", 4.96, 2.18, 3.38, 0.30, size=12.4, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.96, 1.47, 3.68, 1.32, color=PALE_GREEN, outline=TEAL)
    add_text(slide, f"{count_text(fine['included_runs'])} completed calls", 9.20, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "All validated completed queries ≥3 enter Phase 20", 9.20, 2.18, 3.20, 0.30, size=10.4, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_table(
        slide,
        ["Mutually exclusive slot outcome", "Slots", "KDA called?", "Enters Phase 20?"],
        [
            ["Source DEG contrast not validated", count_text(fine["source_not_validated"]), "No", "No"],
            ["Validated source; effective query = 0", count_text(fine["effective_query_zero"]), "No", "No"],
            ["Effective query 1–2", count_text(fine["effective_query_one_two"]), "No", "No"],
            ["Effective query 3–9", count_text(fine["effective_query_three_nine"]), "Yes", "Yes"],
            ["Effective query ≥10", count_text(fine["effective_query_ten_plus"]), "Yes", "Yes"],
            ["TOTAL", count_text(fine["planned_slots"]), f"{count_text(fine['included_runs'])} calls", f"{count_text(fine['included_runs'])} runs"],
        ],
        0.78, 3.17, [6.15, 1.20, 2.05, 2.25], row_h=0.49, header_h=0.47,
        accent=BLUE, font_size=9.5, highlight_last=True,
    )
    finalize_slide(
        slide, slide_no=5,
        source="phase20_source_run_manifest.tsv; phase20_status.tsv",
        goal=f"Show exactly how {count_text(fine['planned_slots'])} planned slots become {count_text(fine['included_runs'])} included KDA runs.",
        walkthrough=f"The rows are mutually exclusive outcomes derived from the canonical source manifest. Phase 12 calls KDA when at least three effective query genes remain. All {count_text(fine['included_runs'])} completed calls enter Phase 20, including {count_text(fine['effective_query_three_nine'])} calls with query size three through nine and {count_text(fine['effective_query_ten_plus'])} with at least ten genes.",
        boundary=f"The {count_text(fine['source_not_validated'])} unavailable slots come from three non-estimable source contrasts times two directions. They are unavailable models, not completed tests with null results. The effective-query-zero row combines absent directional queries and complete loss during network mapping.",
        transition="Zoom out once to see how the counting unit changes through the full fine-cell funnel.",
    )
    register(slide, title, "Fine-cell")

    # 6 — fine funnel overview
    title = f"Fine-cell funnel at a glance: {count_text(fine['comparisons'])} comparisons to {count_text(fine['relaxed_candidates'])} candidate units"
    slide = content_slide(prs, title=title, kicker="Fine funnel overview", slide_no=6, accent=BLUE,
                          subtitle="Non-proportional overview: each arrow can change both the count and the counting unit.")
    add_flow_row(slide, [
        (count_text(fine["comparisons"]), "DEG comparisons"),
        (count_text(fine["planned_slots"]), "directional slots"),
        (count_text(fine["included_runs"]), "KDA calls"),
        (count_text(fine["included_runs"]), "included runs"),
        (count_text(fine["opportunities"]), "gene × run opportunities"),
    ], 1.62, accent=BLUE)
    add_text(slide, "RUN CONSTRUCTION", 0.69, 1.36, 2.1, 0.18, size=8.7, color=BLUE, bold=True)
    add_flow_row(slide, [
        (count_text(fine["non_mt_rows"]), "non-MT run rows"),
        (count_text(fine["aggregate_rows"]), "gene × category units"),
        (count_text(fine["relaxed_coverage_units"]), "coverage ≥ 0.50"),
        (count_text(fine["relaxed_support_units"]), "≥1 supported run"),
        (count_text(fine["relaxed_candidates"]), "relaxed candidate units"),
    ], 4.13, accent=GOLD)
    add_text(slide, "CANDIDATE CONSTRUCTION", 0.69, 3.87, 2.8, 0.18, size=8.7, color=GOLD, bold=True)
    add_rect(slide, 0.82, 3.14, 11.70, 0.58, color=PALE_PURPLE, outline=PURPLE)
    add_text(slide, "Counting unit changes: comparisons → slots → calls/runs → gene × run → gene × category", 1.04, 3.30, 11.25, 0.25, size=12.2, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, f"The final {fine['relaxed_candidates']} units are {fine['relaxed_genes']} distinct genes in {fine['relaxed_categories']} sex/APOE × broad-network categories.", accent=GOLD)
    finalize_slide(
        slide, slide_no=6,
        source="phase20_source_run_manifest.tsv; phase20_filter_funnel.tsv; phase20_status.tsv",
        goal="Provide a complete visual map before explaining each reduction in detail.",
        walkthrough="Read the top row as run construction and the lower row as candidate construction. The largest drop after aggregation is the individual supporting-run gate, but it is applied to gene-by-category units after complete run evidence and coverage have been established.",
        boundary=f"Do not calculate percentages between nodes that use different units. For example, {count_text(fine['included_runs'])} runs and {count_text(fine['aggregate_rows'])} gene-by-category units are not directly comparable denominators.",
        transition="Return to the start of the funnel and inspect how DEG rules construct the mitochondrial query.",
    )
    register(slide, title, "Fine-cell")

    # 7 — fine DEG filtering
    title = "Step 1A: DEG thresholds reduce expression tests to mitochondrial query genes"
    slide = content_slide(prs, title=title, kicker="Fine funnel • DEG query", slide_no=7, accent=BLUE,
                          subtitle="KDA receives directional core-MT paper DEGs, not the complete DEG list.")
    add_rect(slide, 0.66, 1.48, 4.07, 4.95, color=WHITE, outline=LIGHT)
    add_panel_title(slide, "Paper-DEG rule", 0.94, 1.78, 3.48, accent=BLUE)
    add_bullets(slide, [
        "Detected in ≥10% of AD or NCI nuclei.",
        "Within-contrast BH FDR < 0.05.",
        "|log₂ fold change| > log₂(1.3) ≈ 0.379.",
        "Positive and negative results form AD-up and AD-down queries.",
    ], 0.94, 2.34, 3.46, size=12.4, line_h=0.62, accent=BLUE)
    add_rect(slide, 0.94, 5.15, 3.44, 0.86, color=PALE_BLUE, outline=BLUE)
    add_text(slide, "Query = paper DEG ∩ core_mito_protein ∩ requested direction", 1.14, 5.37, 3.05, 0.43, size=11.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_metric(slide, "2,864,117", "tested gene × comparison rows", 5.06, 1.56, 2.26, accent=BLUE)
    add_metric(slide, "118,297", "paper-DEG memberships", 7.50, 1.56, 2.26, accent=BLUE, note="58,112 up • 60,185 down")
    add_metric(slide, "9,262", "core-MT query memberships", 9.94, 1.56, 2.26, accent=GOLD, note="4,258 up • 5,004 down")
    add_rect(slide, 5.06, 3.30, 7.14, 2.62, color=PALE_GREEN, outline=TEAL)
    add_text(slide, "NETWORK-BACKGROUND INTERSECTION", 5.36, 3.60, 3.4, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(slide, "9,262 provisional query memberships", 5.38, 4.04, 6.35, 0.34, size=17.5, color=NAVY, bold=True)
    add_text(slide, "− 1,329 absent from the run-specific induced background", 5.38, 4.52, 5.82, 0.30, size=13.0, color=VERMILION, bold=True)
    add_text(slide, "= 7,933 effective-query memberships", 5.38, 5.06, 6.35, 0.34, size=18, color=TEAL, bold=True)
    add_text(slide, "Membership counts can repeat the same gene across contrasts and directions.", 5.39, 5.55, 5.93, 0.25, size=10.4, color=GRAY)
    add_takeaway(slide, "Query gene ≠ candidate driver: query genes seed enrichment; candidate network genes are tested.", accent=BLUE)
    finalize_slide(
        slide, slide_no=7,
        source="Phase 08 paper-DEG and core-MT audit tables; Phase 12 query-membership tables",
        goal="Explain how the mitochondrial KDA query is derived from the full fine-cell DEG analysis.",
        walkthrough="The paper-DEG filters first act within each fine-cell by group comparison. Only passing core mitochondrial proteins in the requested sign become provisional query memberships. The matching broad network is induced on genes tested in that contrast, so out-of-background query memberships are removed before run eligibility is assessed.",
        boundary="These are membership counts, not unique-gene counts. A query membership also does not mean the gene will be a candidate driver.",
        transition="Next distinguish source failure, an empty effective query, a one- or two-gene query, and the canonical three-gene execution floor.",
    )
    register(slide, title, "Fine-cell")

    # 8 — fine eligibility detail
    title = "Step 1B: source validity and query size decide whether KDA is called and included"
    slide = content_slide(prs, title=title, kicker="Fine funnel • run eligibility", slide_no=8, accent=BLUE,
                          subtitle="Each status answers a different question; none should be relabeled as ‘no biology.’")
    add_rect(slide, 0.66, 1.47, 5.10, 4.90, color=WHITE, outline=LIGHT)
    add_panel_title(slide, "One canonical query-size gate", 0.95, 1.77, 4.48, accent=BLUE)
    add_metric(slide, "≥3", "execution and Phase 20 inclusion", 0.95, 2.31, 2.02, accent=BLUE, bg=PALE_BLUE)
    add_metric(slide, count_text(fine["effective_query_three_nine"]), "included runs with query n=3–9", 3.20, 2.31, 2.02, accent=TEAL, bg=PALE_GREEN)
    add_bullets(slide, [
        f"All {count_text(fine['included_runs'])} validated completed calls enter the aggregate.",
        f"That includes {count_text(fine['effective_query_three_nine'])} size-3–9 calls, retained with Phase 12's small-query warning, and {count_text(fine['effective_query_ten_plus'])} of size ≥10.",
        f"{count_text(fine['completed_empty'])} completed calls returned no original q≤0.05 gene; they still contribute complete evidence.",
    ], 0.96, 3.82, 4.42, size=11.1, line_h=0.66, accent=BLUE)
    add_rect(slide, 6.05, 1.47, 6.60, 4.90, color=PALE_RED, outline=VERMILION)
    add_panel_title(slide, "What ‘source contrast not estimable’ means", 6.37, 1.77, 5.95, accent=VERMILION)
    add_text(slide, "The upstream MAST comparison could not be fitted because one arm had fewer than three cells.", 6.39, 2.35, 5.83, 0.62, size=14.0, color=NAVY, bold=True)
    add_table(slide, ["Fine cell type", "Group", "AD cells", "NCI cells"], [
        ["Fib SLC4A4", "M_e2", "0", "11"],
        ["CAMs", "M_e2", "45", "1"],
        ["Mic MKI67", "M_e2", "2", "9"],
    ], 6.39, 3.17, [2.55, 1.00, 1.15, 1.15], row_h=0.50, header_h=0.47,
       accent=VERMILION, font_size=9.8)
    add_text(slide, "3 unavailable comparisons × 2 directions = 6 unavailable slots", 6.52, 5.34, 5.50, 0.56, size=13.5, color=VERMILION, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "This is not a completed contrast with zero DEGs.", 6.52, 5.96, 5.50, 0.25, size=10.2, color=GRAY, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Never collapse ‘not estimable,’ ‘no query,’ ‘too-small query,’ and ‘empty KDA return’ into one status.", accent=VERMILION)
    finalize_slide(
        slide, slide_no=8,
        source="phase20_source_run_manifest.tsv; phase20_source_status.tsv; source-DEG audit",
        goal="Clarify the eligibility statuses and the canonical minimum of three effective query genes.",
        walkthrough=f"A query of at least three genes allows the Phase 12 call and Phase 20 inclusion. The {count_text(fine['effective_query_three_nine'])} calls with three through nine genes retain Phase 12's below-ten warning but are included alongside the {count_text(fine['effective_query_ten_plus'])} larger-query calls. The right panel gives the three concrete non-estimable comparisons and shows why they create six unavailable direction slots.",
        boundary="Small-query calls use the same KDA and support predicates as larger-query calls, but their smaller queries warrant cautious interpretation. Non-estimability is a data/model condition, not evidence against a biological effect.",
        transition=f"With {count_text(fine['included_runs'])} runs fixed, inspect what one KDA call actually tests.",
    )
    register(slide, title, "Fine-cell")

    # 9 — one KDA call
    title = "Step 2: KDA tests network genes—not only DEGs or mitochondrial genes"
    slide = content_slide(prs, title=title, kicker="Fine funnel • KDA", slide_no=9, accent=BLUE,
                          subtitle="The mitochondrial DEG set is the query; each potential key driver is one network gene.")
    steps = [
        ("1", "Candidate search", "Within 3 undirected hops of the query"),
        ("2", "Directed layers", "Downstream neighborhoods at layers 1–3"),
        ("3", "Enrichment test", "Upper-tail hypergeometric P at each layer"),
        ("4", "Best evidence", "Retain best layer and raw P per gene"),
        ("5", "Stock within-run BH", "original_run_q defines returned subset"),
    ]
    for index, (number, heading, detail) in enumerate(steps):
        x = 0.69 + index * 2.47
        add_rect(slide, x, 1.54, 2.18, 2.16, color=WHITE, outline=LIGHT)
        add_circle(slide, x + 0.16, 1.75, 0.43, BLUE)
        add_text(slide, number, x + 0.16, 1.84, 0.43, 0.18, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, heading, x + 0.18, 2.31, 1.83, 0.37, size=14.1, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x + 0.18, 2.83, 1.83, 0.58, size=10.1, color=GRAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0.69, 4.05, 5.65, 2.17, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "What a candidate can be", 0.98, 4.36, 5.04, accent=BLUE)
    add_bullets(slide, [
        "Need not be a DEG.", "Need not be mitochondrial.",
        "Need not be part of the query.", "Is a network node—not a gene set.",
    ], 1.01, 4.90, 4.82, size=11.8, line_h=0.30, accent=BLUE)
    add_rect(slide, 6.61, 4.05, 5.99, 2.17, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Included-run call outcomes", 6.91, 4.36, 5.37, accent=TEAL)
    add_metric(slide, count_text(fine["explicit_tests"]), "explicit gene × run tests", 6.92, 4.86, 1.65, accent=TEAL, bg=WHITE)
    add_metric(slide, count_text(fine["stock_q05_rows"]), "published original q≤0.05 rows", 8.74, 4.86, 1.65, accent=TEAL, bg=WHITE)
    add_metric(slide, count_text(fine["completed_empty"]), f"of {count_text(fine['included_runs'])} calls returned none", 10.56, 4.86, 1.65, accent=TEAL, bg=WHITE)
    add_takeaway(slide, "Published original_run_q≤0.05 returns are an output subset; the Phase 20 source reconstructs the full pre-significance family.", accent=TEAL)
    finalize_slide(
        slide, slide_no=9,
        source="phase20_source_candidate_tests.tsv.gz; validated Phase 12 KDA returns",
        goal="Define a candidate gene and show how enrichment is calculated within one KDA run.",
        walkthrough="The query is a set of directional core-MT DEGs. Potential driver genes are identified near that query in the network, their directed downstream neighborhoods are tested, the best layer is kept, and stock BH creates original_run_q across the explicit family.",
        boundary="Topology annotations such as root status or global key-driver labels are not Phase 20 selection gates, and a returned gene is not automatically a final cross-run candidate.",
        transition="The next slide explains why Phase 20 needs nonsignificant and implicit evidence as well as stock returns.",
    )
    register(slide, title, "Fine-cell")

    # 10 — complete evidence
    title = "Stock returns are only a subset of the complete evidence family"
    slide = content_slide(prs, title=title, kicker="Fine funnel • complete evidence", slide_no=10, accent=BLUE,
                          subtitle=f"Phase 20 evaluates every assessable network gene across the {count_text(fine['included_runs'])} included runs.")
    add_metric(slide, count_text(fine["explicit_tests"]), "explicit final_raw_p tests", 0.72, 1.50, 2.53, accent=BLUE)
    add_metric(slide, count_text(fine["implicit_rows"]), "implicit P = 1 rows", 3.48, 1.50, 2.53, accent=GOLD)
    add_metric(slide, count_text(fine["absent_rows"]), "absent-background rows", 6.24, 1.50, 2.53, accent=PURPLE)
    add_metric(slide, count_text(fine["opportunities"]), "gene × run opportunities", 9.00, 1.50, 3.52, accent=TEAL)
    add_table(slide, ["Gene state in one run", "Statistical treatment", "Why"], [
            ["Explicitly tested", "Use reconstructed final_raw_p", "Candidate entered the run's explicit family"],
        ["In background, not explicit", "Implicit P = 1", "Zero query overlap outside explicit family"],
        ["Absent from background", "Missing; omit from ACAT", "Gene was not assessable in that run"],
    ], 0.78, 3.18, [3.40, 2.80, 5.55], row_h=0.68, header_h=0.51, accent=BLUE, font_size=10.5)
    add_rect(slide, 0.80, 5.88, 11.74, 0.58, color=PALE_RED, outline=VERMILION)
    add_text(slide, f"The {count_text(fine['stock_q05_rows'])} significant published rows are not the starting table for the cross-run analysis.", 1.04, 6.04, 11.28, 0.24, size=13.1, color=VERMILION, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Implicit P=1 means assessable with zero overlap; missing means not assessable. They are not interchangeable.", accent=PURPLE)
    finalize_slide(
        slide, slide_no=10,
        source="phase20_source_candidate_tests.tsv.gz; phase20_source_checks.tsv",
        goal="Show why reconstruction includes explicit, implicit-zero, and missing evidence states.",
        walkthrough="Explicit tests contribute reconstructed final_raw_p values validated against the published Phase 12 returns. A gene in the background but outside the explicit family contributes P equals one because it has zero overlap. A gene absent from the induced background is missing and is omitted from ACAT while reducing coverage.",
        boundary="Treating missing rows as P equals one would confound lack of network coverage with evidence against enrichment. Using only returned genes would create selection bias.",
        transition="Now remove core-MT candidate drivers and group the retained evidence into category-specific units.",
    )
    register(slide, title, "Fine-cell")

    # 11 — grouping/non-MT
    title = "Step 3: evidence is grouped by gene × sex/APOE × broad network"
    slide = content_slide(prs, title=title, kicker="Fine funnel • aggregation unit", slide_no=11, accent=BLUE,
                          subtitle="Eligible fine-cell runs and either or both query directions can contribute within one category.")
    add_rect(slide, 0.67, 1.48, 5.25, 4.95, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "First restrict candidate drivers", 0.96, 1.78, 4.64, accent=BLUE)
    add_text(slide, count_text(fine["opportunities"]), 0.98, 2.38, 2.05, 0.47, size=28, color=NAVY, bold=True, font=FONT_HEAD)
    add_text(slide, "all-class gene × run rows", 2.86, 2.52, 2.27, 0.25, size=11.1, color=GRAY, bold=True)
    add_text(slide, f"− {count_text(fine['mt_rows'])} core-MT candidate-driver rows", 1.00, 3.20, 4.44, 0.34, size=15.2, color=VERMILION, bold=True)
    add_text(slide, f"= {count_text(fine['non_mt_rows'])} non-MT gene × run rows", 1.00, 3.82, 4.44, 0.38, size=18.2, color=BLUE, bold=True)
    add_text(slide, "Source taxonomy: case_id = non_mt_driver and is_core_mito = FALSE.", 1.00, 4.55, 4.40, 0.58, size=11.5, color=GRAY)
    add_rect(slide, 1.00, 5.40, 4.40, 0.60, color=WHITE, outline=LIGHT)
    add_text(slide, "Core-MT DEGs can seed queries; all core-MT genes are excluded as final non-MT drivers.", 1.20, 5.52, 4.00, 0.34, size=10.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.18, 1.48, 6.46, 4.95, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Then define the category unit", 6.49, 1.78, 5.84, accent=TEAL)
    add_rect(slide, 6.50, 2.32, 5.80, 1.12, color=WHITE, outline=TEAL)
    add_text(slide, "current_symbol + signature_group + broad_network", 6.73, 2.55, 5.34, 0.28, size=14.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "gene + sex/APOE + broad cell type", 6.73, 2.93, 5.34, 0.24, size=11.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_metric(slide, count_text(fine["aggregate_rows"]), "non-MT gene × category units", 6.50, 3.80, 1.76, accent=TEAL, bg=WHITE)
    add_metric(slide, count_text(fine["aggregate_genes"]), "distinct gene symbols", 8.48, 3.80, 1.76, accent=TEAL, bg=WHITE)
    add_metric(slide, f"{fine['analyzable_categories']} / {fine['structural_categories']}", "categories with ≥1 run", 10.46, 3.80, 1.76, accent=TEAL, bg=WHITE)
    add_bullets(slide, [
        "Never combines different sex/APOE groups.",
        "Never combines different broad networks.",
        "Can combine one or more fine types and either or both directions.",
        f"{fine['empty_categories']} structural categories have no included run; they are not estimable.",
    ], 6.55, 5.13, 5.50, size=10.0, line_h=0.29, accent=TEAL)
    add_takeaway(slide, "One gene may generate several final candidate units because its group and broad-network contexts differ.", accent=TEAL)
    finalize_slide(
        slide, slide_no=11,
        source="phase20_source_candidate_tests.tsv.gz; phase20_driver_aggregates.tsv.gz",
        goal="Define the final fine-cell inferential unit and the non-core-MT candidate restriction.",
        walkthrough="Core mitochondrial DEGs can seed a query, but all core mitochondrial genes are removed from the final candidate-driver universe. The remaining rows are grouped by current gene symbol, sex/APOE group, and broad network. One or more eligible fine-cell runs and either or both directions can contribute within a fixed category.",
        boundary="A gene appearing in two categories represents two context-specific units, not an automatic replication or a formal interaction test.",
        transition="Before combining P values, require adequate evidence availability through the coverage gate.",
    )
    register(slide, title, "Fine-cell")

    # 12 — coverage
    title = "Step 4: coverage asks whether a gene is assessable across eligible runs"
    slide = content_slide(prs, title=title, kicker="Fine funnel • coverage", slide_no=12, accent=BLUE,
                          subtitle="Coverage measures availability of usable evidence—not the fraction of significant runs.")
    add_rect(slide, 0.71, 1.53, 4.24, 2.27, color=PALE_BLUE, outline=BLUE)
    add_text(slide, "coverage = usable runs / eligible runs", 1.00, 1.91, 3.66, 0.44, size=14.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(slide, "usable = explicit raw P or implicit P=1", 1.00, 2.55, 3.66, 0.28, size=12.1, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "absent-background rows are not usable", 1.00, 3.02, 3.66, 0.28, size=11.2, color=GRAY, align=PP_ALIGN.CENTER)
    add_metric(slide, count_text(fine["aggregate_rows"]), "input gene × category units", 5.30, 1.57, 2.20, accent=BLUE)
    add_metric(slide, f"−{count_text(fine['zero_usable_units'])}", "zero usable runs", 7.73, 1.57, 2.20, accent=VERMILION)
    add_metric(slide, f"−{count_text(fine['partial_below_relaxed'])}", "some evidence, coverage <0.50", 10.16, 1.57, 2.20, accent=VERMILION)
    add_rect(slide, 5.30, 3.20, 7.06, 0.86, color=PALE_GREEN, outline=TEAL)
    add_text(slide, f"{count_text(fine['relaxed_coverage_units'])} units / {count_text(fine['relaxed_coverage_genes'])} genes pass relaxed coverage ≥ 0.50", 5.55, 3.43, 6.56, 0.34, size=14.2, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.71, 4.24, 11.65, 1.83, color=WHITE, outline=LIGHT)
    add_panel_title(slide, "Coverage does not mean", 1.00, 4.54, 2.65, accent=VERMILION)
    add_bullets(slide, [
        "percentage of significant runs", "percentage of stock-returned runs",
        "percentage of supporting fine cell types", "recurrence across categories",
    ], 3.82, 4.54, 7.90, size=12.3, line_h=0.34, accent=VERMILION)
    add_takeaway(slide, "Coverage protects ACAT from being interpreted when too much of a gene's eligible run evidence is missing.", accent=BLUE)
    finalize_slide(
        slide, slide_no=12,
        source="phase20_driver_aggregates.tsv.gz; phase20_filter_funnel.tsv",
        goal="Explain why a separate availability gate is needed before cross-run combination.",
        walkthrough="Coverage divides runs with usable explicit or implicit evidence by all eligible runs for the gene-category unit. Missing-background rows reduce coverage. The relaxed threshold retains units with at least half of their eligible run evidence available.",
        boundary="Coverage is not a success rate. A unit can have complete coverage and no supporting run, or limited coverage despite a strong result in one available run.",
        transition="Among coverage-qualified units, require at least one independently supported run.",
    )
    register(slide, title, "Fine-cell")

    # 13 — supporting-run gate
    title = "Step 5: a supported unit needs at least one individually convincing run"
    slide = content_slide(prs, title=title, kicker="Fine funnel • support", slide_no=13, accent=BLUE,
                          subtitle="ACAT significance alone is not sufficient for the primary candidate list.")
    add_rect(slide, 0.68, 1.50, 4.21, 4.89, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "Relaxed supporting-run rule", 0.97, 1.82, 3.61, accent=BLUE)
    rules = [
        ("1", "other_query_overlap ≥ 2"),
        ("2", "fold enrichment > 1"),
        ("3", "final_run_q ≤ 0.10"),
    ]
    for index, (number, rule_text) in enumerate(rules):
        y = 2.52 + index * 0.85
        add_circle(slide, 1.02, y, 0.42, BLUE)
        add_text(slide, number, 1.02, y + 0.09, 0.42, 0.17, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.60, y - 0.02, 2.76, 0.49, color=WHITE, outline=LIGHT)
        add_text(slide, rule_text, 1.79, y + 0.08, 2.38, 0.24, size=12.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "All three must hold in the same run.", 1.05, 5.34, 3.45, 0.26, size=11.6, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.22, 1.50, 7.41, 4.89, color=WHITE, outline=LIGHT)
    add_panel_title(slide, "Non-MT explicit-test sub-funnel", 5.53, 1.82, 6.79, accent=TEAL)
    funnel_nodes = [
        (count_text(fine["non_mt_explicit"]), "explicit tests"),
        (count_text(fine["positive_overlap"]), "positive overlap"),
        (count_text(fine["overlap_ge2"]), "other overlap ≥2"),
        (count_text(fine["overlap_ge2_fold_gt1"]), "+ FE >1"),
        (count_text(fine["relaxed_support_events"]), "+ final q≤0.10"),
    ]
    add_flow_row(slide, funnel_nodes, 2.48, accent=TEAL, start_x=5.48, total_w=6.84, node_h=1.08)
    add_rect(slide, 5.50, 4.15, 6.80, 1.64, color=PALE_GREEN, outline=TEAL)
    add_text(slide, f"{count_text(fine['relaxed_support_events'])} supporting gene × run events", 5.77, 4.35, 6.25, 0.30, size=16.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"→ {count_text(fine['support_precoverage_units'])} gene × category units → {count_text(fine['relaxed_support_units'])} after coverage ≥0.50", 5.77, 4.74, 6.25, 0.31, size=13.4, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"{count_text(fine['relaxed_support_units'])} units represent {count_text(fine['relaxed_support_genes'])} genes across {fine['relaxed_support_categories']} categories", 5.77, 5.13, 6.25, 0.27, size=11.3, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"Strict support uses final_run_q≤0.05 and yields {count_text(fine['strict_support_events'])} supporting events.", 5.77, 5.48, 6.25, 0.24, size=9.3, color=GRAY, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "The support gate supplies a concrete run-level anchor for every retained aggregate unit.", accent=TEAL)
    finalize_slide(
        slide, slide_no=13,
        source="phase20_source_candidate_tests.tsv.gz; phase20_driver_aggregates.tsv.gz",
        goal="Explain the individual-run evidence requirement that accompanies cross-run ACAT.",
        walkthrough=f"The other_query_overlap, fold-enrichment, and final_run_q gates must all pass in one explicit run. The sub-funnel starts from all non-MT explicit tests and ends with {count_text(fine['relaxed_support_events'])} supporting gene-by-run events, which collapse to {count_text(fine['support_precoverage_units'])} category units before coverage removes one.",
        boundary="The support field is reconstructed final_run_q from each run's full explicit family and validated against Phase 12 published returns. It is distinct from original_run_q and from the category q obtained after ACAT and category-specific BH correction.",
        transition="Now combine all usable raw P values and explain how the category q is formed.",
    )
    register(slide, title, "Fine-cell")

    # 14 — ACAT/BH
    title = "Step 6A: ACAT combines raw P values; category BH controls FDR within category"
    slide = content_slide(prs, title=title, kicker="Fine funnel • ACAT", slide_no=14, accent=BLUE,
                          subtitle="The shorthand ‘ACAT q’ contains two separate operations.")
    add_rect(slide, 0.67, 1.47, 7.17, 4.98, color=PALE_BLUE, outline=BLUE)
    stages = [
        ("Run-level raw P values", "explicit final_raw_p or implicit P=1"),
        ("Equal-weight ACAT", "missing backgrounds omitted"),
        ("category_acat_p", "one P per gene × category"),
        ("BH within category", "all coverage-qualified genes"),
        ("category q", "relaxed_category_acat_q"),
    ]
    for index, (heading, detail) in enumerate(stages):
        y = 1.78 + index * 0.86
        add_circle(slide, 0.98, y + 0.04, 0.38, BLUE if index < 3 else GOLD)
        add_text(slide, str(index + 1), 0.98, y + 0.12, 0.38, 0.16, size=8.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.52, y, 5.92, 0.52, color=WHITE, outline=LIGHT)
        add_text(slide, heading, 1.73, y + 0.09, 2.57, 0.25, size=12.8, color=NAVY, bold=True)
        add_text(slide, detail, 4.30, y + 0.10, 2.91, 0.24, size=10.3, color=GRAY, align=PP_ALIGN.RIGHT)
    add_rect(slide, 8.10, 1.47, 4.53, 3.17, color=WHITE, outline=LIGHT)
    add_panel_title(slide, "Example unit", 8.39, 1.78, 3.95, accent=TEAL)
    add_text(slide, "Five eligible runs", 8.42, 2.35, 3.87, 0.28, size=14.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "P = .001, .40, 1, 1, NA", 8.42, 2.83, 3.87, 0.35, size=19, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "coverage = 4/5 = 0.80", 8.42, 3.40, 3.87, 0.28, size=13.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "ACAT combines [.001, .40, 1, 1]; NA is omitted", 8.55, 3.86, 3.61, 0.49, size=10.6, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.10, 4.88, 4.53, 1.57, color=PALE_RED, outline=VERMILION)
    add_text(slide, "BH denominator", 8.38, 5.18, 1.58, 0.24, size=10.2, color=VERMILION, bold=True)
    add_text(slide, f"{count_text(fine['relaxed_coverage_units'])} units are partitioned into {fine['analyzable_categories']} category-specific BH families; each uses its own coverage-qualified genes.", 8.38, 5.49, 3.95, 0.70, size=10.4, color=NAVY, bold=True)
    add_takeaway(slide, "ACAT combines evidence; BH supplies the within-category multiple-testing correction across genes.", accent=GOLD)
    finalize_slide(
        slide, slide_no=14,
        source="phase20_methods.md; phase20_driver_aggregates.tsv.gz; phase20_filter_funnel.tsv",
        goal="Separate ACAT combination from the subsequent BH correction that creates the category q value.",
        walkthrough=f"For each gene-category unit, ACAT combines all usable final_raw_p values, including implicit P equals one rows. Missing backgrounds are omitted. The {count_text(fine['relaxed_coverage_units'])} coverage-qualified units are partitioned into {fine['analyzable_categories']} category-specific BH families, and BH is applied across the genes within each category.",
        boundary="ACAT does not combine run-level q values. The BH family is not restricted to genes that already passed the support gate.",
        transition="Apply coverage, support, and category-q thresholds to define relaxed, strict, and exploratory outputs.",
    )
    register(slide, title, "Fine-cell")

    # 15 — threshold tiers
    title = "Step 6B: relaxed, strict, and exploratory tiers answer different questions"
    slide = content_slide(prs, title=title, kicker="Fine funnel • final gates", slide_no=15, accent=BLUE,
                          subtitle="The discovery-oriented main list uses less stringent prespecified thresholds than the strict reference.")
    add_table(slide, ["Gate / output", "Relaxed main", "Strict reference", "Exploratory-only"], [
        ["Coverage", "≥ 0.50", "≥ 0.80", "≥ 0.50"],
        ["Supporting-run q", "≤ 0.10", "≤ 0.05", "≤ 0.10"],
        ["Category q", "≤ 0.10", "≤ 0.05", "0.10 < q ≤ 0.20"],
        ["Candidate units", count_text(fine["relaxed_candidates"]), count_text(fine["strict_candidates"]), f"{count_text(fine['exploratory_candidates'])} additional"],
        ["Distinct genes", count_text(fine["relaxed_genes"]), count_text(fine["strict_genes"]), count_text(fine["exploratory_genes"])],
        ["Categories", count_text(fine["relaxed_categories"]), count_text(fine["strict_categories"]), count_text(fine["exploratory_categories"])],
    ], 0.78, 1.51, [3.25, 2.68, 2.68, 2.68], row_h=0.53, header_h=0.50,
       accent=BLUE, font_size=10.6)
    add_rect(slide, 0.79, 5.28, 3.73, 1.13, color=PALE_BLUE, outline=BLUE)
    add_text(slide, "RELAXED FUNNEL", 1.03, 5.52, 1.40, 0.20, size=8.6, color=BLUE, bold=True)
    add_text(slide, f"{count_text(fine['aggregate_rows'])} → {count_text(fine['relaxed_coverage_units'])} → {count_text(fine['relaxed_support_units'])} → {count_text(fine['relaxed_candidates'])}", 1.03, 5.88, 3.25, 0.26, size=15.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.80, 5.28, 3.73, 1.13, color=PALE_GREEN, outline=TEAL)
    add_text(slide, "STRICT FUNNEL", 5.04, 5.52, 1.35, 0.20, size=8.6, color=TEAL, bold=True)
    add_text(slide, f"{count_text(fine['aggregate_rows'])} → {count_text(fine['strict_coverage_units'])} → {count_text(fine['strict_coverage_support_units'])} → {count_text(fine['strict_candidates'])}", 5.04, 5.88, 3.25, 0.26, size=15.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.81, 5.28, 3.73, 1.13, color=PALE_GOLD, outline=GOLD)
    add_text(slide, "DISPLAY FLAGS", 9.05, 5.52, 1.35, 0.20, size=8.6, color=GOLD, bold=True)
    add_text(slide, "Top 5 / top 10 and stability labels do not filter candidates", 9.05, 5.78, 3.25, 0.47, size=10.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "A final fine-cell candidate is a non-core-MT gene × category unit—not merely one returned KDA gene.", accent=BLUE)
    finalize_slide(
        slide, slide_no=15,
        source="phase20_filter_funnel.tsv; phase20_status.tsv; phase20_driver_aggregates.tsv.gz",
        goal="Define the main, strict-reference, and exploratory thresholds and their yields.",
        walkthrough=f"The relaxed and strict funnels use parallel coverage, supporting-run, and category-q thresholds. All {fine['strict_candidates']} strict units are contained within the {fine['relaxed_candidates']} relaxed units. Exploratory-only leads pass the relaxed availability and support gates but have category q between 0.10 and 0.20.",
        boundary=f"Top-five and top-ten flags change presentation only. Stability labels also characterize evidence but do not remove rows from the {fine['relaxed_candidates']}-candidate main list.",
        transition=f"Summarize what the {fine['relaxed_candidates']} final units represent before showing the two requested figures.",
    )
    register(slide, title, "Fine-cell")

    # 16 — fine output summary
    title = f"Fine-cell output: {fine['relaxed_candidates']} candidate units represent {fine['relaxed_genes']} distinct genes"
    slide = content_slide(prs, title=title, kicker="Fine results • candidate genes", slide_no=16, accent=BLUE,
                          subtitle="Candidate yield is uneven across groups and networks; it should not be read as a formal interaction test.")
    add_metric(slide, count_text(fine["relaxed_candidates"]), "relaxed gene × category units", 0.72, 1.52, 2.12, accent=BLUE)
    add_metric(slide, count_text(fine["relaxed_genes"]), "distinct gene symbols", 3.06, 1.52, 2.12, accent=BLUE)
    add_metric(slide, count_text(fine["relaxed_categories"]), "categories with candidates", 5.40, 1.52, 2.12, accent=BLUE)
    add_metric(slide, count_text(fine["strict_candidates"]), "strict candidate units", 7.74, 1.52, 2.12, accent=TEAL)
    add_metric(slide, count_text(fine["top5_rows"]), "top-five displayed rows", 10.08, 1.52, 2.12, accent=GOLD)
    add_rect(slide, 0.72, 3.14, 5.76, 3.10, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "Where candidate units occur", 1.02, 3.46, 5.16, accent=BLUE)
    add_text(slide, "By sex/APOE group", 1.03, 4.03, 2.05, 0.24, size=10.2, color=BLUE, bold=True)
    group_summary = "  •  ".join(
        f"{group} {fine['candidate_groups'][group]}"
        for group in ("M_e2", "F_e2", "F_e4", "F_e33", "M_e33", "M_e4")
    )
    add_text(slide, group_summary, 1.03, 4.40, 4.97, 0.60, size=11.7, color=NAVY, bold=True)
    add_text(slide, "By broad network", 1.03, 5.16, 2.05, 0.24, size=10.2, color=BLUE, bold=True)
    network_summary = (
        f"Excitatory {fine['candidate_networks']['Excitatory_neurons']} • "
        f"Astrocytes {fine['candidate_networks']['Astrocytes']} • "
        f"OPCs {fine['candidate_networks']['OPCs']} • "
        f"Inhibitory {fine['candidate_networks']['Inhibitory_neurons']}\n"
        f"Oligo {fine['candidate_networks']['Oligodendrocytes']} • "
        f"Microglia {fine['candidate_networks']['Microglia']} • "
        f"Vasculature {fine['candidate_networks']['Vasculature_cells']}"
    )
    add_text(slide, network_summary, 1.03, 5.48, 4.97, 0.62, size=10.3, color=NAVY, bold=True)
    add_rect(slide, 6.76, 3.14, 5.48, 3.10, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Selected recurrent genes", 7.06, 3.46, 4.88, accent=TEAL)
    recurrence = facts["fine_recurrence"]
    recurring = [(gene, recurrence[gene]) for gene in ["RPL11", "RPS15", "RPLP1", "RPL15", "SELENOW", "RPS13"]]
    maximum_recurrence = max(recurrence.values())
    for index, (gene, count) in enumerate(recurring):
        y = 4.02 + index * 0.32
        add_text(slide, gene, 7.08, y, 1.42, 0.22, size=11.2, color=NAVY, bold=True)
        add_rect(slide, 8.63, y + 0.03, 2.62 * count / maximum_recurrence, 0.16, color=TEAL, outline=None, radius=False)
        add_text(slide, str(count), 11.42, y, 0.35, 0.22, size=10.4, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
    tied_at_three = sum(value == 3 for value in recurrence.values())
    add_text(slide, f"Examples shown; {tied_at_three} genes tie at 3 categories. Counts are categories, not runs.", 7.08, 5.99, 4.74, 0.24, size=8.7, color=GRAY, italic=True)
    add_takeaway(slide, "Unequal run availability and category size can affect yield; absence from a list is not a sex/APOE interaction.", accent=VERMILION)
    finalize_slide(
        slide, slide_no=16,
        source="phase20_relaxed_candidates.tsv; phase20_category_manifest.tsv; phase20_top5_summary.tsv",
        goal=f"Translate the {fine['relaxed_candidates']} candidate units into distinct genes, categories, strict support, and display counts.",
        walkthrough=f"A single gene can occupy more than one category, so {fine['relaxed_candidates']} units reduce to {fine['relaxed_genes']} unique symbols. The visible distributions show where candidate units occur, while the recurrence preview shows which genes appear across the largest number of final categories.",
        boundary="Category counts are descriptive and are influenced by unequal included-run availability. They are not tests that one group or network has more biology than another.",
        transition="The recurrence figure now displays the top 20 genes by category count.",
    )
    register(slide, title, "Fine-cell")

    # 17 — fine recurrence figure
    rpl11 = fine["recurrence_details"]["RPL11"]
    rps15 = fine["recurrence_details"]["RPS15"]
    title = f"RPL11 recurs across {rpl11['categories']} supported categories"
    slide = content_slide(prs, title=title, kicker="Fine results • driver recurrence", slide_no=17, accent=BLUE,
                          subtitle="Recurrence is descriptive across final sex/APOE × broad-network categories.")
    add_rect(slide, 0.47, 1.27, 7.16, 5.40, color=WHITE, outline=LIGHT)
    add_picture_contain(slide, FIG["fine_recurrence"], 0.58, 1.37, 6.94, 5.18,
                        alt="Fine-cell Phase 20 horizontal bar chart of the twenty most recurrent relaxed non-MT key drivers across supported categories")
    add_rect(slide, 7.91, 1.42, 4.75, 4.98, color=PALE_BLUE, outline=BLUE)
    add_panel_title(slide, "How to read it", 8.22, 1.74, 4.13, accent=BLUE)
    add_bullets(slide, [
        "Bar length = number of final categories containing the gene.",
        "Fill = how many of those categories also pass the strict reference.",
        f"RPL11: {rpl11['categories']} categories, {rpl11['groups']} groups, {rpl11['networks']} networks; strict in {rpl11['strict']}.",
        f"RPS15: {rps15['categories']} categories, {rps15['groups']} groups, {rps15['networks']} networks; strict in {rps15['strict']}.",
        "The plot shows the top 20 after category count, best q, then symbol ordering.",
    ], 8.23, 2.28, 4.02, size=11.3, line_h=0.61, accent=BLUE)
    add_rect(slide, 8.22, 5.55, 4.10, 0.56, color=WHITE, outline=LIGHT)
    add_text(slide, "Not a formal consistency or sex/APOE-difference test.", 8.41, 5.69, 3.72, 0.27, size=10.1, color=VERMILION, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Recurrence summarizes category presence after all coverage, support, ACAT, and category-q gates.", accent=BLUE)
    finalize_slide(
        slide, slide_no=17,
        source="results/figures/analysis/phase_20_sex_apoe_kda/driver_recurrence figure bundle",
        goal="Interpret the requested fine-cell driver-recurrence figure without confusing categories and runs.",
        walkthrough=f"The bar reports how many final sex/APOE by broad-network categories contain a relaxed candidate unit for the gene. The fill records strict-category count. RPL11 appears in {rpl11['categories']} final categories and RPS15 in {rps15['categories']}; the next tier contains genes recurring in three categories.",
        boundary="The chart does not count significant fine-cell runs and does not test heterogeneity or consistency across sex/APOE groups.",
        transition="The next requested figure shows the ranked presentation subset within supported categories.",
    )
    register(slide, title, "Fine-cell")

    # 18 — fine top5 figure
    title = "Top-five lists retain up to five candidates per supported category"
    slide = content_slide(prs, title=title, kicker="Fine results • top-five candidates", slide_no=18, accent=BLUE)
    add_rect(slide, 0.47, 1.00, 8.22, 6.05, color=WHITE, outline=LIGHT)
    add_picture_contain(slide, FIG["fine_top5"], 0.58, 1.09, 8.00, 5.86,
                        alt="Fine-cell Phase 20 tile chart of up to five relaxed non-MT key drivers per supported sex/APOE and broad-cell category")
    add_rect(slide, 8.94, 1.09, 3.72, 5.74, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "Display contract", 9.23, 1.40, 3.13, accent=GOLD)
    add_metric(slide, count_text(fine["top5_rows"]), "displayed candidate units", 9.23, 1.90, 1.33, accent=GOLD, bg=WHITE)
    add_metric(slide, count_text(fine["relaxed_categories"]), "represented categories", 10.77, 1.90, 1.33, accent=GOLD, bg=WHITE)
    add_bullets(slide, [
        "Order: category q, then ACAT P, then symbol.",
        f"{fine['top5_strict_rows']} displayed rows are strict; {fine['top5_rows'] - fine['top5_strict_rows']} are relaxed-only.",
        f"Only {fine['relaxed_categories']} candidate-bearing categories are shown.",
        f"Omitted: {fine['analyzable_no_candidate']} analyzable with no candidate + {fine['empty_categories']} with no included run.",
        f"Up to five per row; no backfill and no change to the {fine['relaxed_candidates']} candidates.",
    ], 9.23, 3.28, 2.97, size=10.0, line_h=0.53, accent=GOLD)
    add_text(slide, "Presentation view only—no additional statistical filter.", 9.23, 6.24, 3.00, 0.33, size=9.6, color=VERMILION_TEXT, bold=True, align=PP_ALIGN.CENTER)
    finalize_slide(
        slide, slide_no=18,
        source="results/figures/analysis/phase_20_sex_apoe_kda/top5_candidates figure bundle; phase20_top5_summary.tsv",
        goal="Interpret the requested fine-cell top-five tile chart and its ranking rules.",
        walkthrough=f"Each row is one supported sex/APOE by broad-network category. Candidates are ordered using stored relaxed rank, which follows category q, ACAT P, and gene symbol. At most five are displayed, producing {fine['top5_rows']} tiles across {fine['relaxed_categories']} categories.",
        boundary=f"A blank tile is not backfilled with a failing gene. The plot shows {fine['relaxed_categories']} candidate-bearing categories; it omits {fine['analyzable_no_candidate']} analyzable categories with no candidate and {fine['empty_categories']} categories with no included run.",
        transition="Part II repeats the setup and funnel explanation for the direct broad-cell branch.",
    )
    register(slide, title, "Fine-cell")


if __name__ == "__main__":
    main()
