#!/usr/bin/env python3
"""Insert the dataset slide after slide 6 of the genetic-support deck.

The deck has manual edits (a deleted slide and reordered slides), so it must
not be rebuilt from ``build_genetic_support_simple_aggr_deck.py``. This script
surgically appends one slide listing the public datasets of the frozen screen,
categorized by the three evidence sources defined on slide 5, and moves it to
position 7 (directly after "Combining GWAS and QTL"). Every other slide is
left byte-content untouched and verified by text comparison.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).resolve().parent))
import update_phase11_seaad_simple_aggr_part2 as ui  # noqa: E402  (shared styling helpers)


ROOT = Path(__file__).resolve().parents[2]
DECK = (
    ROOT
    / "docs"
    / "presentations"
    / "human_genetic_support_for_key_drivers_simple_aggr_08292026.pptx"
)
AUDIT_PATH = (
    ROOT
    / "results/presentations/human_genetic_support_simple_aggr"
    / "dataset_slide_insert_checks.tsv"
)

EXPECTED_INPUT_SLIDES = 13
INSERT_AFTER = 6  # 1-based slide number the new slide follows
NEW_TITLE = "The public datasets behind the three evidence sources"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DECK)
    parser.add_argument("--output", type=Path, default=DECK)
    parser.add_argument("--audit", type=Path, default=AUDIT_PATH)
    return parser.parse_args()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def build_dataset_slide(prs: Presentation):
    slide = ui.new_slide(prs)
    ui.add_title_block(
        slide,
        NEW_TITLE,
        "The same three categories as the approach — one ready-made source, and the GWAS + QTL pair combined per driver.",
    )
    cards = [
        (
            "1 · PUBLISHED RESULTS",
            ui.PALE_GREEN,
            ui.TEAL,
            [
                "FunGen-xQTL AD integration snapshot: six public summary files.",
                "Contains AD fine-mapping, xQTL, and TWAS gene lists.",
            ],
        ),
        (
            "2 · DISEASE GWAS",
            ui.PALE_SKY,
            ui.BLUE,
            [
                "Clinical AD: Bellenguez 2022 meta-analysis (GCST90027158), complete genome-wide statistics.",
                "CSF biomarkers: amyloid-β42, total tau, p-tau181 (GCST90726396–98), ≈19,000 donors each.",
            ],
        ),
        (
            "3 · BRAIN QTL",
            ui.PALE_GOLD,
            ui.GOLD,
            [
                "NIAGADS NG00184 brain eQTL/sQTL fine-mapping (bulk and single-nucleus).",
                "eQTL Catalogue r7 panels: microglia (Young 2019), neurons (Aygun 2021), neocortex (Walker 2019).",
            ],
        ),
    ]
    for index, (kicker, bg, accent, items) in enumerate(cards):
        x = 0.66 + index * 4.10
        ui.add_rect(slide, x, 1.55, 3.83, 3.95, color=bg, outline=accent)
        ui.add_text(slide, kicker, x + 0.26, 1.84, 3.30, 0.28, size=11, color=ui.readable_accent(accent), bold=True)
        y = 2.35
        for item in items:
            ui.add_rect(slide, x + 0.24, y, 3.35, 1.32, color=ui.WHITE, outline=ui.LIGHT)
            ui.add_text(slide, item, x + 0.42, y + 0.16, 3.02, 1.05, size=10.6, color=ui.DARK)
            y += 1.48
    ui.add_rect(slide, 0.66, 5.75, 12.03, 0.60, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_text(
        slide,
        "Gene positions and symbols from GENCODE v44 and HGNC. All datasets are group-level summary statistics — no individual genotypes.",
        0.90, 5.91, 11.60, 0.30, size=10.2, color=ui.GRAY, align=PP_ALIGN.CENTER,
    )
    ui.add_text(
        slide,
        "Support for a driver needs the pair: category 2 anchors the disease signal, category 3 anchors the gene — joined by the same-variant test.",
        0.88, 6.55, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Attach concrete public datasets to each of the three evidence categories.",
        walkthrough="Category one is the FunGen-xQTL integration snapshot with its ready-made gene-level results. Category two holds the disease-side GWAS: the Bellenguez 2022 clinical AD meta-analysis and the three CSF biomarker studies. Category three holds the gene-side QTL data: NIAGADS NG00184 brain fine-mapping and the three eQTL Catalogue panels.",
        boundary="Gene annotation references (GENCODE, HGNC) support mapping only; everything is summary-level public data.",
        transition="With the data in place: which drivers have screening results.",
    )
    return slide


def main() -> int:
    args = parse_args()
    input_path = args.input.resolve()
    output_path = args.output.resolve()
    if not input_path.is_file():
        raise FileNotFoundError(input_path)
    original_hash = sha256_file(input_path)

    prs = Presentation(str(input_path))
    if len(prs.slides) != EXPECTED_INPUT_SLIDES:
        raise RuntimeError(
            f"Expected the {EXPECTED_INPUT_SLIDES}-slide deck, found {len(prs.slides)}"
        )
    if "Combining GWAS and QTL" not in slide_text(prs.slides[INSERT_AFTER - 1]):
        raise RuntimeError(f"Slide {INSERT_AFTER} does not match the pre-insert contract")
    if NEW_TITLE in "".join(slide_text(s) for s in prs.slides):
        raise RuntimeError("Dataset slide already present; refusing to insert twice")
    before_texts = [slide_text(slide) for slide in prs.slides]

    ui.set_notes_body_template(prs.slides[0].notes_slide.notes_placeholder._element)
    build_dataset_slide(prs)
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    slide_id_list.remove(slide_ids[-1])
    slide_id_list.insert(INSERT_AFTER, slide_ids[-1])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    file_descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{output_path.name}.", suffix=".tmp", dir=output_path.parent
    )
    os.close(file_descriptor)
    temporary = Path(temporary_name)
    try:
        prs.save(str(temporary))
        temporary.replace(output_path)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise

    reloaded = Presentation(str(output_path))
    after_texts = [slide_text(slide) for slide in reloaded.slides]
    others_unchanged = (
        after_texts[: INSERT_AFTER] == before_texts[: INSERT_AFTER]
        and after_texts[INSERT_AFTER + 1 :] == before_texts[INSERT_AFTER:]
    )
    inserted_ok = NEW_TITLE in after_texts[INSERT_AFTER]
    notes_ok = (
        reloaded.slides[INSERT_AFTER].has_notes_slide
        and reloaded.slides[INSERT_AFTER].notes_slide.notes_text_frame is not None
        and reloaded.slides[INSERT_AFTER].notes_slide.notes_text_frame.text.strip() != ""
    )
    checks: list[dict[str, Any]] = [
        {
            "check_id": "output_slide_count",
            "observed": len(reloaded.slides),
            "expected": EXPECTED_INPUT_SLIDES + 1,
            "passed": len(reloaded.slides) == EXPECTED_INPUT_SLIDES + 1,
        },
        {
            "check_id": f"new_slide_at_position_{INSERT_AFTER + 1}",
            "observed": inserted_ok,
            "expected": True,
            "passed": inserted_ok,
        },
        {
            "check_id": "all_other_slides_unchanged",
            "observed": "unchanged" if others_unchanged else "changed",
            "expected": "unchanged",
            "passed": others_unchanged,
        },
        {
            "check_id": "new_slide_has_notes",
            "observed": notes_ok,
            "expected": True,
            "passed": notes_ok,
        },
    ]
    audit = args.audit.resolve()
    audit.parent.mkdir(parents=True, exist_ok=True)
    with audit.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["check_id", "observed", "expected", "passed"],
            delimiter="\t",
            lineterminator="\n",
        )
        writer.writeheader()
        writer.writerows(checks)
    failed = [row["check_id"] for row in checks if not row["passed"]]
    if failed:
        raise RuntimeError("Slide insert failed checks: " + ", ".join(failed))
    print(f"updated={output_path}")
    print(f"inserted_slide_position={INSERT_AFTER + 1}")
    print(f"original_sha256={original_hash}")
    print(f"updated_sha256={sha256_file(output_path)}")
    print(f"audit={audit}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
