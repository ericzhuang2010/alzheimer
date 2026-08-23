#!/usr/bin/env python3
"""Append a multi-dataset mouse-validation recommendation slide.

The target deck is edited in place by default.  Existing slides, speaker notes,
and geometry are preserved; only one new editable slide is appended.  The
scientific content comes from the approved mouse dataset comparison report.
"""

from __future__ import annotations

import argparse
import copy
import csv
import hashlib
import importlib.util
import os
import sys
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches


REPO = Path(__file__).resolve().parents[2]
BASE_HELPER_PATH = REPO / "scripts/presentations/build_mouse_dataset_comparison_deck.py"
TARGET_DECK = REPO / "docs/presentations/mouse_dataset_comparison_08252026.pptx"
DEFAULT_REPORT_DIR = REPO / (
    "results/presentations/validation_mouse/mouse_dataset_comparison_08252026"
)

_SPEC = importlib.util.spec_from_file_location(
    "mouse_dataset_comparison_base_helpers", BASE_HELPER_PATH
)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Cannot load base deck helpers: {BASE_HELPER_PATH}")
base = importlib.util.module_from_spec(_SPEC)
sys.modules[_SPEC.name] = base
_SPEC.loader.exec_module(base)

SOURCE_MD = base.SOURCE_MD
SOURCE_MD_SHA256 = base.SOURCE_MD_SHA256
SLIDE_W = base.SLIDE_W
SLIDE_H = base.SLIDE_H
WHITE = base.WHITE
OFF_WHITE = base.OFF_WHITE
NAVY = base.NAVY
BLUE = base.BLUE
TEAL = base.TEAL
ORANGE = base.ORANGE
VERMILION = base.VERMILION
PURPLE = base.PURPLE
DARK = base.DARK
MID = base.MID
LIGHT = base.LIGHT
PALE_BLUE = base.PALE_BLUE
PALE_GREEN = base.PALE_GREEN
PALE_ORANGE = base.PALE_ORANGE
PALE_RED = base.PALE_RED
PALE_PURPLE = base.PALE_PURPLE
PALE_GRAY = base.PALE_GRAY

EXPECTED_BASE_SLIDES = 14
EXPECTED_FINAL_SLIDES = 15
TITLE = "Recommended mouse validation uses complementary datasets—not one merged cohort"
NOTE_HEADINGS = ("What to point at:", "Main takeaway:", "Boundary / transition:")


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def display_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(REPO))
    except ValueError:
        return str(path.resolve())


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def slide_signature(slide) -> tuple:
    """Stable content/geometry signature used to prove old slides were preserved."""
    shapes = []
    for shape in slide.shapes:
        shapes.append((
            shape.name,
            str(shape.shape_type),
            int(shape.left), int(shape.top), int(shape.width), int(shape.height),
            shape.text if getattr(shape, "has_text_frame", False) else "",
        ))
    return tuple(shapes), slide.notes_slide.notes_text_frame.text


def _add_core_card(slide, x: float, number: str, question: str,
                   dataset: str, role: str, boundary: str, *,
                   accent, bg) -> None:
    base.add_rect(slide, x, 1.62, 3.78, 1.60, color=bg,
                  outline=accent, line_width=1.5)
    base.add_rect(slide, x, 1.62, 3.78, 0.12, color=accent,
                  outline=None, radius=False)
    base.add_text(slide, f"{number} • {question}", x + 0.22, 1.84,
                  3.34, 0.22, size=9.2, color=accent, bold=True)
    base.add_text(slide, dataset, x + 0.22, 2.10, 3.34, 0.33,
                  size=17.5, color=NAVY, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    base.add_text(slide, role, x + 0.22, 2.47, 3.34, 0.36,
                  size=11.2, color=DARK, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    base.add_text(slide, boundary, x + 0.22, 2.88, 3.34, 0.22,
                  size=8.7, color=MID, italic=True,
                  valign=MSO_ANCHOR.MIDDLE)


def _add_disease_chip(slide, x: float, label: str, datasets: str) -> None:
    base.add_rect(slide, x, 3.75, 2.83, 0.62, color=WHITE,
                  outline=ORANGE, line_width=1.0)
    base.add_text(slide, label.upper(), x + 0.12, 3.84, 0.88, 0.18,
                  size=7.6, color=ORANGE, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    base.add_text(slide, datasets, x + 1.06, 3.77, 1.63, 0.42,
                  size=8.8, color=NAVY, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _add_step(slide, x: float, number: str, text: str) -> None:
    base.add_rect(slide, x, 4.86, 2.70, 0.96, color=PALE_GRAY,
                  outline=LIGHT, line_width=1.0)
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x + 0.14), Inches(5.08), Inches(0.42), Inches(0.42),
    )
    base.fill(circle, NAVY)
    circle.line.fill.background()
    base.add_text(slide, number, x + 0.14, 5.12, 0.42, 0.20,
                  size=10.0, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    base.add_text(slide, text, x + 0.68, 4.99, 1.86, 0.50,
                  size=10.4, color=NAVY, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)


def _add_note(slide, template_slide, value: str) -> None:
    """Write speaker notes even when a PowerPoint-resaved notes master is empty."""
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is None:
        template_notes = template_slide.notes_slide
        if template_notes.notes_text_frame is None:
            raise AssertionError("Template slide has no notes body placeholder")
        for shape in template_notes.shapes:
            if not getattr(shape, "is_placeholder", False):
                continue
            notes_slide.shapes._spTree.insert_element_before(
                copy.deepcopy(shape.element), "p:extLst"
            )
        notes_slide.__dict__.pop("placeholders", None)
        notes_slide.__dict__.pop("shapes", None)
    text_frame = notes_slide.notes_text_frame
    if text_frame is None:
        raise AssertionError("Could not create a notes body placeholder")
    text_frame.text = value


def append_recommendation_slide(prs: Presentation) -> None:
    template_slide = prs.slides[0]
    slide = base.new_slide(prs, bg=OFF_WHITE)
    base.add_header(
        slide, "Recommended validation architecture", TITLE, 15, accent=BLUE
    )
    base.add_text(
        slide,
        "Freeze the human PFC result first; then assign one pre-specified role to each independent mouse study.",
        0.57, 1.12, 12.10, 0.27, size=11.3, color=MID,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )

    _add_core_card(
        slide, 0.55, "1", "WHERE?", "GSE185063",
        "Cortical cell-type localization of APOE modules",
        "Cortex, not explicit PFC • no AD model", accent=BLUE, bg=PALE_BLUE,
    )
    _add_core_card(
        slide, 4.78, "2", "MECHANISM + SEX?", "GSE241553",
        "Microglia/CAM apoE3-versus-apoE4 induction under amyloid",
        "Conditional induction • no non-amyloid arm", accent=TEAL, bg=PALE_GREEN,
    )
    _add_core_card(
        slide, 9.00, "3", "FACTORIAL SUPPORT?", "GSE163857",
        "APOE × sex × 5xFAD/control in microglia",
        "Bulk, non-regional • small control groups", accent=PURPLE, bg=PALE_PURPLE,
    )

    base.add_text(
        slide, "ADD ONE DISEASE-DIRECTION LAYER MATCHED TO THE HUMAN CELL TYPE",
        0.60, 3.38, 12.10, 0.20, size=9.2, color=ORANGE,
        bold=True, align=PP_ALIGN.CENTER,
    )
    _add_disease_chip(slide, 0.55, "Microglia", "GSE127892 / GSE140510")
    _add_disease_chip(slide, 3.69, "Astrocytes", "GSE140399 + GSE143758")
    _add_disease_chip(
        slide, 6.83, "Neurons / oligo",
        "GSE140399 / GSE140510\n+ GSE212606",
    )
    _add_disease_chip(slide, 9.97, "Vascular", "GSE140399 / GSE212606")

    base.add_text(
        slide, "COMBINE CONCLUSIONS—NOT RAW COUNT MATRICES",
        0.60, 4.56, 12.10, 0.20, size=9.2, color=NAVY,
        bold=True, align=PP_ALIGN.CENTER,
    )
    steps = (
        (0.55, "1", "Freeze human PFC cell type, contrast, and module"),
        (3.65, "2", "Analyze each study separately at mouse/pool level"),
        (6.75, "3", "Compare orthologous direction and enrichment"),
        (9.85, "4", "Report convergent evidence with boundaries"),
    )
    for index, (x, number, text) in enumerate(steps):
        _add_step(slide, x, number, text)
        if index < len(steps) - 1:
            base.add_text(slide, "→", x + 2.72, 5.12, 0.36, 0.25,
                          size=18.0, color=MID, bold=True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    base.add_ribbon(
        slide,
        "Report “triangulated external validation” or “convergent support”—not complete direct mouse PFC replication.",
        y=6.36, accent=BLUE,
    )
    base.add_source(
        slide,
        "Source: docs/validation_mouse/mouse_dataset_comparison.md, §§5, 8, 10, 13, and 17.",
    )
    _add_note(slide, template_slide, """What to point at:
The recommended minimum is three complementary anchors, not one supposedly complete mouse cohort. GSE185063 asks where the human APOE-associated module appears across cortical cell types. GSE241553 tests sex-aware microglia/CAM apoE induction under amyloid. GSE163857 supplies the direct APOE-by-sex-by-5xFAD factorial microglial comparison. Then add one disease-direction dataset chosen for the human cell type; GSE212606 can supplement neuronal and oligodendroglial questions.

Main takeaway:
Freeze the human PFC cell type, contrast, direction, and module before looking at mouse results. Analyze every study separately at its valid mouse or pooled-sample unit, then compare orthologous effect direction, ranked genes, module enrichment, and pathway concordance.

Boundary / transition:
Do not merge raw count matrices across these studies or erase region, disease-model, APOE, sex, and replication differences with batch correction. Agreement should be reported as triangulated external validation or convergent support—not as one complete, direct mouse PFC replication cohort.""")


def validate_final_deck(path: Path, *,
                        preserved_signatures: tuple | None = None) -> None:
    base.validate_source_contract(SOURCE_MD)
    if not path.exists() or path.stat().st_size < 50_000:
        raise AssertionError(f"Deck missing or unexpectedly small: {path}")
    prs = Presentation(path)
    if len(prs.slides) != EXPECTED_FINAL_SLIDES:
        raise AssertionError(
            f"Expected {EXPECTED_FINAL_SLIDES} slides, found {len(prs.slides)}"
        )
    if (abs(prs.slide_width - SLIDE_W) > 1
            or abs(prs.slide_height - SLIDE_H) > 1):
        raise AssertionError("Deck is not 13.333333 × 7.5 inch widescreen")
    if SOURCE_MD_SHA256 not in (prs.core_properties.comments or ""):
        raise AssertionError("Deck is not bound to the approved Markdown source")

    if preserved_signatures is not None:
        actual = tuple(slide_signature(slide) for slide in list(prs.slides)[:14])
        if actual != preserved_signatures:
            raise AssertionError("An existing slide changed while appending the recommendation")

    slide = prs.slides[-1]
    visible = _slide_text(slide)
    required = (
        TITLE,
        "GSE185063", "GSE241553", "GSE163857",
        "GSE127892 / GSE140510", "GSE140399 + GSE143758", "GSE212606",
        "Freeze human PFC cell type, contrast, and module",
        "Analyze each study separately at mouse/pool level",
        "Compare orthologous direction and enrichment",
        "COMBINE CONCLUSIONS—NOT RAW COUNT MATRICES",
        "triangulated external validation", "convergent support",
        "Source:",
    )
    missing = [text for text in required if text.lower() not in visible.lower()]
    if missing:
        raise AssertionError(f"Recommendation slide missing required claims: {missing}")
    if any(shape.shape_type == 13 for shape in slide.shapes):
        raise AssertionError("Recommendation slide unexpectedly contains a picture")

    note = slide.notes_slide.notes_text_frame.text.strip()
    if len(note.split()) < 80:
        raise AssertionError("Recommendation slide speaker notes are too short")
    for heading in NOTE_HEADINGS:
        if heading not in note:
            raise AssertionError(f"Recommendation note missing heading: {heading}")

    tolerance = Inches(0.02)
    for shape in slide.shapes:
        if shape.left < -tolerance or shape.top < -tolerance:
            raise AssertionError("Recommendation slide has shape above/left of canvas")
        if shape.left + shape.width > SLIDE_W + tolerance:
            raise AssertionError("Recommendation slide has shape beyond right edge")
        if shape.top + shape.height > SLIDE_H + tolerance:
            raise AssertionError("Recommendation slide has shape beyond bottom edge")

    with zipfile.ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise AssertionError("PPTX ZIP integrity check failed")
        slide_xml = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        if len(slide_xml) != EXPECTED_FINAL_SLIDES:
            raise AssertionError("PPTX package slide XML count is incorrect")


def append_to_deck(input_path: Path = TARGET_DECK,
                   output_path: Path = TARGET_DECK) -> tuple[Path, str]:
    base.validate_source_contract(SOURCE_MD)
    input_path = input_path.resolve()
    output_path = output_path.resolve()
    if not input_path.exists():
        raise FileNotFoundError(input_path)

    original_hash = sha256(input_path)
    prs = Presentation(input_path)
    if len(prs.slides) == EXPECTED_FINAL_SLIDES:
        if TITLE.lower() not in _slide_text(prs.slides[-1]).lower():
            raise AssertionError("Deck already has 15 slides, but the final slide is unexpected")
        validate_final_deck(input_path)
        if input_path != output_path:
            raise AssertionError("A finalized deck cannot be copied through append mode")
        return input_path, original_hash
    if len(prs.slides) != EXPECTED_BASE_SLIDES:
        raise AssertionError(
            f"Expected a 14-slide base deck, found {len(prs.slides)} slides"
        )

    preserved = tuple(slide_signature(slide) for slide in prs.slides)
    append_recommendation_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output_path.stem}.", suffix=".tmp.pptx",
        dir=output_path.parent, delete=False,
    ) as handle:
        temp_path = Path(handle.name)
    try:
        prs.save(temp_path)
        validate_final_deck(temp_path, preserved_signatures=preserved)
        os.replace(temp_path, output_path)
    finally:
        if temp_path.exists():
            temp_path.unlink()
    validate_final_deck(output_path, preserved_signatures=preserved)
    return output_path, original_hash


def _write_tsv(path: Path, rows: list[dict[str, object]], fields: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        mode="w", encoding="utf-8", newline="", delete=False,
        prefix=f".{path.name}.", dir=path.parent,
    ) as handle:
        writer = csv.DictWriter(handle, fieldnames=fields, delimiter="\t")
        writer.writeheader()
        writer.writerows(rows)
        temp_path = Path(handle.name)
    os.replace(temp_path, path)


def write_report(deck_path: Path, report_dir: Path, *, base_deck_sha256: str,
                 visual_review_status: str) -> None:
    if visual_review_status not in {"pending", "complete"}:
        raise ValueError("visual_review_status must be pending or complete")
    validate_final_deck(deck_path)
    prs = Presentation(deck_path)
    report_dir.mkdir(parents=True, exist_ok=True)
    stem = "mouse_dataset_comparison_08252026"

    manifest = []
    for role, path in (
        ("source_markdown", SOURCE_MD),
        ("append_script", Path(__file__).resolve()),
        ("deck", deck_path),
    ):
        manifest.append({
            "artifact_role": role,
            "path": display_path(path),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        })
    _write_tsv(report_dir / f"{stem}_input_manifest.tsv", manifest,
               ["artifact_role", "path", "bytes", "sha256"])

    inventory = []
    for index, slide in enumerate(prs.slides, start=1):
        title = next(
            (shape.text for shape in slide.shapes if shape.name == "slide_title"),
            "",
        )
        inventory.append({
            "slide_number": index,
            "title": title,
            "shape_count": len(slide.shapes),
            "picture_count": sum(1 for shape in slide.shapes if shape.shape_type == 13),
            "speaker_note_words": len(slide.notes_slide.notes_text_frame.text.split()),
        })
    _write_tsv(report_dir / f"{stem}_slide_inventory.tsv", inventory,
               ["slide_number", "title", "shape_count", "picture_count",
                "speaker_note_words"])

    checks = (
        ("source_contract", True, f"Markdown SHA-256 {SOURCE_MD_SHA256}"),
        ("pptx_integrity", True, "PPTX ZIP and python-pptx validation passed"),
        ("slide_count", len(prs.slides) == 15, "14 original slides + 1 recommendation"),
        ("recommendation_content", TITLE in _slide_text(prs.slides[-1]),
         "Three core anchors, cell-type disease layer, and separate-study workflow"),
        ("native_editable", all(shape.shape_type != 13 for shape in prs.slides[-1].shapes),
         "Recommendation uses editable native PowerPoint shapes"),
        ("speaker_notes", len(prs.slides[-1].notes_slide.notes_text_frame.text.split()) >= 80,
         "Three-section recommendation speaker notes"),
        ("visual_review", visual_review_status == "complete",
         "Reviewed in color and grayscale" if visual_review_status == "complete"
         else "Awaiting color and grayscale review"),
    )
    blocking_failures = [
        check_id for check_id, passed, _ in checks
        if check_id != "visual_review" and not passed
    ]
    if blocking_failures:
        raise AssertionError(
            f"Blocking report checks failed: {blocking_failures}"
        )
    _write_tsv(
        report_dir / f"{stem}_checks.tsv",
        [{"check_id": cid, "passed": str(bool(passed)), "detail": detail}
         for cid, passed, detail in checks],
        ["check_id", "passed", "detail"],
    )

    status = [{
        "schema_version": "mouse_dataset_comparison_recommendation_append_v1",
        "validation_status": (
            "validated_complete" if visual_review_status == "complete"
            else "awaiting_visual_review"
        ),
        "visual_review_status": visual_review_status,
        "slides": len(prs.slides),
        "base_deck_sha256": base_deck_sha256,
        "deck_path": display_path(deck_path),
        "deck_bytes": deck_path.stat().st_size,
        "deck_sha256": sha256(deck_path),
        "source_markdown_sha256": SOURCE_MD_SHA256,
        "completed_utc": datetime.now(timezone.utc).isoformat(),
    }]
    _write_tsv(
        report_dir / f"{stem}_status.tsv", status,
        ["schema_version", "validation_status", "visual_review_status", "slides",
         "base_deck_sha256", "deck_path", "deck_bytes", "deck_sha256",
         "source_markdown_sha256", "completed_utc"],
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=TARGET_DECK)
    parser.add_argument("--output", type=Path, default=TARGET_DECK)
    parser.add_argument("--report-dir", type=Path, default=DEFAULT_REPORT_DIR)
    parser.add_argument("--visual-review-status", choices=("pending", "complete"),
                        default="pending")
    parser.add_argument("--validate-only", action="store_true")
    parser.add_argument(
        "--report-only", action="store_true",
        help="Refresh reports for an already-appended deck without changing it",
    )
    return parser.parse_args()


def _existing_base_hash(report_dir: Path) -> str:
    status_path = report_dir / "mouse_dataset_comparison_08252026_status.tsv"
    if not status_path.exists():
        raise FileNotFoundError(
            f"Cannot refresh reports without the prior append status: {status_path}"
        )
    with status_path.open(encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle, delimiter="\t"))
    if len(rows) != 1 or not rows[0].get("base_deck_sha256"):
        raise AssertionError(f"Invalid prior append status: {status_path}")
    return rows[0]["base_deck_sha256"]


def main() -> None:
    args = parse_args()
    output = args.output.resolve()
    report_dir = args.report_dir.resolve()
    if args.validate_only and args.report_only:
        raise ValueError("--validate-only and --report-only are mutually exclusive")
    if args.validate_only:
        validate_final_deck(output)
        print(f"Validated: {output}")
        return
    if args.report_only:
        base_hash = _existing_base_hash(report_dir)
        write_report(output, report_dir, base_deck_sha256=base_hash,
                     visual_review_status=args.visual_review_status)
        print(f"Reports refreshed: {report_dir}")
        return
    built, base_hash = append_to_deck(args.input, output)
    if base_hash == sha256(built):
        status_path = report_dir / "mouse_dataset_comparison_08252026_status.tsv"
        if status_path.exists():
            base_hash = _existing_base_hash(report_dir)
    write_report(built, report_dir, base_deck_sha256=base_hash,
                 visual_review_status=args.visual_review_status)
    print(f"Appended and validated: {built}")


if __name__ == "__main__":
    main()
