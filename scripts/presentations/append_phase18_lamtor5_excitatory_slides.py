#!/usr/bin/env python3
"""Append the LAMTOR5 excitatory-neuron pathway and STRING slides."""

from __future__ import annotations

import argparse
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN


REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO))

from scripts.presentations.append_phase18_cox7c_astrocyte_slides import (  # noqa: E402
    replace_text_preserve_style,
)
from scripts.presentations.build_phase18_key_driver_selection_deck import (  # noqa: E402
    BLUE,
    GOLD,
    GRAY,
    LIGHT,
    NAVY,
    PALE_BLUE,
    PURPLE,
    TEAL,
    VERMILION,
    WHITE,
    add_bullets,
    add_header,
    add_metric,
    add_panel_title,
    add_picture_contain,
    add_rect,
    add_source,
    add_text,
    new_slide,
    trim_white,
)


DEFAULT_DECK = REPO / "docs/presentations/key_driver_selection_analysis 08182026.pptx"
FIG_ROOT = REPO / "results/figures/analysis/phase_18_key_driver_selection/LAMTOR5/excitatory"
PATHWAY_FIG = FIG_ROOT / "phase18_lamtor5_excitatory_consensus_network_pathways.png"
STRING_FIG = FIG_ROOT / "stringdb_full_medium_conf.png"


def append_slides(input_path: Path, output_path: Path) -> Path:
    for path in (input_path, PATHWAY_FIG, STRING_FIG):
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation(input_path)
    if len(prs.slides) != 31:
        raise RuntimeError(f"Expected a 31-slide source deck, found {len(prs.slides)}")

    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    if "LAMTOR5 • EXCITATORY NEURONS" in all_text:
        raise RuntimeError("LAMTOR5 excitatory-neuron slides already appear to be present")

    replace_text_preserve_style(prs.slides[1], "Slides 16–31", "Slides 16–33")
    replace_text_preserve_style(
        prs.slides[1],
        "Place APOE, RPL11, COX7C, and SELENOW on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
        "Place APOE, RPL11, COX7C, SELENOW, and LAMTOR5 on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "Next: localize candidate evidence across sex/APOE strata, then examine APOE, RPL11, COX7C, and SELENOW in depth.",
        "Next: localize candidate evidence across sex/APOE strata, then examine APOE, RPL11, COX7C, SELENOW, and LAMTOR5 in depth.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–31 show APOE, RPL11, COX7C, and SELENOW network/protein examples",
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–33 show APOE, RPL11, COX7C, SELENOW, and LAMTOR5 network/protein examples",
    )

    with tempfile.TemporaryDirectory(prefix="lamtor5_excitatory_deck_assets_") as temp_dir:
        assets = Path(temp_dir)
        pathway = trim_white(PATHWAY_FIG, assets)
        string = trim_white(STRING_FIG, assets)

        # 32 — LAMTOR5 excitatory-neuron directed pathway figure
        slide = new_slide(prs)
        add_header(
            slide,
            "LAMTOR5 • excitatory neurons",
            "LAMTOR5 anchors a stable lysosome-to-respiration hypothesis",
            32,
            accent=PURPLE,
            subtitle="Twelve conservative-support runs reconstruct the excitatory neighborhood; LAMTOR5 also reproduces in inhibitory neurons with complete omission retention in both networks.",
        )
        add_picture_contain(
            slide,
            pathway,
            0.10,
            1.22,
            8.92,
            5.92,
            alt="LAMTOR5-centered excitatory-neuron consensus network with directed Bayesian-network edges and contextual pathway outlines",
        )
        add_rect(slide, 9.18, 1.42, 3.62, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.45, 1.70, 3.07, accent=PURPLE)
        add_metric(slide, "12", "conservative-support excitatory runs", 9.45, 2.18, 1.39, accent=PURPLE)
        add_metric(slide, "2.59×10⁻³", "aggregate ACAT q", 10.99, 2.18, 1.54, accent=GOLD)
        add_panel_title(slide, "Network readout", 9.45, 3.48, 3.07, accent=GOLD)
        add_bullets(
            slide,
            [
                "17 nodes and 16 directed edges retain all eight observed mitochondrial query hits at the 3/12 display threshold.",
                "UQCR10 → LAMTOR5 supplies upstream context; nine direct outputs include ATP5IF1, NDUFA6, TMEM11, and UQCRHL.",
                "ETC/OXPHOS: 5 genes, FDR 8.23×10⁻⁴; complex-I biogenesis: 3 genes, FDR 0.0449.",
                "Cristae and mTORC1 outlines are contextual, not significant; 4/12 would remove only TMEM126A.",
            ],
            9.45,
            3.96,
            2.94,
            size=10.3,
            line_h=0.58,
            accent=GOLD,
        )
        add_text(
            slide,
            "Interpretation: recurrence and respiratory enrichment support a focused regulatory hypothesis; neither the arrows nor pathway labels prove causality or pathway activity.",
            9.48,
            6.42,
            2.88,
            0.34,
            size=9.1,
            color=NAVY,
            bold=True,
        )
        add_source(
            slide,
            "Source: phase18_lamtor5_excitatory_consensus_network_pathways.png, caption, methods, ORA table, and gene-by-gene analysis",
        )

        # 33 — LAMTOR5 excitatory-neuron STRING figure and mechanistic test
        slide = new_slide(prs)
        add_header(
            slide,
            "LAMTOR5 • excitatory neurons",
            "STRING validates the respiratory module—but leaves LAMTOR5 isolated",
            33,
            accent=PURPLE,
            subtitle="Medium-confidence STRING associations are undirected and not excitatory-neuron- or AD-specific.",
        )
        add_rect(slide, 0.42, 1.39, 6.37, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(
            slide,
            string,
            0.64,
            1.62,
            5.93,
            4.98,
            alt="STRING medium-confidence functional association network for LAMTOR5 and excitatory-neuron consensus-neighborhood proteins",
        )
        add_rect(slide, 7.04, 1.39, 5.76, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 7.34, 1.72, 5.15, accent=PURPLE)
        add_bullets(
            slide,
            [
                "A dense target component connects complex I/III/V proteins, including NDUFA6, NDUFB6, UQCR10/UQCRHL, and ATP5 subunits.",
                "LAMTOR5 has no medium-confidence edge to this module, so STRING does not corroborate the modeled LAMTOR5→target chain.",
                "CHCHD10, MRPS36, POP7, and MAGEH1 are also isolated; ELOC–POMP forms a separate pair, showing uneven protein-level support.",
            ],
            7.34,
            2.20,
            5.05,
            size=10.8,
            line_h=0.59,
            accent=PURPLE,
        )
        add_panel_title(slide, "Mechanism worth testing", 7.34, 4.28, 5.15, accent=TEAL)
        add_text(
            slide,
            "External biology supplies a plausible bridge: LAMTOR5 is a Ragulator subunit that recruits mTORC1 to lysosomes; mTORC1 regulates mitochondrial output, Aβ can disrupt lysosome-to-mitochondria signaling, and defective LAMTOR5 can impair V-ATPase assembly and acidification. Direct LAMTOR5 perturbation in AD remains untested.",
            7.35,
            4.74,
            5.02,
            0.78,
            size=10.0,
            color=GRAY,
        )
        add_rect(slide, 7.35, 5.64, 4.98, 0.78, color=PALE_BLUE, outline=BLUE)
        add_text(
            slide,
            "In APOE-isogenic excitatory neurons, combine CRISPRi/CRISPRa with amino-acid withdrawal/re-feeding; measure lysosomal pH, mTORC1 localization, p-S6K/4EBP1, mitophagy, OCR, and the frozen target module, then perform LAMTOR5 rescue.",
            7.56,
            5.76,
            4.57,
            0.54,
            size=9.0,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "Female ε2 AD-up versus female ε4/male ε2 AD-down support suggests context dependence—not a universal direction.",
            7.42,
            6.58,
            4.84,
            0.20,
            size=8.5,
            color=VERMILION,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_source(
            slide,
            "Sources: stringdb_full_medium_conf.png • phase18_key_driver_gene_by_gene_initial_analysis.md • Bar-Peled 2012; Morita 2013; Norambuena 2018; Zhang 2024",
        )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    validate_output(output_path)
    return output_path


def validate_output(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 33:
        raise RuntimeError(f"Expected 33 slides, found {len(prs.slides)}")
    image_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    if image_shapes != 22:
        raise RuntimeError(f"Expected 22 figure images, found {image_shapes}")

    endings = []
    for slide_index in range(len(prs.slides) - 2, len(prs.slides)):
        slide = prs.slides[slide_index]
        endings.append(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
        )
    if not all("LAMTOR5 • EXCITATORY NEURONS" in text for text in endings):
        raise RuntimeError("The final two slides are not the expected LAMTOR5 excitatory-neuron slides")

    with zipfile.ZipFile(path) as archive:
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels"))
        )
    if "filter_attrition" in package_text:
        raise RuntimeError("Deprecated filter_attrition content found in the PPTX package")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    args = parser.parse_args()
    result = append_slides(args.input, args.output)
    print(result)
