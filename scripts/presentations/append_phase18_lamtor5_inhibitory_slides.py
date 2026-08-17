#!/usr/bin/env python3
"""Append the LAMTOR5 inhibitory-neuron pathway and STRING slides."""

from __future__ import annotations

import argparse
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN


REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO))

from scripts.presentations.append_phase18_cox7c_astrocyte_slides import (  # noqa: E402
    replace_text_preserve_style,
)
from scripts.presentations.build_phase18_key_driver_selection_deck import (  # noqa: E402
    BLUE,
    GOLD,
    GRAY,
    LIGHT,
    NAVY,
    PALE_BLUE,
    PURPLE,
    TEAL,
    VERMILION,
    WHITE,
    add_bullets,
    add_header,
    add_metric,
    add_panel_title,
    add_picture_contain,
    add_rect,
    add_source,
    add_text,
    new_slide,
    trim_white,
)


DEFAULT_DECK = REPO / "docs/presentations/key_driver_selection_analysis 08182026.pptx"
FIG_ROOT = REPO / "results/figures/analysis/phase_18_key_driver_selection/LAMTOR5/inhibitory"
PATHWAY_FIG = FIG_ROOT / "phase18_lamtor5_inhibitory_consensus_network_pathways.png"
STRING_FIG = FIG_ROOT / "stringdb_full_medium_conf.png"


def append_slides(input_path: Path, output_path: Path) -> Path:
    for path in (input_path, PATHWAY_FIG, STRING_FIG):
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation(input_path)
    if len(prs.slides) != 33:
        raise RuntimeError(f"Expected a 33-slide source deck, found {len(prs.slides)}")

    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    if "LAMTOR5 • INHIBITORY NEURONS" in all_text:
        raise RuntimeError("LAMTOR5 inhibitory-neuron slides already appear to be present")

    replace_text_preserve_style(prs.slides[1], "Slides 16–33", "Slides 16–35")
    replace_text_preserve_style(
        prs.slides[8],
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–33 show APOE, RPL11, COX7C, SELENOW, and LAMTOR5 network/protein examples",
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–35 show APOE, RPL11, COX7C, SELENOW, and LAMTOR5 network/protein examples",
    )

    with tempfile.TemporaryDirectory(prefix="lamtor5_inhibitory_deck_assets_") as temp_dir:
        assets = Path(temp_dir)
        pathway = trim_white(PATHWAY_FIG, assets)
        string = trim_white(STRING_FIG, assets)

        # 34 — LAMTOR5 inhibitory-neuron directed pathway figure
        slide = new_slide(prs)
        add_header(
            slide,
            "LAMTOR5 • inhibitory neurons",
            "LAMTOR5 expands into a respiratory/proteostasis program",
            34,
            accent=PURPLE,
            subtitle="Five inhibitory runs show complete omission retention; every supporting context uses an AD-down mitochondrial query.",
        )
        add_picture_contain(
            slide,
            pathway,
            0.10,
            1.22,
            8.92,
            5.92,
            alt="LAMTOR5-centered inhibitory-neuron consensus network with directed Bayesian-network edges and contextual pathway outlines",
        )
        add_rect(slide, 9.18, 1.42, 3.62, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.45, 1.70, 3.07, accent=PURPLE)
        add_metric(slide, "5", "conservative-support inhibitory runs", 9.45, 2.18, 1.39, accent=PURPLE)
        add_metric(slide, "4.14×10⁻³", "aggregate ACAT q", 10.99, 2.18, 1.54, accent=GOLD)
        add_panel_title(slide, "Network readout", 9.45, 3.48, 3.07, accent=GOLD)
        add_bullets(
            slide,
            [
                "28 nodes and 28 directed edges retain all 14 observed mitochondrial query hits at 1/5; 2/5 retains nine, leaving five single-run branches exploratory.",
                "Upstream chain: RPS15 → COX7C → LAMTOR5; nine direct outputs include ATP5PF, which fans into respiratory and redox targets.",
                "ETC/OXPHOS: 11 genes, FDR 9.95×10⁻¹²; complex-I biogenesis: 7 genes, FDR 4.69×10⁻⁷; protein degradation: 4 genes, FDR 0.0253.",
            ],
            9.45,
            3.96,
            2.94,
            size=10.0,
            line_h=0.80,
            accent=GOLD,
        )
        add_text(
            slide,
            "Interpretation: stronger pathway coherence but shallower recurrence than the excitatory network; translation and mTORC1 outlines remain contextual.",
            9.48,
            6.42,
            2.88,
            0.34,
            size=9.1,
            color=NAVY,
            bold=True,
        )
        add_source(
            slide,
            "Source: phase18_lamtor5_inhibitory_consensus_network_pathways.png, caption, methods, ORA table, and gene-by-gene analysis",
        )

        # 35 — LAMTOR5 inhibitory-neuron STRING figure and paired validation
        slide = new_slide(prs)
        add_header(
            slide,
            "LAMTOR5 • inhibitory neurons",
            "STRING supports the targets—not a LAMTOR5 connection",
            35,
            accent=PURPLE,
            subtitle="Medium-confidence STRING associations are undirected and not inhibitory-neuron- or AD-specific.",
        )
        add_rect(slide, 0.42, 1.39, 6.37, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(
            slide,
            string,
            0.64,
            1.62,
            5.93,
            4.98,
            alt="STRING medium-confidence functional association network for LAMTOR5 and inhibitory-neuron consensus-neighborhood proteins",
        )
        add_rect(slide, 7.04, 1.39, 5.76, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 7.34, 1.72, 5.15, accent=PURPLE)
        add_bullets(
            slide,
            [
                "A dense respiratory component spans complex I, IV, and V proteins plus SOD1 and CISD1, strongly supporting target-module coherence.",
                "LAMTOR5 remains isolated at medium confidence, so STRING does not provide the missing LAMTOR5→respiratory protein bridge.",
                "A separate SRP14/RPS15/SPCS2/DAD1 branch and several isolated proteins show that support is modular, not uniform across modeled edges.",
            ],
            7.34,
            2.20,
            5.05,
            size=10.8,
            line_h=0.59,
            accent=PURPLE,
        )
        add_panel_title(slide, "Cross-neuronal test", 7.34, 4.28, 5.15, accent=TEAL)
        add_text(
            slide,
            "Across excitatory and inhibitory neurons, LAMTOR5 totals 17 conservative-support runs across 13 fine cell types with complete omission retention. Ragulator–mTORC1 and V-ATPase/acidification biology make the route plausible, but no direct gene-specific AD perturbation study was identified.",
            7.35,
            4.74,
            5.02,
            0.72,
            size=10.1,
            color=GRAY,
        )
        add_rect(slide, 7.35, 5.60, 4.98, 0.82, color=PALE_BLUE, outline=BLUE)
        add_text(
            slide,
            "Compare APOE-isogenic inhibitory and excitatory neurons after LAMTOR5 CRISPRi/CRISPRa plus amino-acid withdrawal/re-feeding: lysosomal pH, mTORC1 localization, p-S6K/4EBP1, mitophagy, OCR, and each frozen target module; require sgRNA-resistant rescue.",
            7.56,
            5.72,
            4.57,
            0.58,
            size=8.9,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "All five inhibitory runs are AD-down—four male ε2 and one female ε4—but this is descriptive, not an interaction test.",
            7.42,
            6.58,
            4.84,
            0.20,
            size=8.5,
            color=VERMILION,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_source(
            slide,
            "Sources: stringdb_full_medium_conf.png • phase18_key_driver_gene_by_gene_initial_analysis.md • Bar-Peled 2012; Morita 2013; Norambuena 2018; Zhang 2024",
        )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    validate_output(output_path)
    return output_path


def validate_output(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 35:
        raise RuntimeError(f"Expected 35 slides, found {len(prs.slides)}")
    image_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    if image_shapes != 24:
        raise RuntimeError(f"Expected 24 figure images, found {image_shapes}")

    endings = []
    for slide_index in range(len(prs.slides) - 2, len(prs.slides)):
        slide = prs.slides[slide_index]
        endings.append(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
        )
    if not all("LAMTOR5 • INHIBITORY NEURONS" in text for text in endings):
        raise RuntimeError("The final two slides are not the expected LAMTOR5 inhibitory-neuron slides")

    with zipfile.ZipFile(path) as archive:
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels"))
        )
    if "filter_attrition" in package_text:
        raise RuntimeError("Deprecated filter_attrition content found in the PPTX package")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    args = parser.parse_args()
    result = append_slides(args.input, args.output)
    print(result)
