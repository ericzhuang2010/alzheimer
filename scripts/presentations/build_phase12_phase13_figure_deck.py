#!/usr/bin/env python3
"""Build a scientific slide deck for the requested Phase 12/13 figures.

The deck is intentionally generated from project-local figures and supporting
documentation. It does not refit models or recalculate the source analyses.
"""

from __future__ import annotations

import csv
import os
import tempfile
from collections import Counter
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
OUT_DIR = REPO / "docs" / "presentations"
OUT_PATH = OUT_DIR / "phase12_kda_phase13_modifier_figures.pptx"

FIG = {
    "circular": REPO / "results/figures/analysis/phase12_kda/reduced_circular_figure/phase12_kda_reduced_circular.png",
    "sex_apoe": REPO / "results/figures/analysis/phase12_kda/sex_apoe_figure/phase12_kda_sex_apoe.png",
    "wang": REPO / "results/figures/analysis/phase12_kda/network_figures/phase12_kda_wang_subnetworks.png",
    "sex_reversal": REPO / "results/figures/analysis/phase12_kda/network_figures/phase12_kda_sex_reversal_networks.png",
    "atp": REPO / "results/figures/analysis/phase12_kda/network_figures/phase12_kda_atp_convergence.png",
    "connectivity": REPO / "results/figures/analysis/phase12_kda/network_figures/phase12_kda_connectivity_evidence.png",
    "modifier": REPO / "results/figures/analysis/phase13_respiratory_modifier/modifier_landscape/phase13_modifier_landscape.png",
}

TABLE = {
    "circular": REPO / "results/figures/analysis/phase12_kda/reduced_circular_figure/phase12_kda_reduced_circular_plotted_data.tsv",
    "sex_apoe": REPO / "results/figures/analysis/phase12_kda/sex_apoe_figure/phase12_kda_sex_apoe_plotted_data.tsv",
    "atp": REPO / "results/figures/analysis/phase12_kda/network_figures/phase12_kda_atp_convergence_pairs.tsv",
    "corr": REPO / "results/figures/analysis/phase12_kda/network_figures/phase12_kda_connectivity_evidence_correlations.tsv",
    "modifier": REPO / "results/figures/analysis/phase13_respiratory_modifier/modifier_landscape/phase13_modifier_landscape_plotted_data.tsv",
}

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 35, 61)
NAVY_2 = RGBColor(30, 59, 91)
BLUE = RGBColor(46, 124, 171)
CYAN = RGBColor(64, 176, 196)
TEAL = RGBColor(0, 158, 115)
ORANGE = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
MAGENTA = RGBColor(204, 121, 167)
YELLOW = RGBColor(240, 228, 66)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(232, 238, 244)
MID = RGBColor(118, 132, 148)
DARK = RGBColor(31, 39, 48)
GRAY = RGBColor(85, 96, 108)
PALE_BLUE = RGBColor(225, 238, 247)
PALE_ORANGE = RGBColor(252, 237, 211)
PALE_PURPLE = RGBColor(239, 232, 246)
PALE_GREEN = RGBColor(225, 243, 237)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh, delimiter="\t"))


def finite_number(value: str | None) -> float | None:
    if value is None or value.strip() in {"", "NA", "NaN", "nan"}:
        return None
    return float(value)


def shape_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def shape_line(shape, color: RGBColor, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if transparency:
        shape.line.transparency = transparency


def set_run(run, *, size: float, color: RGBColor, bold: bool = False,
            font: str = FONT_BODY, italic: bool = False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(slide, text: str, x: float, y: float, w: float, h: float,
             *, size: float = 18, color: RGBColor = DARK, bold: bool = False,
             font: str = FONT_BODY, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.04,
             italic: bool = False, rotate: float | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if rotate is not None:
        box.rotation = rotate
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, font=font, italic=italic)
    return box


def add_rich_text(slide, spans: list[tuple[str, dict]], x: float, y: float,
                  w: float, h: float, *, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, margin: float = 0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    for text, style in spans:
        run = p.add_run()
        run.text = text
        set_run(run, size=style.get("size", 18), color=style.get("color", DARK),
                bold=style.get("bold", False), font=style.get("font", FONT_BODY),
                italic=style.get("italic", False))
    return box


def add_rect(slide, x: float, y: float, w: float, h: float,
             *, fill: RGBColor = WHITE, line: RGBColor | None = LIGHT,
             radius: bool = True, line_width: float = 1.0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape_fill(shape, fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape_line(shape, line, line_width)
    return shape


def add_circle(slide, x: float, y: float, d: float, *, fill: RGBColor,
               line: RGBColor | None = None, line_width: float = 1.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape_fill(shape, fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape_line(shape, line, line_width)
    return shape


def add_header(slide, title: str, kicker: str, *, accent: RGBColor = CYAN,
               subtitle: str | None = None) -> None:
    add_text(slide, kicker.upper(), 0.62, 0.28, 3.5, 0.28, size=11, color=accent, bold=True)
    add_text(slide, title, 0.62, 0.58, 12.0, 0.62, size=28, color=NAVY, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.18, 11.9, 0.42, size=13.5, color=GRAY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.58), Inches(12.08), Inches(0.025))
    shape_fill(line, LIGHT)
    line.line.fill.background()


def add_footer(slide, slide_no: int, source: str | None = None) -> None:
    if source:
        add_text(slide, source, 0.62, 7.18, 11.6, 0.18, size=7.4, color=MID)
    add_text(slide, str(slide_no), 12.45, 7.15, 0.35, 0.2, size=8.5, color=MID, align=PP_ALIGN.RIGHT)


def add_bullet_list(slide, items: list[str], x: float, y: float, w: float,
                    *, size: float = 16.5, color: RGBColor = DARK,
                    bullet_color: RGBColor = CYAN, line_h: float = 0.58,
                    bullet_d: float = 0.10, bold_leads: bool = False) -> None:
    cy = y
    for item in items:
        add_circle(slide, x, cy + 0.14, bullet_d, fill=bullet_color)
        if bold_leads and ":" in item:
            lead, rest = item.split(":", 1)
            add_rich_text(slide, [
                (lead + ":", {"size": size, "color": color, "bold": True}),
                (rest, {"size": size, "color": color}),
            ], x + 0.22, cy, w - 0.22, line_h)
        else:
            add_text(slide, item, x + 0.22, cy, w - 0.22, line_h, size=size, color=color)
        cy += line_h


def add_metric(slide, value: str, label: str, x: float, y: float, w: float,
               *, accent: RGBColor = CYAN, fill: RGBColor = WHITE,
               note: str | None = None, label_color: RGBColor = GRAY) -> None:
    add_rect(slide, x, y, w, 1.12 if not note else 1.35, fill=fill, line=LIGHT)
    add_text(slide, value, x + 0.18, y + 0.12, w - 0.36, 0.48, size=25, color=accent, bold=True, font=FONT_HEAD)
    add_text(slide, label, x + 0.18, y + 0.59, w - 0.36, 0.33, size=11.5, color=label_color, bold=True)
    if note:
        add_text(slide, note, x + 0.18, y + 0.94, w - 0.36, 0.26, size=8.8, color=MID)


def add_step(slide, num: str, title: str, body: str, x: float, y: float, w: float,
             *, accent: RGBColor = CYAN, fill: RGBColor = WHITE) -> None:
    add_rect(slide, x, y, w, 1.20, fill=fill, line=LIGHT)
    add_circle(slide, x + 0.18, y + 0.18, 0.44, fill=accent)
    add_text(slide, num, x + 0.18, y + 0.20, 0.44, 0.34, size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, x + 0.75, y + 0.14, w - 0.92, 0.30, size=15, color=NAVY, bold=True)
    add_text(slide, body, x + 0.75, y + 0.48, w - 0.92, 0.58, size=10.7, color=GRAY)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float,
              *, color: RGBColor = MID, width: float = 1.5, dashed: bool = False) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape_line(line, color, width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    line.line.end_arrowhead = True


def set_alt_text(shape, title: str, description: str) -> None:
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        c_nv_pr.set("name", title)
        c_nv_pr.set("descr", description)
    except Exception:
        pass


def trim_white(path: Path, target_dir: Path) -> Path:
    img = Image.open(path).convert("RGB")
    # Treat near-white as background while retaining light-gray plot elements.
    mask = Image.new("L", img.size)
    pix = img.load()
    out = mask.load()
    for yy in range(img.height):
        for xx in range(img.width):
            r, g, b = pix[xx, yy]
            out[xx, yy] = 255 if min(r, g, b) < 246 else 0
    bbox = mask.getbbox()
    if bbox:
        left, top, right, bottom = bbox
        pad = max(8, int(min(img.size) * 0.008))
        bbox = (max(0, left - pad), max(0, top - pad), min(img.width, right + pad), min(img.height, bottom + pad))
        img = img.crop(bbox)
    target = target_dir / path.name
    img.save(target, quality=95)
    return target


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float,
                        *, alt: str) -> None:
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    pic = slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))
    set_alt_text(pic, alt, alt)


def add_figure_slide(prs: Presentation, slide_no: int, img_path: Path,
                     figure_id: str, source: str, *, background: RGBColor = WHITE) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = background
    add_text(slide, figure_id.upper(), 0.35, 0.16, 5.5, 0.24, size=9, color=NAVY, bold=True)
    add_picture_contain(slide, img_path, 0.18, 0.42, 12.97, 6.68, alt=figure_id)
    add_footer(slide, slide_no, source)


def add_table(slide, data: list[list[str]], x: float, y: float, w: float, h: float,
              *, col_widths: list[float] | None = None, header_fill: RGBColor = NAVY,
              header_color: RGBColor = WHITE, body_size: float = 10.5,
              first_col_bold: bool = False):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else (WHITE if r % 2 else OFF_WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                set_run(run, size=10.2 if r == 0 else body_size,
                        color=header_color if r == 0 else DARK,
                        bold=(r == 0 or (first_col_bold and c == 0)))
    return table_shape


def new_slide(prs: Presentation, title: str, kicker: str, slide_no: int,
              *, accent: RGBColor = CYAN, subtitle: str | None = None,
              source: str | None = None, bg: RGBColor = OFF_WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = bg
    add_header(slide, title, kicker, accent=accent, subtitle=subtitle)
    add_footer(slide, slide_no, source)
    return slide


def build_deck() -> Path:
    for path in list(FIG.values()) + list(TABLE.values()):
        if not path.exists():
            raise FileNotFoundError(path)

    circular = read_tsv(TABLE["circular"])
    sex_apoe = read_tsv(TABLE["sex_apoe"])
    atp_pairs = read_tsv(TABLE["atp"])
    corr = read_tsv(TABLE["corr"])
    modifier = read_tsv(TABLE["modifier"])

    circular_unique = len({r["key_driver"] for r in circular})
    circular_mtdna = sum(r["mtDNA_encoded"].upper() == "TRUE" for r in circular)
    circular_recurrence = Counter(r["key_driver"] for r in circular)
    sex_status = Counter(r["display_status"] for r in sex_apoe)
    focused_driver_targets = 27
    all_driver_targets = len(atp_pairs)
    modifier_status = Counter(r["scientific_status"] for r in modifier)
    nominal_score = sum((p := finite_number(r["p_value_score"])) is not None and p < 0.05 for r in modifier)
    min_q_score = min(q for r in modifier if (q := finite_number(r["q_value_score"])) is not None)
    min_q_camera = min(q for r in modifier if (q := finite_number(r["q_value_camera"])) is not None)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Phase 12 KDA and Phase 13 Respiratory-Modifier Figures"
    prs.core_properties.subject = "Methods, figures, findings, and integrated interpretation"
    prs.core_properties.author = "Alzheimer project analysis team"
    prs.core_properties.keywords = "Alzheimer, mitochondria, KDA, APOE, sex, respiratory modifier"

    with tempfile.TemporaryDirectory(prefix="alzheimer_deck_assets_") as td:
        asset_dir = Path(td)
        trimmed = {key: trim_white(path, asset_dir) for key, path in FIG.items()}

        # 1 — Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
        add_rect(slide, 0, 0, 13.333, 7.5, fill=NAVY, line=None, radius=False)
        for x, y, d, col in [(9.1, 0.7, 1.7, TEAL), (10.7, 1.5, 1.05, ORANGE),
                             (11.55, 0.55, 0.56, CYAN), (9.85, 3.05, 0.76, PURPLE),
                             (11.35, 3.05, 1.35, BLUE), (10.45, 4.75, 0.48, MAGENTA)]:
            add_circle(slide, x, y, d, fill=col)
        for x1, y1, x2, y2 in [(9.85,1.55,10.95,2.0),(11.2,2.15,11.8,3.45),(10.25,3.45,11.4,3.65),
                               (10.65,4.05,10.7,4.95),(11.8,3.9,10.75,4.98)]:
            add_arrow(slide, x1, y1, x2, y2, color=RGBColor(151, 178, 204), width=1.2)
        add_text(slide, "MITOCHONDRIAL NETWORK EVIDENCE", 0.72, 0.72, 6.4, 0.3,
                 size=11, color=CYAN, bold=True)
        add_text(slide, "Phase 12 KDA\nand Phase 13\nmodifier analysis", 0.72, 1.20, 7.5, 2.40,
                 size=36, color=WHITE, bold=True, font=FONT_HEAD)
        add_text(slide, "How the figures were generated, what they show, and where the conclusions stop",
                 0.75, 3.95, 6.95, 0.75, size=18, color=RGBColor(207, 220, 234))
        add_rect(slide, 0.75, 5.30, 6.80, 0.06, fill=CYAN, line=None, radius=False)
        add_text(slide, "Phase 12: cross-network KDA ranking, stratified support, and topology\n"
                        "Phase 13: donor-aware respiratory modifier tests across the complete prespecified family",
                 0.75, 5.62, 7.2, 0.92, size=12.5, color=RGBColor(207, 220, 234))
        add_text(slide, "Prepared from validated project outputs • 10 August 2026",
                 0.75, 6.91, 6.8, 0.24, size=9, color=RGBColor(151, 178, 204))

        # 2 — Story map
        slide = new_slide(prs, "One story, three levels of evidence", "Deck map", 2,
                          subtitle="The figures answer related—but not interchangeable—questions.")
        cards = [
            ("1", "Rank & localize", "Which candidate drivers recur, and in which sex/APOE and cell-network contexts?", TEAL, PALE_GREEN),
            ("2", "Connect & synthesize", "How do selected drivers connect to mitochondrial and Complex V genes across directed networks?", BLUE, PALE_BLUE),
            ("3", "Formally test modifiers", "Do donor-level respiratory module effects differ by sex or APOE after complete-family correction?", ORANGE, PALE_ORANGE),
        ]
        for i, (num, title, body, accent, fill) in enumerate(cards):
            x = 0.75 + i * 4.15
            add_rect(slide, x, 2.05, 3.72, 3.65, fill=fill, line=accent)
            add_circle(slide, x + 0.25, 2.30, 0.48, fill=accent)
            add_text(slide, num, x + 0.25, 2.32, 0.48, 0.34, size=15, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, title, x + 0.25, 3.05, 3.12, 0.52, size=20, color=NAVY, bold=True, font=FONT_HEAD)
            add_text(slide, body, x + 0.25, 3.72, 3.13, 1.25, size=14, color=GRAY)
            add_text(slide, ["Phase 12 overview + heatmap", "Four network figures", "Phase 13 landscape"][i],
                     x + 0.25, 5.17, 3.05, 0.3, size=10.5, color=accent, bold=True)
        add_arrow(slide, 4.47, 3.86, 4.84, 3.86, color=MID, width=1.6)
        add_arrow(slide, 8.62, 3.86, 8.99, 3.86, color=MID, width=1.6)
        add_rect(slide, 1.20, 6.17, 10.90, 0.60, fill=WHITE, line=LIGHT)
        add_text(slide, "Interpretive rule", 1.43, 6.31, 1.30, 0.24, size=11, color=VERMILION, bold=True)
        add_text(slide, "Network enrichment generates mechanistic hypotheses; Phase 13 provides the donor-aware interaction test.",
                 2.78, 6.25, 8.98, 0.34, size=13.2, color=NAVY, bold=True)

        # 3 — Phase 12 shared setup
        slide = new_slide(prs, "Phase 12 KDA: common analytical foundation", "Shared setup", 3,
                          subtitle="All Phase 12 figures use eligible primary directional KDA runs and preserve the full candidate matrix.",
                          source="Sources: Phase 12 design documents; prepare_phase12_kda_figure_data.R")
        steps = [
            ("1", "Directional signatures", "AD-up and AD-down mitochondrial DEG signatures are evaluated separately; pooled and combined-direction signatures are excluded."),
            ("2", "Directed neighborhoods", "KDA asks whether a signature is enriched downstream of a candidate in a broad-cell-type Bayesian network."),
            ("3", "Complete test matrix", "Nonsignificant and zero-overlap tests remain. A truly untested candidate stays missing, preventing significance-only ranking bias."),
            ("4", "Two evidence summaries", "Mean −log10(P) ranks candidates; within-run BH-adjusted P≤0.05 counts recurrence. Neither is a causal effect estimate."),
        ]
        for i, (n, t, b) in enumerate(steps):
            add_step(slide, n, t, b, 0.82 + (i % 2) * 6.12, 1.95 + (i // 2) * 1.55, 5.67,
                     accent=[TEAL, BLUE, PURPLE, ORANGE][i], fill=WHITE)
        add_rect(slide, 0.84, 5.35, 11.80, 1.13, fill=NAVY, line=None)
        add_rich_text(slide, [
            ("Key boundary: ", {"size": 16, "color": CYAN, "bold": True}),
            ("KDA significance means a candidate’s downstream neighborhood is enriched for the query signature. ", {"size": 16, "color": WHITE}),
            ("It does not prove a direct gene-to-gene effect or causality.", {"size": 16, "color": WHITE, "bold": True}),
        ], 1.14, 5.63, 11.12, 0.62, valign=MSO_ANCHOR.MIDDLE)

        # 4 — Circular setup
        slide = new_slide(prs, "Reduced circular overview: setup and generation", "Figure 1 setup", 4,
                          subtitle="A deliberately reduced view ranks each network independently and shows only cross-network recurrence.",
                          source="Source: reduced_circular_figure_design.md; visualize_phase12_kda_reduced_circular.R")
        add_metric(slide, "7", "broad networks with eligible KDA runs", 0.80, 1.92, 2.75, accent=TEAL)
        add_metric(slide, "3", "top drivers retained per network", 3.73, 1.92, 2.75, accent=BLUE)
        add_metric(slide, "21", "network–driver sectors plotted", 6.66, 1.92, 2.75, accent=PURPLE)
        add_metric(slide, "295", "eligible primary directional runs reconciled", 9.59, 1.92, 2.75, accent=ORANGE)
        add_text(slide, "Ranking statistic", 0.86, 3.55, 2.2, 0.33, size=15, color=NAVY, bold=True)
        add_rect(slide, 0.82, 3.91, 5.63, 1.66, fill=WHITE, line=LIGHT)
        add_rich_text(slide, [
            ("S", {"size": 22, "color": NAVY, "bold": True, "italic": True}),
            ("n,g", {"size": 12, "color": NAVY, "bold": True}),
            (" = mean across eligible runs of  −log", {"size": 17, "color": DARK}),
            ("10", {"size": 11, "color": DARK}),
            ("(P", {"size": 17, "color": DARK}),
            ("KDA", {"size": 11, "color": DARK}),
            (")", {"size": 17, "color": DARK}),
        ], 1.18, 4.22, 4.88, 0.45, align=PP_ALIGN.CENTER)
        add_text(slide, "Standardized to the maximum within each broad network; coverage is explicit and used in deterministic tie-breaking.",
                 1.10, 4.82, 5.10, 0.48, size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, "Visual encoding", 6.94, 3.55, 2.2, 0.33, size=15, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Outer color band = broad network",
            "Radial bar = standardized MeanOfLog",
            "Center link = same gene selected in >1 network",
            "Gray label marker = mtDNA-encoded sentinel",
        ], 6.96, 3.96, 5.25, size=14.2, bullet_color=TEAL, line_h=0.48)
        add_rect(slide, 0.84, 6.05, 11.80, 0.64, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "CAMs and T cells had no eligible KDA runs and therefore receive no empty sector; absence is not a negative KDA result.",
                 1.08, 6.21, 11.25, 0.28, size=12.2, color=NAVY, bold=True)

        # 5 — Circular figure
        add_figure_slide(prs, 5, trimmed["circular"], "Figure 1 • Recurrent mitochondrial KDA evidence across networks",
                         "Source: phase12_kda_reduced_circular.png")

        # 6 — Circular findings
        slide = new_slide(prs, "Circular overview: findings", "Figure 1 conclusions", 6,
                          subtitle="The overview is dominated by recurrent respiratory-chain sentinels, with a small number of nuclear candidates.",
                          source="Source: phase12_kda_reduced_circular_plotted_data.tsv")
        add_metric(slide, f"{circular_unique}", "unique drivers among 21 sectors", 0.80, 1.95, 2.75, accent=PURPLE)
        add_metric(slide, f"{circular_mtdna}/21", "selected sectors are mtDNA-encoded", 3.73, 1.95, 2.75, accent=VERMILION)
        add_metric(slide, f"{circular_recurrence['MT-CO2']}", "networks select MT-CO2 in the top three", 6.66, 1.95, 2.75, accent=TEAL)
        add_metric(slide, "100%", "ranking coverage for all displayed rows", 9.59, 1.95, 2.75, accent=BLUE)
        add_rect(slide, 0.82, 3.55, 5.72, 2.45, fill=PALE_GREEN, line=TEAL)
        add_text(slide, "What the figure supports", 1.12, 3.84, 4.95, 0.38, size=18, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "MT-CO2 is the dominant recurrent candidate across broad networks.",
            "MT-CO3, MT-ATP6, MT-CYB, and MT-ND4 also recur across network top-three lists.",
            "OPCs elevate nuclear candidates RPS15 and FTL; excitatory and oligodendrocyte lists include COX6B1 and COX4I1.",
        ], 1.12, 4.32, 5.04, size=12.5, bullet_color=TEAL, line_h=0.52)
        add_rect(slide, 6.80, 3.55, 5.72, 2.45, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "What the figure does not establish", 7.10, 3.84, 4.98, 0.38, size=18, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "High rank is not proof that an mtDNA gene is an upstream causal regulator.",
            "Center links show repeated selection—not Bayesian-network edges.",
            "The overview is a sentinel map; mechanistically tractable nuclear candidates are evaluated in the next panels.",
        ], 7.10, 4.32, 5.04, size=12.5, bullet_color=ORANGE, line_h=0.52)

        # 7 — sex/APOE setup
        slide = new_slide(prs, "Sex/APOE evidence map: setup and generation", "Figure 2 setup", 7,
                          subtitle="A conservative candidate screen focuses the display on nuclear, query-independent driver hypotheses.",
                          source="Sources: sex_apoe_dot_heatmap_design.md; visualize_phase12_kda_sex_apoe.R")
        add_step(slide, "1", "Select candidate rows", "Primary directional result; non-mtDNA; driver outside the query; overlap ≥2; signature size ≥10; sufficient ranking coverage.",
                 0.80, 1.92, 5.80, accent=TEAL)
        add_step(slide, "2", "Aggregate within each stratum", "For each row × sex/APOE × direction cell, summarize every tested run—not only significant rows.",
                 6.73, 1.92, 5.80, accent=BLUE)
        add_step(slide, "3", "Encode strength and recurrence", "Color = mean −log10(P); area = fraction of tested runs significant after within-run BH correction.",
                 0.80, 3.40, 5.80, accent=PURPLE)
        add_step(slide, "4", "Keep missingness visible", "Outlined dot = tested, none significant. Gray × = no eligible/tested run. Right bar = within-network overall rank.",
                 6.73, 3.40, 5.80, accent=ORANGE)
        add_metric(slide, "16", "prioritized network–driver rows", 0.84, 5.15, 2.68, accent=TEAL)
        add_metric(slide, "192", "direction × sex/APOE display cells", 3.63, 5.15, 2.68, accent=BLUE)
        add_metric(slide, str(sex_status["tested_significant"]), "cells with ≥1 significant run", 6.42, 5.15, 2.68, accent=PURPLE)
        add_metric(slide, str(sex_status["no_eligible_or_tested_run"]), "cells with no eligible/tested run", 9.21, 5.15, 3.28, accent=ORANGE)
        add_text(slide, "The shared color cap is the prespecified 95th percentile (≈5.00); uncapped values remain in the TSV.",
                 0.96, 6.61, 11.5, 0.26, size=10.5, color=GRAY, italic=True)

        # 8 — sex/APOE figure
        add_figure_slide(prs, 8, trimmed["sex_apoe"], "Figure 2 • Sex- and APOE-stratified support for prioritized drivers",
                         "Source: phase12_kda_sex_apoe.png")

        # 9 — sex/APOE findings
        slide = new_slide(prs, "Sex/APOE map: findings", "Figure 2 conclusions", 9,
                          subtitle="Two ε2-associated directional patterns are broadest; other signals are more localized and often sparsely replicated.",
                          source="Source: sex_apoe_explanation.md; phase12_kda_sex_apoe_plotted_data.tsv")
        add_rect(slide, 0.78, 1.92, 5.86, 3.82, fill=PALE_BLUE, line=BLUE)
        add_text(slide, "Male ε2 • AD-down", 1.10, 2.22, 4.9, 0.46, size=22, color=BLUE, bold=True, font=FONT_HEAD)
        add_text(slide, "Broadest recurring pattern", 1.10, 2.67, 4.8, 0.30, size=12, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Excitatory RPL11: 8/14 significant runs",
            "Inhibitory RPS15: 9/10; LAMTOR5: 5/9",
            "Astrocytic APOE: 2/3; excitatory TMEM147: 4/10",
            "OPC and oligodendrocyte signals are 1/1 and therefore strong but minimally replicated",
        ], 1.10, 3.10, 5.05, size=13.1, bullet_color=BLUE, line_h=0.53)
        add_rect(slide, 6.81, 1.92, 5.74, 3.82, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Female ε2 • AD-up", 7.13, 2.22, 4.9, 0.46, size=22, color=VERMILION, bold=True, font=FONT_HEAD)
        add_text(slide, "Complementary recurring pattern", 7.13, 2.67, 4.8, 0.30, size=12, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Excitatory RPL11: 6/12; TMEM147: 6/11",
            "Excitatory SELENOW: 5/12",
            "Astrocytic RPL11, APOE, and RPS15: each 1/3",
            "Additional female ε4 AD-down and male ε4 signals are network-specific",
        ], 7.13, 3.10, 4.95, size=13.1, bullet_color=ORANGE, line_h=0.53)
        add_rect(slide, 1.32, 6.05, 10.64, 0.66, fill=NAVY, line=None)
        add_text(slide, "Conclusion: the strata generate hypotheses for donor-aware interaction tests; visual differences are not formal sex/APOE interactions.",
                 1.58, 6.23, 10.15, 0.30, size=13.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 10 — shared network setup
        slide = new_slide(prs, "Network figures: shared inputs and visual language", "Shared network setup", 10,
                          subtitle="The same validated production bundle, network backgrounds, degree mapping, and Complex V definition are reused across all four figures.",
                          source="Sources: network_figures_creation_plan.md; phase12_kda_network_figure_common.py")
        add_rect(slide, 0.80, 1.92, 4.03, 4.55, fill=WHITE, line=LIGHT)
        add_text(slide, "Inputs", 1.10, 2.22, 3.3, 0.42, size=20, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Phase 12 run-level KDA and complete candidate tests",
            "Broad-cell-type directed Bayesian networks",
            "Matched Phase 08 AD vs NCI log fold-change",
            "Run-specific effective backgrounds and selected layers",
            "Fixed 26-gene MitoCarta Complex V definition",
        ], 1.10, 2.82, 3.30, size=12.5, bullet_color=TEAL, line_h=0.58)
        add_rect(slide, 4.99, 1.92, 3.55, 4.55, fill=WHITE, line=LIGHT)
        add_text(slide, "Visual grammar", 5.29, 2.22, 2.9, 0.42, size=20, color=NAVY, bold=True)
        # Mini legend
        diamond = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(5.36), Inches(2.93), Inches(0.34), Inches(0.34))
        shape_fill(diamond, CYAN); shape_line(diamond, DARK, 1)
        add_text(slide, "key driver", 5.83, 2.93, 2.0, 0.27, size=11.5, color=GRAY)
        add_circle(slide, 5.36, 3.48, 0.34, fill=PALE_BLUE, line=DARK, line_width=2)
        add_text(slide, "black ring = KDA overlap", 5.83, 3.49, 2.25, 0.27, size=11.5, color=GRAY)
        add_circle(slide, 5.36, 4.03, 0.34, fill=WHITE, line=PURPLE, line_width=2.2)
        add_text(slide, "purple ring = Complex V", 5.83, 4.04, 2.25, 0.27, size=11.5, color=GRAY)
        add_arrow(slide, 5.36, 4.78, 6.10, 4.78, color=DARK, width=2.1)
        add_text(slide, "dark = highlighted path", 6.22, 4.66, 1.95, 0.30, size=11.5, color=GRAY)
        add_arrow(slide, 5.36, 5.32, 6.10, 5.32, color=LIGHT, width=1.2)
        add_text(slide, "pale = neighborhood context", 6.22, 5.20, 1.95, 0.30, size=11.5, color=GRAY)
        add_text(slide, "Node fill: blue = lower in AD; orange = higher in AD", 5.29, 5.85, 2.8, 0.42, size=10.8, color=GRAY, bold=True)
        add_rect(slide, 8.70, 1.92, 3.82, 4.55, fill=NAVY, line=None)
        add_text(slide, "Reproducibility", 9.02, 2.22, 3.0, 0.42, size=20, color=WHITE, bold=True)
        add_bullet_list(slide, [
            "Deterministic layouts and path tie-breaking",
            "Auditable node, edge, path, and plotted-point TSVs",
            "Input/output SHA-256 hashes and UTC generation log",
            "10 automated tests passed",
            "300-dpi-scale PNG plus SVG/PDF vectors",
        ], 9.02, 2.82, 3.02, size=12.3, color=WHITE, bullet_color=CYAN, line_h=0.58)

        # 11 — Wang setup
        slide = new_slide(prs, "Wang-style neighborhoods: setup and generation", "Figure 3 setup", 11,
                          subtitle="Three prespecified representative runs illustrate directed routes from nominated drivers to mitochondrial targets.",
                          source="Sources: wang_subnetworks_explanation.md; plot_phase12_kda_wang_subnetworks.py")
        data = [
            ["Panel", "Driver / context", "Selected run", "Neighborhood", "Why selected"],
            ["A", "APOE / astrocyte", "Ast GRM3 • male ε2 • AD-down", "L2 • n=19 • overlap=5", "Mechanistic ATP/TUFM example"],
            ["B", "LAMTOR5 / excitatory", "Exc L3-4 RORB CUX2 • male ε2 • AD-down", "L3 • n=27 • overlap=8", "Smallest LAMTOR5 q"],
            ["C", "GABARAPL2 / excitatory", "Exc L4-5 RORB GABRG1 • male ε2 • AD-down", "L3 • n=45 • overlap=9", "Smallest GABARAPL2 q"],
        ]
        add_table(slide, data, 0.78, 2.02, 11.78, 2.10, col_widths=[0.68, 2.25, 3.45, 2.13, 3.27], body_size=10.3, first_col_bold=True)
        add_text(slide, "How the panels were built", 0.86, 4.48, 3.0, 0.36, size=17, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Reconstruct the exact directed KDA neighborhood inside the run-specific effective background.",
            "Truncate at the result’s selected layer; preserve all neighborhood edges.",
            "Highlight prespecified shortest routes to Complex V, translation, and stress-control genes.",
            "Label only focal drivers, overlaps, targets, and path intermediates to keep 19–45-node panels readable.",
        ], 0.88, 4.95, 7.15, size=12.7, bullet_color=TEAL, line_h=0.48)
        add_rect(slide, 8.45, 4.45, 4.06, 1.65, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Selection boundary", 8.76, 4.72, 3.35, 0.36, size=16.5, color=NAVY, bold=True)
        add_text(slide, "These are illustrative runs, not the only supporting runs. Cross-run recurrence is evaluated separately in the convergence figure.",
                 8.76, 5.15, 3.34, 0.66, size=11.8, color=GRAY)

        # 12 — Wang figure
        add_figure_slide(prs, 12, trimmed["wang"], "Figure 3 • Directed KDA neighborhoods connect drivers to mitochondrial genes",
                         "Source: phase12_kda_wang_subnetworks.png")

        # 13 — Wang findings
        slide = new_slide(prs, "Wang-style neighborhoods: findings", "Figure 3 conclusions", 13,
                          subtitle="All three selected driver systems contain directed routes to mitochondrial and ATP-synthase genes.",
                          source="Source: network_figures_explained.md; path TSVs")
        findings = [
            ("APOE • astrocyte", "APOE → TUFM; APOE → ATP5PB; APOE → LDHB → ATP5F1A", "FE 17.1 • q=0.00041", TEAL),
            ("LAMTOR5 • excitatory", "LAMTOR5 → ATP5IF1; LAMTOR5 → POP7 → ATP5MC2", "FE 14.0 • q≈1×10⁻⁵", ORANGE),
            ("GABARAPL2 • excitatory", "GABARAPL2 → CHCHD2 → ATP5MC3; route through MAGEF1/SNAPC5 to PARK7", "FE 9.9 • q≈3.4×10⁻⁵", BLUE),
        ]
        for i, (title, route, stat, accent) in enumerate(findings):
            x = 0.80 + i * 4.14
            add_rect(slide, x, 1.98, 3.76, 3.52, fill=WHITE, line=accent)
            add_rect(slide, x, 1.98, 3.76, 0.12, fill=accent, line=None, radius=False)
            add_text(slide, title, x + 0.26, 2.32, 3.18, 0.54, size=18, color=NAVY, bold=True, font=FONT_HEAD)
            add_text(slide, route, x + 0.26, 3.05, 3.20, 1.25, size=13.2, color=GRAY)
            add_text(slide, stat, x + 0.26, 4.72, 3.18, 0.36, size=12.5, color=accent, bold=True)
        add_rect(slide, 1.10, 5.92, 11.12, 0.75, fill=NAVY, line=None)
        add_text(slide, "Conclusion: the selected candidate systems connect mitochondrial translation and stress-control genes to Complex V components—within directed model neighborhoods, not as experimentally proven regulatory chains.",
                 1.42, 6.10, 10.48, 0.37, size=12.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 14 — sex reversal setup
        slide = new_slide(prs, "Sex-reversal networks: setup and generation", "Figure 4 setup", 14,
                          subtitle="Female ε2 AD-up and male ε2 AD-down neighborhoods are aligned to separate topology, membership, and expression color.",
                          source="Sources: sex_reversal_networks_explanation.md; plot_phase12_kda_sex_reversal_networks.py")
        add_rect(slide, 0.80, 1.95, 5.68, 3.82, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Female ε2", 1.10, 2.25, 4.95, 0.43, size=23, color=VERMILION, bold=True, font=FONT_HEAD)
        add_text(slide, "AD-up mitochondrial signature", 1.10, 2.72, 4.95, 0.34, size=14, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Condition-specific neighborhood reconstructed for APOE, LAMTOR5, and GABARAPL2",
            "Matched Phase 08 AD-vs-NCI logFC colors the active nodes",
            "KDA overlap rings are condition-specific",
        ], 1.10, 3.26, 4.88, size=13.0, bullet_color=ORANGE, line_h=0.62)
        add_rect(slide, 6.85, 1.95, 5.68, 3.82, fill=PALE_BLUE, line=BLUE)
        add_text(slide, "Male ε2", 7.15, 2.25, 4.95, 0.43, size=23, color=BLUE, bold=True, font=FONT_HEAD)
        add_text(slide, "AD-down mitochondrial signature", 7.15, 2.72, 4.95, 0.34, size=14, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "The same three driver systems and matched cell contexts",
            "The row’s union graph is computed once and uses identical coordinates",
            "Elements absent in one condition are faded rather than removed",
        ], 7.15, 3.26, 4.88, size=13.0, bullet_color=BLUE, line_h=0.62)
        add_arrow(slide, 5.98, 3.80, 6.66, 3.80, color=PURPLE, width=2.0)
        add_rect(slide, 1.55, 6.08, 10.25, 0.58, fill=PALE_PURPLE, line=PURPLE)
        add_text(slide, "Shared coordinates make orange-to-blue expression reversals directly visible; they do not test whether network topology differs by sex.",
                 1.83, 6.23, 9.70, 0.27, size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

        # 15 — sex reversal figure
        add_figure_slide(prs, 15, trimmed["sex_reversal"], "Figure 4 • Sex-reversed KDA signals retain shared driver-centered structure",
                         "Source: phase12_kda_sex_reversal_networks.png")

        # 16 — sex reversal findings
        slide = new_slide(prs, "Sex-reversal networks: findings", "Figure 4 conclusions", 16,
                          subtitle="The same driver-centered framework accompanies opposite descriptive expression directions in ε2 female and male strata.",
                          source="Source: sex_reversal_networks_explanation.md; node TSV")
        values = [
            ["System / gene", "Female ε2 AD-up logFC", "Male ε2 AD-down logFC"],
            ["Astrocyte APOE", "+0.785", "−0.506"],
            ["TUFM", "+0.506", "−0.939"],
            ["ATP5PB", "+0.153", "−1.151"],
            ["Excitatory LAMTOR5", "+0.505", "−0.928"],
            ["ATP5IF1", "+0.420", "−0.497"],
            ["Excitatory GABARAPL2", "+0.463", "−0.573"],
            ["ATP5MC3", "+0.430", "−0.491"],
        ]
        add_table(slide, values, 0.80, 1.95, 6.62, 4.63, col_widths=[2.86, 1.88, 1.88], body_size=10.8, first_col_bold=True)
        add_rect(slide, 7.75, 1.95, 4.77, 2.08, fill=PALE_GREEN, line=TEAL)
        add_text(slide, "Supported interpretation", 8.05, 2.25, 4.12, 0.37, size=18, color=NAVY, bold=True)
        add_text(slide, "APOE, LAMTOR5, and GABARAPL2 remain linked to mitochondrial/Complex V genes in both ε2 strata, while matched AD effects reverse direction.",
                 8.05, 2.78, 4.05, 0.92, size=13.1, color=GRAY)
        add_rect(slide, 7.75, 4.34, 4.77, 2.24, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Required caution", 8.05, 4.64, 4.12, 0.37, size=18, color=NAVY, bold=True)
        add_text(slide, "No formal AD×sex or AD×APOE interaction was fitted here. Networks are reused across sex, donor groups are small/unequal, and the two strata are not independent replications.",
                 8.05, 5.17, 4.05, 1.06, size=12.5, color=GRAY)

        # 17 — ATP setup
        slide = new_slide(prs, "ATP convergence map: setup and generation", "Figure 5 setup", 17,
                          subtitle="Unlike the selected neighborhoods, this bipartite map synthesizes qualifying primary directional calls across runs.",
                          source="Sources: atp_convergence_explanation.md; plot_phase12_kda_atp_convergence.py")
        filters = [
            "Primary AD-up or AD-down KDA run",
            "Driver is nuclear and outside its overlap",
            "Overall overlap ≥2 genes",
            "Mitochondrial signature ≥10 genes",
            "At least one overlap gene in fixed 26-gene Complex V set",
            "Deduplicate by KDA run ID for each network–driver–target",
        ]
        add_rect(slide, 0.80, 1.92, 5.78, 4.72, fill=WHITE, line=LIGHT)
        add_text(slide, "Qualifying-call filter", 1.10, 2.23, 4.95, 0.42, size=20, color=NAVY, bold=True)
        add_bullet_list(slide, filters, 1.10, 2.82, 4.95, size=12.7, bullet_color=TEAL, line_h=0.55)
        add_rect(slide, 6.83, 1.92, 5.68, 4.72, fill=NAVY, line=None)
        add_text(slide, "Encoding", 7.15, 2.23, 4.8, 0.42, size=20, color=WHITE, bold=True)
        add_bullet_list(slide, [
            "Left nodes: selected driver × broad network; color = network",
            "Driver size: full-network degree using the shared Wang mapping",
            "Right nodes: Complex V genes; area = total supporting calls",
            "Edge width: number of qualifying runs",
            "Line style: directed distance 1, 2, or 3",
        ], 7.15, 2.84, 4.78, size=12.7, color=WHITE, bullet_color=CYAN, line_h=0.58)
        add_metric(slide, str(focused_driver_targets), "focused driver–target relationships", 7.18, 5.67, 2.42,
                   accent=CYAN, fill=NAVY_2, label_color=RGBColor(207, 220, 234))
        add_metric(slide, str(all_driver_targets), "complete qualifying combinations", 9.78, 5.67, 2.42,
                   accent=ORANGE, fill=NAVY_2, label_color=RGBColor(207, 220, 234))

        # 18 — ATP figure
        add_figure_slide(prs, 18, trimmed["atp"], "Figure 5 • Recurrent convergence on ATP synthase / Complex V genes",
                         "Source: phase12_kda_atp_convergence.png")

        # 19 — ATP findings
        slide = new_slide(prs, "ATP convergence map: findings", "Figure 5 conclusions", 19,
                          subtitle="Complex V convergence recurs across fine-cell-type, sex, and APOE contexts—especially in neuronal networks.",
                          source="Source: phase12_kda_atp_convergence_pairs.tsv")
        data = [
            ["Driver / network", "Complex V target", "Qualifying calls", "Direction mix", "Distance"],
            ["GABARAPL2 / excitatory", "ATP5MC3", "15", "4 up / 11 down", "2"],
            ["RPL11 / excitatory", "ATP5PF", "14", "4 up / 10 down", "2"],
            ["LAMTOR5 / excitatory", "ATP5IF1", "12", "4 up / 8 down", "1"],
            ["LAMTOR5 / excitatory", "ATP5MC2", "12", "4 up / 8 down", "2"],
            ["RPL11 / excitatory", "ATP5ME", "12", "1 up / 11 down", "3"],
            ["RPS15 / inhibitory", "ATP5F1E", "6", "1 up / 5 down", "2"],
        ]
        add_table(slide, data, 0.80, 1.95, 8.10, 4.40, col_widths=[2.40, 1.56, 1.45, 1.70, 0.99], body_size=10.8, first_col_bold=True)
        add_rect(slide, 9.17, 1.95, 3.35, 4.40, fill=PALE_PURPLE, line=PURPLE)
        add_text(slide, "Takeaway", 9.48, 2.25, 2.72, 0.40, size=20, color=NAVY, bold=True)
        add_text(slide, "The ATP finding is not confined to one illustrative run.", 9.48, 2.87, 2.72, 0.78, size=17, color=PURPLE, bold=True)
        add_text(slide, "Several driver–Complex V relationships recur across analysis contexts. Edge width pools AD-up and AD-down calls, so recurrence does not imply one universal disease direction.",
                 9.48, 3.88, 2.72, 1.60, size=12.9, color=GRAY)
        add_text(slide, "Strongest family: excitatory GABARAPL2 → ATP5MC3 (15 calls)", 9.48, 5.66, 2.72, 0.45, size=11.2, color=NAVY, bold=True)

        # 20 — connectivity setup
        slide = new_slide(prs, "Connectivity diagnostic: setup and generation", "Figure 6 setup", 20,
                          subtitle="The question is whether aggregate KDA evidence rises with full-network connectivity—and how strongly that varies by network.",
                          source="Sources: connectivity_evidence_explanation.md; plot_phase12_kda_connectivity_evidence.py")
        add_rect(slide, 0.80, 1.92, 3.70, 4.64, fill=WHITE, line=LIGHT)
        add_text(slide, "One plotted record", 1.10, 2.23, 3.05, 0.40, size=20, color=NAVY, bold=True)
        add_text(slide, "candidate gene × broad network", 1.10, 2.72, 3.05, 0.38, size=14, color=TEAL, bold=True)
        add_metric(slide, "50,165", "records with ≥1 ranking test", 1.10, 3.28, 2.92, accent=TEAL)
        add_text(slide, "A gene can appear in multiple networks because degree and KDA ranking are network-specific.",
                 1.10, 4.75, 3.02, 0.85, size=12.6, color=GRAY)
        add_rect(slide, 4.78, 1.92, 3.70, 4.64, fill=PALE_BLUE, line=BLUE)
        add_text(slide, "Panel A", 5.08, 2.23, 3.05, 0.40, size=20, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "x = total-degree percentile within network",
            "y = standardized MeanOfLog KDA evidence",
            "area = significant primary directional calls",
            "color = broad network",
            "labels = prespecified candidates",
        ], 5.08, 2.86, 3.00, size=12.4, bullet_color=BLUE, line_h=0.58)
        add_rect(slide, 8.75, 1.92, 3.77, 4.64, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Panel B", 9.05, 2.23, 3.05, 0.40, size=20, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Separate Spearman correlation within each network",
            "Rank-based; no linearity or normality assumption",
            "Effect magnitude (rho) matters more than tiny P values at n≈5k–10k",
            "Raw-degree facets retained as diagnostic supplement",
        ], 9.05, 2.86, 3.04, size=12.4, bullet_color=ORANGE, line_h=0.66)

        # 21 — connectivity figure
        add_figure_slide(prs, 21, trimmed["connectivity"], "Figure 6 • Network connectivity is contextual evidence",
                         "Source: phase12_kda_connectivity_evidence.png")

        # 22 — connectivity findings
        slide = new_slide(prs, "Connectivity diagnostic: findings", "Figure 6 conclusions", 22,
                          subtitle="Connectivity and aggregate KDA evidence are positively associated in every network, but connectivity is not determinative.",
                          source="Source: phase12_kda_connectivity_evidence_correlations.tsv")
        corr_table = [["Network", "Candidate records", "Spearman ρ"]]
        pretty = {
            "Astrocytes": "Astrocytes", "Excitatory_neurons": "Excitatory neurons",
            "Inhibitory_neurons": "Inhibitory neurons", "Microglia": "Microglia",
            "OPCs": "OPCs", "Oligodendrocytes": "Oligodendrocytes", "Vasculature_cells": "Vasculature",
        }
        for r in corr:
            corr_table.append([pretty[r["broad_network"]], f"{int(r['n_candidates']):,}", f"{float(r['spearman_rho']):.3f}"])
        add_table(slide, corr_table, 0.80, 1.95, 5.35, 4.65, col_widths=[2.45, 1.65, 1.25], body_size=10.8, first_col_bold=True)
        add_rect(slide, 6.45, 1.95, 6.08, 2.08, fill=PALE_GREEN, line=TEAL)
        add_text(slide, "Signal", 6.78, 2.25, 1.0, 0.38, size=18, color=TEAL, bold=True)
        add_text(slide, "Highest associations occur in excitatory neurons (ρ=0.559) and astrocyte/inhibitory networks (ρ≈0.525).",
                 7.74, 2.20, 4.38, 0.85, size=14.4, color=NAVY, bold=True)
        add_text(slide, "All seven correlations are positive; the weakest is vasculature (ρ=0.170).",
                 7.74, 3.12, 4.38, 0.50, size=12.4, color=GRAY)
        add_rect(slide, 6.45, 4.34, 6.08, 2.26, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Boundary", 6.78, 4.64, 1.0, 0.38, size=18, color=VERMILION, bold=True)
        add_text(slide, "High-degree candidates still span a wide range of KDA scores. Connectivity can shape opportunity for overlap but does not replace KDA enrichment, recurrence, or causal validation.",
                 7.74, 4.57, 4.38, 1.18, size=14.0, color=NAVY, bold=True)
        add_text(slide, "Use connectivity as context—not as an independent significance test.",
                 7.74, 5.86, 4.38, 0.42, size=12.5, color=GRAY, italic=True)

        # 23 — network synthesis
        slide = new_slide(prs, "What the four network figures establish together", "Cross-figure synthesis", 23,
                          subtitle="Selected mechanisms, aligned sex contrasts, cross-run recurrence, and a network-bias diagnostic converge on one cautious interpretation.",
                          source="Sources: Phase 12 network figure explanations and plotted-data tables")
        rows = [
            ("Selected neighborhoods", "APOE, LAMTOR5, and GABARAPL2 have directed routes to Complex V and mitochondrial stress/translation genes.", TEAL),
            ("Sex-reversal alignment", "The same driver-centered framework accompanies female ε2 AD-up and male ε2 AD-down expression directions.", ORANGE),
            ("Cross-run convergence", "Complex V targets recur across contexts; the strongest focused link is GABARAPL2→ATP5MC3.", PURPLE),
            ("Connectivity diagnostic", "Hubs tend to rank more strongly, but degree explains only part of the KDA evidence and varies by network.", BLUE),
        ]
        for i, (t, b, accent) in enumerate(rows):
            y = 1.88 + i * 1.13
            add_rect(slide, 0.82, y, 11.72, 0.91, fill=WHITE, line=LIGHT)
            add_rect(slide, 0.82, y, 0.14, 0.91, fill=accent, line=None, radius=False)
            add_text(slide, t, 1.20, y + 0.17, 2.53, 0.34, size=15, color=NAVY, bold=True)
            add_text(slide, b, 3.80, y + 0.13, 8.32, 0.50, size=13.2, color=GRAY)
        add_rect(slide, 1.18, 6.52, 11.02, 0.44, fill=NAVY, line=None)
        add_text(slide, "Integrated Phase 12 conclusion: recurrent, cell-contextual network convergence on mitochondrial/Complex V biology—plausible and auditable, but still hypothesis-generating.",
                 1.42, 6.59, 10.54, 0.24, size=12.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 24 — Phase 13 setup
        slide = new_slide(prs, "Phase 13: donor-aware respiratory-modifier design", "Figure 7 shared setup", 24,
                          subtitle="The formal question is whether the adjusted AD effect on a respiratory module differs between two sex/APOE groups.",
                          source="Sources: phase_13_workflow_explained.md; phase_13_respiratory_modifier_plan.md")
        # Flow
        x_positions = [0.70, 3.20, 5.70, 8.20, 10.70]
        flow = [
            ("Nucleus counts", "Sum within donor × cell context", TEAL),
            ("Donor profiles", "TMM + logCPM using all genes", BLUE),
            ("Module scores", "NCI-reference gene z-scores averaged", PURPLE),
            ("Six AD effects", "Adjusted AD−NCI per sex/APOE stratum", ORANGE),
            ("Modifier test", "Difference of two AD effects", VERMILION),
        ]
        for i, (t, b, accent) in enumerate(flow):
            x = x_positions[i]
            add_rect(slide, x, 2.10, 1.95, 1.48, fill=WHITE, line=accent)
            add_rect(slide, x, 2.10, 1.95, 0.10, fill=accent, line=None, radius=False)
            add_text(slide, t, x + 0.17, 2.39, 1.61, 0.38, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, b, x + 0.14, 2.82, 1.67, 0.54, size=9.6, color=GRAY, align=PP_ALIGN.CENTER)
            if i < 4:
                add_arrow(slide, x + 2.00, 2.84, x + 2.38, 2.84, color=MID, width=1.4)
        add_rect(slide, 0.78, 4.10, 5.74, 2.15, fill=PALE_BLUE, line=BLUE)
        add_text(slide, "Difference-of-differences", 1.08, 4.40, 5.10, 0.40, size=19, color=NAVY, bold=True)
        add_rich_text(slide, [
            ("estimate = Δ(group 1) − Δ(group 2)", {"size": 20, "color": BLUE, "bold": True}),
        ], 1.12, 5.05, 5.02, 0.48, align=PP_ALIGN.CENTER)
        add_text(slide, "Δ = model-adjusted AD module-score mean − NCI module-score mean", 1.10, 5.62, 5.10, 0.33, size=11.2, color=GRAY, align=PP_ALIGN.CENTER)
        add_rect(slide, 6.82, 4.10, 5.70, 2.15, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Prespecified test family", 7.12, 4.40, 5.05, 0.40, size=19, color=NAVY, bold=True)
        add_text(slide, "7 contexts × 7 contrasts × 4 modules = 196 tests", 7.12, 5.07, 5.05, 0.38,
                 size=16.5, color=VERMILION, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Effect unit: NCI-reference donor-level module-score SD", 7.12, 5.58, 5.05, 0.33,
                 size=11.2, color=GRAY, align=PP_ALIGN.CENTER)

        # 25 — Phase13 landscape generation
        slide = new_slide(prs, "Modifier landscape: setup and generation", "Figure 7 setup", 25,
                          subtitle="The renderer reads the immutable validated gate table; it does not refit models or alter q values.",
                          source="Sources: phase13_modifier_landscape_methods.md; plot_phase13_respiratory_modifier_figures.R")
        modules = [
            ("mtDNA OXPHOS", "13 genes • direct respiratory", TEAL),
            ("Nuclear OXPHOS", "86 structural genes • direct respiratory", BLUE),
            ("Mitochondrial translation", "155 genes • supporting program", PURPLE),
            ("MIB/MICOS membrane", "19 genes • supporting program", ORANGE),
        ]
        for i, (t, b, accent) in enumerate(modules):
            x = 0.78 + i * 3.13
            add_rect(slide, x, 1.92, 2.85, 1.30, fill=WHITE, line=accent)
            add_text(slide, t, x + 0.18, 2.16, 2.50, 0.38, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, b, x + 0.18, 2.65, 2.50, 0.28, size=9.5, color=GRAY, align=PP_ALIGN.CENTER)
        add_rect(slide, 0.78, 3.62, 5.72, 2.68, fill=PALE_PURPLE, line=PURPLE)
        add_text(slide, "Heatmap encoding", 1.08, 3.92, 5.10, 0.40, size=19, color=NAVY, bold=True)
        add_bullet_list(slide, [
            "Rows = 7 frozen broad cell contexts",
            "Columns = 3 sex + 4 APOE contrasts",
            "Color = signed modifier estimate (−1.5 to +1.5 SD)",
            "Crossed gray cell = not testable, never zero",
            "Outlines/symbols redundantly encode frozen status",
        ], 1.08, 4.40, 4.98, size=12.2, bullet_color=PURPLE, line_h=0.36)
        add_rect(slide, 6.82, 3.62, 5.70, 2.68, fill=NAVY, line=None)
        add_text(slide, "Scientific gate", 7.12, 3.92, 5.05, 0.40, size=19, color=WHITE, bold=True)
        add_bullet_list(slide, [
            "Complete-family module-score FDR",
            "Independent CAMERA gene-set support",
            "Donor-count and module-coverage thresholds",
            "Confidence interval and meaningful-effect rule",
            "Bootstrap, leave-one-donor-out, and sensitivity stability",
        ], 7.12, 4.42, 4.95, size=12.0, color=WHITE, bullet_color=CYAN, line_h=0.36)
        add_text(slide, "No significance stars; exact gate status controls the conclusion.",
                 7.18, 6.43, 4.82, 0.28, size=10.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

        # 26 — Phase13 figure
        add_figure_slide(prs, 26, trimmed["modifier"], "Figure 7 • Respiratory-modifier landscape: complete prespecified family",
                         "Source: phase13_modifier_landscape.png")

        # 27 — Phase13 findings
        slide = new_slide(prs, "Modifier landscape: findings", "Figure 7 conclusions", 27,
                          subtitle="The production run is technically complete, but the complete-family scientific result is inconclusive.",
                          source="Sources: phase13_modifier_landscape_caption.md; respiratory gate decisions")
        add_metric(slide, str(len(modifier)), "prespecified tests rendered", 0.80, 1.92, 2.72, accent=BLUE)
        add_metric(slide, str(modifier_status["inconclusive"]), "estimable but inconclusive", 3.67, 1.92, 2.72, accent=PURPLE)
        add_metric(slide, str(modifier_status["not_testable"]), "not testable (all vasculature)", 6.54, 1.92, 2.72, accent=ORANGE)
        add_metric(slide, str(modifier_status["supported"]), "supported results", 9.41, 1.92, 2.72, accent=VERMILION)
        add_rect(slide, 0.80, 3.57, 7.40, 2.78, fill=WHITE, line=LIGHT)
        add_text(slide, "Complete-family evidence", 1.10, 3.87, 6.78, 0.40, size=19, color=NAVY, bold=True)
        add_bullet_list(slide, [
            f"{nominal_score} module-score tests had nominal P<0.05, but none passed FDR.",
            f"Minimum module-score q = {min_q_score:.3f}; minimum CAMERA q = {min_q_camera:.3f}.",
            "The direct-respiratory nominal highlight is OPC nuclear OXPHOS, female−male within ε3/3: estimate +0.815 SD, 95% CI [0.130, 1.501], P=0.0199, q=0.876.",
            "Heatmap colors show effect estimates—not support status; visually strong cells can remain inconclusive.",
        ], 1.10, 4.42, 6.62, size=12.4, bullet_color=BLUE, line_h=0.47)
        add_rect(slide, 8.48, 3.57, 4.05, 2.78, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, "Permitted conclusion", 8.79, 3.87, 3.42, 0.40, size=19, color=NAVY, bold=True)
        add_text(slide, "No respiratory module showed a supported sex/APOE modifier after the prespecified complete-family gate.",
                 8.79, 4.48, 3.42, 0.98, size=16.2, color=VERMILION, bold=True)
        add_text(slide, "This is evidence of unresolved uncertainty—not proof that all modifier effects are exactly zero.",
                 8.79, 5.63, 3.42, 0.50, size=11.8, color=GRAY, italic=True)

        # 28 — Integrated interpretation
        slide = new_slide(prs, "Integrated interpretation: Phase 12 and Phase 13", "Cross-phase synthesis", 28,
                          subtitle="The two phases operate at different evidentiary levels and therefore answer different questions.",
                          source="Sources: validated Phase 12 and Phase 13 figure packages")
        data = [
            ["Dimension", "Phase 12 KDA", "Phase 13 modifier analysis"],
            ["Unit of analysis", "Directed candidate neighborhood × run", "Donor-level respiratory module score"],
            ["Primary question", "Which candidates’ neighborhoods are enriched?", "Does the AD effect differ by sex/APOE?"],
            ["Strength", "Localization, topology, recurrence", "Formal difference-of-differences with FDR/stability gate"],
            ["Result", "Recurrent cell-contextual Complex V convergence", "0/196 supported; 180 inconclusive; 16 not testable"],
            ["Conclusion level", "Mechanistic hypotheses", "No confirmed respiratory modifier in this test family"],
        ]
        add_table(slide, data, 0.78, 1.92, 11.78, 3.98, col_widths=[2.20, 4.68, 4.90], body_size=11.4, first_col_bold=True)
        add_rect(slide, 1.03, 6.20, 11.28, 0.64, fill=NAVY, line=None)
        add_text(slide, "There is no contradiction: Phase 12 supports where and how mitochondrial network convergence may occur; Phase 13 did not confirm a broad donor-level sex/APOE modifier after complete-family correction.",
                 1.32, 6.34, 10.70, 0.32, size=12.1, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 29 — Final implications
        slide = new_slide(prs, "Implications and next analytical steps", "Closing", 29,
                          subtitle="Preserve the mechanistic hypotheses while raising the evidentiary bar for sex/APOE claims.",
                          source="Synthesis of project-local figure methods, tables, and validation notes")
        cards = [
            ("Prioritize", "Use APOE, LAMTOR5, GABARAPL2, RPL11, and RPS15–Complex V systems as focused hypotheses, not a validated causal ranking.", TEAL, PALE_GREEN),
            ("Validate", "Test nominated driver→Complex V relationships with perturbation or orthogonal regulatory evidence; retain cell-network context.", BLUE, PALE_BLUE),
            ("Re-test modifiers", "Use donor-aware interaction models in larger/independent cohorts, especially ε2 and vasculature strata with limited testability.", ORANGE, PALE_ORANGE),
            ("Report transparently", "Keep full test families, not-testable states, q values, donor support, and stability checks visible in figures and text.", PURPLE, PALE_PURPLE),
        ]
        for i, (t, b, accent, fill) in enumerate(cards):
            x = 0.80 + (i % 2) * 6.05
            y = 1.95 + (i // 2) * 2.25
            add_rect(slide, x, y, 5.68, 1.82, fill=fill, line=accent)
            add_circle(slide, x + 0.28, y + 0.30, 0.44, fill=accent)
            add_text(slide, str(i + 1), x + 0.28, y + 0.32, 0.44, 0.32, size=14, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, t, x + 0.90, y + 0.25, 4.30, 0.38, size=19, color=NAVY, bold=True)
            add_text(slide, b, x + 0.90, y + 0.75, 4.30, 0.78, size=12.6, color=GRAY)
        add_rect(slide, 1.40, 6.42, 10.53, 0.54, fill=NAVY, line=None)
        add_text(slide, "Bottom line: reproducible network evidence points to mitochondrial/Complex V systems; formal sex/APOE respiratory modification remains unconfirmed.",
                 1.66, 6.48, 10.00, 0.38, size=11.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)

        OUT_DIR.mkdir(parents=True, exist_ok=True)
        prs.save(OUT_PATH)

    return OUT_PATH


if __name__ == "__main__":
    output = build_deck()
    print(output)
