#!/usr/bin/env python3
"""Build a concise PowerPoint explaining the fit limits of GSE143758.

The deck is intentionally a dataset-fit assessment, not an analysis-results
deck.  Its sole scientific source is the user-provided Word document
``GSE143758_dataset_does_not_work.docx``.  The presentation preserves the
important distinction that GSE143758 is useful for narrower 5xFAD disease and
astrocyte questions while being unsuitable for the intended APOE-by-sex
validation endpoint.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import os
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
SOURCE_DOC = REPO / "docs/validation_mouse/GSE143758_dataset_does_not_work.docx"
SOURCE_DOC_SHA256 = "5bfdc1563178dba3aa33cb54e1b325728189890ffef02afc8350041311c6934d"
DEFAULT_OUT = REPO / (
    "docs/presentations/gse143758_mouse_validation_limitations_08232026.pptx"
)
DEFAULT_REPORT_DIR = REPO / (
    "results/presentations/validation_mouse/gse143758_mouse_validation_limitations"
)

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)
FONT = "Arial"
EXPECTED_SLIDES = 5
NOTE_HEADINGS = ("What to point at:", "Main takeaway:", "Boundary / transition:")

NAVY = RGBColor(15, 35, 61)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
ORANGE = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(220, 228, 237)
MID = RGBColor(96, 108, 123)
DARK = RGBColor(35, 42, 51)
PALE_BLUE = RGBColor(224, 239, 248)
PALE_GREEN = RGBColor(224, 244, 238)
PALE_ORANGE = RGBColor(255, 244, 216)
PALE_RED = RGBColor(253, 235, 228)
PALE_GRAY = RGBColor(239, 243, 247)

TITLES = [
    "GSE143758 is valuable, but it cannot validate APOE-by-sex effects",
    "Thirty-seven GEO records reduce to eight independent mice in the main comparison",
    "The intended sex and APOE endpoints are not testable",
    "The target is cortex/PFC, but GSE143758 is primarily hippocampal",
    "Use GSE143758 only as an auxiliary 5xFAD disease and cell-state reference",
]


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def display_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(REPO))
    except ValueError:
        return str(path.resolve())


def read_docx_text(path: Path) -> str:
    """Return visible DOCX text in document order without extra dependencies."""
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("word/document.xml")
    root = ET.fromstring(xml)
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    values = [node.text or "" for node in root.iter(f"{ns}t")]
    return " ".join(value.strip() for value in values if value.strip())


def validate_source_contract(path: Path = SOURCE_DOC) -> str:
    if not path.exists():
        raise FileNotFoundError(path)
    if sha256(path) != SOURCE_DOC_SHA256:
        raise AssertionError("Source DOCX does not match the approved regional-fit revision")
    text = read_docx_text(path)
    required = [
        "GSE143758",
        "37 GSM records",
        "10 sample/library preparations",
        "8 independent mice",
        "4 WT and 4 5xFAD",
        "54,769 high-quality nuclei",
        "one female WT and one female 5xFAD",
        "does not compare human APOE2, APOE3, or APOE4 genotypes",
        "APOE x sex x disease interaction",
        "not a suitable primary dataset for a human prefrontal-cortex (PFC) APOE-by-sex validation study",
        "Main all-cell atlas and age course: hippocampus",
        "A smaller astrocyte-only cortex/PFC subset",
        "Region-matched broad cell types",
        "cortex resources are primarily astrocyte-focused",
        "the same mice also contributed hippocampus",
        "Descriptive cross-region astrocyte-state check",
        "direct PFC, APOE, or sex-interaction validation",
    ]
    missing = [phrase for phrase in required if phrase.lower() not in text.lower()]
    if missing:
        raise AssertionError(f"Source document contract changed; missing: {missing}")
    return text


def fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def stroke(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_run(run, *, size: float, color: RGBColor, bold: bool = False,
            italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             color: RGBColor = WHITE, outline: RGBColor | None = LIGHT,
             radius: bool = True, line_width: float = 1.0):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        kind, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline, line_width)
    return shape


def add_circle(slide, x: float, y: float, d: float, *, color: RGBColor,
               outline: RGBColor | None = None, line_width: float = 1.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline, line_width)
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *,
             size: float = 16, color: RGBColor = DARK, bold: bool = False,
             italic: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.03):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic)
    return box


def add_rich_text(slide, spans: list[tuple[str, dict]], x: float, y: float,
                  w: float, h: float, *, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    frame.margin_top = frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    for value, style in spans:
        run = paragraph.add_run()
        run.text = value
        set_run(
            run,
            size=style.get("size", 16),
            color=style.get("color", DARK),
            bold=style.get("bold", False),
            italic=style.get("italic", False),
        )
    return box


def add_bullets(slide, items: list[str], x: float, y: float, w: float, *,
                size: float = 15.5, accent: RGBColor = BLUE,
                line_h: float = 0.72) -> None:
    for index, item in enumerate(items):
        cy = y + index * line_h
        add_circle(slide, x, cy + 0.18, 0.11, color=accent)
        add_text(
            slide, item, x + 0.24, cy, w - 0.24, line_h,
            size=size, color=DARK, valign=MSO_ANCHOR.MIDDLE,
        )


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, *,
                  color: RGBColor = NAVY, width: float = 2.0):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    stroke(connector, color, width)
    connector.line.end_arrowhead = True
    return connector


def new_slide(prs: Presentation, *, bg: RGBColor = WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_header(slide, kicker: str, title: str, page_no: int, *,
               accent: RGBColor = BLUE, subtitle: str | None = None) -> None:
    add_text(
        slide, kicker.upper(), 0.55, 0.18, 5.9, 0.24,
        size=9.5, color=accent, bold=True,
    )
    title_size = 25.0 if len(title) < 72 else 22.5
    add_text(
        slide, title, 0.55, 0.44, 11.95, 0.72,
        size=title_size, color=NAVY, bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_text(slide, subtitle, 0.57, 1.10, 11.70, 0.28,
                 size=11.2, color=MID)
    add_text(
        slide, f"{page_no:02d}", 12.42, 0.21, 0.36, 0.20,
        size=9.0, color=MID, bold=True, align=PP_ALIGN.RIGHT,
    )


def add_source(slide, value: str) -> None:
    add_text(slide, value, 0.55, 7.23, 12.20, 0.16, size=6.8, color=MID)


def add_note(slide, value: str) -> None:
    slide.notes_slide.notes_text_frame.text = value


def add_ribbon(slide, value: str, *, y: float = 6.30,
               fill_color: RGBColor = NAVY, accent: RGBColor = ORANGE) -> None:
    add_rect(slide, 0.65, y, 12.03, 0.58, color=fill_color,
             outline=None, radius=False)
    add_rect(slide, 0.65, y, 0.10, 0.58, color=accent,
             outline=None, radius=False)
    add_text(
        slide, value, 0.90, y + 0.10, 11.48, 0.34,
        size=13.0, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )


def add_status_card(slide, x: float, y: float, w: float, h: float, *,
                    status: str, title: str, body: str,
                    accent: RGBColor, bg: RGBColor) -> None:
    add_rect(slide, x, y, w, h, color=bg, outline=accent, line_width=1.6)
    add_rect(slide, x, y, w, 0.13, color=accent, outline=None, radius=False)
    add_text(slide, status, x + 0.20, y + 0.28, 0.92, 0.42,
             size=15.0, color=accent, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, x + 1.18, y + 0.24, w - 1.42, 0.50,
             size=16.5, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, body, x + 0.25, y + 0.88, w - 0.50, h - 1.08,
             size=13.7, color=DARK)


def _slide_text(slide) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            values.append(shape.text)
    return "\n".join(values)


def build_deck(output_path: Path = DEFAULT_OUT) -> Path:
    validate_source_contract(SOURCE_DOC)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "GSE143758 Mouse Validation Dataset-Fit Assessment"
    prs.core_properties.subject = (
        "Why GSE143758 cannot answer the intended APOE-by-sex and cortex/PFC validation question"
    )
    prs.core_properties.author = "Alzheimer project analysis team"
    prs.core_properties.keywords = (
        "GSE143758, 5xFAD, snRNA-seq, APOE, sex, cortex, PFC, mouse validation, dataset fit"
    )
    prs.core_properties.comments = (
        "Generated from docs/validation_mouse/GSE143758_dataset_does_not_work.docx. "
        "This is a dataset-fit assessment, not a new analysis of expression data. "
        f"Source DOCX SHA-256: {SOURCE_DOC_SHA256}."
    )

    # Slide 1 — decision
    slide = new_slide(prs, bg=OFF_WHITE)
    add_header(
        slide, "Dataset-fit decision", TITLES[0], 1, accent=VERMILION,
        subtitle="The dataset is useful for narrower 5xFAD questions; the intended factorial endpoint is the mismatch.",
    )
    add_rect(slide, 0.70, 1.62, 3.35, 4.30, color=PALE_RED,
             outline=VERMILION, line_width=2.0)
    add_text(slide, "NO-GO", 1.05, 2.05, 2.65, 0.72,
             size=32.0, color=VERMILION, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "for APOE × sex × disease\nvalidation", 1.02, 2.88,
             2.70, 0.96, size=19.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 1.05, 4.23, 2.65, 0.05, color=VERMILION,
             outline=None, radius=False)
    add_text(slide, "Design mismatch—not a low-quality dataset", 1.02, 4.55,
             2.70, 0.68, size=14.0, color=DARK, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_status_card(slide, 4.45, 1.62, 3.80, 1.78,
                    status="YES", title="AD-related mouse snRNA-seq",
                    body="5xFAD versus WT with broad brain cell types.",
                    accent=TEAL, bg=PALE_GREEN)
    add_status_card(slide, 8.55, 1.62, 3.80, 1.78,
                    status="YES", title="Disease and cell-state biology",
                    body="Useful for astrocyte and age questions, mainly in hippocampus.",
                    accent=BLUE, bg=PALE_BLUE)
    add_status_card(slide, 4.45, 3.72, 3.80, 1.78,
                    status="NO", title="Human APOE isoforms",
                    body="No human APOE3 or APOE4 genotype groups.",
                    accent=VERMILION, bg=PALE_RED)
    add_status_card(slide, 8.55, 3.72, 3.80, 1.78,
                    status="NO", title="Replicated female groups",
                    body="Only one reported female WT and one female 5xFAD mouse.",
                    accent=ORANGE, bg=PALE_ORANGE)
    add_ribbon(slide, "The APOE × sex × disease interaction is not testable in this design.",
               y=6.22, accent=VERMILION)
    add_source(slide, "Source: GSE143758 assessment DOCX, Executive summary and Tables 1–2.")
    add_note(slide, """What to point at:
Start with the red no-go card, then contrast the two green/blue strengths with the two missing design factors. The dataset is genuine mouse brain single-nucleus RNA-seq and is useful for 5xFAD disease and astrocyte questions. The problem is alignment with the requested APOE-by-sex endpoint.

Main takeaway:
GSE143758 is not intrinsically a bad dataset. It is the wrong experimental design for estimating APOE3-versus-APOE4, sex effects, or their interaction with disease. Regional fit is also limited: our intended validation needs cortex/PFC evidence, whereas the principal all-cell atlas is hippocampal.

Boundary / transition:
The next slide explains why the apparent size of the GEO series does not repair those missing biological factors.""")

    # Slide 2 — counting hierarchy
    slide = new_slide(prs, bg=WHITE)
    add_header(slide, "Biological replication", TITLES[1], 2, accent=BLUE)
    xs = [0.60, 3.76, 6.92, 10.08]
    widths = [2.58, 2.58, 2.58, 2.58]
    colors = [PALE_BLUE, PALE_GREEN, PALE_ORANGE, PALE_RED]
    accents = [BLUE, TEAL, ORANGE, VERMILION]
    counts = ["37", "10", "8", "54,769"]
    labels = [
        "GSM records",
        "sample/library\npreparations",
        "independent male mice",
        "high-quality nuclei",
    ]
    details = [
        "Whole GEO series:\nages, regions, and validations",
        "Main 7-month hippocampal\natlas—not PFC",
        "4 WT + 4 5xFAD",
        "Measurements nested\ninside those mice",
    ]
    for index in range(3):
        add_connector(slide, xs[index] + widths[index], 3.48,
                      xs[index + 1] - 0.12, 3.48, color=MID, width=2.2)
    for x, w, bg, accent, count, label, detail in zip(
        xs, widths, colors, accents, counts, labels, details
    ):
        add_rect(slide, x, 1.68, w, 3.62, color=bg, outline=accent,
                 line_width=1.8)
        add_text(slide, count, x + 0.18, 2.02, w - 0.36, 0.70,
                 size=30.0, color=accent, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, x + 0.18, 2.78, w - 0.36, 0.72,
                 size=17.0, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, detail, x + 0.20, 3.82, w - 0.40, 0.94,
                 size=13.0, color=DARK,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_ribbon(
        slide,
        "The mouse is the independent unit; extra libraries and nuclei improve measurement, not animal replication.",
        y=5.92, accent=BLUE,
    )
    add_source(slide, "Source: GSE143758 assessment DOCX, Figure 1 and Table 4.")
    add_note(slide, """What to point at:
Walk from left to right. Thirty-seven is the count of GEO sample records for the umbrella series, not thirty-seven independent mice. The principal seven-month male hippocampal atlas contains ten library preparations arising from eight mice, four WT and four 5xFAD, and those mice contribute 54,769 retained nuclei.

Main takeaway:
Biological n for the main disease comparison is four mice per genotype. Thousands of nuclei cannot be counted as thousands of independent replicates.

Boundary / transition:
Even that valid male disease contrast does not supply the replicated female or human APOE groups needed for the planned interaction. It is also a hippocampal atlas; Slide 4 explains why the smaller PFC component does not provide direct region-matched replication.""")

    # Slide 3 — sex and APOE design
    slide = new_slide(prs, bg=OFF_WHITE)
    add_header(
        slide, "Missing design factors", TITLES[2], 3, accent=VERMILION,
        subtitle="Reported study components are shown below; they do not form a balanced APOE-by-sex experiment.",
    )
    add_text(slide, "WT", 2.52, 1.50, 1.70, 0.36, size=17.0,
             color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "5xFAD", 4.88, 1.50, 1.70, 0.36, size=17.0,
             color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Male\nmain atlas", 0.72, 2.10, 1.50, 0.76,
             size=16.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Female\nreported validation", 0.72, 3.74, 1.50, 0.82,
             size=15.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    for x, value in ((2.30, "n = 4"), (4.66, "n = 4")):
        add_rect(slide, x, 1.98, 2.14, 1.22, color=PALE_GREEN,
                 outline=TEAL, line_width=1.8)
        add_text(slide, value, x + 0.15, 2.22, 1.84, 0.48,
                 size=25.0, color=TEAL, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    for x in (2.30, 4.66):
        add_rect(slide, x, 3.62, 2.14, 1.22, color=PALE_RED,
                 outline=VERMILION, line_width=1.8)
        add_text(slide, "n = 1", x + 0.15, 3.86, 1.84, 0.48,
                 size=25.0, color=VERMILION, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Qualitative only: no within-group female variance",
             2.20, 5.00, 4.72, 0.38, size=14.0, color=VERMILION,
             bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 7.28, 1.64, 5.35, 3.96, color=PALE_ORANGE,
             outline=ORANGE, line_width=1.8)
    add_text(slide, "APOE factor absent", 7.72, 2.00, 4.48, 0.52,
             size=23.0, color=ORANGE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "The dataset has endogenous mouse Apoe expression,\nbut no human APOE2, APOE3, or APOE4 genotype groups.",
             7.72, 2.78, 4.48, 1.12, size=16.0, color=NAVY,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Without an APOE factor, APOE3-vs-APOE4 and\nAPOE × sex × disease effects cannot be estimated.",
             7.72, 4.22, 4.48, 0.82, size=15.0, color=DARK, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_ribbon(slide, "More nuclei cannot replace missing or unreplicated animal groups.",
               y=6.18, accent=VERMILION)
    add_source(slide, "Source: GSE143758 assessment DOCX, Tables 1–2 and 5–7.")
    add_note(slide, """What to point at:
The left matrix shows the strongest inferential component: four male WT and four male 5xFAD mice. The reported female validation has one mouse in each disease group, so it can show that a state is observable but cannot estimate female variability or support a reliable disease-by-sex interaction. The right card shows the separate structural problem: there is no human APOE isoform factor.

Main takeaway:
Both defining dimensions of the planned validation are unavailable: sex is unreplicated and APOE3/APOE4 is absent.

Boundary / transition:
The next slide explains why counting every GEO record together would not create the missing factorial design.""")

    # Slide 4 — mixed subexperiments
    slide = new_slide(prs, bg=WHITE)
    add_header(
        slide, "Regional fit and series structure", TITLES[3], 4,
        accent=ORANGE,
        subtitle="PFC is a small astrocyte-focused subset; the 37 records also mix ages, repeat preparations, and protocols.",
    )
    add_rect(slide, 0.75, 1.52, 5.18, 1.20, color=PALE_BLUE,
             outline=BLUE, line_width=1.8)
    add_text(slide, "HUMAN VALIDATION TARGET", 1.05, 1.70, 4.58, 0.24,
             size=10.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Cortex / PFC", 1.05, 2.00, 4.58, 0.48,
             size=23.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "≠", 6.12, 1.82, 1.10, 0.56,
             size=30.0, color=VERMILION, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 7.40, 1.52, 5.18, 1.20, color=PALE_ORANGE,
             outline=ORANGE, line_width=1.8)
    add_text(slide, "MAIN GSE143758 ALL-CELL ATLAS", 7.70, 1.70, 4.58, 0.24,
             size=10.5, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "7-month hippocampus", 7.70, 2.00, 4.58, 0.48,
             size=22.0, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    card_x = [0.60, 3.76, 6.92, 10.08]
    card_titles = ["Main atlas", "PFC subset", "Not independent", "Other heterogeneity"]
    card_bodies = [
        "Broad all-cell analysis\n4 WT + 4 5xFAD males",
        "7- and 10-month astrocytes\n1 WT + 1 5xFAD per age",
        "Some PFC samples reuse\nhippocampus-profiled mice",
        "Age, lysis, chemistry,\nand batch also vary",
    ]
    card_accents = [BLUE, ORANGE, VERMILION, MID]
    card_bgs = [PALE_BLUE, PALE_ORANGE, PALE_RED, PALE_GRAY]
    add_text(slide, "The 37-record umbrella also mixes:",
             4.55, 2.88, 4.24, 0.24,
             size=12.2, color=MID, bold=True, align=PP_ALIGN.CENTER)
    for x, title, body, accent, bg in zip(
        card_x, card_titles, card_bodies, card_accents, card_bgs
    ):
        add_rect(slide, x, 3.18, 2.58, 2.16, color=bg,
                 outline=accent, line_width=1.5)
        add_rect(slide, x, 3.18, 2.58, 0.13, color=accent,
                 outline=None, radius=False)
        add_text(slide, title, x + 0.18, 3.46, 2.22, 0.50,
                 size=15.8, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x + 0.20, 4.16, 2.18, 0.84,
                 size=12.8, color=DARK,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_ribbon(
        slide,
        "This supports cross-region concordance—not direct PFC replication of the human result.",
        y=5.88, accent=ORANGE,
    )
    add_source(slide, "Source: GSE143758 assessment DOCX, Sections 1, 4–6, 8, and 12.")
    add_note(slide, """What to point at:
Start with the regional mismatch. Our intended validation needs cortex or prefrontal cortex, but the main all-cell atlas and the age-course data are hippocampal. The PFC component is astrocyte-focused, with one WT and one 5xFAD mouse at each of two ages, and those mice were also profiled in hippocampus. The 37 GEO records additionally mix ages, repeated preparations, and protocols.

Main takeaway:
GSE143758 does not provide a replicated, all-cell PFC disease comparison. Combining all 37 records would also confound biological and technical differences and would not increase independent animal replication.

Boundary / transition:
The final slide converts these limitations into a practical decision about what GSE143758 can still contribute.""")

    # Slide 5 — recommendation
    slide = new_slide(prs, bg=OFF_WHITE)
    add_header(slide, "Recommendation", TITLES[4], 5, accent=TEAL)
    add_rect(slide, 0.70, 1.46, 5.82, 4.55, color=WHITE,
             outline=TEAL, line_width=1.8)
    add_rect(slide, 0.70, 1.46, 5.82, 0.14, color=TEAL,
             outline=None, radius=False)
    add_text(slide, "Defensible auxiliary uses", 1.04, 1.84, 5.12, 0.50,
             size=22.0, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_bullets(slide, [
        "Mouse-level 7-month male 5xFAD-versus-WT disease direction",
        "Disease-associated astrocyte and age-progression questions",
        "Descriptive cortex/PFC astrocyte concordance",
        "Pre-specified pathway or module concordance",
    ], 1.05, 2.58, 5.08, size=15.0, accent=TEAL, line_h=0.74)

    add_rect(slide, 6.82, 1.46, 5.82, 4.55, color=WHITE,
             outline=VERMILION, line_width=1.8)
    add_rect(slide, 6.82, 1.46, 5.82, 0.14, color=VERMILION,
             outline=None, radius=False)
    add_text(slide, "Claims this dataset cannot support", 7.16, 1.84,
             5.12, 0.50, size=22.0, color=VERMILION, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_bullets(slide, [
        "Human APOE3-versus-APOE4 effects",
        "Reliable female-versus-male effects or interactions",
        "APOE × sex × disease inference",
        "Replicated, all-cell cortex/PFC validation",
    ], 7.17, 2.58, 5.08, size=15.0, accent=VERMILION, line_h=0.74)
    add_ribbon(
        slide,
        "Decision: no-go for the original endpoint and direct PFC replication; cross-region analysis only.",
        y=6.22, accent=TEAL,
    )
    add_source(slide, "Source: GSE143758 assessment DOCX, Sections 8, 12, and 14.")
    add_note(slide, """What to point at:
Use the two columns to separate a defensible auxiliary role from unsupported claims. The dataset can test male 5xFAD-versus-WT disease direction at the mouse level and examine disease-associated astrocyte or age-related programs. Its smaller cortex/PFC component can support descriptive astrocyte concordance, but not broad, independently replicated cortex/PFC validation. It also cannot supply human APOE isoform or replicated sex inference.

Main takeaway:
The correct project decision is no-go for the intended APOE-by-sex and cortex/PFC validation endpoint, while preserving a narrower, separately labeled disease-reference use if that changed question is valuable.

Boundary / transition:
A report-ready statement is: GSE143758 was considered as a 5xFAD-versus-WT snRNA-seq reference but was not used for the primary validation because human APOE groups, replicated female groups, and a replicated all-cell PFC cohort were unavailable.""")

    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output_path.stem}.", suffix=".tmp.pptx",
        dir=output_path.parent, delete=False,
    ) as handle:
        temp_path = Path(handle.name)
    try:
        prs.save(temp_path)
        validate_deck(temp_path)
        os.replace(temp_path, output_path)
    finally:
        if temp_path.exists():
            temp_path.unlink()
    validate_deck(output_path)
    return output_path


def validate_deck(path: Path) -> None:
    validate_source_contract(SOURCE_DOC)
    if not path.exists() or path.stat().st_size < 50_000:
        raise AssertionError(f"Deck missing or unexpectedly small: {path}")
    prs = Presentation(path)
    if len(prs.slides) != EXPECTED_SLIDES:
        raise AssertionError(f"Expected {EXPECTED_SLIDES} slides, found {len(prs.slides)}")
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        raise AssertionError("Deck is not 13.333333 × 7.5 inch widescreen")
    if SOURCE_DOC_SHA256 not in (prs.core_properties.comments or ""):
        raise AssertionError("Deck is not bound to the approved source DOCX revision")

    combined: list[str] = []
    for index, (slide, title) in enumerate(zip(prs.slides, TITLES), start=1):
        visible = _slide_text(slide)
        if " ".join(title.split()) not in " ".join(visible.split()):
            raise AssertionError(f"Slide {index} missing expected title: {title}")
        if "Source:" not in visible:
            raise AssertionError(f"Slide {index} has no source line")
        note = slide.notes_slide.notes_text_frame.text.strip()
        if len(note.split()) < 55:
            raise AssertionError(f"Slide {index} note is too short: {len(note.split())} words")
        for heading in NOTE_HEADINGS:
            if heading not in note:
                raise AssertionError(f"Slide {index} note missing heading: {heading}")
        combined.extend([visible, note])
        tolerance = Inches(0.02)
        for shape in slide.shapes:
            if shape.left < -tolerance or shape.top < -tolerance:
                raise AssertionError(f"Slide {index} has shape above/left of canvas")
            if shape.left + shape.width > SLIDE_W + tolerance:
                raise AssertionError(f"Slide {index} has shape beyond right edge")
            if shape.top + shape.height > SLIDE_H + tolerance:
                raise AssertionError(f"Slide {index} has shape beyond bottom edge")

    joined = "\n".join(combined)
    required = [
        "37",
        "GSM records",
        "10",
        "sample/library",
        "8",
        "independent male mice",
        "4 WT + 4 5xFAD",
        "54,769",
        "n = 1",
        "no human APOE",
        "not testable",
        "mouse is the independent unit",
        "no-go for the original endpoint",
        "Cortex / PFC",
        "7-month hippocampus",
        "cross-region concordance",
        "not direct PFC replication",
        "Some PFC samples reuse",
        "Replicated, all-cell cortex/PFC validation",
    ]
    for phrase in required:
        if phrase.lower() not in joined.lower():
            raise AssertionError(f"Required presentation claim missing: {phrase}")
    forbidden = [
        "37 mice",
        "powered sex comparison",
        "dataset failed",
        "dataset is useless",
        "nuclei are independent replicates",
        "APOE3-versus-APOE4 validation",
    ]
    for phrase in forbidden:
        if phrase.lower() in joined.lower():
            raise AssertionError(f"Forbidden presentation wording found: {phrase}")

    with zipfile.ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise AssertionError("PPTX ZIP integrity check failed")
        slide_xml = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        if len(slide_xml) != EXPECTED_SLIDES:
            raise AssertionError("PPTX package slide XML count is incorrect")


def _write_tsv(path: Path, rows: list[dict[str, object]], fields: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        mode="w", encoding="utf-8", newline="", delete=False,
        prefix=f".{path.name}.", dir=path.parent,
    ) as handle:
        writer = csv.DictWriter(handle, fieldnames=fields, delimiter="\t")
        writer.writeheader()
        writer.writerows(rows)
        temp = Path(handle.name)
    os.replace(temp, path)


def write_report(deck_path: Path, report_dir: Path, *,
                 visual_review_status: str) -> None:
    if visual_review_status not in {"pending", "complete"}:
        raise ValueError("visual_review_status must be pending or complete")
    validate_deck(deck_path)
    prs = Presentation(deck_path)
    report_dir.mkdir(parents=True, exist_ok=True)
    stem = "gse143758_mouse_validation_limitations"

    inputs = []
    for role, path in (
        ("source_docx", SOURCE_DOC),
        ("builder", Path(__file__).resolve()),
        ("deck", deck_path),
    ):
        inputs.append({
            "artifact_role": role,
            "path": display_path(path),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        })
    _write_tsv(
        report_dir / f"{stem}_input_manifest.tsv", inputs,
        ["artifact_role", "path", "bytes", "sha256"],
    )

    inventory = []
    for index, slide in enumerate(prs.slides, start=1):
        inventory.append({
            "slide_number": index,
            "title": TITLES[index - 1],
            "shape_count": len(slide.shapes),
            "picture_count": sum(1 for shape in slide.shapes if shape.shape_type == 13),
            "speaker_note_words": len(slide.notes_slide.notes_text_frame.text.split()),
        })
    _write_tsv(
        report_dir / f"{stem}_slide_inventory.tsv", inventory,
        ["slide_number", "title", "shape_count", "picture_count", "speaker_note_words"],
    )

    checks = [
        ("source_contract", True,
         f"Approved regional-fit DOCX SHA-256: {SOURCE_DOC_SHA256}"),
        ("pptx_integrity", True, "PPTX ZIP and python-pptx validation passed"),
        ("slide_count", len(prs.slides) == 5, "Five slides; no appendix"),
        ("widescreen", prs.slide_width == SLIDE_W and prs.slide_height == SLIDE_H,
         "13.333333 × 7.5 inches"),
        ("editable_native_graphics", all(row["picture_count"] == 0 for row in inventory),
         "All scientific diagrams use editable PowerPoint shapes"),
        ("speaker_notes", all(row["speaker_note_words"] >= 55 for row in inventory),
         "Every slide has three-section speaker notes"),
        ("sample_hierarchy", True, "37 GSM → 10 libraries → 8 mice → 54,769 nuclei"),
        ("sex_replication", True, "Male 4/4; reported female 1/1"),
        ("apoe_factor", True, "No human APOE isoform manipulation"),
        ("regional_fit", True,
         "Primary all-cell atlas is hippocampal; the smaller PFC component is astrocyte-focused and not independent"),
        ("decision_scope", True, "No-go for intended endpoint; auxiliary use retained"),
        ("visual_review", visual_review_status == "complete",
         "PPTX native-shape previews reviewed in color and grayscale"
         if visual_review_status == "complete"
         else "Awaiting PPTX native-shape color/grayscale review"),
    ]
    check_rows = [
        {"check_id": check_id, "passed": str(bool(passed)), "detail": detail}
        for check_id, passed, detail in checks
    ]
    _write_tsv(
        report_dir / f"{stem}_checks.tsv", check_rows,
        ["check_id", "passed", "detail"],
    )

    completed = datetime.now(timezone.utc).isoformat()
    status = [{
        "schema_version": "gse143758_mouse_validation_limitations_deck_v2",
        "deck_id": "gse143758_mouse_validation_limitations",
        "validation_status": (
            "validated_complete" if visual_review_status == "complete"
            else "awaiting_visual_review"
        ),
        "visual_review_status": visual_review_status,
        "slides": len(prs.slides),
        "input_files": len(inputs),
        "checks": len(checks),
        "deck_path": display_path(deck_path),
        "deck_bytes": deck_path.stat().st_size,
        "deck_sha256": sha256(deck_path),
        "source_doc_sha256": SOURCE_DOC_SHA256,
        "completed_utc": completed,
    }]
    _write_tsv(
        report_dir / f"{stem}_status.tsv", status,
        [
            "schema_version", "deck_id", "validation_status",
            "visual_review_status", "slides", "input_files", "checks",
            "deck_path", "deck_bytes", "deck_sha256", "source_doc_sha256",
            "completed_utc",
        ],
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--report-dir", type=Path, default=DEFAULT_REPORT_DIR)
    parser.add_argument(
        "--visual-review-status", choices=("pending", "complete"), default="pending",
    )
    parser.add_argument("--validate-only", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = args.output.resolve()
    if args.validate_only:
        validate_deck(output)
        print(f"Validated: {output}")
        return
    built = build_deck(output)
    write_report(
        built, args.report_dir.resolve(),
        visual_review_status=args.visual_review_status,
    )
    print(f"Built and validated: {built}")


if __name__ == "__main__":
    main()
