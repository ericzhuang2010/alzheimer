#!/usr/bin/env python3
"""Build the mouse public-dataset comparison PowerPoint.

The presentation is a faithful, editable summary of
``docs/validation_mouse/mouse_dataset_comparison.md``.  The first two slides
define the eight evaluation criteria; each remaining slide applies those same
criteria to one dataset profile.  The deck does not combine studies into one
cohort and does not present newly calculated expression results.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import os
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
SOURCE_MD = REPO / "docs/validation_mouse/mouse_dataset_comparison.md"
SOURCE_MD_SHA256 = "1d6178dee8b717c5baa43148027110600f92621d4c327145626815868a8385aa"
DEFAULT_OUT = REPO / "docs/presentations/mouse_dataset_comparison_08232026.pptx"
DEFAULT_REPORT_DIR = REPO / "results/presentations/validation_mouse/mouse_dataset_comparison"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)
FONT = "Arial"
EXPECTED_SLIDES = 14
NOTE_HEADINGS = ("What to point at:", "Main takeaway:", "Boundary / transition:")

NAVY = RGBColor(15, 35, 61)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
ORANGE = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(220, 228, 237)
MID = RGBColor(96, 108, 123)
DARK = RGBColor(35, 42, 51)
PALE_BLUE = RGBColor(224, 239, 248)
PALE_GREEN = RGBColor(224, 244, 238)
PALE_ORANGE = RGBColor(255, 244, 216)
PALE_RED = RGBColor(253, 235, 228)
PALE_PURPLE = RGBColor(240, 233, 247)
PALE_GRAY = RGBColor(239, 243, 247)

CRITERIA = (
    "Region", "APOE", "Sex", "Disease",
    "Modality", "Replication", "Cell breadth", "Public usability",
)
STATE_STYLE = {
    "STRONG": (TEAL, PALE_GREEN),
    "PARTIAL": (ORANGE, PALE_ORANGE),
    "NOT MATCHED": (VERMILION, PALE_RED),
}
ROLE_STYLE = {
    "CORE ANCHOR": (BLUE, PALE_BLUE),
    "DISEASE DIRECTION": (TEAL, PALE_GREEN),
    "SPECIALIZED SUPPORT": (PURPLE, PALE_PURPLE),
}


@dataclass(frozen=True)
class CriterionEvidence:
    state: str
    fact: str


@dataclass(frozen=True)
class DatasetProfile:
    accessions: str
    title: str
    role: str
    section: str
    source: str
    snapshot: str
    criteria: tuple[CriterionEvidence, ...]
    best_use: str
    limitation: str
    safe_claim: str
    note_focus: str


def E(state: str, fact: str) -> CriterionEvidence:
    if state not in STATE_STYLE:
        raise ValueError(f"Unknown criterion state: {state}")
    return CriterionEvidence(state, fact)


DATASETS: tuple[DatasetProfile, ...] = (
    DatasetProfile(
        "GSE185063",
        "Among audited candidates, GSE185063 is the strongest broad cortical APOE anchor, but lacks disease and usable sex metadata",
        "CORE ANCHOR", "§6.1", "GEO GSE185063; PMID 36040482",
        "Cortex • snRNA-seq • 16 mice • 4/APOE×age group",
        (
            E("PARTIAL", "Cortex/cortical mantle; not explicit PFC."),
            E("STRONG", "Human APOE3 versus APOE4 knock-in."),
            E("PARTIAL", "Study used both sexes; snRNA mouse map unresolved."),
            E("NOT MATCHED", "No AD model or matched disease contrast."),
            E("STRONG", "snRNA-seq with mouse-level samples."),
            E("STRONG", "16 mice; four per APOE-by-age group."),
            E("STRONG", "Neurons, glia, OPCs, and vascular cells."),
            E("PARTIAL", "Counts available; sex mapping still needs audit."),
        ),
        "Localize human PFC APOE-associated modules across cortical cell types.",
        "No disease factor; sex-specific inference is unsafe without a verified mouse map.",
        "Use the label “cortical cross-species APOE concordance”—not direct mouse PFC replication.",
        "The letter F in names such as E3F means floxed, not female.",
    ),
    DatasetProfile(
        "GSE241553",
        "GSE241553 adds sex-balanced cortical APOE induction under amyloid, but not an AD-control contrast",
        "CORE ANCHOR", "§6.2", "GEO GSE241553; PMID 37857825",
        "Cortex • scRNA-seq • 24 mice • 3F+3M/group",
        (
            E("PARTIAL", "Cortex; not an explicit PFC subdivision."),
            E("PARTIAL", "Conditional apoE3/apoE4 induction in microglia/CAMs."),
            E("STRONG", "3 females + 3 males in each of four groups."),
            E("PARTIAL", "All mice are on an amyloid-model background; no non-amyloid arm."),
            E("STRONG", "scRNA-seq with mouse-resolved samples."),
            E("STRONG", "24 mice; six per line-by-induction group."),
            E("PARTIAL", "Multiple populations; direct perturbation is microglia/CAM."),
            E("STRONG", "Per-sample files and group identities are available."),
        ),
        "Test sex-aware microglia/CAM apoE3-versus-apoE4 induction responses under amyloid.",
        "C versus T is non-induced versus induced—not control versus AD; only three mice per sex×line×induction group.",
        "Call this a microglia/CAM-specific APOE induction response under amyloid.",
        "The preferred mechanistic contrast is (E4-T − E4-C) − (E3-T − E3-C).",
    ),
    DatasetProfile(
        "GSE163857",
        "Among audited candidates, GSE163857 is the only direct APOE × sex × disease design, but only in non-regional bulk microglia",
        "CORE ANCHOR", "§6.3", "GEO GSE163857; PMID 34746703",
        "Brain microglia • bulk RNA-seq • 30 mice • all 8 APOE×sex×disease groups",
        (
            E("NOT MATCHED", "Sorted brain microglia; no cortex/PFC dissection."),
            E("STRONG", "Human APOE3/APOE4 targeted replacement."),
            E("STRONG", "Female and male mice are explicitly identified."),
            E("STRONG", "5xFAD versus matched targeted-replacement controls."),
            E("PARTIAL", "Bulk RNA-seq of sorted microglia."),
            E("PARTIAL", "Independent mice, but control groups have n=1–3."),
            E("NOT MATCHED", "Microglia only."),
            E("STRONG", "Compact processed mouse-level count matrix."),
        ),
        "Directly benchmark APOE-by-sex-by-5xFAD effects in microglia.",
        "Not regional or broad-cell; small, imbalanced controls make the three-way interaction underpowered.",
        "Report factorial microglial support with effect estimates and uncertainty—not PFC replication.",
        "The 30-mouse design contains all eight groups, but the male APOE3 control group has one mouse.",
    ),
    DatasetProfile(
        "GSE140399",
        "GSE140399 supplies broad-cell cortical disease direction, but pooled mice make it descriptive",
        "DISEASE DIRECTION", "§6.4", "GEO GSE140399; PMID 31932797",
        "15-month cortex/hippocampus • snRNA-seq • one pool/condition",
        (
            E("PARTIAL", "Cortex is separate from hippocampus; not explicit PFC."),
            E("NOT MATCHED", "No human APOE isoform; TREM2 is the modifier."),
            E("NOT MATCHED", "Sex is not a usable public study factor."),
            E("STRONG", "WT, 5xFAD, Trem2KO, and Trem2KO×5xFAD."),
            E("PARTIAL", "snRNA-seq, but animals were pooled before profiling."),
            E("PARTIAL", "Three mice form one pool per region×genotype."),
            E("STRONG", "Broad neural and glial cell coverage."),
            E("PARTIAL", "Counts are usable; pooling blocks ordinary inference."),
        ),
        "Check the direction of frozen human modules in homologous cortical cell types.",
        "One expression pool per condition means no mouse-level variance or ordinary inferential P values.",
        "Use “descriptive cortical concordance”—never “replicated in three mice per group.”",
        "Pooling three mice improves the pool estimate but still produces one inferential profile.",
    ),
    DatasetProfile(
        "GSE140510",
        "GSE140510 provides replicated broad-cell disease direction, but tissue is not explicitly cortex/PFC",
        "DISEASE DIRECTION", "§6.5", "GEO GSE140510; PMID 31932797",
        "7-month mouse brain • snRNA-seq • 12 mice • 3/genotype",
        (
            E("NOT MATCHED", "Public label is mouse brain, not cortex/PFC."),
            E("NOT MATCHED", "No human APOE isoform; TREM2 is the modifier."),
            E("NOT MATCHED", "Sex is not a usable public study factor."),
            E("STRONG", "WT, 5xFAD, Trem2KO, and Trem2KO×5xFAD."),
            E("STRONG", "snRNA-seq with separately processed mice."),
            E("STRONG", "12 independent mice; three per genotype."),
            E("STRONG", "Broad neural, glial, and vascular coverage."),
            E("STRONG", "Mouse-level files support pseudobulk contrasts."),
        ),
        "Provide replicated broad-cell AD-direction support across homologous cell types.",
        "Regional mismatch and no APOE or sex factor; keep separate from cortical GSE140399.",
        "Use when inferential replication matters more than exact cortical localization.",
        "This study trades stronger mouse-level replication for weaker regional matching.",
    ),
    DatasetProfile(
        "GSE143758",
        "GSE143758 supports hippocampal and astrocyte disease biology, not broad PFC validation",
        "DISEASE DIRECTION", "§6.6", "GEO GSE143758; PMID 32341542",
        "Main: hippocampus, 4 WT+4 5xFAD • Cortical subset: 1+1 per age",
        (
            E("PARTIAL", "Cortex data exist, but the broad atlas is hippocampal."),
            E("NOT MATCHED", "No human APOE isoform factor."),
            E("PARTIAL", "Main atlas is male; only two females total, one/genotype."),
            E("STRONG", "5xFAD versus WT disease contrast."),
            E("STRONG", "snRNA-seq; mouse-level main atlas."),
            E("PARTIAL", "Main n=4+4; cortex has n=1 per age×genotype group."),
            E("PARTIAL", "Main atlas broad; cortex matrix is astrocyte-focused."),
            E("PARTIAL", "Thirty-seven records mix components; IDs need care."),
        ),
        "Support hippocampal astrocyte, aging, and descriptive cross-region pathways.",
        "The cortical component is astrocyte-focused, with one record per age×genotype group; it is not a replicated broad-cell PFC cohort.",
        "Cortical data exist, but use them only for descriptive astrocyte concordance—not replicated all-cell PFC validation.",
        "Most of the study and its main all-cell atlas are not cortex/PFC; the cortical arm has one WT and one 5xFAD mouse at each of two ages.",
    ),
    DatasetProfile(
        "GSE225503 / GSE239999",
        "GSE225503/GSE239999 add APOE2/3/4 immune aging under 5xFAD, but lack AD controls and usable sex metadata",
        "SPECIALIZED SUPPORT", "§6.7", "GEO GSE225503/GSE239999; PMID 38159571",
        "Mixed cortex + hippocampus CD45+ cells • scRNA/multiome • multiplexed",
        (
            E("PARTIAL", "Cortical and hippocampal regions; separation needs audit."),
            E("STRONG", "Human APOE2, APOE3, and APOE4 knock-in."),
            E("NOT MATCHED", "Public sex metadata are insufficient."),
            E("PARTIAL", "All animals are 5xFAD; no non-AD group."),
            E("PARTIAL", "Immune scRNA-seq/multiome with HTO demultiplexing."),
            E("PARTIAL", "Mouse identities must be recovered from metadata."),
            E("PARTIAL", "CD45-positive immune compartment only."),
            E("PARTIAL", "Large RDS objects; region/mouse fields need audit."),
        ),
        "Study APOE2/3/4 immune aging, late microglial states, and chromatin under 5xFAD.",
        "No non-AD arm, usable sex factor, broad cell coverage, or direct PFC match.",
        "Call this APOE2/3/4 immune-cell support under 5xFAD—not an APOE-by-sex AD-control test.",
        "Eight GEO records are multiplexed experiment files, not eight independent mice.",
    ),
    DatasetProfile(
        "GSE127892",
        "GSE127892 captures sex-by-amyloid microglial trajectories, but pooled mice prevent mouse-level inference",
        "SPECIALIZED SUPPORT", "§6.8", "GEO GSE127892; PMID 31018141",
        "Cortex+hippocampus • plate scRNA • 32 condition pools",
        (
            E("PARTIAL", "Cortex and hippocampus; not explicit PFC."),
            E("NOT MATCHED", "Endogenous mouse Apoe; no human isoform contrast."),
            E("PARTIAL", "Female and male conditions, but animals are pooled."),
            E("STRONG", "AppNL-G-F versus WT across four ages."),
            E("PARTIAL", "Plate scRNA-seq; GEO records are cells/wells."),
            E("PARTIAL", "Two mice pooled into one profile per condition."),
            E("NOT MATCHED", "Microglia only."),
            E("PARTIAL", "12,288 cell records require pool reconstruction."),
        ),
        "Describe sex-by-amyloid microglial state trajectories across age and region.",
        "Pooling removes mouse-level variance; human APOE isoforms and broad cells are absent.",
        "Treat as condition-level descriptive support—not a replicated mouse-level interaction.",
        "The 32-condition design is rich, but each condition has one pooled expression profile.",
    ),
    DatasetProfile(
        "GSE212606",
        "GSE212606 offers exceptional sex-aware whole-brain breadth, but APOE4 is confounded with TREM2 R47H",
        "SPECIALIZED SUPPORT", "§6.9", "GEO GSE212606; PMID 37774676",
        "Whole brain • EasySci RNA+ATAC • >300 subtypes • animal manifest needed",
        (
            E("NOT MATCHED", "Whole brain; not cortex or PFC."),
            E("PARTIAL", "LOAD model combines APOE4 with TREM2 R47H."),
            E("PARTIAL", "Both sexes reported; animal-level map is required."),
            E("STRONG", "WT, 5xFAD, and LOAD-model comparisons."),
            E("PARTIAL", "Single-cell RNA/ATAC; GEO objects are aggregate."),
            E("PARTIAL", "Replication becomes auditable only with a manifest."),
            E("STRONG", ">300 reported subtypes, including rare populations."),
            E("PARTIAL", "Very large archive; mouse identities must be linked."),
        ),
        "Assess broad sex/disease convergence, rare populations, and chromatin support.",
        "APOE4 cannot be separated from TREM2 R47H; not region matched or directly APOE3-versus-APOE4.",
        "Use as broad whole-brain convergence and chromatin evidence—not direct PFC/APOE validation.",
        "The scale is a strength for cell breadth, but it does not repair the genetic or regional confounding.",
    ),
    DatasetProfile(
        "GSE213446",
        "GSE213446 provides APOE-by-tau support in hippocampus, but mouse-level sex and replication are unresolved",
        "SPECIALIZED SUPPORT", "§6.10", "GEO GSE213446",
        "Hippocampus • snRNA-seq • 10 public libraries • APOE×tau×treatment",
        (
            E("NOT MATCHED", "Hippocampus; not cortex/PFC."),
            E("STRONG", "Human APOE3/APOE4 with and without P301S."),
            E("PARTIAL", "Sex biology reported; public library map absent."),
            E("PARTIAL", "Tauopathy versus non-tau, not the target amyloid design."),
            E("PARTIAL", "snRNA-seq, but animals per library are unclear."),
            E("PARTIAL", "One listed library/condition; mouse n is unaudited."),
            E("STRONG", "Broad hippocampal cell populations."),
            E("PARTIAL", "Counts available; mouse/sex metadata need recovery."),
        ),
        "Provide pathway-level, cross-region and cross-pathology APOE-by-tau support.",
        "Hippocampal tau model; sex mapping and independent mouse replication are unresolved.",
        "Keep conclusions descriptive until animal-level metadata are recovered.",
        "Antibiotic versus water treatment adds another biological factor that must remain explicit.",
    ),
    DatasetProfile(
        "GSE212317 / GSE213391",
        "GSE212317/GSE213391 support female APOE immunometabolism, but pooling and age/disease confounding limit inference",
        "SPECIALIZED SUPPORT", "§6.11", "GEO GSE212317/GSE213391; PMID 36871219",
        "Whole left hemisphere • glia-enriched scRNA/spatial/bulk • female only",
        (
            E("NOT MATCHED", "Broad left hemisphere; not region-specific PFC."),
            E("STRONG", "Human APOE3 and APOE4 contexts."),
            E("NOT MATCHED", "Female only; no sex comparison."),
            E("PARTIAL", "Age and 5xFAD status are not balanced factorially."),
            E("PARTIAL", "Glia-enriched scRNA plus spatial/bulk resources."),
            E("PARTIAL", "Three mice pooled into one library per group."),
            E("PARTIAL", "Glia enriched rather than broad all-cell coverage."),
            E("PARTIAL", "Processed resources usable; pooling limits inference."),
        ),
        "Explore female APOE immunometabolism, glycolysis, TCA, and lipid handling.",
        "One pooled library/group, female-only design, and age/disease confounding prevent interaction claims.",
        "Use for descriptive female APOE immunometabolism—not sex or clean APOE-by-disease replication.",
        "Middle-age groups are 5xFAD whereas young and aged groups represent different pathology contexts.",
    ),
    DatasetProfile(
        "GSE127884",
        "GSE127884 tests endogenous Apoe dependence in male amyloid microglia—not human APOE isoforms or sex effects",
        "SPECIALIZED SUPPORT", "§6.12", "GEO GSE127884; PMID 31018141",
        "Cortex+hippocampus • plate scRNA • male microglia • cells/wells in GEO",
        (
            E("PARTIAL", "Cortex and hippocampus; not explicit PFC."),
            E("NOT MATCHED", "Mouse Apoe versus Apoe knockout; no human E3/E4."),
            E("NOT MATCHED", "Male only."),
            E("STRONG", "APP/PS1 and controls with/without mouse Apoe."),
            E("PARTIAL", "Plate single-cell microglial RNA-seq."),
            E("NOT MATCHED", "Animal-level replication is not exposed cleanly."),
            E("NOT MATCHED", "Microglia only."),
            E("NOT MATCHED", "3,072 GEO records are cells/wells, not mice."),
        ),
        "Test whether endogenous mouse Apoe is required for amyloid-associated microglial states.",
        "Male-only, microglia-only design with no human isoform comparison and unclear mouse replication.",
        "Use as mechanistic Apoe-dependence support—not human APOE or sex validation.",
        "The 3,072 GEO records must never be used as the biological sample count.",
    ),
)

TITLES = (
    "Evaluation criteria used in this report: biological fit",
    "Evaluation criteria used in this report: evidence quality",
    *(profile.title for profile in DATASETS),
)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def display_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(REPO))
    except ValueError:
        return str(path.resolve())


def validate_source_contract(path: Path = SOURCE_MD) -> str:
    if not path.exists():
        raise FileNotFoundError(path)
    if sha256(path) != SOURCE_MD_SHA256:
        raise AssertionError("mouse_dataset_comparison.md changed from the approved source")
    text = path.read_text(encoding="utf-8")
    required = [
        "Evaluation criteria used in this report",
        "No openly auditable public dataset identified in this re-audit satisfies the complete design",
        "The independent biological unit is usually the **mouse**",
        "Pooling two or three mice before sequencing does not create independent mouse-level profiles",
        "A batch-correction method cannot convert hippocampus into PFC",
        *(profile.accessions.split(" / ")[0] for profile in DATASETS),
        "GSE239999", "GSE213391",
    ]
    missing = [value for value in required if value.lower() not in text.lower()]
    if missing:
        raise AssertionError(f"Source report contract changed; missing: {missing}")
    return text


def fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def stroke(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_run(run, *, size: float, color: RGBColor, bold: bool = False,
            italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             color: RGBColor = WHITE, outline: RGBColor | None = LIGHT,
             radius: bool = True, line_width: float = 1.0, name: str | None = None):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        stroke(shape, outline, line_width)
    if name:
        shape.name = name
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *,
             size: float = 16, color: RGBColor = DARK, bold: bool = False,
             italic: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.03,
             name: str | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic)
    return box


def new_slide(prs: Presentation, *, bg: RGBColor = WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_header(slide, kicker: str, title: str, page_no: int,
               *, accent: RGBColor = BLUE) -> None:
    add_text(slide, kicker.upper(), 0.55, 0.16, 6.1, 0.22,
             size=9.2, color=accent, bold=True)
    if len(title) > 112:
        title_size = 18.2
    elif len(title) > 92:
        title_size = 19.5
    elif len(title) > 72:
        title_size = 21.0
    else:
        title_size = 24.0
    add_text(slide, title, 0.55, 0.39, 12.05, 0.78,
             size=title_size, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE, name="slide_title")
    add_text(slide, f"{page_no:02d}", 12.47, 0.17, 0.30, 0.18,
             size=8.8, color=MID, bold=True, align=PP_ALIGN.RIGHT)


def add_source(slide, value: str) -> None:
    add_text(slide, value, 0.55, 7.24, 12.20, 0.14, size=6.7, color=MID,
             name="source_line")


def add_note(slide, value: str) -> None:
    slide.notes_slide.notes_text_frame.text = value


def add_ribbon(slide, text: str, *, y: float = 6.38,
               accent: RGBColor = BLUE) -> None:
    add_rect(slide, 0.55, y, 12.23, 0.56, color=NAVY, outline=None,
             radius=False, name="takeaway_ribbon")
    add_rect(slide, 0.55, y, 0.10, 0.56, color=accent, outline=None,
             radius=False)
    add_text(slide, text, 0.82, y + 0.08, 11.65, 0.36, size=12.0,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def add_definition_card(slide, x: float, y: float, w: float, h: float,
                        title: str, strong: str, partial: str, not_match: str,
                        *, accent: RGBColor) -> None:
    add_rect(slide, x, y, w, h, color=WHITE, outline=accent, line_width=1.5)
    add_rect(slide, x, y, w, 0.11, color=accent, outline=None, radius=False)
    add_text(slide, title, x + 0.20, y + 0.23, w - 0.40, 0.36,
             size=17.0, color=NAVY, bold=True)
    rows = (("STRONG", strong, TEAL), ("PARTIAL", partial, ORANGE),
            ("NOT MATCHED", not_match, VERMILION))
    for index, (label, body, color) in enumerate(rows):
        cy = y + 0.72 + index * 0.39
        add_text(slide, label, x + 0.20, cy, 1.02, 0.23,
                 size=8.1, color=DARK, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x + 1.30, cy - 0.02, w - 1.50, 0.29,
                 size=10.2, color=DARK, valign=MSO_ANCHOR.MIDDLE)


def add_criterion_card(slide, criterion: str, evidence: CriterionEvidence,
                       x: float, y: float, w: float, h: float,
                       *, index: int) -> None:
    accent, bg = STATE_STYLE[evidence.state]
    add_rect(slide, x, y, w, h, color=WHITE, outline=LIGHT,
             line_width=1.0, name=f"criterion_card_{index}_{criterion}")
    add_rect(slide, x, y, 0.08, h, color=accent, outline=None,
             radius=False)
    criterion_size = 8.7 if criterion == "Public usability" else 10.6
    add_text(slide, criterion, x + 0.18, y + 0.10, 1.18, 0.22,
             size=criterion_size, color=NAVY, bold=True)
    add_rect(slide, x + 0.16, y + 0.40, 1.08, 0.29, color=bg,
             outline=accent, radius=True, line_width=0.8,
             name=f"criterion_state_{index}")
    add_text(slide, evidence.state, x + 0.19, y + 0.43, 1.02, 0.19,
             size=8.2 if evidence.state != "NOT MATCHED" else 7.4,
             color=DARK, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, evidence.fact, x + 1.38, y + 0.14, w - 1.56, h - 0.27,
             size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE)


def add_right_card(slide, x: float, y: float, w: float, h: float,
                   label: str, body: str, *, accent: RGBColor,
                   bg: RGBColor, body_size: float = 12.2) -> None:
    add_rect(slide, x, y, w, h, color=bg, outline=accent, line_width=1.2)
    add_text(slide, label.upper(), x + 0.20, y + 0.15, w - 0.40, 0.20,
             size=8.7, color=accent, bold=True)
    add_text(slide, body, x + 0.20, y + 0.43, w - 0.40, h - 0.55,
             size=body_size, color=NAVY, bold=label == "Study snapshot",
             valign=MSO_ANCHOR.MIDDLE)


def build_criteria_slides(prs: Presentation) -> None:
    slide = new_slide(prs, bg=OFF_WHITE)
    add_header(slide, "Mouse dataset comparison", TITLES[0], 1, accent=BLUE)
    add_definition_card(slide, 0.65, 1.35, 5.95, 1.92, "Region",
                        "Explicit frontal/PFC tissue.",
                        "Cortex, or mixed cortex + hippocampus.",
                        "Hippocampus, whole brain, or unclear anatomy.", accent=BLUE)
    add_definition_card(slide, 6.75, 1.35, 5.95, 1.92, "APOE",
                        "Human APOE3 versus APOE4 in the intended context.",
                        "Conditional APOE, APOE4+other risk, or mouse Apoe.",
                        "No APOE factor.", accent=PURPLE)
    add_definition_card(slide, 0.65, 3.50, 5.95, 1.92, "Sex",
                        "Female/male labels for each mouse with replication.",
                        "Both sexes mentioned but mapping is missing or pooled.",
                        "One sex or unusable sex metadata.", accent=ORANGE)
    add_definition_card(slide, 6.75, 3.50, 5.95, 1.92, "Disease",
                        "AD model and matched non-AD animals in one study.",
                        "All animals diseased, or pathology differs from target.",
                        "No disease contrast.", accent=VERMILION)
    add_ribbon(
        slide,
        "No exact dataset was identified in the openly auditable set; complementary datasets answer different parts of the question.",
        y=5.86, accent=BLUE,
    )
    add_source(slide, "Source: docs/validation_mouse/mouse_dataset_comparison.md, §§2–3 (report dated Aug 23, 2026).")
    add_note(slide, """What to point at:
The first four criteria ask whether the biology matches the human prefrontal-cortex question. Strong means a direct match, partial means useful but qualified support, and not matched means that dimension cannot answer the intended claim. Cortex is only a partial regional match to human PFC unless a frontal subdivision is documented.

Main takeaway:
No exact dataset was identified in the openly auditable candidate set. That does not mean no such dataset exists; it means the reviewed public studies contribute complementary pieces rather than one complete PFC × APOE × sex × disease design.

Boundary / transition:
The next slide defines evidence quality and the biological replicate rule used throughout the dataset profiles.""")

    slide = new_slide(prs, bg=WHITE)
    add_header(slide, "Mouse dataset comparison", TITLES[1], 2, accent=TEAL)
    definitions = (
        ("Modality", "sc/snRNA with mouse identity", "bulk-sorted or pooled single-cell", "no usable expression"),
        ("Replication", "independent mice per group", "pooling or very small groups", "cells/wells only; no animal n"),
        ("Cell breadth", "neurons + glia + vascular", "immune/glia enriched", "one sorted cell type"),
        ("Public usability", "counts + animal metadata", "metadata reconstruction needed", "key identities unavailable"),
    )
    for index, (title, strong, partial, not_match) in enumerate(definitions):
        x = 0.55 + index * 3.18
        add_definition_card(slide, x, 1.30, 3.00, 2.02, title,
                            strong, partial, not_match,
                            accent=(BLUE, TEAL, PURPLE, ORANGE)[index])
    add_right_card(slide, 0.65, 3.68, 5.92, 1.65, "Biological replicate",
                   "3 mice × 5,000 cells/mouse\nInferential n = 3 mice—not 15,000 cells",
                   accent=BLUE, bg=PALE_BLUE, body_size=17.0)
    add_right_card(slide, 6.77, 3.68, 5.92, 1.65, "Two non-negotiable rules",
                   "Pooling mice creates one pooled profile.\nBatch correction cannot convert hippocampus into PFC.",
                   accent=VERMILION, bg=PALE_RED, body_size=15.2)
    add_ribbon(slide, "The eight states are not summed into a score: one critical mismatch can still invalidate a specific claim.",
               y=5.86, accent=TEAL)
    add_source(slide, "Source: docs/validation_mouse/mouse_dataset_comparison.md, §§2–3; biological replicate rule.")
    add_note(slide, """What to point at:
The second four criteria determine whether a study can support mouse-level inference and whether the public files can be audited. The worked example shows why many nuclei do not increase biological n. A pooled library is one profile even if several animals contributed tissue.

Main takeaway:
These criteria are noncompensatory. Exceptional cell count or cell-type breadth cannot replace missing independent mice, a missing APOE factor, or the wrong anatomy. Batch correction can reduce technical variation, but it cannot transform hippocampus into prefrontal cortex.

Boundary / transition:
Every remaining slide applies the same eight labels and factual evidence lines, followed by a defensible use and a claim boundary.""")


def dataset_note(profile: DatasetProfile, next_accession: str | None) -> str:
    transition = (
        f"The next profile applies the same criteria to {next_accession}."
        if next_accession else
        "This is the final dataset profile; compare roles across slides rather than treating these studies as one merged cohort."
    )
    return f"""What to point at:
Start with the role and study snapshot, then scan the eight criterion states. {profile.note_focus} The colored state is always paired with a factual line, so the judgment does not depend on color alone.

Main takeaway:
Best use: {profile.best_use} Main limitation: {profile.limitation}

Boundary / transition:
{profile.safe_claim} The biological unit must remain the mouse—or the pre-sequencing pool when mice were pooled—and raw matrices from separate studies should not be merged as one cohort. {transition}"""


def build_dataset_slide(prs: Presentation, profile: DatasetProfile,
                        page_no: int, next_accession: str | None) -> None:
    role_color, role_bg = ROLE_STYLE[profile.role]
    slide = new_slide(prs, bg=OFF_WHITE if page_no % 2 else WHITE)
    add_header(slide, f"Dataset profile • {profile.role}", profile.title,
               page_no, accent=role_color)

    positions = []
    for col in range(2):
        for row in range(4):
            positions.append((0.55 + col * 4.13, 1.34 + row * 1.08))
    for index, (criterion, evidence, (x, y)) in enumerate(
        zip(CRITERIA, profile.criteria, positions), start=1
    ):
        add_criterion_card(slide, criterion, evidence, x, y, 3.94, 0.93,
                           index=index)

    add_rect(slide, 8.91, 1.32, 3.87, 0.48, color=role_bg,
             outline=role_color, line_width=1.2)
    add_text(slide, profile.role, 9.10, 1.42, 3.49, 0.24,
             size=10.3, color=role_color, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_right_card(slide, 8.91, 1.94, 3.87, 1.02, "Study snapshot",
                   profile.snapshot, accent=role_color, bg=role_bg,
                   body_size=11.4)
    add_right_card(slide, 8.91, 3.10, 3.87, 1.36, "Best use",
                   profile.best_use, accent=TEAL, bg=PALE_GREEN,
                   body_size=12.0)
    add_right_card(slide, 8.91, 4.60, 3.87, 1.36, "Do not claim",
                   profile.limitation, accent=VERMILION, bg=PALE_RED,
                   body_size=11.7)
    add_ribbon(slide, profile.safe_claim, y=6.36, accent=role_color)
    add_source(slide, f"Source: mouse_dataset_comparison.md {profile.section}; {profile.source}. Suitability assessment, not new matrix analysis.")
    add_note(slide, dataset_note(profile, next_accession))


def build_deck(output_path: Path = DEFAULT_OUT) -> Path:
    validate_source_contract(SOURCE_MD)
    if len(DATASETS) != 12 or len(TITLES) != EXPECTED_SLIDES:
        raise AssertionError("Deck data contract must be two criteria + twelve dataset slides")
    for profile in DATASETS:
        if len(profile.criteria) != len(CRITERIA):
            raise AssertionError(f"{profile.accessions} does not have eight criteria")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Mouse Public-Dataset Comparison for Human PFC Validation"
    prs.core_properties.subject = "Eight-criterion assessment of 12 mouse dataset profiles"
    prs.core_properties.author = "Alzheimer project analysis team"
    prs.core_properties.keywords = (
        "mouse validation, PFC, APOE, sex, Alzheimer disease, scRNA-seq, snRNA-seq, GEO"
    )
    prs.core_properties.comments = (
        "Generated from docs/validation_mouse/mouse_dataset_comparison.md. "
        "This deck summarizes dataset suitability and does not report new expression analysis. "
        f"Source Markdown SHA-256: {SOURCE_MD_SHA256}."
    )

    build_criteria_slides(prs)
    for index, profile in enumerate(DATASETS):
        next_accession = DATASETS[index + 1].accessions if index + 1 < len(DATASETS) else None
        build_dataset_slide(prs, profile, index + 3, next_accession)

    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output_path.stem}.", suffix=".tmp.pptx",
        dir=output_path.parent, delete=False,
    ) as handle:
        temp_path = Path(handle.name)
    try:
        prs.save(temp_path)
        validate_deck(temp_path)
        os.replace(temp_path, output_path)
    finally:
        if temp_path.exists():
            temp_path.unlink()
    validate_deck(output_path)
    return output_path


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def validate_deck(path: Path) -> None:
    validate_source_contract(SOURCE_MD)
    if not path.exists() or path.stat().st_size < 50_000:
        raise AssertionError(f"Deck missing or unexpectedly small: {path}")
    prs = Presentation(path)
    if len(prs.slides) != EXPECTED_SLIDES:
        raise AssertionError(f"Expected {EXPECTED_SLIDES} slides, found {len(prs.slides)}")
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        raise AssertionError("Deck is not 13.333333 × 7.5 inch widescreen")
    if SOURCE_MD_SHA256 not in (prs.core_properties.comments or ""):
        raise AssertionError("Deck is not bound to the approved Markdown source")

    combined: list[str] = []
    tolerance = Inches(0.02)
    for index, (slide, expected_title) in enumerate(zip(prs.slides, TITLES), start=1):
        visible = _slide_text(slide)
        normalized = " ".join(visible.split())
        if " ".join(expected_title.split()) not in normalized:
            raise AssertionError(f"Slide {index} missing title: {expected_title}")
        if "Source:" not in visible:
            raise AssertionError(f"Slide {index} has no visible source line")
        note = slide.notes_slide.notes_text_frame.text.strip()
        if len(note.split()) < 55:
            raise AssertionError(f"Slide {index} note is too short: {len(note.split())} words")
        for heading in NOTE_HEADINGS:
            if heading not in note:
                raise AssertionError(f"Slide {index} note missing heading: {heading}")
        if any(shape.shape_type == 13 for shape in slide.shapes):
            raise AssertionError(f"Slide {index} unexpectedly contains a raster picture")
        for shape in slide.shapes:
            if shape.left < -tolerance or shape.top < -tolerance:
                raise AssertionError(f"Slide {index} has shape above/left of canvas")
            if shape.left + shape.width > SLIDE_W + tolerance:
                raise AssertionError(f"Slide {index} has shape beyond right edge")
            if shape.top + shape.height > SLIDE_H + tolerance:
                raise AssertionError(f"Slide {index} has shape beyond bottom edge")
        if index >= 3:
            for criterion in CRITERIA:
                if criterion.lower() not in visible.lower():
                    raise AssertionError(f"Slide {index} missing criterion: {criterion}")
            if "best use" not in visible.lower() or "do not claim" not in visible.lower():
                raise AssertionError(f"Slide {index} missing use/boundary cards")
            state_badges = sum(
                visible.count(state) for state in STATE_STYLE
            )
            if state_badges != 8:
                raise AssertionError(f"Slide {index} has {state_badges}, not 8, state badges")
        combined.extend([visible, note])

    joined = "\n".join(combined)
    required = [
        "Evaluation criteria used in this report",
        "No exact dataset was identified in the openly auditable set",
        "Inferential n = 3 mice—not 15,000 cells",
        "Pooling mice creates one pooled profile",
        "Batch correction cannot convert hippocampus into PFC",
        "most of the study and its main all-cell atlas are not cortex/PFC",
        "Cortical data exist, but use them only for descriptive astrocyte concordance",
        "all 8 APOE×sex×disease groups",
        *(profile.accessions.split(" / ")[0] for profile in DATASETS),
        "GSE239999", "GSE213391",
    ]
    for phrase in required:
        if phrase.lower() not in joined.lower():
            raise AssertionError(f"Required deck claim missing: {phrase}")
    forbidden = [
        "no public dataset exists",
        "cells are biological replicates",
        "nuclei are biological replicates",
        "37 mice",
        "12,288 mice",
        "3,072 mice",
        "Cortex/PFC:",
        "cortex/PFC mice",
        "n=1/cell",
        "factorial cells",
        "control cells have n=1–3",
        "direct mouse PFC replication" + " confirmed",
        "one integrated mouse replication cohort",
    ]
    for phrase in forbidden:
        if phrase.lower() in joined.lower():
            raise AssertionError(f"Forbidden deck wording found: {phrase}")

    with zipfile.ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise AssertionError("PPTX ZIP integrity check failed")
        slide_xml = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        if len(slide_xml) != EXPECTED_SLIDES:
            raise AssertionError("PPTX package slide XML count is incorrect")


def _write_tsv(path: Path, rows: list[dict[str, object]], fields: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        mode="w", encoding="utf-8", newline="", delete=False,
        prefix=f".{path.name}.", dir=path.parent,
    ) as handle:
        writer = csv.DictWriter(handle, fieldnames=fields, delimiter="\t")
        writer.writeheader()
        writer.writerows(rows)
        temp_path = Path(handle.name)
    os.replace(temp_path, path)


def criterion_matrix_rows() -> list[dict[str, object]]:
    rows: list[dict[str, object]] = []
    for profile in DATASETS:
        for order, (criterion, evidence) in enumerate(zip(CRITERIA, profile.criteria), start=1):
            rows.append({
                "accession_group": profile.accessions,
                "role": profile.role,
                "criterion_order": order,
                "criterion": criterion,
                "state": evidence.state,
                "evidence": evidence.fact,
                "source_section": profile.section,
            })
    return rows


def write_report(deck_path: Path, report_dir: Path, *,
                 visual_review_status: str) -> None:
    if visual_review_status not in {"pending", "complete"}:
        raise ValueError("visual_review_status must be pending or complete")
    validate_deck(deck_path)
    prs = Presentation(deck_path)
    report_dir.mkdir(parents=True, exist_ok=True)
    stem = "mouse_dataset_comparison"

    inputs = []
    for role, path in (
        ("source_markdown", SOURCE_MD),
        ("builder", Path(__file__).resolve()),
        ("deck", deck_path),
    ):
        inputs.append({
            "artifact_role": role,
            "path": display_path(path),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        })
    _write_tsv(report_dir / f"{stem}_input_manifest.tsv", inputs,
               ["artifact_role", "path", "bytes", "sha256"])

    inventory = []
    for index, slide in enumerate(prs.slides, start=1):
        inventory.append({
            "slide_number": index,
            "slide_type": "criteria" if index <= 2 else "dataset_profile",
            "title": TITLES[index - 1],
            "shape_count": len(slide.shapes),
            "picture_count": sum(1 for shape in slide.shapes if shape.shape_type == 13),
            "speaker_note_words": len(slide.notes_slide.notes_text_frame.text.split()),
        })
    _write_tsv(report_dir / f"{stem}_slide_inventory.tsv", inventory,
               ["slide_number", "slide_type", "title", "shape_count",
                "picture_count", "speaker_note_words"])

    matrix = criterion_matrix_rows()
    _write_tsv(report_dir / f"{stem}_criterion_matrix.tsv", matrix,
               ["accession_group", "role", "criterion_order", "criterion",
                "state", "evidence", "source_section"])

    deck_text = "\n".join(
        _slide_text(slide) + "\n" + slide.notes_slide.notes_text_frame.text
        for slide in prs.slides
    )
    regional_guardrail = (
        "Cortical data exist, but use them only for descriptive astrocyte concordance" in deck_text
        and "Cortex/PFC:" not in deck_text
        and "cortex/PFC mice" not in deck_text
    )
    replicate_guardrail = (
        "Inferential n = 3 mice—not 15,000 cells" in deck_text
        and "Pooling mice creates one pooled profile" in deck_text
        and "n=1/cell" not in deck_text
        and "factorial cells" not in deck_text
    )
    no_overall_score = "not summed into a score" in deck_text
    gse143758_context = (
        "broad atlas is hippocampal" in deck_text
        and "cortex matrix is astrocyte-focused" in deck_text
        and "one record per age×genotype group" in deck_text
    )

    checks = [
        ("source_contract", True, f"Markdown SHA-256 {SOURCE_MD_SHA256}"),
        ("pptx_integrity", True, "PPTX ZIP and python-pptx validation passed"),
        ("slide_structure", len(prs.slides) == 14,
         "Two criteria slides followed by twelve dataset profiles"),
        ("criterion_matrix", len(matrix) == 96,
         "12 profiles × 8 noncompensatory criteria"),
        ("all_native_editable", all(row["picture_count"] == 0 for row in inventory),
         "No raster pictures; diagrams and scorecards are editable shapes"),
        ("speaker_notes", all(row["speaker_note_words"] >= 55 for row in inventory),
         "Every slide has three-section speaker notes"),
        ("regional_guardrail", regional_guardrail,
         "Cortex is partial PFC support; hippocampus/whole brain remain distinct"),
        ("replicate_guardrail", replicate_guardrail,
         "Mouse or pre-sequencing pool—not cell/nucleus/well—is the inferential unit"),
        ("no_overall_score", no_overall_score,
         "Criterion states are not summed into a universal dataset rank"),
        ("gse143758_context", gse143758_context,
         "Cortex exists but main broad-cell evidence is hippocampal and cortex is astrocyte-focused"),
        ("visual_review", visual_review_status == "complete",
         "Reviewed in color and grayscale" if visual_review_status == "complete"
         else "Awaiting color and grayscale review"),
    ]
    blocking_failures = [
        check_id for check_id, passed, _ in checks
        if check_id != "visual_review" and not passed
    ]
    if blocking_failures:
        raise AssertionError(f"Presentation report checks failed: {blocking_failures}")
    _write_tsv(
        report_dir / f"{stem}_checks.tsv",
        [{"check_id": cid, "passed": str(bool(passed)), "detail": detail}
         for cid, passed, detail in checks],
        ["check_id", "passed", "detail"],
    )

    status = [{
        "schema_version": "mouse_dataset_comparison_deck_v1",
        "deck_id": "mouse_dataset_comparison",
        "validation_status": (
            "validated_complete" if visual_review_status == "complete"
            else "awaiting_visual_review"
        ),
        "visual_review_status": visual_review_status,
        "slides": len(prs.slides),
        "criteria_slides": 2,
        "dataset_slides": 12,
        "criterion_rows": len(matrix),
        "checks": len(checks),
        "deck_path": display_path(deck_path),
        "deck_bytes": deck_path.stat().st_size,
        "deck_sha256": sha256(deck_path),
        "source_markdown_sha256": SOURCE_MD_SHA256,
        "completed_utc": datetime.now(timezone.utc).isoformat(),
    }]
    _write_tsv(
        report_dir / f"{stem}_status.tsv", status,
        ["schema_version", "deck_id", "validation_status",
         "visual_review_status", "slides", "criteria_slides",
         "dataset_slides", "criterion_rows", "checks", "deck_path",
         "deck_bytes", "deck_sha256", "source_markdown_sha256",
         "completed_utc"],
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--report-dir", type=Path, default=DEFAULT_REPORT_DIR)
    parser.add_argument("--visual-review-status", choices=("pending", "complete"),
                        default="pending")
    parser.add_argument("--validate-only", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = args.output.resolve()
    if args.validate_only:
        validate_deck(output)
        print(f"Validated: {output}")
        return
    built = build_deck(output)
    write_report(built, args.report_dir.resolve(),
                 visual_review_status=args.visual_review_status)
    print(f"Built and validated: {built}")


if __name__ == "__main__":
    main()
