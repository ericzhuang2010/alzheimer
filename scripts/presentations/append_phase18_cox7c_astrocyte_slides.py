#!/usr/bin/env python3
"""Append the COX7C astrocyte pathway and STRING slides to the Phase 18 deck."""

from __future__ import annotations

import argparse
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN


REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO))

from scripts.presentations.build_phase18_key_driver_selection_deck import (  # noqa: E402
    BLUE,
    GRAY,
    LIGHT,
    NAVY,
    PALE_BLUE,
    TEAL,
    VERMILION,
    WHITE,
    add_bullets,
    add_header,
    add_metric,
    add_panel_title,
    add_picture_contain,
    add_rect,
    add_source,
    add_text,
    new_slide,
    trim_white,
)


DEFAULT_DECK = REPO / "docs/presentations/key_driver_selection_analysis 08182026.pptx"
FIG_ROOT = REPO / "results/figures/analysis/phase_18_key_driver_selection/COX7C/astrocytes"
PATHWAY_FIG = FIG_ROOT / "phase18_cox7c_astrocyte_consensus_network_pathways.png"
STRING_FIG = FIG_ROOT / "stringdb_full_medium_conf.png"


def replace_text_preserve_style(slide, old: str, new: str) -> None:
    """Replace exact paragraph text while retaining the first run's formatting."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            current = "".join(run.text for run in paragraph.runs)
            if current != old:
                continue
            if paragraph.runs:
                paragraph.runs[0].text = new
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = new
            return
    raise RuntimeError(f"Expected deck text was not found: {old!r}")


def append_slides(input_path: Path, output_path: Path) -> Path:
    for path in (input_path, PATHWAY_FIG, STRING_FIG):
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation(input_path)
    if len(prs.slides) != 25:
        raise RuntimeError(f"Expected a 25-slide source deck, found {len(prs.slides)}")

    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    if "COX7C • ASTROCYTES" in all_text:
        raise RuntimeError("COX7C astrocyte slides already appear to be present")

    # Keep the agenda and synthesis handoff consistent with the appended slides.
    replace_text_preserve_style(prs.slides[1], "Slides 16–25", "Slides 16–27")
    replace_text_preserve_style(
        prs.slides[1],
        "Place RPL11 and APOE on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
        "Place APOE, RPL11, and COX7C on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "The broad atlas establishes which candidates emerged; the next slides localize support before the RPL11 case study.",
        "The broad atlas establishes which candidates emerged; the next slides localize support before the candidate case studies.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "Next: localize candidate evidence across sex/APOE strata, then examine RPL11 in depth.",
        "Next: localize candidate evidence across sex/APOE strata, then examine APOE, RPL11, and COX7C in depth.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–25 show RPL11 and APOE network/protein examples",
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–27 show APOE, RPL11, and COX7C network/protein examples",
    )

    with tempfile.TemporaryDirectory(prefix="cox7c_deck_assets_") as temp_dir:
        assets = Path(temp_dir)
        pathway = trim_white(PATHWAY_FIG, assets)
        string = trim_white(STRING_FIG, assets)

        # 26 — COX7C astrocyte directed pathway figure
        slide = new_slide(prs)
        add_header(
            slide,
            "COX7C • astrocytes",
            "A compact astrocyte signal maps to a coherent respiratory module",
            26,
            accent=BLUE,
            subtitle="COX7C was self-excluded from both mitochondrial queries; the retained hits therefore reflect other genes in its modeled neighborhood.",
        )
        add_picture_contain(
            slide,
            pathway,
            0.11,
            1.22,
            8.91,
            5.92,
            alt="COX7C-centered astrocyte consensus network with directed Bayesian-network edges and significant pathway outlines",
        )
        add_rect(slide, 9.18, 1.42, 3.62, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.45, 1.70, 3.07, accent=BLUE)
        add_metric(slide, "2", "conservative-support astrocyte runs", 9.45, 2.18, 1.39, accent=BLUE)
        add_metric(slide, "5.14×10⁻⁴", "aggregate ACAT q", 10.99, 2.18, 1.54, accent=TEAL)
        add_panel_title(slide, "Network readout", 9.45, 3.48, 3.07, accent=TEAL)
        add_bullets(
            slide,
            [
                "23 nodes and 23 directed edges retain ten self-excluded mitochondrial query hits.",
                "Upstream chain: RPLP1 → RPL11 → COX7C; COX7C has 15 direct modeled outputs.",
                "ETC/OXPHOS: 10 genes, FDR 3.27×10⁻¹¹; cristae: 4 genes, FDR 1.45×10⁻⁴; cytosolic ribosome: 4 genes, FDR 0.00722.",
                "The 1/2 display threshold preserves ten hits; requiring 2/2 would retain only four.",
            ],
            9.45,
            3.96,
            2.94,
            size=10.5,
            line_h=0.58,
            accent=TEAL,
        )
        add_text(
            slide,
            "Interpretation: strong respiratory-module coherence, but the two-run astrocyte result needs independent replication.",
            9.48,
            6.46,
            2.88,
            0.28,
            size=9.4,
            color=NAVY,
            bold=True,
        )
        add_source(
            slide,
            "Source: phase18_cox7c_astrocyte_consensus_network_pathways.png, caption, methods, ORA table, and gene-by-gene analysis",
        )

        # 27 — COX7C astrocyte STRING figure and validation interpretation
        slide = new_slide(prs)
        add_header(
            slide,
            "COX7C • astrocytes",
            "STRING places COX7C in a dense respiratory core",
            27,
            accent=BLUE,
            subtitle="Medium-confidence STRING associations are undirected and not astrocyte- or AD-specific; dense complex membership is expected.",
        )
        add_rect(slide, 0.42, 1.39, 6.37, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(
            slide,
            string,
            0.64,
            1.62,
            5.93,
            4.98,
            alt="STRING medium-confidence functional association network for COX7C and astrocyte consensus-neighborhood proteins",
        )
        add_rect(slide, 7.04, 1.39, 5.76, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 7.34, 1.72, 5.15, accent=BLUE)
        add_bullets(
            slide,
            [
                "COX7C sits in a dense association core spanning complex I, III, IV, and V proteins plus TOMM7 and SLIRP.",
                "A highly connected ribosomal branch echoes the modeled RPLP1 → RPL11 → COX7C upstream chain.",
                "Several Bayesian-neighborhood proteins remain weakly connected or isolated, so STRING does not validate every modeled edge.",
            ],
            7.34,
            2.20,
            5.05,
            size=11.3,
            line_h=0.61,
            accent=BLUE,
        )
        add_panel_title(slide, "Driver or respiratory sentinel?", 7.34, 4.21, 5.15, accent=TEAL)
        add_text(
            slide,
            "COX7C is nuclear encoded and is mostly AD-down in earlier DEG results, contrasting with many AD-up mtDNA-encoded subunits. That pattern supports a mitonuclear-stoichiometry hypothesis, but STRING density alone cannot make COX7C an upstream regulator.",
            7.35,
            4.68,
            5.02,
            0.82,
            size=10.9,
            color=GRAY,
        )
        add_rect(slide, 7.35, 5.66, 4.98, 0.80, color=PALE_BLUE, outline=BLUE)
        add_text(
            slide,
            "Compare mild COX7C and COX4I1 perturbation: exact target module, complex-IV assembly/activity, protein stoichiometry, respiration, and mitochondrial mass. Replicate the astrocyte result independently.",
            7.57,
            5.80,
            4.55,
            0.50,
            size=9.8,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "The weak direct bulk-sQTL supports prioritization; one source is not replication of both network contexts.",
            7.41,
            6.57,
            4.86,
            0.20,
            size=8.9,
            color=VERMILION,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_source(
            slide,
            "Sources: stringdb_full_medium_conf.png • phase18_key_driver_gene_by_gene_initial_analysis.md • human-genetic-support review",
        )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    validate_output(output_path)
    return output_path


def validate_output(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 27:
        raise RuntimeError(f"Expected 27 slides, found {len(prs.slides)}")
    image_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    if image_shapes != 16:
        raise RuntimeError(f"Expected 16 figure images, found {image_shapes}")

    endings = []
    for slide_index in range(len(prs.slides) - 2, len(prs.slides)):
        slide = prs.slides[slide_index]
        endings.append(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
        )
    if not all("COX7C • ASTROCYTES" in text for text in endings):
        raise RuntimeError("The final two slides are not the expected COX7C astrocyte slides")

    with zipfile.ZipFile(path) as archive:
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels"))
        )
    if "filter_attrition" in package_text:
        raise RuntimeError("Deprecated filter_attrition content found in the PPTX package")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    args = parser.parse_args()
    result = append_slides(args.input, args.output)
    print(result)
