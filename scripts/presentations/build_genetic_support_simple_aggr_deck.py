#!/usr/bin/env python3
"""Build the human-genetic-support deck for the current (simple) driver list.

This standalone deck re-anchors the frozen Phase 19 public-data genetic
evidence to the authoritative returned-only simple-aggregation key drivers.
It replaces the outdated 2026-08-25 genetic-support deck, which described the
previous 25-gene candidate freeze; that file is left untouched.

Driver-membership facts are recomputed from the authoritative ROSMAP and
SEA-AD aggregation tables, regional AD association values are read from the
frozen Phase 19 recovery summary, and the qualitative evidence grades follow
the consolidated Phase 19 results summary. The build fails if any recomputed
fact drifts from the frozen expected values.

WARNING (2026-08-29): the published deck has since received manual edits (a
deleted slide, reordered slides) and a surgically inserted dataset slide
(``insert_genetic_support_dataset_slide.py``). Rerunning this builder would
overwrite those changes. Do not rebuild without first reconciling the edits.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import update_phase11_seaad_simple_aggr_part2 as ui  # noqa: E402  (shared styling helpers)


ROOT = Path(__file__).resolve().parents[2]
DEFAULT_OUT = (
    ROOT
    / "docs"
    / "presentations"
    / "human_genetic_support_for_key_drivers_simple_aggr_08292026.pptx"
)
ROS_CATEGORIES = (
    ROOT
    / "results/minerva_production/20_sex_apoe_kda_simple_aggr"
    / "simple_category_gene_aggregates.tsv"
)
SEA_CATEGORIES = (
    ROOT
    / "results/validation_human/11_sex_apoe_kda_simple_aggr"
    / "simple_category_gene_aggregates.tsv"
)
REGIONAL_GWAS = (
    ROOT
    / "results/minerva_production/19_genetic_support_tier2_recovery"
    / "recovery_regional_gwas_summary.tsv"
)
AUDIT_PATH = (
    ROOT
    / "results/presentations/human_genetic_support_simple_aggr"
    / "build_checks.tsv"
)

SCREENED_NUCLEAR = [
    "APOE", "COX7C", "COX4I1", "COX6B1", "UQCR10", "RPL11", "RPLP1", "RPL15",
    "LAPTM4A", "SELENOW", "RPS13", "LAMTOR5", "DYNLT1", "RPS15", "ATP6V1F",
    "RPL38", "FTL", "ANKRD11", "NCOA1",
]
SCREENED_MTDNA = ["MT-ATP6", "MT-CO2", "MT-CO3", "MT-CYB", "MT-ND4", "MT-ND5"]
MT_CLASS_NUCLEAR = ["COX7C", "COX4I1", "COX6B1", "UQCR10"]
NO_SIGNAL_DRIVERS = [
    "RPL11", "RPLP1", "RPL15", "LAMTOR5", "RPS13", "DYNLT1", "ATP6V1F",
    "FTL", "LAPTM4A", "RPL38", "NCOA1",
]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--audit", type=Path, default=AUDIT_PATH)
    return parser.parse_args()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def load_facts() -> dict[str, Any]:
    for path in (ROS_CATEGORIES, SEA_CATEGORIES, REGIONAL_GWAS):
        if not path.is_file():
            raise FileNotFoundError(path)

    full = pd.read_csv(
        ROS_CATEGORIES, sep="\t", dtype=str, keep_default_na=False, na_values=["NA"]
    )
    ros = full[full["case_id"].eq("non_mt_driver")].copy()
    ros["score"] = ros["returned_run_q_acat_score"].astype(float)
    ros = ros.sort_values(
        ["signature_group", "broad_network", "score", "current_symbol"],
        kind="mergesort",
    )
    ros["rank"] = ros.groupby(["signature_group", "broad_network"]).cumcount() + 1

    sea = pd.read_csv(
        SEA_CATEGORIES, sep="\t", dtype=str, keep_default_na=False, na_values=["NA"]
    )
    sea_genes = set(sea.loc[sea["case_id"].eq("non_mt_driver"), "current_symbol"])

    regional = pd.read_csv(REGIONAL_GWAS, sep="\t", dtype=str)
    if len(regional) != 19:
        raise RuntimeError("Expected 19 regional GWAS rows")
    regional_significant = set(
        regional.loc[regional["regional_gwas_signal"].eq("TRUE"), "gene"]
    )
    if regional_significant != {"APOE", "RPS15", "ANKRD11", "COX7C"}:
        raise RuntimeError(f"Unexpected significant regions: {regional_significant}")
    min_p = dict(zip(regional["gene"], regional["regional_min_p"]))
    lead = dict(zip(regional["gene"], regional["regional_lead_variant"]))

    ros_genes = set(ros["current_symbol"])
    per_gene: dict[str, dict[str, Any]] = {}
    for gene in SCREENED_NUCLEAR:
        sub = ros[ros["current_symbol"].eq(gene)]
        per_gene[gene] = {
            "is_driver": len(sub) > 0,
            "categories": len(sub),
            "best_rank": int(sub["rank"].min()) if len(sub) else None,
            "best_q": float(sub["score"].min()) if len(sub) else None,
            "top5": bool((sub["rank"] <= 5).any()) if len(sub) else False,
            "in_seaad": gene in sea_genes,
        }
    carried = [g for g in SCREENED_NUCLEAR if per_gene[g]["is_driver"]]
    if len(carried) != 15 or set(SCREENED_NUCLEAR) - set(carried) != set(MT_CLASS_NUCLEAR):
        raise RuntimeError("Carried-driver contract failed")
    if not all(per_gene[g]["top5"] for g in carried):
        raise RuntimeError("Expected every carried driver to have a top-five appearance")
    for gene in MT_CLASS_NUCLEAR:
        if not (full["current_symbol"].eq(gene) & full["case_id"].eq("mt_driver")).any():
            raise RuntimeError(f"{gene} missing from the mitochondrial driver class")

    apoe = ros[ros["current_symbol"].eq("APOE")]
    apoe_categories = sorted(
        f"{row.signature_group}" for row in apoe.itertuples()
    )
    if set(apoe["broad_network"]) != {"Astrocytes"} or len(apoe) != 3:
        raise RuntimeError("APOE driver-context contract failed")
    rps15 = per_gene["RPS15"]
    if rps15["categories"] != 12 or not rps15["in_seaad"]:
        raise RuntimeError("RPS15 contract failed")
    if per_gene["SELENOW"]["categories"] != 8 or per_gene["ANKRD11"]["categories"] != 2:
        raise RuntimeError("SELENOW/ANKRD11 contract failed")

    unscreened = ros_genes - set(SCREENED_NUCLEAR) - set(SCREENED_MTDNA)
    cross_cohort_unscreened = sorted(
        (ros_genes & sea_genes) - set(SCREENED_NUCLEAR) - set(SCREENED_MTDNA)
    )
    priority = [g for g in ["WDR82", "HGSNAT", "TTC8", "BEX3"] if g in cross_cohort_unscreened]
    if len(priority) != 4:
        raise RuntimeError("Cross-cohort unscreened priority contract failed")

    ros["calls"] = ros["returned_call_count"].astype(int)
    gene_depth = ros.groupby("current_symbol").agg(
        cats=("current_symbol", "size"), calls=("calls", "sum")
    )
    top5_data = pd.read_csv(
        ROOT
        / "results/figures/analysis/phase_20_sex_apoe_simple_aggr"
        / "top5_candidates/phase20_simple_aggr_top5_candidates_plot_data.tsv",
        sep="\t",
    )
    depth = {
        "unit_count": len(ros),
        "one_off_genes": int(((gene_depth.cats == 1) & (gene_depth.calls == 1)).sum()),
        "multi_category_genes": int((gene_depth.cats >= 2).sum()),
        "calls3_genes": int((gene_depth.calls >= 3).sum()),
        "top5_genes": int(top5_data["current_symbol"].nunique()),
        "seaad_shared_genes": len(ros_genes & sea_genes),
    }
    expected_depth = {
        "unit_count": 689,
        "one_off_genes": 285,
        "multi_category_genes": 125,
        "calls3_genes": 78,
        "top5_genes": 105,
        "seaad_shared_genes": 35,
    }
    if depth != expected_depth:
        raise RuntimeError(f"Evidence-depth contract failed: {depth}")

    return {
        "depth": depth,
        "ros_gene_count": len(ros_genes),
        "per_gene": per_gene,
        "carried": carried,
        "apoe_groups": apoe_categories,
        "min_p": min_p,
        "lead": lead,
        "unscreened_count": len(unscreened),
        "cross_cohort_unscreened": cross_cohort_unscreened,
        "priority": priority,
    }


def q_text(value: float) -> str:
    return f"{value:.1e}".replace("e-0", "e-")


def add_flow(slide, nodes, y, *, accent, start_x=0.75, total_w=11.85, node_h=1.25):
    gap = 0.30
    node_w = (total_w - gap * (len(nodes) - 1)) / len(nodes)
    for index, (heading, detail) in enumerate(nodes):
        x = start_x + index * (node_w + gap)
        ui.add_rect(slide, x, y, node_w, node_h, color=ui.WHITE, outline=ui.LIGHT)
        ui.add_text(
            slide, heading, x + 0.12, y + 0.14, node_w - 0.24, 0.52,
            size=12.6, color=ui.NAVY, bold=True, align=PP_ALIGN.CENTER,
        )
        ui.add_text(
            slide, detail, x + 0.12, y + 0.66, node_w - 0.24, 0.50,
            size=9.4, color=ui.GRAY, align=PP_ALIGN.CENTER,
        )
        if index < len(nodes) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(x + node_w + 0.06),
                Inches(y + 0.47),
                Inches(0.18),
                Inches(0.32),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = accent
            arrow.line.fill.background()


def build_deck(output_path: Path, facts: dict[str, Any]) -> list[str]:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    properties = prs.core_properties
    properties.title = "Human genetic support for the key drivers"
    properties.subject = (
        "Public-data genetic evidence for the key drivers from the network "
        "key-driver (KDA) analysis"
    )
    properties.author = "Alzheimer project analysis team"
    properties.language = "en-US"
    generated_at = datetime.now(timezone.utc).replace(tzinfo=None)
    properties.created = generated_at
    properties.modified = generated_at
    titles: list[str] = []

    per_gene = facts["per_gene"]
    min_p = facts["min_p"]

    # 1 — title
    slide = ui.new_slide(prs, bg=ui.NAVY)
    ui.add_text(slide, "HUMAN GENETIC SUPPORT", 0.78, 0.67, 4.0, 0.28, size=11, color=ui.SKY, bold=True)
    title = "Do the key drivers carry inherited Alzheimer's risk?"
    ui.add_text(
        slide, title, 0.78, 1.20, 11.60, 1.38, size=33, color=ui.WHITE, bold=True,
        font=ui.FONT_HEAD, valign=MSO_ANCHOR.MIDDLE,
    )
    ui.add_text(
        slide,
        "Public genetic evidence for the key drivers from the network key-driver (KDA) analysis",
        0.82, 2.72, 11.2, 0.48, size=16, color=ui.RGBColor(204, 219, 234),
    )
    cards = [
        ("4", "drivers with genetic evidence"),
        ("APOE + RPS15", "lead the genetic evidence"),
        (f"15 of {facts['ros_gene_count']}", "drivers screened so far"),
    ]
    for index, (value, label) in enumerate(cards):
        x = 0.82 + index * 4.02
        ui.add_rect(slide, x, 3.75, 3.72, 1.60, color=ui.NAVY_2, outline=ui.BLUE)
        ui.add_text(slide, value, x + 0.25, 4.02, 3.22, 0.50, size=22, color=ui.SKY, bold=True, font=ui.FONT_HEAD, align=PP_ALIGN.CENTER)
        ui.add_text(slide, label, x + 0.25, 4.62, 3.22, 0.55, size=11.5, color=ui.RGBColor(210, 224, 238), align=PP_ALIGN.CENTER)
    ui.add_text(slide, "Public-data study • 29 August 2026", 0.82, 6.55, 7.0, 0.24, size=10, color=ui.RGBColor(157, 176, 195))
    ui.add_notes(
        slide,
        goal="Set the question: does inherited human variation support the network key drivers?",
        walkthrough="Public genetic evidence exists for fifteen of the driver genes so far. Four of them show genetic evidence, led by APOE and RPS15; the rest of the driver list has not yet been screened.",
        boundary="No new association analysis was run; all genetic values are from the frozen public-data workstreams.",
        transition="First, how the two kinds of evidence relate.",
    )
    titles.append(title)

    # 2 — framing: two complementary questions
    title = "Network drivers and inherited risk are complementary evidence"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "The KDA analysis finds genes central in diseased brain cells; genetics asks whether DNA variation near them alters risk.")
    ui.add_rect(slide, 0.64, 1.52, 5.86, 4.35, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_text(slide, "NETWORK KEY DRIVERS", 0.95, 1.84, 4.6, 0.27, size=10, color=ui.BLUE, bold=True)
    ui.add_text(slide, "Which genes are central in diseased brain-cell networks?", 0.95, 2.18, 4.95, 0.62, size=18, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
    ui.add_metric(slide, str(facts["ros_gene_count"]), "driver genes", 0.95, 3.10, 2.30, accent=ui.BLUE)
    ui.add_metric(slide, "32", "sex/APOE · cell-type categories", 3.50, 3.10, 2.30, accent=ui.BLUE)
    ui.add_text(
        slide,
        "The strongest drivers also replicate in an independent human cohort.",
        0.96, 4.55, 5.25, 0.55, size=11.6, color=ui.GRAY,
    )
    ui.add_rect(slide, 6.83, 1.52, 5.86, 4.35, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_text(slide, "INHERITED RISK", 7.14, 1.84, 4.6, 0.27, size=10, color=ui.TEAL_TEXT, bold=True)
    ui.add_text(slide, "Which genes are linked to inherited AD risk?", 7.14, 2.18, 4.95, 0.62, size=18, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
    ui.add_bullets(slide, [
        "Uses public AD and biomarker genetics, independent of the networks.",
        "Evidence is per gene: disease association plus gene-activity links.",
        "Goal: find drivers where both kinds of evidence converge.",
    ], 7.15, 3.20, 5.25, size=11.6, line_h=0.62, accent=ui.TEAL)
    ui.add_text(
        slide,
        "Genetic support strengthens a driver; its absence does not refute one — the two ask different causal questions.",
        0.88, 6.30, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Frame genetics as an independent, complementary line of evidence for the driver list.",
        walkthrough="The left card summarizes the driver list under test: 433 genes across 32 sex/APOE-by-cell-type categories, with the strongest drivers replicating in an independent cohort. The right card states the genetic question asked of each gene.",
        boundary="Network centrality does not imply a germline risk variant, and vice versa; convergence is the interesting outcome.",
        transition="First, what exactly the driver list contains.",
    )
    titles.append(title)

    # 3 — where the 433 come from
    depth = facts["depth"]
    title = f"Where the {facts['ros_gene_count']} driver genes come from"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Every gene the network analysis returned as significant is on the list — there are no extra selection gates.")
    add_flow(slide, [
        ("295 KDA calls", "network key-driver calls across the ROSMAP categories"),
        ("Significant returns", "genes returned at q ≤ 0.05 within a call"),
        ("Non-MT focus", "mitochondrial-protein genes set aside"),
        (f"{depth['unit_count']} units", "gene × category pairs"),
        (f"{facts['ros_gene_count']} genes", "unique driver symbols"),
    ], 1.85, accent=ui.TEAL)
    ui.add_rect(slide, 0.75, 3.60, 11.85, 1.15, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_text(slide, "Why units and genes differ", 1.00, 3.82, 4.0, 0.28, size=11.5, color=ui.BLUE, bold=True)
    ui.add_text(
        slide,
        "A gene counts once per sex/APOE × cell-type category it returns in, so recurrent genes hold several of the "
        f"{depth['unit_count']} units. Collapsing to unique symbols leaves {facts['ros_gene_count']} genes.",
        1.00, 4.14, 11.30, 0.50, size=11.5, color=ui.DARK,
    )
    ui.add_rect(slide, 0.75, 5.05, 11.85, 1.15, color=ui.PALE_GOLD, outline=ui.GOLD)
    ui.add_text(slide, "A census, not a shortlist", 1.00, 5.27, 4.0, 0.28, size=11.5, color=ui.GOLD_TEXT, bold=True)
    ui.add_text(
        slide,
        "One significant return anywhere is enough for membership. The list is deliberately permissive, so the depth "
        "of network evidence behind each gene varies widely — shown on the next slide.",
        1.00, 5.59, 11.30, 0.50, size=11.5, color=ui.DARK,
    )
    ui.add_notes(
        slide,
        goal="Define exactly what the 433 number counts before any genetic claims reference it.",
        walkthrough="The 295 ROSMAP network key-driver calls return significant genes at within-call q at most 0.05. After setting aside mitochondrial-protein genes, the returns form 689 gene-by-category units, which collapse to 433 unique driver genes.",
        boundary="Membership requires only one significant return; the within-call q values are exploratory and not cross-call FDR-controlled.",
        transition="How deep the evidence is behind each of the 433.",
    )
    titles.append(title)

    # 4 — evidence depth
    title = f"Evidence depth inside the {facts['ros_gene_count']} is very uneven"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "One significant return is enough for membership; deeper network evidence is concentrated in a minority.")
    bars = [
        (depth["one_off_genes"], "One-off genes: a single category, a single returned call", ui.VERMILION),
        (depth["multi_category_genes"], "Recur in ≥2 sex/APOE × cell-type categories", ui.BLUE),
        (depth["top5_genes"], "Appear in a category top-five display", ui.BLUE),
        (depth["calls3_genes"], "Backed by ≥3 returned calls", ui.TEAL),
        (depth["seaad_shared_genes"], "Also return in the independent SEA-AD cohort", ui.TEAL),
    ]
    total = facts["ros_gene_count"]
    for index, (count, label, color) in enumerate(bars):
        y = 1.70 + index * 0.88
        ui.add_text(slide, label, 0.85, y, 5.75, 0.55, size=12.0, color=ui.NAVY, bold=True)
        ui.add_rect(slide, 6.75, y + 0.05, 4.60 * count / total, 0.34, color=color, outline=None, radius=False)
        ui.add_text(
            slide,
            f"{count}  ({round(100 * count / total)}%)",
            6.75 + 4.60 * count / total + 0.12, y + 0.06, 1.60, 0.30,
            size=12.5, color=color, bold=True,
        )
    ui.add_text(slide, "Groups overlap; each bar is counted against all 433 genes.", 6.75, 6.12, 5.60, 0.26, size=9.5, color=ui.GRAY, italic=True)
    ui.add_text(
        slide,
        "Follow-up work should weight the deep-evidence minority — not treat all 433 genes equally.",
        0.88, 6.55, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Show that the permissive list has a small core of well-supported genes and a large one-off tail.",
        walkthrough="Two hundred eighty-five of the four hundred thirty-three genes, about two-thirds, have exactly one category and one returned call, often with q just under 0.05. One hundred twenty-five recur in at least two categories, one hundred five reach a top-five display, seventy-eight have at least three returned calls, and thirty-five also return in SEA-AD.",
        boundary="The groups overlap, so the bars do not sum to 433; each is a share of the full list.",
        transition="The high-level approach: three sources of genetic evidence and how they combine.",
    )
    titles.append(title)

    # 5 — three sources of evidence
    title = "High-level approach: three sources of public genetic evidence"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "One source is ready-made gene-level evidence; the other two each tell half the story.")
    cards = [
        (
            "1 · PUBLISHED RESULTS",
            "FunGen AD integration",
            "Already-computed variant-to-gene mappings: AD fine-mapping, TWAS, and gene lists.",
            "Directly supports an AD–gene relationship (this is how APOE is established).",
            "Covers only what has been published.",
            ui.PALE_GREEN, ui.TEAL,
        ),
        (
            "2 · DISEASE GWAS",
            "Clinical AD + CSF biomarkers",
            "DNA variants associated with an AD diagnosis or with amyloid/tau spinal-fluid levels.",
            "Anchors a genomic region to disease.",
            "Points to a region — not to a specific gene.",
            ui.PALE_SKY, ui.BLUE,
        ),
        (
            "3 · BRAIN QTL",
            "Expression + splicing panels",
            "DNA variants that change a gene's RNA amount or splicing in brain cells and tissue.",
            "Anchors a variant to the driver gene itself.",
            "Says nothing about disease on its own.",
            ui.PALE_GOLD, ui.GOLD,
        ),
    ]
    for index, (kicker, heading, what, gives, limit, bg, accent) in enumerate(cards):
        x = 0.66 + index * 4.10
        ui.add_rect(slide, x, 1.55, 3.83, 4.55, color=bg, outline=accent)
        ui.add_text(slide, kicker, x + 0.26, 1.82, 3.30, 0.26, size=10, color=ui.readable_accent(accent), bold=True)
        ui.add_text(slide, heading, x + 0.26, 2.14, 3.32, 0.60, size=16.5, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
        ui.add_text(slide, what, x + 0.26, 2.90, 3.32, 0.95, size=10.8, color=ui.DARK)
        ui.add_rect(slide, x + 0.24, 3.95, 3.35, 0.90, color=ui.WHITE, outline=ui.LIGHT)
        ui.add_text(slide, gives, x + 0.40, 4.12, 3.05, 0.60, size=10.6, color=ui.readable_accent(accent), bold=True)
        ui.add_text(slide, f"Limit: {limit}", x + 0.26, 5.05, 3.32, 0.75, size=10.2, color=ui.GRAY, italic=True)
    ui.add_text(
        slide,
        "Genetic support for a driver comes from source 1 directly — or from combining sources 2 and 3.",
        0.88, 6.45, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Reduce the data inventory to three roles: ready-made gene-level results, disease anchors, and gene anchors.",
        walkthrough="Published AD integration results from the FunGen consortium already map variants to genes and directly support AD-gene relationships. Disease GWAS — the clinical AD meta-analysis and the three CSF biomarker studies — tie variants to disease but only at region resolution. Brain QTL panels tie variants to a specific gene's activity but carry no disease information.",
        boundary="All three are summary-level public data; no individual genotypes were used.",
        transition="How sources two and three combine into gene-level support.",
    )
    titles.append(title)

    # 6 — combining GWAS and QTL
    title = "Combining GWAS and QTL: one variant, two effects"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "A driver gains genetic support when the disease signal and the gene-activity signal trace to the same DNA variant.")
    logic = [
        ("GWAS says", "a DNA variant near the driver raises AD risk", ui.PALE_SKY, ui.BLUE),
        ("QTL says", "the same variant changes the driver's activity in brain", ui.PALE_GOLD, ui.GOLD),
        ("Together", "the driver's activity is a plausible path from inherited DNA to disease", ui.PALE_GREEN, ui.TEAL),
    ]
    for index, (label, statement, bg, accent) in enumerate(logic):
        y = 1.60 + index * 0.98
        ui.add_rect(slide, 0.80, y, 2.20, 0.80, color=accent, outline=None)
        ui.add_text(slide, label, 0.92, y + 0.24, 1.96, 0.32, size=13.5, color=ui.WHITE, bold=True, align=PP_ALIGN.CENTER)
        ui.add_rect(slide, 3.15, y, 9.40, 0.80, color=bg, outline=accent)
        ui.add_text(slide, statement, 3.42, y + 0.22, 8.90, 0.38, size=13.0, color=ui.NAVY, bold=True)
    add_flow(slide, [
        ("Step 1", "find an AD-associated region near the driver (GWAS)"),
        ("Step 2", "confirm the driver has a QTL signal in that region"),
        ("Step 3", "same-variant test: do the two signals share one variant?"),
    ], 4.75, accent=ui.TEAL)
    ui.add_rect(slide, 0.80, 6.20, 11.85, 0.72, color=ui.PALE_RED, outline=ui.VERMILION)
    ui.add_text(
        slide,
        "Proximity is not enough — an AD region can hold many genes and variants. The same-variant test is the decisive step, "
        "and it needs complete statistics that are not yet publicly available for any driver.",
        1.04, 6.34, 11.35, 0.48, size=10.8, color=ui.VERMILION_TEXT, bold=True,
    )
    ui.add_notes(
        slide,
        goal="Explain the combination logic that turns two region-level signals into gene-level support.",
        walkthrough="GWAS ties a variant to disease; QTL ties a variant to the driver's activity. If a formal test shows both signals share one variant, the driver's regulation becomes a plausible causal path. Operationally: find the AD region, confirm the driver's QTL, then run the same-variant (colocalization) test.",
        boundary="The same-variant test requires complete regional statistics, fitted multi-signal models, and matched reference panels; those inputs are missing publicly, so no driver has completed step three yet.",
        transition="The staged pipeline implements exactly this order.",
    )
    titles.append(title)

    # 4 — staged design
    title = "Evidence was checked in stages; a route stops at the first failed gate"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Stopping early is a data statement, not proof that a gene has no genetic role.")
    add_flow(slide, [
        ("Nearby AD signal", "region association below the genome-wide threshold"),
        ("Gene-activity link", "significant brain QTL for the gene itself"),
        ("Compatible models", "complete statistics and matched reference data"),
        ("Same-variant test", "do both signals point to one variant?"),
    ], 1.85, accent=ui.BLUE)
    ui.add_rect(slide, 0.75, 3.60, 11.85, 1.10, color=ui.PALE_GOLD, outline=ui.GOLD)
    ui.add_text(slide, "Where the screen currently ends", 1.00, 3.82, 4.0, 0.28, size=11.5, color=ui.GOLD_TEXT, bold=True)
    ui.add_text(
        slide,
        "No gene completed the final same-variant test: the required complete gene-activity models and matched reference data "
        "are not yet publicly available. All current conclusions rest on the first two stages.",
        1.00, 4.14, 11.30, 0.48, size=11.5, color=ui.DARK,
    )
    ui.add_rect(slide, 0.75, 5.00, 11.85, 1.10, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_text(slide, "Reading a stopped route", 1.00, 5.22, 4.0, 0.28, size=11.5, color=ui.VERMILION_TEXT, bold=True)
    ui.add_text(
        slide,
        "\u201CNo signal\u201D means the tested region truly had no qualifying association. \u201CNot assessable\u201D means the required "
        "measurement or model was unavailable — the gene was not tested and proven negative.",
        1.00, 5.54, 11.30, 0.48, size=11.5, color=ui.DARK,
    )
    ui.add_notes(
        slide,
        goal="Prevent over-reading: separate signal-negative results from unassessable routes and note that zero same-variant tests completed.",
        walkthrough="Four stages: regional AD association, gene-level activity link, model compatibility, and the same-variant comparison. Every route stopped at or before stage three because complete fitted models and matched reference panels were unavailable publicly.",
        boundary="Genome-wide threshold five-times-ten-to-the-minus-eight; QTL gates were gene-specific corrected thresholds.",
        transition="Which drivers have screening results so far.",
    )
    titles.append(title)

    # 5 — screening coverage
    title = f"15 of the {facts['ros_gene_count']} drivers have genetic screening results so far"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Screening has focused on prominent driver genes; every screened gene holds a top-five category position.")
    ui.add_metric(slide, "15", "drivers screened", 1.55, 1.50, 3.00, accent=ui.TEAL)
    ui.add_metric(slide, "4", "with genetic evidence", 5.15, 1.50, 3.00, accent=ui.BLUE)
    ui.add_metric(slide, str(facts["unscreened_count"]), "drivers awaiting screening", 8.75, 1.50, 3.00, accent=ui.GOLD)
    ui.add_rect(slide, 0.72, 3.10, 5.90, 3.05, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "The 15 screened drivers", 1.02, 3.42, 5.30, accent=ui.TEAL)
    ui.add_text(
        slide,
        "APOE, RPS15, SELENOW, ANKRD11,\nRPL11, RPLP1, RPL15, RPS13, RPL38,\nLAMTOR5, DYNLT1, ATP6V1F, FTL, LAPTM4A, NCOA1",
        1.04, 4.00, 5.30, 1.30, size=13.0, color=ui.NAVY, bold=True,
    )
    ui.add_text(slide, "Each appears in a top-five category display of the driver list.", 1.04, 5.55, 5.30, 0.30, size=10.5, color=ui.GRAY)
    ui.add_rect(slide, 6.90, 3.10, 5.75, 3.05, color=ui.PALE_SKY, outline=ui.SKY)
    ui.add_panel_title(slide, "What screening covered per gene", 7.20, 3.42, 5.15, accent=ui.BLUE)
    ui.add_bullets(slide, [
        "Nearby DNA association with clinical AD.",
        "Brain gene-activity links (RNA amount and splicing).",
        "Spinal-fluid biomarker association and gene-based tests.",
    ], 7.22, 4.00, 5.20, size=11.5, line_h=0.62, accent=ui.BLUE)
    ui.add_notes(
        slide,
        goal="State the screening coverage plainly before showing any evidence.",
        walkthrough="Fifteen driver genes have public screening results; four of them show genetic evidence, and four hundred eighteen drivers await screening. Each screened gene was checked for a nearby clinical-AD association, brain gene-activity links, and spinal-fluid biomarker evidence.",
        boundary="Coverage reflects which genes were screened, not any ranking by genetic promise.",
        transition="What the evidence says about the fifteen.",
    )
    titles.append(title)

    # 6 — evidence map
    title = "Genetic evidence for the 15 screened drivers, at a glance"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Grades from the public-data screen; driver contexts from the KDA results.")
    ui.add_table(
        slide,
        ["Gene(s)", "Driver contexts (KDA)", "Genetic evidence"],
        [
            ["APOE", "Astrocytes in 3 sex/APOE groups", "Strong: direct fine-mapped AD coding variant + all 3 CSF biomarkers"],
            ["RPS15", "12 categories — most recurrent driver", "Promising, unresolved: AD region P = 4.1e-30 + brain activity links"],
            ["SELENOW", "8 categories", "Weak: appears in a public TWAS gene list, no model detail"],
            ["ANKRD11", "2 OPC categories", "Region only: AD region P = 1.3e-11, no gene-level link passed"],
            ["11 other drivers", "1–10 categories each", "No qualifying AD or CSF signal in the tested design"],
        ],
        0.80, 1.55, [1.95, 3.75, 6.10], row_h=0.62, header_h=0.50, font_size=10.6,
    )
    ui.add_text(
        slide,
        "Evidence strength does not track network recurrence: RPL11 (10 categories) has no signal, while APOE (3 categories) is the strongest.",
        0.88, 5.55, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Give the complete evidence-by-driver map in one view before the two spotlights.",
        walkthrough="APOE is the only strong result. RPS15 combines a very strong regional AD association with positive brain activity links but no completed same-variant test. SELENOW has list-membership only; ANKRD11 has a significant region without a gene-level link. The other eleven drivers had no qualifying signal.",
        boundary="A significant nearby region does not assign the association to the candidate gene.",
        transition="Spotlight one: RPS15.",
    )
    titles.append(title)

    # 7 — RPS15 spotlight
    rps15_q = per_gene["RPS15"]["best_q"]
    title = "RPS15: network centrality and genetic proximity converge"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "The most recurrent driver is also the most interesting unresolved genetic candidate.")
    ui.add_rect(slide, 0.66, 1.52, 5.90, 4.55, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_panel_title(slide, "In the KDA networks", 0.96, 1.84, 5.30, accent=ui.BLUE)
    ui.add_bullets(slide, [
        "Most recurrent driver: 12 of 32 categories.",
        f"Best category score q = {q_text(rps15_q)} (top rank).",
        "Recurs across both sexes and five cell types.",
        "Also returns in the independent SEA-AD cohort.",
    ], 0.98, 2.45, 5.35, size=12.0, line_h=0.72, accent=ui.BLUE)
    ui.add_rect(slide, 6.85, 1.52, 5.82, 4.55, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "In human genetics", 7.15, 1.84, 5.22, accent=ui.TEAL)
    ui.add_bullets(slide, [
        f"AD region strongly associated: P = {min_p['RPS15']} ({facts['lead']['RPS15']}).",
        "Brain-tissue RNA-amount link passed its gate.",
        "Three further brain activity tracks are signal-positive.",
        "Same-variant test unresolved: complete models unavailable.",
    ], 7.17, 2.45, 5.30, size=12.0, line_h=0.72, accent=ui.TEAL)
    ui.add_rect(slide, 0.66, 6.25, 12.0, 0.62, color=ui.PALE_RED, outline=ui.VERMILION)
    ui.add_text(
        slide,
        "Caution: the AD signal is near RPS15 but has not been assigned to it; cell-type-exact activity data were negative or missing.",
        0.90, 6.42, 11.55, 0.30, size=11.0, color=ui.VERMILION_TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Make RPS15 the flagship follow-up target of the deck.",
        walkthrough="RPS15 is the single most recurrent driver and replicates in SEA-AD. Its candidate region carries a strong AD association, and bulk-brain activity links are positive, but no complete fitted model plus matched reference existed to run the same-variant test, and exact OPC or inhibitory-neuron data were negative or unavailable.",
        boundary="Positive brain tracks partially reuse the discovery cohort, so they support mechanism rather than independent validation.",
        transition="Spotlight two: APOE.",
    )
    titles.append(title)

    # 8 — APOE spotlight
    title = "APOE: the strongest genetic result is an astrocyte driver"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Strong gene-level genetic support; the exact astrocyte mechanism is still open.")
    ui.add_rect(slide, 0.66, 1.52, 5.90, 4.55, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "Genetic evidence", 0.96, 1.84, 5.30, accent=ui.TEAL)
    ui.add_bullets(slide, [
        "AD coding variant mapped directly to APOE (P ≈ 1.9e-155).",
        "Region significant for all three spinal-fluid biomarkers.",
        "Gene-based biomarker tests also pass for all three.",
        "Only screened gene passing every first-stage gate.",
    ], 0.98, 2.45, 5.35, size=12.0, line_h=0.72, accent=ui.TEAL)
    ui.add_rect(slide, 6.85, 1.52, 5.82, 4.55, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_panel_title(slide, "In the KDA networks", 7.15, 1.84, 5.22, accent=ui.BLUE)
    ui.add_bullets(slide, [
        "An astrocyte driver in three sex/APOE groups"
        f" ({', '.join(facts['apoe_groups'])}).",
        "Best within-category rank: 3.",
        "Modest recurrence compared with RPS15 (3 vs 12 categories).",
        "Does not return in the SEA-AD validation cohort.",
    ], 7.17, 2.45, 5.30, size=12.0, line_h=0.72, accent=ui.BLUE)
    ui.add_rect(slide, 0.66, 6.25, 12.0, 0.62, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_text(
        slide,
        "Genetics confirms APOE as an AD gene; whether the risk acts through APOE activity in astrocytes remains untested.",
        0.90, 6.42, 11.55, 0.30, size=11.0, color=ui.GRAY, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Keep APOE's strong result in view while separating gene-level support from mechanism-level proof.",
        walkthrough="APOE has the direct fine-mapped coding variant and passes both gates for CSF amyloid-beta 42, total tau, and p-tau181. In the networks it is an astrocyte driver in three sex/APOE groups, though it does not return in SEA-AD.",
        boundary="A single-signal sensitivity analysis suggested a shared CSF protein signal, but it is explicitly sensitivity-only; the primary multi-signal test was not met.",
        transition="How to read the eleven negatives.",
    )
    titles.append(title)

    # 9 — reading negatives
    title = "Eleven drivers had no qualifying signal — read this carefully"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "\u201CNo genetic support found\u201D has three different origins with different meanings.")
    cards = [
        ("Truly signal-negative", "The complete tested region had no qualifying AD or biomarker association under the frozen thresholds.", ui.PALE_SKY, ui.BLUE),
        ("Could not be evaluated", "Required measurements, models, or reference data were missing — stopped routes, not negative results.", ui.PALE_GOLD, ui.GOLD),
        ("Outside the design", "Rare variants, sex/APOE interactions, and protein-level effects were not tested.", ui.PALE_RED, ui.VERMILION),
    ]
    for index, (heading, body, bg, accent) in enumerate(cards):
        x = 0.66 + index * 4.10
        ui.add_rect(slide, x, 1.55, 3.83, 2.85, color=bg, outline=accent)
        ui.add_text(slide, heading, x + 0.26, 1.85, 3.30, 0.40, size=15.5, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
        ui.add_text(slide, body, x + 0.26, 2.45, 3.32, 1.80, size=11.0, color=ui.DARK)
    ui.add_rect(slide, 0.66, 4.75, 12.0, 1.35, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_panel_title(slide, "Example: replication without genetic support", 0.96, 4.98, 8.0, accent=ui.PURPLE)
    ui.add_text(
        slide,
        "DYNLT1 is a driver in four categories and also returns in SEA-AD — yet its region had no qualifying AD signal. "
        "Network recurrence and inherited risk are different questions; a gene can be central to disease networks without a common risk variant nearby.",
        0.98, 5.42, 11.40, 0.60, size=11.5, color=ui.DARK,
    )
    ui.add_notes(
        slide,
        goal="Prevent the negatives from being read as refutation of the network drivers.",
        walkthrough="The three origins: genuinely signal-negative regions, routes the data could not evaluate, and mechanisms the common-variant design does not cover. DYNLT1 illustrates the distinction: cross-cohort network replication with no nearby common AD variant.",
        boundary="Lower confidence in a simple inherited cis-risk mechanism is warranted; refutation of a functional role is not.",
        transition="The biggest limitation of the current evidence: coverage of the new list.",
    )
    titles.append(title)

    # 10 — the gap
    title = f"{facts['unscreened_count']} of {facts['ros_gene_count']} drivers have never been screened"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Screening coverage is now the main gap in the genetic evidence.")
    ui.add_rect(slide, 0.72, 1.55, 5.90, 4.50, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_panel_title(slide, "First-priority screening targets", 1.02, 1.87, 5.30, accent=ui.BLUE)
    ui.add_text(
        slide,
        "Drivers that replicate in the independent SEA-AD cohort but were never genetically screened:",
        1.04, 2.42, 5.30, 0.55, size=11.5, color=ui.DARK,
    )
    ui.add_text(
        slide,
        ", ".join(facts["priority"]),
        1.04, 3.10, 5.30, 0.42, size=17, color=ui.BLUE, bold=True, font=ui.FONT_HEAD,
    )
    ui.add_text(
        slide,
        "WDR82 and HGSNAT lead the cross-cohort convergence; TTC8 and BEX3 recur broadly in both cohorts.",
        1.04, 3.70, 5.30, 0.80, size=11.0, color=ui.GRAY,
    )
    ui.add_rect(slide, 6.90, 1.55, 5.75, 4.50, color=ui.PALE_GRAY, outline=ui.LIGHT)
    ui.add_panel_title(slide, "Also unscreened", 7.20, 1.87, 5.15, accent=ui.GRAY)
    ui.add_bullets(slide, [
        "Most of the recurrence leaders of the driver list (e.g., SELENOM, GABARAPL2, PGAM1, GAPDH).",
        "Most genes shown in the category top-five displays.",
        "The same public-data pipeline can screen them; the gap is coverage, not method.",
    ], 7.22, 2.45, 5.20, size=11.5, line_h=0.95, accent=ui.GRAY)
    ui.add_notes(
        slide,
        goal="Turn the coverage gap into a concrete, prioritized screening request.",
        walkthrough="Only fifteen drivers have been screened. The natural first extension is the cross-cohort replicated set — WDR82, HGSNAT, TTC8, BEX3 — followed by the recurrence leaders of the driver list. The same public-data pipeline applies; the gap is coverage rather than method.",
        boundary="Prioritization here uses network recurrence and replication, not any genetic pre-screening.",
        transition="Close with what we know and what to do next.",
    )
    titles.append(title)

    # 11 — take-home
    title = "What we know now — and the shortest path to stronger claims"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title)
    ui.add_rect(slide, 0.66, 1.30, 5.95, 5.05, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "Current state of evidence", 0.97, 1.62, 5.35, accent=ui.TEAL)
    ui.add_bullets(slide, [
        "APOE: strong genetic support; an astrocyte driver in three groups.",
        "RPS15: most recurrent driver + strong nearby AD region; the single best follow-up target.",
        "SELENOW and ANKRD11: partial genetic signals worth tracking.",
        "No same-variant test has been completed for any gene.",
    ], 0.99, 2.20, 5.40, size=11.8, line_h=0.85, accent=ui.TEAL)
    ui.add_rect(slide, 6.90, 1.30, 5.78, 5.05, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_panel_title(slide, "Next steps, in order", 7.21, 1.62, 5.18, accent=ui.BLUE)
    ui.add_bullets(slide, [
        "Obtain complete gene-activity models and matched references to finish the RPS15 and APOE same-variant tests.",
        "Extend the screen to the unscreened cross-cohort drivers (WDR82, HGSNAT, TTC8, BEX3).",
        "Then cover the recurrence leaders and category top-five genes.",
        "Add protein-level, rare-variant, and sex/APOE-interaction designs.",
    ], 7.23, 2.20, 5.25, size=11.8, line_h=0.95, accent=ui.BLUE)
    ui.add_text(
        slide,
        "Bottom line: genetics already supports a handful of drivers — led by APOE and RPS15 — but most of the driver list is still unscreened.",
        0.88, 6.60, 11.60, 0.30, size=12.0, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Leave a precise summary and an ordered request list.",
        walkthrough="The left panel is the current evidence state; the right panel orders the next steps from cheapest-decisive to most ambitious.",
        boundary="All conclusions remain screening-level until a valid same-variant analysis is completed.",
        transition="End of deck.",
    )
    titles.append(title)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return titles


def slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def main() -> int:
    args = parse_args()
    output_path = args.output.resolve()
    facts = load_facts()
    titles = build_deck(output_path, facts)

    reloaded = Presentation(str(output_path))
    checks: list[dict[str, Any]] = [
        {
            "check_id": "slide_count",
            "observed": len(reloaded.slides),
            "expected": len(titles),
            "passed": len(reloaded.slides) == len(titles),
        }
    ]
    for index, title in enumerate(titles, start=1):
        present = title in slide_text(reloaded.slides[index - 1])
        checks.append(
            {
                "check_id": f"slide{index}_title",
                "observed": present,
                "expected": True,
                "passed": present,
            }
        )
    notes_ok = all(
        slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip()
        for slide in reloaded.slides
    )
    checks.append(
        {"check_id": "all_slides_have_notes", "observed": notes_ok, "expected": True, "passed": notes_ok}
    )
    old_deck = ROOT / "docs" / "presentations" / "human_genetic_support_for_key_drivers_08252026.pptx"
    checks.append(
        {
            "check_id": "old_deck_untouched_exists",
            "observed": old_deck.is_file(),
            "expected": True,
            "passed": old_deck.is_file(),
        }
    )

    audit = args.audit.resolve()
    audit.parent.mkdir(parents=True, exist_ok=True)
    with audit.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["check_id", "observed", "expected", "passed"],
            delimiter="\t",
            lineterminator="\n",
        )
        writer.writeheader()
        writer.writerows(checks)
    failed = [row["check_id"] for row in checks if not row["passed"]]
    if failed:
        raise RuntimeError("Deck build failed checks: " + ", ".join(failed))
    print(f"built={output_path}")
    print(f"slides={len(titles)}")
    print(f"sha256={sha256_file(output_path)}")
    print(f"audit={audit}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
