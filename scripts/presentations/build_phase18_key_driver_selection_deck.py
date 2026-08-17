#!/usr/bin/env python3
"""Build a presentation-ready key-driver selection deck.

The deck uses the requested selection and human-genetic-support figures.
The deprecated ``filter_attrition`` directory is deliberately not referenced.
"""

from __future__ import annotations

import argparse
import csv
import tempfile
import zipfile
from collections import Counter, defaultdict
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
FIG_ROOT = REPO / "results/figures/analysis/phase_18_key_driver_selection"
GENETIC_SUPPORT_FIG_ROOT = REPO / "results/figures/analysis/phase_19_genetic_support"
DEFAULT_OUT = REPO / "docs/presentations/key_driver_selection_analysis 08192026.pptx"

FIG = {
    "selection": FIG_ROOT / "key_driver_selection_process/phase18_key_driver_selection_process.png",
    "circular_mt": FIG_ROOT / "two_case_circular/phase18_mt_driver_circular.png",
    "circular_non_mt": FIG_ROOT / "two_case_circular/phase18_non_mt_driver_circular.png",
    "atlas_mt": FIG_ROOT / "evidence_atlas_mt/phase18_evidence_atlas_mt.png",
    "atlas_non_mt": FIG_ROOT / "evidence_atlas_non_mt/phase18_evidence_atlas_non_mt.png",
    "sex_mt": FIG_ROOT / "sex_apoe_mt/phase18_sex_apoe_mt.png",
    "sex_non_mt": FIG_ROOT / "sex_apoe_non_mt/phase18_sex_apoe_non_mt.png",
    "rpl11_astro_pathway": FIG_ROOT / "RPL11/astrocyte/phase18_rpl11_astrocyte_consensus_network_pathways.png",
    "rpl11_astro_string": FIG_ROOT / "RPL11/astrocyte/stringdb_full_medium_conf.png",
    "rpl11_exc_pathway": FIG_ROOT / "RPL11/excitatory/phase18_rpl11_excitatory_consensus_network_pathways.png",
    "rpl11_exc_string": FIG_ROOT / "RPL11/excitatory/stringdb_full_medium_conf.png",
    "apoe_astro_pathway": FIG_ROOT / "APOE/astrocytes/phase18_apoe_astrocyte_consensus_network_pathways.png",
    "apoe_astro_string": FIG_ROOT / "APOE/astrocytes/stringdb_full_medium_conf.png",
    "genetic_support": GENETIC_SUPPORT_FIG_ROOT / "genetic_support_slide_summary.png",
}

DATA = {
    "circular": FIG_ROOT / "two_case_circular/phase18_two_case_circular_plot_data.tsv",
    "atlas_mt": FIG_ROOT / "evidence_atlas_mt/phase18_evidence_atlas_mt_gene_summary.tsv",
    "atlas_non_mt": FIG_ROOT / "evidence_atlas_non_mt/phase18_evidence_atlas_non_mt_gene_summary.tsv",
    "sex_mt": FIG_ROOT / "sex_apoe_mt/phase18_sex_apoe_mt_plot_data.tsv",
    "sex_non_mt": FIG_ROOT / "sex_apoe_non_mt/phase18_sex_apoe_non_mt_plot_data.tsv",
}

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 35, 61)
NAVY_2 = RGBColor(30, 59, 91)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
GOLD = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(226, 233, 240)
MID = RGBColor(116, 130, 145)
DARK = RGBColor(31, 39, 48)
GRAY = RGBColor(82, 94, 108)
PALE_BLUE = RGBColor(225, 238, 247)
PALE_GREEN = RGBColor(225, 243, 237)
PALE_GOLD = RGBColor(252, 240, 211)

# Arial is installed on the project Mac and is broadly portable across
# PowerPoint environments.  Using it avoids silent serif substitutions in
# Quick Look/Keynote when Aptos is not installed system-wide.
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def read_tsv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh, delimiter="\t"))


def fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def line(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             color: RGBColor = WHITE, outline: RGBColor | None = LIGHT,
             radius: bool = True, transparency: int = 0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color, transparency)
    if outline is None:
        shape.line.fill.background()
    else:
        line(shape, outline)
    return shape


def add_circle(slide, x: float, y: float, d: float, color: RGBColor,
               *, outline: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        line(shape, outline)
    return shape


def set_run(run, *, size: float, color: RGBColor, bold: bool = False,
            font: str = FONT_BODY, italic: bool = False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *,
             size: float = 16, color: RGBColor = DARK, bold: bool = False,
             font: str = FONT_BODY, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.03,
             italic: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, font=font, italic=italic)
    return box


def add_rich_text(slide, spans: list[tuple[str, dict]], x: float, y: float,
                  w: float, h: float, *, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    for text, style in spans:
        run = p.add_run()
        run.text = text
        set_run(run, size=style.get("size", 16), color=style.get("color", DARK),
                bold=style.get("bold", False), italic=style.get("italic", False),
                font=style.get("font", FONT_BODY))
    return box


def add_bullets(slide, items: list[str], x: float, y: float, w: float, *,
                size: float = 14, color: RGBColor = DARK,
                accent: RGBColor = TEAL, line_h: float = 0.56) -> None:
    cy = y
    for item in items:
        add_circle(slide, x, cy + 0.14, 0.09, accent)
        add_text(slide, item, x + 0.19, cy, w - 0.19, line_h,
                 size=size, color=color)
        cy += line_h


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, *,
               accent: RGBColor = TEAL, bg: RGBColor = WHITE) -> None:
    add_rect(slide, x, y, w, 1.02, color=bg, outline=LIGHT)
    value_size = 24 if len(value) <= 7 else 16.5
    add_text(slide, value, x + 0.16, y + 0.10, w - 0.32, 0.48,
             size=value_size, color=accent, bold=True, font=FONT_HEAD,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, label, x + 0.16, y + 0.61, w - 0.32, 0.24,
             size=10.2, color=GRAY, bold=True)


def set_alt_text(shape, title: str, description: str) -> None:
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        c_nv_pr.set("name", title)
        c_nv_pr.set("descr", description)
    except Exception:
        pass


def trim_white(path: Path, target_dir: Path) -> Path:
    """Trim only empty white margins; preserve the rendered scientific figure."""
    img = Image.open(path).convert("RGB")
    bg = Image.new("RGB", img.size, (255, 255, 255))
    diff = ImageChops.difference(img, bg).convert("L")
    # Ignore faint antialiasing/noise in the outer whitespace.
    diff = diff.point(lambda value: 255 if value > 10 else 0)
    bbox = diff.getbbox()
    if bbox:
        pad = max(10, int(min(img.size) * 0.006))
        left, top, right, bottom = bbox
        img = img.crop((max(0, left - pad), max(0, top - pad),
                        min(img.width, right + pad), min(img.height, bottom + pad)))
    target = target_dir / path.name
    img.save(target, optimize=True)
    return target


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float,
                        *, alt: str):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    pic = slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))
    set_alt_text(pic, alt, alt)
    return pic


def add_header(slide, kicker: str, title: str, slide_no: int, *,
               accent: RGBColor = TEAL, subtitle: str | None = None) -> None:
    add_text(slide, kicker.upper(), 0.55, 0.22, 4.8, 0.24,
             size=9.5, color=accent, bold=True)
    add_text(slide, title, 0.55, 0.47, 11.95, 0.50,
             size=25, color=NAVY, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, 0.57, 0.99, 11.65, 0.30,
                 size=11.5, color=GRAY)
    add_text(slide, f"{slide_no:02d}", 12.38, 0.24, 0.42, 0.22,
             size=9, color=MID, bold=True, align=PP_ALIGN.RIGHT)


def add_source(slide, source: str) -> None:
    add_text(slide, source, 0.55, 7.25, 12.0, 0.14, size=6.9, color=MID)


def add_takeaway_ribbon(slide, text: str, *, accent: RGBColor = TEAL,
                        y: float = 6.78) -> None:
    add_rect(slide, 0.55, y, 12.23, 0.48, color=NAVY, outline=None,
             radius=False, transparency=2)
    add_rect(slide, 0.55, y, 0.10, 0.48, color=accent, outline=None, radius=False)
    add_text(slide, text, 0.80, y + 0.09, 11.70, 0.27,
             size=11.3, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_panel_title(slide, text: str, x: float, y: float, w: float,
                    *, accent: RGBColor = TEAL) -> None:
    add_rect(slide, x, y + 0.03, 0.07, 0.34, color=accent, outline=None, radius=False)
    add_text(slide, text, x + 0.18, y, w - 0.18, 0.38,
             size=16.5, color=NAVY, bold=True, font=FONT_HEAD)


def add_connector(slide, x1: float, y1: float, x2: float, y2: float,
                  color: RGBColor = SKY) -> None:
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line(shape, color, 1.6)
    shape.line.end_arrowhead = True


def new_slide(prs: Presentation, *, bg: RGBColor = OFF_WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_section_divider(prs: Presentation, section_no: int, title: str,
                        subtitle: str, topics: list[str], page_no: int,
                        *, accent: RGBColor) -> None:
    """Add a concise dark section divider using the deck's visual language."""
    slide = new_slide(prs, bg=NAVY)
    add_rect(slide, 0, 0, 13.333, 7.5, color=NAVY, outline=None, radius=False)
    add_text(slide, f"SECTION {section_no:02d}", 0.76, 0.63, 3.30, 0.28,
             size=11, color=accent, bold=True)
    add_rect(slide, 0.76, 1.28, 0.11, 2.26, color=accent, outline=None, radius=False)
    add_text(slide, title, 1.16, 1.36, 8.35, 1.50,
             size=32, color=WHITE, bold=True, font=FONT_HEAD,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, subtitle, 1.18, 3.18, 7.80, 0.92,
             size=15.3, color=RGBColor(204, 219, 234))

    add_text(slide, f"{section_no:02d}", 9.55, 0.88, 2.52, 1.92,
             size=104, color=NAVY_2, bold=True, font=FONT_HEAD,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_circle(slide, 11.63, 1.06, 0.56, accent)
    add_circle(slide, 10.98, 2.93, 0.28, SKY)
    add_connector(slide, 10.85, 2.20, 11.91, 1.34, RGBColor(105, 137, 169))
    add_connector(slide, 10.84, 2.23, 11.12, 3.05, RGBColor(105, 137, 169))

    add_text(slide, "IN THIS SECTION", 0.78, 4.67, 2.10, 0.25,
             size=9.8, color=SKY, bold=True)
    topic_x = [0.78, 4.20, 7.62]
    topic_w = 3.10
    for index, topic in enumerate(topics):
        x = topic_x[index]
        add_rect(slide, x, 5.12, topic_w, 0.76,
                 color=NAVY_2, outline=accent)
        add_text(slide, topic, x + 0.18, 5.30, topic_w - 0.36, 0.32,
                 size=11.4, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_text(slide, f"{page_no:02d}", 12.38, 0.34, 0.42, 0.22,
             size=9, color=RGBColor(145, 171, 197), bold=True,
             align=PP_ALIGN.RIGHT)


def summarize() -> dict:
    circular = read_tsv(DATA["circular"])
    displayed = [r for r in circular if r["slot_status"] == "ranked_candidate"]
    by_class = defaultdict(list)
    for row in displayed:
        by_class[row["case_id"]].append(row)

    mt_gene_counts = Counter(r["current_symbol"] for r in by_class["mt_driver"])
    non_mt_gene_counts = Counter(r["current_symbol"] for r in by_class["non_mt_driver"])
    mt = read_tsv(DATA["atlas_mt"])
    non_mt = read_tsv(DATA["atlas_non_mt"])

    def sex_totals(path: Path) -> list[tuple[str, str, int, int]]:
        rows = read_tsv(path)
        groups: dict[tuple[str, str], dict] = defaultdict(lambda: {"runs": 0, "contexts": set()})
        for row in rows:
            key = (row["signature_direction"], row["signature_group"])
            support = int(row["conservative_support_count"])
            groups[key]["runs"] += support
            if support > 0:
                groups[key]["contexts"].add((row["current_symbol"], row["broad_network"]))
        ranked = sorted(((direction, group, values["runs"], len(values["contexts"]))
                         for (direction, group), values in groups.items()),
                        key=lambda item: (item[2], item[3]), reverse=True)
        return ranked

    return {
        "mt_positions": len(by_class["mt_driver"]),
        "non_mt_positions": len(by_class["non_mt_driver"]),
        "mt_genes": len(mt_gene_counts),
        "non_mt_genes": len(non_mt_gene_counts),
        "mt_gene_counts": mt_gene_counts,
        "non_mt_gene_counts": non_mt_gene_counts,
        "mt_summary": {r["current_symbol"]: r for r in mt},
        "non_mt_summary": {r["current_symbol"]: r for r in non_mt},
        "sex_mt": sex_totals(DATA["sex_mt"]),
        "sex_non_mt": sex_totals(DATA["sex_non_mt"]),
    }


def build_deck(output_path: Path = DEFAULT_OUT) -> Path:
    for path in [*FIG.values(), *DATA.values()]:
        if not path.exists():
            raise FileNotFoundError(path)
    if any("filter_attrition" in str(path) for path in [*FIG.values(), *DATA.values()]):
        raise RuntimeError("Deprecated filter_attrition content must not enter this deck")

    summary = summarize()
    mt_co2 = summary["mt_summary"]["MT-CO2"]
    rpl11 = summary["non_mt_summary"]["RPL11"]
    rps15 = summary["non_mt_summary"]["RPS15"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Phase 18 Key-Driver Selection"
    prs.core_properties.subject = "Selection, cross-network evidence, and descriptive sex/APOE support"
    prs.core_properties.author = "Alzheimer project analysis team"
    prs.core_properties.keywords = "Alzheimer, mitochondria, KDA, key drivers, APOE, sex"

    with tempfile.TemporaryDirectory(prefix="phase18_deck_assets_") as td:
        assets = Path(td)
        trimmed = {key: trim_white(path, assets) for key, path in FIG.items()}

        # 1 — title and narrative frame
        slide = new_slide(prs, bg=NAVY)
        add_rect(slide, 0, 0, 13.333, 7.5, color=NAVY, outline=None, radius=False)
        for x, y, d, c in [
            (9.10, 0.78, 1.18, TEAL), (10.63, 1.18, 0.75, GOLD),
            (11.78, 0.62, 0.42, SKY), (9.62, 2.82, 0.58, PURPLE),
            (11.12, 3.02, 1.02, BLUE), (10.18, 4.76, 0.42, VERMILION),
            (11.83, 5.12, 0.70, TEAL),
        ]:
            add_circle(slide, x, y, d, c)
        for x1, y1, x2, y2 in [
            (9.70, 1.35, 10.98, 1.56), (11.02, 1.72, 11.58, 3.24),
            (9.95, 3.10, 11.35, 3.50), (10.26, 3.40, 10.39, 4.94),
            (11.63, 4.02, 12.17, 5.42), (10.57, 4.98, 11.91, 5.48),
        ]:
            add_connector(slide, x1, y1, x2, y2, RGBColor(137, 167, 195))
        add_text(slide, "PHASE 18 • KEY-DRIVER PRIORITIZATION", 0.72, 0.72, 7.4, 0.28,
                 size=11, color=SKY, bold=True)
        add_text(slide, "From run-level KDA\nevidence to a focused\nkey-driver atlas",
                 0.72, 1.25, 7.85, 2.35, size=35, color=WHITE,
                 bold=True, font=FONT_HEAD)
        add_text(slide,
                 "Selection logic • cross-network recurrence • evidence breadth • descriptive sex/APOE patterns",
                 0.75, 4.04, 7.28, 0.72, size=17, color=RGBColor(210, 222, 235))
        add_rect(slide, 0.75, 5.14, 6.95, 0.055, color=TEAL, outline=None, radius=False)
        add_text(slide,
                 "MT drivers and non-MT drivers are ranked separately within seven broad cell-network models.",
                 0.75, 5.45, 7.20, 0.70, size=12.8, color=RGBColor(190, 207, 225))
        add_text(slide, "Prepared from validated Phase 18 outputs • 16 August 2026",
                 0.75, 6.90, 7.10, 0.23, size=9, color=RGBColor(145, 171, 197))

        # 2 — agenda
        slide = new_slide(prs)
        add_header(
            slide,
            "Presentation roadmap",
            "Agenda",
            2,
            accent=BLUE,
            subtitle="Three sections move from candidate selection to independent genetic support and then to network and proteomics validation.",
        )
        agenda_rows = [
            (
                "01", "Key-driver selection & evidence landscape", "Slides 03–11",
                "Selection logic, MT/non-MT recurrence, evidence breadth, synthesis, and descriptive sex/APOE patterns.",
                TEAL, PALE_GREEN,
            ),
            (
                "02", "Human genetic support", "Slides 12–15",
                "Public summary data, direct-mapping rules, context-aware grading, and the candidate-level result profile.",
                BLUE, PALE_BLUE,
            ),
            (
                "03", "Key-driver networks & proteomics validation", "Slides 16–25",
                "Place RPL11 and APOE on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
                GOLD, PALE_GOLD,
            ),
        ]
        for index, (number, title, slide_range, body, accent, bg) in enumerate(agenda_rows):
            y = 1.55 + index * 1.63
            add_rect(slide, 0.74, y, 11.86, 1.36, color=WHITE, outline=LIGHT)
            add_rect(slide, 0.74, y, 0.10, 1.36, color=accent, outline=None, radius=False)
            add_circle(slide, 1.10, y + 0.34, 0.68, accent)
            add_text(slide, number, 1.10, y + 0.53, 0.68, 0.23,
                     size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                     valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, title, 2.04, y + 0.22, 6.75, 0.36,
                     size=18.3, color=NAVY, bold=True, font=FONT_HEAD)
            add_text(slide, slide_range, 9.75, y + 0.23, 2.31, 0.27,
                     size=10.2, color=accent, bold=True, align=PP_ALIGN.RIGHT)
            add_text(slide, body, 2.05, y + 0.72, 9.90, 0.40,
                     size=11.5, color=GRAY)
        add_takeaway_ribbon(
            slide,
            "Narrative arc: prioritize candidates → seek independent human evidence → place key drivers on networks and validate at the protein level.",
            accent=BLUE,
            y=6.62,
        )

        # 3 — section divider: selection and evidence
        add_section_divider(
            prs,
            1,
            "Key-driver selection &\nevidence landscape",
            "How 25 candidates emerged across 47 displayed network contexts—and where support recurs.",
            ["Selection logic", "MT / non-MT evidence", "Sex / APOE patterns"],
            3,
            accent=TEAL,
        )

        # 4 — selection process
        slide = new_slide(prs)
        add_header(slide, "Selection process", "A strict intersection of coverage, support, and aggregate evidence", 4,
                   subtitle="The pipeline changes counting units as run-level tests are aggregated into network × driver-class candidates.")
        add_picture_contain(slide, trimmed["selection"], 0.35, 1.28, 8.55, 5.84,
                            alt="Phase 18 key-driver selection process from 95,557 run-level tests to 47 displayed positions")
        add_rect(slide, 9.08, 1.42, 3.72, 5.38, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What earns a place", 9.36, 1.72, 3.05, accent=TEAL)
        add_bullets(slide, [
            "All three gates are required: ≥80% usable-run coverage, ≥1 conservative-support run, and aggregate ACAT q ≤ 0.05.",
            "The intersection retains 78 of 10,433 represented candidate units (0.75%).",
            "Candidates are ranked independently in 14 lists: 7 broad networks × 2 driver classes.",
            "Ranks 1–5 are displayed with no backfilling when fewer candidates pass.",
        ], 9.36, 2.28, 3.06, size=12.3, line_h=0.76)
        add_metric(slide, "47", "displayed gene × network positions", 9.36, 5.50, 1.45, accent=BLUE)
        add_metric(slide, "25", "unique displayed genes", 10.97, 5.50, 1.45, accent=GOLD)
        add_text(slide, "Prioritization—not proof of causal regulation.", 9.39, 6.63, 3.02, 0.25,
                 size=10.5, color=VERMILION, bold=True)
        add_source(slide, "Source: phase18_key_driver_selection_process.png and accompanying validated caption/methods")

        # 5 — MT circular
        slide = new_slide(prs, bg=WHITE)
        add_picture_contain(slide, trimmed["circular_mt"], 0.05, 0.02, 13.23, 6.92,
                            alt="Circular view of selected mitochondrial key drivers across seven broad cell networks")
        add_takeaway_ribbon(
            slide,
            f"MT view: {summary['mt_genes']} genes occupy {summary['mt_positions']} displayed positions; "
            f"MT-CO2 recurs in {summary['mt_gene_counts']['MT-CO2']} of 7 networks, while gray slots preserve lists with fewer than five passers.",
            accent=TEAL,
        )
        add_source(slide, "Source: phase18_mt_driver_circular.png • center curves show recurrence, not network edges")

        # 6 — non-MT circular
        slide = new_slide(prs, bg=WHITE)
        add_picture_contain(slide, trimmed["circular_non_mt"], 0.05, 0.02, 13.23, 6.92,
                            alt="Circular view of selected non-mitochondrial key drivers across seven broad cell networks")
        add_takeaway_ribbon(
            slide,
            f"Non-MT view: {summary['non_mt_genes']} genes occupy {summary['non_mt_positions']} displayed positions; "
            f"RPL11 recurs in {summary['non_mt_gene_counts']['RPL11']} networks, and vasculature has no passing non-MT candidate.",
            accent=GOLD,
        )
        add_source(slide, "Source: phase18_non_mt_driver_circular.png • non-MT means outside core MitoCarta, not absence of mitochondrial function")

        # 7 — MT evidence atlas
        slide = new_slide(prs, bg=WHITE)
        add_picture_contain(slide, trimmed["atlas_mt"], 0.05, 0.06, 13.23, 6.94,
                            alt="Evidence atlas for 10 selected mitochondrial key-driver genes")
        add_takeaway_ribbon(
            slide,
            f"MT evidence is broadest for MT-CO2: {mt_co2['conservative_supporting_run_count']}/{mt_co2['explicitly_tested_run_count']} conservative-support runs, "
            f"{mt_co2['supporting_fine_cell_type_count']} fine cell types, all 7 networks, all 6 primary groups, and both AD directions.",
            accent=BLUE,
        )
        add_source(slide, "Source: phase18_evidence_atlas_mt.png • stability is diagnostic and did not determine selection")

        # 8 — non-MT evidence atlas
        slide = new_slide(prs)
        add_header(slide, "Evidence atlas • non-MT", "Ribosomal candidates carry the broadest non-MT evidence", 8,
                   accent=GOLD, subtitle="Network evidence, recurrence, breadth, and leave-one-fine-cell-type stability are complementary summaries.")
        add_picture_contain(slide, trimmed["atlas_non_mt"], 0.25, 1.28, 9.17, 5.86,
                            alt="Evidence atlas for 15 selected non-mitochondrial key-driver genes")
        add_rect(slide, 9.62, 1.45, 3.18, 5.39, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Presentation readout", 9.89, 1.72, 2.64, accent=GOLD)
        add_bullets(slide, [
            f"RPL11 passes in 4 networks with {rpl11['conservative_supporting_run_count']}/{rpl11['explicitly_tested_run_count']} conservative-support runs across {rpl11['supporting_fine_cell_type_count']} fine cell types.",
            f"RPS15 has the widest non-MT fine-cell breadth ({rps15['supporting_fine_cell_type_count']}) and passes in 3 networks.",
            "The dashed RPS15 excitatory-neuron tile passed selection but ranked below that network's five-gene display cap.",
            "RPL/RPS genes recur across networks; LAMTOR5 spans excitatory and inhibitory networks; APOE is astrocyte-localized.",
        ], 9.89, 2.28, 2.62, size=11.8, line_h=0.79, accent=GOLD)
        add_rect(slide, 9.89, 5.69, 2.62, 0.70, color=PALE_GOLD, outline=GOLD)
        add_text(slide, "The top-five rule is a display cap—not an evidence threshold.",
                 10.08, 5.84, 2.24, 0.38, size=10.3, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_source(slide, "Source: phase18_evidence_atlas_non_mt.png and gene summary TSV")

        # 10 — MT sex/APOE
        slide = new_slide(prs)
        add_header(slide, "Sex/APOE evidence • MT", "MT support is widespread but uneven across descriptive strata", 10,
                   accent=BLUE, subtitle="Each row is a gene × broad-network context; columns partition AD-up and AD-down queries by sex and APOE group.")
        add_picture_contain(slide, trimmed["sex_mt"], 0.18, 1.28, 7.05, 5.98,
                            alt="Sex- and APOE-stratified support for selected mitochondrial key drivers")
        add_rect(slide, 7.44, 1.44, 5.37, 5.55, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "How to read the matrix", 7.75, 1.72, 4.72, accent=BLUE)
        add_bullets(slide, [
            "Dot area = fraction of usable queries with conservative support; fill = capped −log10(stratum ACAT P).",
            "Small open circles mean usable tests with zero conservative support; dashes mean no eligible query.",
            "The right tracks report official network-level q, support/usable, coverage, evidence tier, and rank.",
        ], 7.75, 2.25, 4.67, size=12.2, line_h=0.68, accent=BLUE)
        add_panel_title(slide, "Descriptive pattern", 7.75, 4.32, 4.72, accent=TEAL)
        mt_top = summary["sex_mt"][:3]
        group_label = {"F_e2": "female ε2", "F_e33": "female ε3/ε3", "F_e4": "female ε4",
                       "M_e2": "male ε2", "M_e33": "male ε3/ε3", "M_e4": "male ε4"}
        direction_label = {"AD_up_mito": "AD-up", "AD_down_mito": "AD-down"}
        pattern_lines = [
            f"{direction_label[d]} {group_label[g]}: {runs} supporting runs across {contexts} contexts"
            for d, g, runs, contexts in mt_top
        ]
        add_bullets(slide, pattern_lines, 7.75, 4.82, 4.60, size=12.0,
                    line_h=0.50, accent=TEAL)
        add_text(slide, "These are not formal sex, APOE, or interaction tests.",
                 7.78, 6.48, 4.54, 0.25, size=11.0, color=VERMILION, bold=True)
        add_source(slide, "Source: phase18_sex_apoe_mt.png • supporting-run totals are descriptive across displayed contexts")

        # 11 — non-MT sex/APOE
        slide = new_slide(prs)
        add_header(slide, "Sex/APOE evidence • non-MT", "Non-MT support concentrates in a smaller set of strata", 11,
                   accent=GOLD, subtitle="The same visual grammar reveals broad recurrence for some drivers and sharply localized evidence for others.")
        add_picture_contain(slide, trimmed["sex_non_mt"], 0.18, 1.30, 8.62, 5.86,
                            alt="Sex- and APOE-stratified support for selected non-mitochondrial key drivers")
        add_rect(slide, 9.01, 1.46, 3.79, 5.39, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Leading patterns", 9.30, 1.73, 3.21, accent=GOLD)
        non_top = summary["sex_non_mt"][:4]
        non_pattern_lines = [
            f"{direction_label[d]} {group_label[g]}: {runs} supporting runs in {contexts} contexts"
            for d, g, runs, contexts in non_top
        ]
        add_bullets(slide, non_pattern_lines, 9.30, 2.25, 3.15, size=11.7,
                    line_h=0.62, accent=GOLD)
        add_panel_title(slide, "Candidate examples", 9.30, 4.75, 3.21, accent=TEAL)
        add_text(slide,
                 "RPL11 is recurrent across four networks; RPS15 is broad across neuronal/OPC contexts; APOE is localized to astrocytes. "
                 "Large dots with small denominators should not be mistaken for broad replication.",
                 9.31, 5.24, 3.09, 1.06, size=11.3, color=GRAY)
        add_text(slide, "Patterned evidence generates hypotheses; it does not establish specificity.",
                 9.31, 6.39, 3.08, 0.34, size=10.4, color=VERMILION, bold=True)
        add_source(slide, "Source: phase18_sex_apoe_non_mt.png • descriptive totals match the validated figure guide")

        # Candidate synthesis; moved to position 9 before saving.
        slide = new_slide(prs)
        synthesis_slide_id = prs.slides._sldIdLst[-1]
        add_header(slide, "Synthesis", "What Phase 18 establishes—and what it leaves open", 9,
                   subtitle="The broad atlas establishes which candidates emerged; the next slides localize support before the RPL11 case study.")
        cards = [
            ("Selection", "78 passing units become 47 displayed network positions across two separately ranked driver classes.", TEAL, PALE_GREEN),
            ("MT landscape", "Respiratory-chain genes dominate recurrence; MT-CO2 is the clearest cross-network sentinel.", BLUE, PALE_BLUE),
            ("Non-MT landscape", "Ribosomal candidates lead breadth and recurrence, with additional signaling, trafficking, and APOE-linked hypotheses.", GOLD, PALE_GOLD),
        ]
        for i, (title, body, accent, bg) in enumerate(cards):
            x = 0.72 + i * 4.18
            add_rect(slide, x, 1.72, 3.78, 2.45, color=bg, outline=accent)
            add_circle(slide, x + 0.28, 2.04, 0.34, accent)
            add_text(slide, title, x + 0.76, 1.99, 2.62, 0.39,
                     size=19, color=NAVY, bold=True, font=FONT_HEAD)
            add_text(slide, body, x + 0.28, 2.73, 3.17, 1.02,
                     size=13.0, color=GRAY)
        add_rect(slide, 0.92, 4.62, 11.48, 1.37, color=NAVY, outline=None)
        add_rich_text(slide, [
            ("Interpretive boundary: ", {"size": 17, "color": SKY, "bold": True}),
            ("the figures prioritize network-associated key-driver candidates and show where evidence recurs. ",
             {"size": 17, "color": WHITE}),
            ("They do not prove causality or formally test sex/APOE interactions.",
             {"size": 17, "color": WHITE, "bold": True}),
        ], 1.28, 4.93, 10.75, 0.74, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide,
                 "Next: localize candidate evidence across sex/APOE strata, then examine RPL11 in depth.",
                 1.28, 6.35, 10.76, 0.42, size=13.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_source(slide, "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–25 show RPL11 and APOE network/protein examples")

        # 12 — section divider: human genetic support
        add_section_divider(
            prs,
            2,
            "Human genetic support",
            "An independent check of whether public Alzheimer’s disease genetic signals map directly to the selected candidate genes.",
            ["Public summary data", "Direct-mapping rules", "Candidate-level results"],
            12,
            accent=BLUE,
        )

        # 13 — human genetic support: data
        slide = new_slide(prs)
        add_header(
            slide,
            "Human genetic support • data",
            "A compact public summary bundle tests all 25 candidates",
            13,
            accent=BLUE,
            subtitle="This is an independent evidence line across every displayed gene × broad-network context—not a rerun of the network model.",
        )

        add_rect(slide, 0.55, 1.45, 7.42, 4.94, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What entered the analysis", 0.86, 1.72, 6.82, accent=BLUE)

        source_rows = [
            (
                "01", "Selected key drivers",
                "25 unique genes evaluated in all 47 displayed broad cell-type/network contexts.",
                TEAL, PALE_GREEN,
            ),
            (
                "02", "Public human-genetics summaries",
                "Official FunGen-xQTL Alzheimer’s disease snapshot: GWAS/fine-mapping, xQTL, TWAS and gene-level result tables.",
                BLUE, PALE_BLUE,
            ),
            (
                "03", "Gene identity references",
                "GENCODE v44 (GRCh38) coordinates plus the HGNC 2026-06-05 symbol table keep locus and gene mappings explicit.",
                GOLD, PALE_GOLD,
            ),
        ]
        for index, (number, title, body, accent, bg) in enumerate(source_rows):
            y = 2.22 + index * 1.15
            add_rect(slide, 0.86, y, 6.78, 0.94, color=bg, outline=accent)
            add_circle(slide, 1.08, y + 0.22, 0.49, accent)
            add_text(slide, number, 1.08, y + 0.31, 0.49, 0.20,
                     size=9.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                     valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, title, 1.78, y + 0.15, 2.20, 0.26,
                     size=13.8, color=NAVY, bold=True, font=FONT_HEAD)
            add_text(slide, body, 4.00, y + 0.14, 3.35, 0.57,
                     size=10.8, color=GRAY, valign=MSO_ANCHOR.MIDDLE)

        add_rect(slide, 0.86, 5.82, 6.78, 0.35, color=OFF_WHITE, outline=None, radius=False)
        add_text(slide, "Summary-level only: no individual genotypes, phenotypes, or protected participant data.",
                 1.02, 5.91, 6.45, 0.18, size=9.7, color=VERMILION,
                 bold=True, align=PP_ALIGN.CENTER)

        add_rect(slide, 8.20, 1.45, 4.58, 4.94, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Scale and portability", 8.50, 1.72, 3.97, accent=TEAL)
        add_metric(slide, "25", "unique candidate genes", 8.50, 2.22, 1.74, accent=BLUE)
        add_metric(slide, "47", "gene × network contexts", 10.44, 2.22, 1.74, accent=TEAL)
        add_metric(slide, "6", "public source files", 8.50, 3.49, 1.74, accent=GOLD)
        add_metric(slide, "~8.8 MiB", "total downloaded bundle", 10.44, 3.49, 1.74, accent=PURPLE)
        add_rect(slide, 8.50, 4.83, 3.68, 1.16, color=NAVY, outline=None)
        add_text(slide, "LOCAL WORKSTATION READY", 8.76, 5.04, 3.17, 0.24,
                 size=10.2, color=SKY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Small, tabular inputs; no cluster-scale compute required.",
                 8.76, 5.38, 3.17, 0.39, size=11.5, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)

        add_takeaway_ribbon(
            slide,
            "Human genetics asks whether an independently measured AD-associated signal maps to the same candidate gene in a relevant context.",
            accent=BLUE,
            y=6.62,
        )
        add_source(slide, "Sources: FunGen-xQTL public snapshot (2026-07-29) • GENCODE v44, GRCh38 • HGNC 2026-06-05")

        # 14 — human genetic support: method
        slide = new_slide(prs)
        add_header(
            slide,
            "Human genetic support • method",
            "Direct gene mapping—not proximity—determines the grade",
            14,
            accent=TEAL,
            subtitle="The workflow screens broadly, then prevents a nearby association from being mistaken for support of a specific candidate gene.",
        )

        flow = [
            ("25 genes", "47 candidate contexts", BLUE, PALE_BLUE),
            ("19 nuclear genes", "screen ±1 Mb per locus", TEAL, PALE_GREEN),
            ("356 records", "candidate-locus variant records", GOLD, PALE_GOLD),
            ("Direct mapping", "target gene + evidence context", PURPLE, RGBColor(239, 232, 246)),
        ]
        flow_x = [0.55, 3.60, 6.65, 9.70]
        for index, (title, body, accent, bg) in enumerate(flow):
            x = flow_x[index]
            add_rect(slide, x, 1.53, 2.52, 1.13, color=bg, outline=accent)
            add_text(slide, title, x + 0.16, 1.72, 2.20, 0.31,
                     size=16.5, color=NAVY, bold=True, font=FONT_HEAD,
                     align=PP_ALIGN.CENTER)
            add_text(slide, body, x + 0.16, 2.15, 2.20, 0.24,
                     size=9.6, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
            if index < len(flow) - 1:
                add_connector(slide, x + 2.55, 2.10, flow_x[index + 1] - 0.08, 2.10, SKY)

        add_rect(slide, 0.55, 2.98, 4.00, 2.73, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Set aside", 0.84, 3.26, 3.40, accent=MID)
        add_text(slide, "352", 0.85, 3.79, 1.18, 0.55,
                 size=30, color=MID, bold=True, font=FONT_HEAD)
        add_text(slide, "regional proximity-only records", 1.86, 3.84, 2.32, 0.39,
                 size=13.3, color=NAVY, bold=True)
        add_text(slide,
                 "They help locate candidate regions but do not name the candidate as the molecular target, so they are not counted as gene support.",
                 0.86, 4.56, 3.32, 0.73, size=11.4, color=GRAY)

        add_rect(slide, 4.67, 2.98, 4.00, 2.73, color=PALE_GREEN, outline=TEAL)
        add_panel_title(slide, "Grade direct evidence", 4.96, 3.26, 3.40, accent=TEAL)
        add_text(slide, "4 + 1", 4.97, 3.79, 1.34, 0.55,
                 size=30, color=TEAL, bold=True, font=FONT_HEAD)
        add_text(slide, "direct xQTL records + TWAS-list row", 6.25, 3.84, 2.02, 0.50,
                 size=12.3, color=NAVY, bold=True)
        add_text(slide,
                 "Grade strong / moderate / weak from the source’s direct gene mapping, evidence type, and match to the broad cell context.",
                 4.98, 4.56, 3.31, 0.73, size=11.4, color=GRAY)

        add_rect(slide, 8.79, 2.98, 4.00, 2.73, color=PALE_GOLD, outline=GOLD)
        add_panel_title(slide, "Mitochondrial route", 9.08, 3.26, 3.40, accent=GOLD)
        add_text(slide, "6", 9.09, 3.79, 0.66, 0.55,
                 size=30, color=GOLD, bold=True, font=FONT_HEAD)
        add_text(slide, "mtDNA genes • 20 contexts", 9.83, 3.84, 2.56, 0.39,
                 size=13.0, color=NAVY, bold=True)
        add_text(slide,
                 "The available source lacks mtDNA-specific heteroplasmy, haplogroup, copy-number and NUMT controls; these contexts are not assessable.",
                 9.10, 4.56, 3.32, 0.73, size=11.4, color=GRAY)

        add_rect(slide, 0.78, 5.95, 12.00, 0.42, color=OFF_WHITE, outline=LIGHT, radius=False)
        add_text(slide,
                 "Colocalization boundary: the public source provides prefiltered summaries, not H0–H4 posterior probabilities; missing statistics were not reconstructed or relabeled.",
                 1.00, 6.05, 11.57, 0.22, size=10.3, color=VERMILION,
                 bold=True, align=PP_ALIGN.CENTER)
        add_takeaway_ribbon(
            slide,
            "Decision rule: direct gene mapping can support a candidate; proximity alone cannot. Missing source coverage remains unresolved—not negative.",
            accent=TEAL,
            y=6.67,
        )
        add_source(slide, "Method: registered candidate list • ±1 Mb nuclear-locus screen • direct-target review • context-aware evidence grading")

        # 15 — human genetic support: results
        slide = new_slide(prs)
        add_header(
            slide,
            "Human genetic support • results",
            "APOE stands out; COX7C and SELENOW remain limited",
            15,
            accent=BLUE,
            subtitle="Most candidates remain unresolved because direct mapping or appropriate source coverage was unavailable; the grades are not causal probabilities.",
        )
        add_picture_contain(
            slide,
            trimmed["genetic_support"],
            0.50,
            1.37,
            12.33,
            5.70,
            alt="Human genetic support summary across 47 candidate gene-by-network contexts, highlighting APOE, COX7C, SELENOW, unresolved nuclear candidates, and mitochondrial contexts not assessable in the source",
        )
        add_source(slide, "Source: validated human-genetic-support summary • FunGen-xQTL public snapshot • GENCODE v44 • HGNC 2026-06-05")

        # 16 — section divider: network graphs and proteomics validation
        add_section_divider(
            prs,
            3,
            "Key-driver network graphs &\nproteomics validation",
            "Place prioritized key drivers in their cell-type network context, then test whether predicted modules are supported at the protein level.",
            ["Network graph overlays", "Protein-level evidence", "Proteomics validation"],
            16,
            accent=GOLD,
        )

        # 17 — first-batch candidate panel
        slide = new_slide(prs)
        add_header(slide, "First-batch validation panel", "Why these seven candidates go first", 17,
                   subtitle="Every listed gene × broad-network pairing is a separate validation unit; support in one network does not validate another.")

        candidate_rows = [
            (
                "APOE", "Astrocytes", "CALIBRATOR",
                "Known-AD control; reproduce its frozen astrocyte mitochondrial target module.",
                BLUE, PALE_BLUE,
            ),
            (
                "COX7C", "Astrocytes • Inhibitory neurons", "GENETICS-INFORMED",
                "Weak direct bulk-sQTL mapping; test network targets, complex-IV assembly, and protein stoichiometry.",
                SKY, PALE_BLUE,
            ),
            (
                "SELENOW", "Excitatory neurons", "REINFORCED",
                "Stable KDA plus tau/redox evidence; the source TWAS flag remains limited.",
                TEAL, PALE_GREEN,
            ),
            (
                "LAMTOR5", "Excitatory neurons • Inhibitory neurons", "CROSS-NETWORK",
                "Two neuronal networks test lysosome–mTOR/mitochondria recurrence.",
                PURPLE, RGBColor(239, 232, 246),
            ),
            (
                "RPL11", "Astrocytes • Excitatory neurons\nMicroglia • Oligodendrocytes", "BIAS TEST",
                "Four selected networks; require expression-, degree-, and ribosomal-matched nulls.",
                GOLD, PALE_GOLD,
            ),
            (
                "FTL", "OPCs", "EXPLORATORY",
                "Sparse OPC signal; localize the iron/ferroptosis mechanism before promotion.",
                VERMILION, RGBColor(250, 232, 224),
            ),
            (
                "ANKRD11", "OPCs", "NOVELTY TEST",
                "Sparse OPC novelty test for chromatin-linked regulation.",
                NAVY_2, RGBColor(234, 239, 245),
            ),
        ]

        add_rect(slide, 0.55, 1.40, 12.23, 0.43, color=NAVY, outline=None, radius=False)
        add_text(slide, "GENE", 0.74, 1.51, 1.03, 0.18,
                 size=8.8, color=SKY, bold=True)
        add_text(slide, "BROAD NETWORK(S) TO VALIDATE", 1.92, 1.51, 3.10, 0.18,
                 size=8.8, color=SKY, bold=True)
        add_text(slide, "ROLE", 5.19, 1.51, 1.80, 0.18,
                 size=8.8, color=SKY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "WHY / PRIMARY TEST", 7.12, 1.51, 5.12, 0.18,
                 size=8.8, color=SKY, bold=True)

        for index, (gene, networks, role, body, accent, role_bg) in enumerate(candidate_rows):
            y = 1.88 + index * 0.63
            row_bg = WHITE if index % 2 == 0 else OFF_WHITE
            add_rect(slide, 0.55, y, 12.23, 0.56, color=row_bg, outline=LIGHT, radius=False)
            add_rect(slide, 0.55, y, 0.08, 0.56, color=accent, outline=None, radius=False)
            add_text(slide, gene, 0.74, y + 0.12, 1.05, 0.28,
                     size=14.0, color=NAVY, bold=True, font=FONT_HEAD,
                     valign=MSO_ANCHOR.MIDDLE)
            network_size = 8.8 if gene in {"LAMTOR5", "RPL11"} else 9.5
            add_text(slide, networks, 1.92, y + 0.08, 3.10, 0.40,
                     size=network_size, color=NAVY, bold=True,
                     valign=MSO_ANCHOR.MIDDLE)
            add_rect(slide, 5.16, y + 0.13, 1.84, 0.30, color=role_bg, outline=accent)
            role_size = 7.1 if role == "GENETICS-INFORMED" else 7.7
            add_text(slide, role, 5.21, y + 0.21, 1.74, 0.13,
                     size=role_size, color=accent, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, body, 7.12, y + 0.08, 5.38, 0.40,
                     size=9.5, color=GRAY, valign=MSO_ANCHOR.MIDDLE)

        add_takeaway_ribbon(
            slide,
            "Validate every listed gene × network unit separately; COX7C enters both astrocyte and inhibitory-neuron workstreams, but one bulk-sQTL result is not two replications.",
            accent=TEAL, y=6.55,
        )
        add_source(slide, "Sources: Phase 18 call_key_driver_returns.tsv • cross-validation guide, Section 2 • Tier 1 human-genetic-support report")

        # 18 — RPL11 rationale
        slide = new_slide(prs, bg=NAVY)
        add_rect(slide, 0, 0, 13.333, 7.5, color=NAVY, outline=None, radius=False)
        add_text(slide, "RPL11 DEEP DIVE", 0.72, 0.54, 4.0, 0.28,
                 size=10.5, color=SKY, bold=True)
        add_text(slide, "A ribosomal-stress candidate at the\ntranslation–mitochondria interface",
                 0.72, 0.98, 8.25, 1.26, size=31, color=WHITE,
                 bold=True, font=FONT_HEAD)
        add_text(slide,
                 "RPL11 is the broadest non-MT network candidate, but the mechanistic claim depends on distinguishing a true stress transducer from a highly connected module marker.",
                 0.75, 2.53, 8.10, 0.80, size=15.0, color=RGBColor(204, 219, 234))

        metrics = [
            ("4", "passing broad networks", TEAL),
            ("29/84", "significant / tested runs", SKY),
            ("25", "conservative-support runs", GOLD),
            ("15", "supporting fine cell types", PURPLE),
        ]
        for i, (value, label, accent) in enumerate(metrics):
            x = 0.75 + i * 2.06
            add_rect(slide, x, 3.57, 1.86, 1.22, color=NAVY_2, outline=accent)
            add_text(slide, value, x + 0.16, 3.75, 1.54, 0.46,
                     size=24, color=accent, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
            add_text(slide, label, x + 0.15, 4.26, 1.56, 0.30,
                     size=9.6, color=RGBColor(207, 220, 234), bold=True, align=PP_ALIGN.CENTER)

        add_rect(slide, 9.22, 0.73, 3.43, 2.63, color=PALE_GREEN, outline=TEAL)
        add_text(slide, "Mechanistic model", 9.50, 1.03, 2.86, 0.36,
                 size=18, color=NAVY, bold=True, font=FONT_HEAD)
        add_text(slide,
                 "Free RPL11 can inhibit MDM2 and activate p53; neuronal ribosomal stress can require RPL11 for p53-linked cell death.",
                 9.50, 1.56, 2.80, 1.17, size=12.8, color=GRAY)
        add_text(slide, "RPL11 transmits stress into mitochondrial and survival programs.",
                 9.50, 2.78, 2.82, 0.34, size=10.7, color=TEAL, bold=True)

        add_rect(slide, 9.22, 3.61, 3.43, 2.64, color=PALE_GOLD, outline=GOLD)
        add_text(slide, "Module-marker model", 9.50, 3.91, 2.86, 0.36,
                 size=18, color=NAVY, bold=True, font=FONT_HEAD)
        add_text(slide,
                 "RPL11 may inherit strong KDA and STRING signals from high expression, ribosomal connectivity, or a general disease/quality response.",
                 9.50, 4.44, 2.80, 1.16, size=12.8, color=GRAY)
        add_text(slide, "Target specificity and matched ribosomal nulls decide between the models.",
                 9.50, 5.66, 2.82, 0.40, size=10.7, color=VERMILION, bold=True)

        add_rect(slide, 0.75, 5.28, 8.05, 1.06, color=NAVY_2, outline=None)
        add_text(slide, "Human context", 1.02, 5.52, 1.40, 0.28,
                 size=12, color=SKY, bold=True)
        add_text(slide,
                 "AD capillary proteomics reported increased RPL11/RPL15 protein, while earlier AD work supports ribosomal dysfunction. The capillary signal does not establish neuronal or astrocytic localization.",
                 2.34, 5.42, 6.07, 0.63, size=11.4, color=WHITE)
        add_text(slide, "18", 12.37, 0.34, 0.42, 0.22,
                 size=9, color=RGBColor(145, 171, 197), bold=True, align=PP_ALIGN.RIGHT)
        add_source(slide, "Biological discussion: phase18_key_driver_gene_by_gene_initial_analysis.md • Zhang 2003; Slomnicki 2018; Suzuki 2022; Ding 2005")

        # 19 — astrocyte directed pathway figure
        slide = new_slide(prs)
        add_header(slide, "RPL11 • astrocytes", "A focused respiratory branch emerges from three supporting runs", 19,
                   accent=TEAL, subtitle="Directed Bayesian-network edges are model-derived regulatory hypotheses; pathway outlines summarize the displayed genes.")
        add_picture_contain(slide, trimmed["rpl11_astro_pathway"], 0.18, 1.25, 8.48, 5.85,
                            alt="RPL11-centered astrocyte consensus network with pathway annotations")
        add_rect(slide, 8.87, 1.43, 3.93, 5.46, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.16, 1.70, 3.35, accent=TEAL)
        add_metric(slide, "3/20", "supporting / usable astrocyte runs", 9.16, 2.18, 1.54, accent=TEAL)
        add_metric(slide, "3.44×10⁻⁵", "aggregate ACAT q", 10.88, 2.18, 1.57, accent=BLUE)
        add_panel_title(slide, "Network readout", 9.16, 3.48, 3.35, accent=BLUE)
        add_bullets(slide, [
            "18 nodes, 17 directed edges, and 12 mitochondrial query hits.",
            "Upstream chain: RPS25 → RPLP1 → RPL11. Direct outputs: COX7C, CWC15, and PRDX1.",
            "ETC/OXPHOS: 8 genes, BH FDR 1.98×10⁻⁸. Cristae formation: 3 genes, FDR 0.00475.",
            "The ribosome outline is contextual (FDR 0.0805), not significant at 0.05.",
        ], 9.16, 3.96, 3.27, size=11.4, line_h=0.61, accent=BLUE)
        add_text(slide, "Interpretation: a reproducible but relatively small astrocyte signal connects RPL11 to a coherent respiratory/cristae branch.",
                 9.18, 6.46, 3.20, 0.25, size=9.8, color=NAVY, bold=True)
        add_source(slide, "Source: phase18_rpl11_astrocyte_consensus_network_pathways.png, caption, methods, and ORA table")

        # 20 — astrocyte STRING figure
        slide = new_slide(prs)
        add_header(slide, "RPL11 • astrocytes", "STRING supports target-module coherence more clearly than direct regulation", 20,
                   accent=TEAL, subtitle="Medium-confidence STRING functional associations are undirected and are not astrocyte- or AD-specific.")
        add_rect(slide, 0.43, 1.39, 5.38, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(slide, trimmed["rpl11_astro_string"], 0.70, 1.62, 4.84, 4.98,
                            alt="STRING medium-confidence network for RPL11 and astrocyte mitochondrial target proteins")
        add_rect(slide, 6.08, 1.39, 6.72, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 6.42, 1.74, 5.98, accent=TEAL)
        add_bullets(slide, [
            "The respiratory targets form a dense protein-association core spanning complex I, III, IV, and V proteins.",
            "RPL11 is sparsely connected relative to the target-to-target core; PSAP and CYB5R3 remain visually isolated in this view.",
            "The visual density therefore supports a coherent mitochondrial module, but it cannot by itself show that RPL11 regulates that module.",
        ], 6.42, 2.27, 5.91, size=13.0, line_h=0.69, accent=TEAL)
        add_panel_title(slide, "Claim boundary", 6.42, 4.63, 5.98, accent=GOLD)
        add_text(slide,
                 "A direct RPL11 claim requires the input-only physical network, edge-level experimental/curated evidence, a ≥0.700 threshold, and matched ribosomal/topology controls. Until then, the medium-confidence figure is suggestive module-level cross-validation.",
                 6.43, 5.12, 5.80, 1.02, size=12.2, color=GRAY)
        add_rect(slide, 6.43, 6.23, 5.76, 0.42, color=PALE_GOLD, outline=GOLD)
        add_text(slide, "Allowed conclusion: protein-level coherence of the target module; direct RPL11 support remains provisional.",
                 6.66, 6.31, 5.30, 0.22, size=10.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_source(slide, "Source: astrocyte/stringdb_full_medium_conf.png • interpretation follows rpl11_astrocyte_string_analysis_guide.md")

        # 21 — excitatory-neuron directed pathway figure
        slide = new_slide(prs)
        add_header(slide, "RPL11 • excitatory neurons", "A broader translation–mitochondria interface recurs across 20 runs", 21,
                   accent=GOLD, subtitle="The excitatory context has stronger aggregate evidence and complete leave-one-fine-cell-type candidate retention.")
        add_picture_contain(slide, trimmed["rpl11_exc_pathway"], 0.14, 1.24, 8.63, 5.89,
                            alt="RPL11-centered excitatory-neuron consensus network with pathway annotations")
        add_rect(slide, 8.96, 1.42, 3.84, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.23, 1.70, 3.28, accent=GOLD)
        add_metric(slide, "20/97", "supporting / usable excitatory runs", 9.23, 2.18, 1.55, accent=GOLD)
        add_metric(slide, "1.84×10⁻⁹", "aggregate ACAT q", 10.96, 2.18, 1.55, accent=BLUE)
        add_panel_title(slide, "Network readout", 9.23, 3.48, 3.28, accent=BLUE)
        add_bullets(slide, [
            "35 nodes, 34 directed edges, 21 mitochondrial query hits, and 9 direct RPL11 outputs.",
            "Direct outputs include ribosomal proteins plus COX7C and SMDT1, creating a model-derived bridge to mitochondrial functions.",
            "Cytosolic ribosome: 14 genes, FDR 1.35×10⁻¹⁶; ETC/OXPHOS: 12 genes, FDR 3.20×10⁻¹³.",
            "Mitochondrial protein degradation (FDR 0.00180) and cristae formation (FDR 0.00762) add proteostasis/structure hypotheses.",
        ], 9.23, 3.96, 3.19, size=11.2, line_h=0.61, accent=BLUE)
        add_text(slide, "Interpretation: the excitatory signal is broader and more recurrent, but its ribosomal density makes specificity controls essential.",
                 9.25, 6.47, 3.12, 0.25, size=9.8, color=NAVY, bold=True)
        add_source(slide, "Source: phase18_rpl11_excitatory_consensus_network_pathways.png, node/edge tables, and ORA table")

        # 22 — excitatory-neuron STRING figure
        slide = new_slide(prs)
        add_header(slide, "RPL11 • excitatory neurons", "STRING resolves ribosomal and respiratory protein neighborhoods", 22,
                   accent=GOLD, subtitle="The displayed functional network is useful for architecture; it does not preserve Bayesian edge direction.")
        add_rect(slide, 0.35, 1.40, 7.63, 5.43, color=WHITE, outline=LIGHT)
        add_picture_contain(slide, trimmed["rpl11_exc_string"], 0.55, 1.62, 7.25, 4.99,
                            alt="STRING medium-confidence network for RPL11 and excitatory-neuron target proteins")
        add_rect(slide, 8.22, 1.40, 4.58, 5.43, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Architecture", 8.51, 1.74, 3.98, accent=GOLD)
        add_bullets(slide, [
            "RPL11 sits with cytosolic/mitochondrial ribosomal proteins, while the respiratory proteins form a dense OXPHOS core.",
            "The view is consistent with a ribosome–mitochondria interface rather than one simple RPL11→target protein chain.",
            "Respiratory-protein density is expected and can dominate whole-network PPI enrichment.",
        ], 8.51, 2.25, 3.86, size=12.2, line_h=0.70, accent=GOLD)
        add_panel_title(slide, "Exploratory elements", 8.51, 4.52, 3.98, accent=TEAL)
        add_text(slide,
                 "MRPS18C, COX17, COX20, and PINK1 appear as STRING-side connectors rather than prespecified Bayesian-consensus nodes. They are mechanism hypotheses—not validation endpoints.",
                 8.52, 5.03, 3.79, 0.84, size=11.5, color=GRAY)
        add_rect(slide, 8.51, 6.05, 3.84, 0.56, color=PALE_GOLD, outline=GOLD)
        add_text(slide, "Test RPL11-to-target edges separately from target-to-target connectivity.",
                 8.72, 6.17, 3.43, 0.28, size=10.3, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_source(slide, "Source: excitatory/stringdb_full_medium_conf.png • medium-confidence functional STRING view")

        # 23 — cross-cell-type synthesis and validation plan
        slide = new_slide(prs)
        add_header(slide, "RPL11 synthesis", "The excitatory signal is broader; the astrocyte signal is more focused", 23,
                   subtitle="Both contexts support a translation/stress-to-mitochondria hypothesis, while leaving target specificity unresolved.")

        add_rect(slide, 0.69, 1.55, 5.89, 3.63, color=PALE_GREEN, outline=TEAL)
        add_text(slide, "ASTROCYTES", 1.00, 1.83, 2.20, 0.30, size=11, color=TEAL, bold=True)
        add_text(slide, "Focused respiratory branch", 1.00, 2.19, 4.93, 0.43,
                 size=21, color=NAVY, bold=True, font=FONT_HEAD)
        add_bullets(slide, [
            "3/20 supporting runs; aggregate q 3.44×10⁻⁵.",
            "18-node directed consensus with ETC/OXPHOS and cristae enrichment.",
            "STRING shows a dense target core but comparatively sparse RPL11 connectivity.",
            "Best current framing: focused, plausible, and still in need of protein-edge/null validation.",
        ], 1.00, 2.78, 5.04, size=12.5, line_h=0.50, accent=TEAL)

        add_rect(slide, 6.79, 1.55, 5.86, 3.63, color=PALE_GOLD, outline=GOLD)
        add_text(slide, "EXCITATORY NEURONS", 7.10, 1.83, 2.75, 0.30, size=11, color=GOLD, bold=True)
        add_text(slide, "Broad translation–mitochondria interface", 7.10, 2.19, 5.02, 0.43,
                 size=21, color=NAVY, bold=True, font=FONT_HEAD)
        add_bullets(slide, [
            "20/97 supporting runs; aggregate q 1.84×10⁻⁹.",
            "35-node consensus combines ribosome, OXPHOS, cristae, and mitochondrial proteostasis.",
            "STRING separates a ribosomal neighborhood and dense respiratory core with exploratory connectors.",
            "Best current framing: stronger recurrence, but higher risk of generic ribosomal-module effects.",
        ], 7.10, 2.78, 4.97, size=12.5, line_h=0.50, accent=GOLD)

        add_rect(slide, 0.91, 5.51, 11.50, 1.28, color=NAVY, outline=None)
        add_text(slide, "DECISIVE VALIDATION", 1.18, 5.72, 1.78, 0.28, size=10.5, color=SKY, bold=True)
        add_text(slide,
                 "Match for expression, degree, neighborhood size, and ribosomal class → use partial CRISPRi/CRISPRa in APOE-isogenic astrocytes and excitatory neurons → measure predicted MT targets, p53, nascent translation, respiration, and viability → add p53 inhibition/rescue and spatial/vascular proteomics.",
                 2.82, 5.68, 9.15, 0.78, size=11.6, color=WHITE)
        add_source(slide, "Interpretation and validation priorities synthesized from phase18_key_driver_gene_by_gene_initial_analysis.md")

        # 24 — APOE astrocyte directed pathway figure
        slide = new_slide(prs)
        add_header(
            slide,
            "APOE • astrocytes",
            "APOE anchors a recurrent astrocyte mitochondrial neighborhood",
            24,
            accent=BLUE,
            subtitle="APOE is an established AD gene; the Phase 18 contribution is its placement within a specific astrocyte mitochondrial network.",
        )
        add_picture_contain(
            slide,
            trimmed["apoe_astro_pathway"],
            0.12,
            1.23,
            8.93,
            5.91,
            alt="APOE-centered astrocyte consensus network with directed Bayesian-network edges and contextual pathway outlines",
        )
        add_rect(slide, 9.20, 1.42, 3.60, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.47, 1.70, 3.05, accent=BLUE)
        add_metric(slide, "4/14", "supporting / tested astrocyte runs", 9.47, 2.18, 1.38, accent=BLUE)
        add_metric(slide, "1.27×10⁻²", "aggregate ACAT q", 10.99, 2.18, 1.53, accent=TEAL)
        add_panel_title(slide, "Network readout", 9.47, 3.48, 3.05, accent=TEAL)
        add_bullets(slide, [
            "19 nodes, 18 directed edges, and seven mitochondrial query hits across four supporting runs.",
            "APOE-linked paths reach ATP5F1A, ATP5PB, CHCHD10, LDHB, NME3, TUFM, and AGT.",
            "Amyloid, cholesterol/efflux, and cristae outlines are contextual; none has BH FDR < 0.05.",
            "Support spans three fine cell types and both mitochondrial directions; leave-one-fine-cell-type retention is 2/3.",
        ], 9.47, 3.96, 2.91, size=10.7, line_h=0.58, accent=TEAL)
        add_text(
            slide,
            "Interpretation: credible network placement, but less stable than the strongest neuronal candidates and not proof of direct APOE regulation.",
            9.49,
            6.43,
            2.87,
            0.31,
            size=9.5,
            color=NAVY,
            bold=True,
        )
        add_source(slide, "Source: phase18_apoe_astrocyte_consensus_network_pathways.png, caption, methods, and gene-by-gene analysis")

        # 25 — APOE astrocyte STRING figure and validation interpretation
        slide = new_slide(prs)
        add_header(
            slide,
            "APOE • astrocytes",
            "STRING reveals two distinct APOE-neighborhood modules",
            25,
            accent=BLUE,
            subtitle="Medium-confidence STRING associations are undirected and context-agnostic; they complement, but do not validate, Bayesian edge direction.",
        )
        add_rect(slide, 0.43, 1.39, 6.04, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(
            slide,
            trimmed["apoe_astro_string"],
            0.70,
            1.69,
            5.50,
            4.66,
            alt="STRING medium-confidence functional association network for APOE and the astrocyte consensus-neighborhood proteins",
        )
        add_rect(slide, 6.72, 1.39, 6.08, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 7.03, 1.72, 5.47, accent=BLUE)
        add_bullets(slide, [
            "A mitochondrial protein cluster contains TUFM, ATP5F1A, ATP5PB, CHCHD10, and NME3.",
            "APOE instead connects with a lipid/amyloid/stress neighborhood containing PLTP, CST3, AGT, ITM2B, and GPX4.",
            "Their separation argues against reading this figure as a simple direct APOE→mitochondrial protein chain.",
        ], 7.03, 2.20, 5.42, size=11.8, line_h=0.60, accent=BLUE)
        add_panel_title(slide, "Mechanism and decisive test", 7.03, 4.15, 5.47, accent=TEAL)
        add_text(
            slide,
            "Prior cell studies connect APOE4 to lysosomal cholesterol accumulation, impaired mitophagy, altered mitochondrial dynamics, and disrupted fatty-acid metabolism. Test whether that biology reproduces the exact Phase 18 target neighborhood in APOE2/3/4-isogenic astrocytes.",
            7.04,
            4.62,
            5.35,
            0.92,
            size=11.4,
            color=GRAY,
        )
        add_rect(slide, 7.04, 5.72, 5.31, 0.73, color=PALE_GREEN, outline=TEAL)
        add_text(
            slide,
            "Measure respiration, substrate use, cholesterol distribution, mitophagy, lipid-droplet transfer, and the frozen target module; fit donor-level disease × sex × APOE interactions separately.",
            7.24,
            5.85,
            4.90,
            0.44,
            size=9.9,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "Observed female ε2, male ε2, and male ε4 support is descriptive—not evidence of an interaction.",
            7.08,
            6.55,
            5.24,
            0.22,
            size=9.4,
            color=VERMILION,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_source(slide, "Sources: stringdb_full_medium_conf.png • phase18_key_driver_gene_by_gene_initial_analysis.md • Lee 2023; Schmukler 2020; Qi 2021; Williams 2020")

        # Narrative order: agenda/section → selection/circular/atlas → synthesis
        # → sex/APOE → human genetics → validation/RPL11. The synthesis slide is
        # authored after the sex/APOE slides, so move its saved slide-id entry to
        # position nine, then update the small upper-right slide numerals.
        slide_ids = prs.slides._sldIdLst
        slide_ids.remove(synthesis_slide_id)
        slide_ids.insert(8, synthesis_slide_id)
        for position, ordered_slide in enumerate(prs.slides, start=1):
            for shape in ordered_slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if shape.left <= Inches(12.0) or shape.top >= Inches(0.7):
                    continue
                current = shape.text.strip()
                if len(current) == 2 and current.isdigit():
                    runs = shape.text_frame.paragraphs[0].runs
                    if runs:
                        runs[0].text = f"{position:02d}"

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    validate_deck(output_path)
    return output_path


def validate_deck(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 25:
        raise RuntimeError(f"Expected 25 slides, found {len(prs.slides)}")
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        raise RuntimeError("Deck is not 16:9 widescreen")
    image_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    if image_shapes != 14:
        raise RuntimeError(f"Expected fourteen requested figure images, found {image_shapes}")
    ordered_text = [
        "\n".join(shape.text for shape in slide.shapes
                  if getattr(shape, "has_text_frame", False))
        for slide in prs.slides
    ]
    expected_sequence = {
        1: "PRESENTATION ROADMAP",
        2: "SECTION 01",
        3: "SELECTION PROCESS",
        8: "SYNTHESIS",
        9: "SEX/APOE EVIDENCE • MT",
        10: "SEX/APOE EVIDENCE • NON-MT",
        11: "SECTION 02",
        12: "HUMAN GENETIC SUPPORT • DATA",
        13: "HUMAN GENETIC SUPPORT • METHOD",
        14: "HUMAN GENETIC SUPPORT • RESULTS",
        15: "SECTION 03",
        16: "FIRST-BATCH VALIDATION PANEL",
        17: "RPL11 DEEP DIVE",
        22: "RPL11 SYNTHESIS",
        23: "APOE • ASTROCYTES",
        24: "APOE • ASTROCYTES",
    }
    for zero_based_index, marker in expected_sequence.items():
        if marker not in ordered_text[zero_based_index]:
            raise RuntimeError(f"Expected {marker!r} on slide {zero_based_index + 1}")
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        if len([name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]) != 25:
            raise RuntimeError("PPTX package has an unexpected slide count")
        package_text = "\n".join(
            zf.read(name).decode("utf-8", errors="ignore")
            for name in names if name.endswith((".xml", ".rels"))
        )
        if "filter_attrition" in package_text:
            raise RuntimeError("Deprecated filter_attrition reference found in PPTX package")
        visible_genetics_text = "\n".join(ordered_text[12:15]).lower()
        if "phase 19" in visible_genetics_text or "phase19" in visible_genetics_text:
            raise RuntimeError("Internal phase identifier found on a human-genetic-support slide")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUT,
        help=f"Output PPTX path (default: {DEFAULT_OUT})",
    )
    args = parser.parse_args()
    output = build_deck(args.output)
    print(output)
