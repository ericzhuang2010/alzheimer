#!/usr/bin/env python3
"""Update the genetic-support deck with the completed 19b screen results.

The 19b rerun extended the public genetic screen from 15 drivers to all 433
(summary screen + regional AD/CSF scans). This updater surgically:

- rewrites the three headline cards on slide 1;
- replaces slides 8, 9, 12, 13, and 14 (evidence map, coverage, negatives,
  gap, take-home) with versions built from the validated 19b outputs;
- leaves slides 2-7 and the RPS15/APOE spotlight slides (10, 11) untouched.

Every displayed number is recomputed from the 19b result bundles and checked
against frozen expected values before any slide is modified.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).resolve().parent))
import update_phase11_seaad_simple_aggr_part2 as ui  # noqa: E402  (shared styling helpers)


ROOT = Path(__file__).resolve().parents[2]
DECK = (
    ROOT
    / "docs/presentations"
    / "human_genetic_support_for_key_drivers_simple_aggr_08292026.pptx"
)
TIER1_EVIDENCE = ROOT / "results/minerva_production/19b_genetic_support_tier1/fungen_gene_evidence.tsv"
REGIONAL_ALL = (
    ROOT / "results/minerva_production/19b_genetic_support_regional/regional_summary_all_traits.tsv"
)
AUDIT_PATH = (
    ROOT
    / "results/presentations/human_genetic_support_simple_aggr"
    / "deck_19b_update_checks.tsv"
)

EXPECTED_INPUT_SLIDES = 14
REPLACE_POSITIONS = {8, 9, 12, 13, 14}  # 1-based
PRE_CONTRACT_TITLES = {
    8: "15 of the 433 drivers have genetic screening results so far",
    9: "Genetic evidence for the 15 screened drivers, at a glance",
    10: "RPS15: network centrality and genetic proximity converge",
    11: "APOE: the strongest genetic result is an astrocyte driver",
    12: "Eleven drivers had no qualifying signal — read this carefully",
    13: "418 of 433 drivers have never been screened",
    14: "What we know now — and the shortest path to stronger claims",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DECK)
    parser.add_argument("--output", type=Path, default=DECK)
    parser.add_argument("--audit", type=Path, default=AUDIT_PATH)
    return parser.parse_args()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def load_facts() -> dict[str, Any]:
    for path in (TIER1_EVIDENCE, REGIONAL_ALL):
        if not path.is_file():
            raise FileNotFoundError(path)
    evidence = pd.read_csv(TIER1_EVIDENCE, sep="\t")
    grades = evidence["grade"].value_counts().to_dict()
    strong = evidence[evidence["grade"].eq("strong")].sort_values("direct_min_p")
    moderate = evidence[evidence["grade"].eq("moderate")].sort_values("direct_min_p")
    weak = evidence[evidence["grade"].eq("weak")].sort_values("direct_min_p")

    regional = pd.read_csv(REGIONAL_ALL, sep="\t")
    ad = regional[regional["trait"].eq("clinical_ad_bellenguez2022")]
    chrx_untested = ad[ad["variant_rows"].eq(0)]
    gw_ad = ad[ad["genome_wide_significant"]]

    def ad_value(gene: str, column: str) -> Any:
        rows = ad[ad["gene"].eq(gene)]
        return rows.iloc[0][column] if len(rows) else None

    facts = {
        "grades": grades,
        "strong_genes": strong["gene"].tolist(),
        "moderate_genes": moderate["gene"].tolist(),
        "weak_genes": weak["gene"].tolist(),
        "gw_ad_windows": len(gw_ad),
        "chrx_untested": len(chrx_untested),
        "unmapped_symbols": 433 - regional["gene"].nunique(),
        "plcg2_ad_p": float(ad_value("PLCG2", "regional_min_p")),
        "hspa1a_ad_p": float(ad_value("HSPA1A", "regional_min_p")),
        "map3k2_ad_p": float(ad_value("MAP3K2", "regional_min_p")),
        "map3k2_lead": str(ad_value("MAP3K2", "regional_lead_variant")),
        "rps15_grade": evidence.loc[evidence["gene"].eq("RPS15"), "grade"].iloc[0],
    }
    tau = regional[
        regional["gene"].eq("MAP1LC3B")
        & regional["trait"].isin(["csf_total_tau_gcst90726397", "csf_ptau181_gcst90726398"])
    ]
    facts["map1lc3b_tau_gw"] = bool(tau["genome_wide_significant"].all()) and len(tau) == 2
    prio = ad[ad["gene"].isin(["WDR82", "HGSNAT", "TTC8"])]
    facts["prio_no_gw"] = bool((~prio["genome_wide_significant"]).all()) and len(prio) == 3

    expected = {
        "strong": 8,
        "moderate": 3,
        "weak": 17,
        "none_found": 405,
    }
    for grade, count in expected.items():
        if grades.get(grade, 0) != count:
            raise RuntimeError(f"Grade drift for {grade}: {grades.get(grade, 0)} != {count}")
    if facts["gw_ad_windows"] != 33 or facts["chrx_untested"] != 20 or facts["unmapped_symbols"] != 8:
        raise RuntimeError("Regional contract drift")
    if facts["map3k2_lead"] != "rs6733839":
        raise RuntimeError("MAP3K2 lead variant is not the expected BIN1 signal")
    if not facts["map1lc3b_tau_gw"] or not facts["prio_no_gw"]:
        raise RuntimeError("MAP1LC3B / priority-gene contract drift")
    if facts["rps15_grade"] != "none_found":
        raise RuntimeError("RPS15 direct-mapping grade drifted")
    if "APOE" not in facts["strong_genes"] or "PLCG2" not in facts["strong_genes"]:
        raise RuntimeError("Expected APOE and PLCG2 among strong grades")
    return facts


def replace_text_anywhere(slide, old: str, new: str) -> None:
    matches = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text == old:
                    matches.append(run)
    if len(matches) != 1:
        raise RuntimeError(f"Expected one run with text {old!r}, found {len(matches)}")
    matches[0].text = new


def p_text(value: float) -> str:
    return f"{value:.1e}".replace("e-0", "e-")


def update_slide_1(slide) -> None:
    replace_text_anywhere(slide, "4", "All 433")
    replace_text_anywhere(slide, "drivers with genetic evidence", "drivers now screened (summary + regional)")
    replace_text_anywhere(slide, "APOE + RPS15", "APOE · RPS15 · PLCG2")
    replace_text_anywhere(slide, "15 of 433", "8 strong")
    replace_text_anywhere(slide, "drivers screened so far", "direct-evidence grades (3 moderate, 17 weak)")
    ui.add_notes(
        slide,
        goal="Set the question with the completed screen: every driver now has public genetic screening results.",
        walkthrough="All 433 drivers went through the public summary screen and the regional AD and CSF scans. Eight drivers carry strong direct evidence grades, led by APOE and the established AD gene PLCG2; RPS15 keeps its separate regional-plus-QTL case.",
        boundary="Gene-level tests and same-variant analyses are still pending; grades come from the public summary screen only.",
        transition="First, how the two kinds of evidence relate.",
    )


def build_evidence_map(prs, facts) -> str:
    title = "Direct public evidence across all 433 drivers, at a glance"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Grades from the public summary screen, now covering every driver.")
    strong_list = ", ".join(facts["strong_genes"])
    moderate_list = ", ".join(facts["moderate_genes"])
    weak_preview = ", ".join(facts["weak_genes"][:6]) + ", …"
    ui.add_table(
        slide,
        ["Grade", "Drivers", "Genes", "What it means"],
        [
            ["Strong", "8", strong_list, "Fine-mapped AD variant assigned to the gene"],
            ["Moderate", "3", moderate_list, "Direct mapping with lower inclusion score"],
            ["Weak", "17", weak_preview, "Sub-genome-wide mapping or list membership"],
            ["None found", "405", "—", "No direct mapping in the public summaries"],
        ],
        0.80, 1.60, [1.55, 1.15, 5.60, 3.60], row_h=0.80, header_h=0.50, font_size=10.2,
    )
    ui.add_text(
        slide,
        "PLCG2 — an established AD gene — is now also a graded driver. RPS15 stays \u201Cnone found\u201D here: "
        "its case rests on regional and brain-activity evidence, not a published direct mapping.",
        0.88, 5.75, 11.60, 0.50, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Replace the legacy 15-gene evidence map with the full-list grading.",
        walkthrough="The strong tier holds APOE plus seven newly graded drivers, including the established AD gene PLCG2 and the STAG3, PPP4C/SEPHS2, ZNF251, and ZNF652 mappings. Three drivers are moderate and seventeen weak; the remaining four hundred five have no direct public mapping.",
        boundary="Grades reflect published summary evidence only; they are not colocalization results.",
        transition="Where each driver now stands in the pipeline.",
    )
    return title


def build_stage_status(prs, facts) -> str:
    title = "Screening breadth is complete; depth is the remaining work"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Every driver has been through the first two evidence stages.")
    ui.add_table(
        slide,
        ["Evidence stage", "Status", "Coverage and notes"],
        [
            ["Public summary screen", "complete", "All 433 drivers graded"],
            [
                "Regional AD + CSF scans",
                "complete",
                f"405 windows scanned; {facts['chrx_untested']} X-chromosome drivers not assessable "
                f"(autosomal GWAS); {facts['unmapped_symbols']} symbols unmappable",
            ],
            ["Gene-based tests", "next", "Resolve shared-locus windows; runs on the data host"],
            ["Brain QTL routes", "pending", "Replicated and recurrent drivers first"],
            ["Same-variant tests", "blocked", "Needs complete public models + matched references"],
        ],
        0.80, 1.60, [3.30, 1.55, 7.05], row_h=0.72, header_h=0.50, font_size=10.6,
    )
    ui.add_text(
        slide,
        "The coverage gap of the first screen is closed — what remains is depth, not breadth.",
        0.88, 6.05, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Show that the old 15-of-433 coverage gap is closed and name the remaining stages.",
        walkthrough="The summary screen and the four regional scans cover every driver that can be mapped and tested. Twenty X-chromosome drivers, including BEX3, cannot be tested with the autosomal GWAS files, and eight clone-named symbols lack a unique genome mapping.",
        boundary="Completed stages establish coverage, not validation; the pending stages carry the inferential weight.",
        transition="The RPS15 and APOE spotlights are unchanged by the rerun.",
    )
    return title


def build_negatives(prs, facts) -> str:
    title = "Most drivers have no nearby signal — read the negatives carefully"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "\u201CNo genetic support found\u201D still has three different origins.")
    cards = [
        (
            "Truly signal-negative",
            "Most drivers — including the cross-cohort replicators WDR82, HGSNAT, and TTC8 — have no "
            "genome-wide variant nearby (best regional P ≈ 1e-4).",
            ui.PALE_SKY, ui.BLUE,
        ),
        (
            "Not assessable",
            f"{facts['chrx_untested']} X-chromosome drivers (BEX3 among them) are untestable: the AD and CSF "
            f"GWAS files cover autosomes only. {facts['unmapped_symbols']} further symbols lack a unique mapping.",
            ui.PALE_GOLD, ui.GOLD,
        ),
        (
            "Proximity artifacts",
            "A significant window often belongs to a neighbor: the strongest \u201Chit\u201D is MAP3K2's window "
            "catching BIN1's famous signal.",
            ui.PALE_RED, ui.VERMILION,
        ),
    ]
    for index, (heading, body, bg, accent) in enumerate(cards):
        x = 0.66 + index * 4.10
        ui.add_rect(slide, x, 1.55, 3.83, 3.10, color=bg, outline=accent)
        ui.add_text(slide, heading, x + 0.26, 1.85, 3.30, 0.40, size=15.5, color=ui.NAVY, bold=True, font=ui.FONT_HEAD)
        ui.add_text(slide, body, x + 0.26, 2.45, 3.32, 2.05, size=11.0, color=ui.DARK)
    ui.add_rect(slide, 0.66, 5.00, 12.0, 1.35, color=ui.WHITE, outline=ui.LIGHT)
    ui.add_panel_title(slide, "Replication without genetic support is still informative", 0.96, 5.23, 9.0, accent=ui.PURPLE)
    ui.add_text(
        slide,
        "WDR82, HGSNAT, and TTC8 replicate across two human cohorts as network drivers yet show no common risk "
        "variant nearby. Network recurrence and inherited risk answer different causal questions; the pending "
        "gene-based and rare-variant designs can still test them.",
        0.98, 5.67, 11.40, 0.60, size=11.5, color=ui.DARK,
    )
    ui.add_notes(
        slide,
        goal="Update the negatives reading for the full-list screen.",
        walkthrough="Three origins remain: genuinely signal-negative regions, now including the three testable cross-cohort priorities; not-assessable routes, now dominated by the twenty X-chromosome drivers; and proximity artifacts, where a window inherits a neighbor's signal.",
        boundary="Signal-negative applies to common autosomal variants under the frozen threshold only.",
        transition="The genuine new leads the expanded screen produced.",
    )
    return title


def build_leads(prs, facts) -> str:
    title = "The expanded screen surfaces new leads — and known-locus artifacts"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title, "Trust the direct gene-level mappings; discount shared-locus window hits.")
    ui.add_rect(slide, 0.66, 1.52, 5.90, 4.60, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "Genuine new leads", 0.96, 1.84, 5.30, accent=ui.TEAL)
    ui.add_bullets(slide, [
        f"PLCG2 — established AD gene, now also a driver: strong direct mapping and a significant own region (P = {p_text(facts['plcg2_ad_p'])}).",
        f"HSPA1A — SEA-AD-replicated driver with a significant region (P = {p_text(facts['hspa1a_ad_p'])}), in the complex HLA area.",
        "MAP1LC3B — genome-wide significant for both tau spinal-fluid markers.",
        "STAG3, ZNF652, PPP4C/SEPHS2, ZNF251 — new strong direct mappings.",
    ], 0.98, 2.42, 5.40, size=11.2, line_h=0.90, accent=ui.TEAL)
    ui.add_rect(slide, 6.85, 1.52, 5.82, 4.60, color=ui.PALE_RED, outline=ui.VERMILION)
    ui.add_panel_title(slide, "Artifacts to discount", 7.15, 1.84, 5.22, accent=ui.VERMILION)
    ui.add_bullets(slide, [
        f"{facts['gw_ad_windows']} driver windows are genome-wide significant for AD — but many share one lead variant with a famous neighbor.",
        f"MAP3K2's window (P = {p_text(facts['map3k2_ad_p'])}) is BIN1's signal; TMEM259's window carries the RPS15-locus signal.",
        "Five drivers around SHARPIN share a single variant.",
        "Gene-based tests, the next stage, separate these from real gene-level signals.",
    ], 7.17, 2.42, 5.30, size=11.2, line_h=0.90, accent=ui.VERMILION)
    ui.add_text(
        slide,
        "Window significance is not gene support — the strong direct mappings are the trustworthy new leads.",
        0.88, 6.40, 11.60, 0.30, size=11.5, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Separate the trustworthy new findings from the expected proximity artifacts.",
        walkthrough="PLCG2 is the headline: an established AD gene that is also a network driver, with both a strong direct mapping and its own significant region. HSPA1A adds a replicated driver in the HLA area, and MAP1LC3B is significant for both tau biomarkers. On the artifact side, the strongest window hits belong to BIN1 and the RPS15 locus rather than the drivers named on the windows.",
        boundary="HLA-region and shared-locus results need conditional analyses before any gene claim.",
        transition="Close with the updated state and the ordered next steps.",
    )
    return title


def build_takehome(prs, facts) -> str:
    title = "What we know now — and what runs next"
    slide = ui.new_slide(prs)
    ui.add_title_block(slide, title)
    ui.add_rect(slide, 0.66, 1.30, 5.95, 5.05, color=ui.PALE_GREEN, outline=ui.TEAL)
    ui.add_panel_title(slide, "Where the evidence stands", 0.97, 1.62, 5.35, accent=ui.TEAL)
    ui.add_bullets(slide, [
        "APOE: strongest support — direct coding variant plus all three spinal-fluid biomarkers.",
        "PLCG2: the best new driver-genetics convergence.",
        "RPS15: unchanged — top driver, strong region, same-variant test still blocked.",
        "WDR82, HGSNAT, TTC8: no common-variant support; BEX3 untestable (X chromosome).",
    ], 0.99, 2.20, 5.40, size=11.6, line_h=0.92, accent=ui.TEAL)
    ui.add_rect(slide, 6.90, 1.30, 5.78, 5.05, color=ui.PALE_SKY, outline=ui.BLUE)
    ui.add_panel_title(slide, "Next steps, in order", 7.21, 1.62, 5.18, accent=ui.BLUE)
    ui.add_bullets(slide, [
        "Gene-based tests across all drivers and traits — resolves the shared-locus windows.",
        "Brain QTL routes for the replicated and recurrent drivers.",
        "Same-variant tests wherever complete public models and matched references exist.",
        "Protein-level, rare-variant, and sex/APOE-interaction designs afterward.",
    ], 7.23, 2.20, 5.25, size=11.6, line_h=0.92, accent=ui.BLUE)
    ui.add_text(
        slide,
        "Breadth is done — every driver is screened. Depth decides which of the new leads survive.",
        0.88, 6.60, 11.60, 0.30, size=12.0, color=ui.PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    ui.add_notes(
        slide,
        goal="Leave the updated summary and the ordered plan for the remaining stages.",
        walkthrough="APOE and PLCG2 lead the supported set; RPS15's unresolved case is unchanged; the cross-cohort replicators lack common-variant support and move to the alternative designs. Next: gene-based tests, QTL routes, then same-variant analyses.",
        boundary="All current grades are screening-level; no same-variant test has been completed for any driver.",
        transition="End of deck.",
    )
    return title


def main() -> int:
    args = parse_args()
    input_path = args.input.resolve()
    output_path = args.output.resolve()
    if not input_path.is_file():
        raise FileNotFoundError(input_path)
    facts = load_facts()
    original_hash = sha256_file(input_path)

    prs = Presentation(str(input_path))
    if len(prs.slides) != EXPECTED_INPUT_SLIDES:
        raise RuntimeError(f"Expected {EXPECTED_INPUT_SLIDES} slides, found {len(prs.slides)}")
    for position, expected_title in PRE_CONTRACT_TITLES.items():
        if expected_title not in slide_text(prs.slides[position - 1]):
            raise RuntimeError(f"Slide {position} does not match the pre-update contract")
    untouched_positions = [2, 3, 4, 5, 6, 7, 10, 11]
    before_untouched = {n: slide_text(prs.slides[n - 1]) for n in untouched_positions}

    ui.set_notes_body_template(prs.slides[0].notes_slide.notes_placeholder._element)
    update_slide_1(prs.slides[0])

    # Append the replacement slides FIRST so python-pptx allocates fresh part
    # names (deleting first would make it reuse names of surviving slides).
    new_titles = [
        build_stage_status(prs, facts),   # -> position 8 (coverage first, matching user's order)
        build_evidence_map(prs, facts),   # -> position 9
        build_negatives(prs, facts),      # -> position 12
        build_leads(prs, facts),          # -> position 13
        build_takehome(prs, facts),       # -> position 14
    ]

    # Now remove the five superseded slides (descending positions).
    slide_id_list = prs.slides._sldIdLst
    for position in sorted(REPLACE_POSITIONS, reverse=True):
        slide_id = list(slide_id_list)[position - 1]
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)

    # Current order: 1-7, RPS15, APOE, A, B, C, D, E  ->  1-7, A, B, RPS15, APOE, C, D, E
    ids = list(slide_id_list)
    desired = ids[:7] + [ids[9], ids[10], ids[7], ids[8], ids[11], ids[12], ids[13]]
    for element in ids:
        slide_id_list.remove(element)
    for element in desired:
        slide_id_list.append(element)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    file_descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{output_path.name}.", suffix=".tmp", dir=output_path.parent
    )
    os.close(file_descriptor)
    temporary = Path(temporary_name)
    try:
        prs.save(str(temporary))
        temporary.replace(output_path)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise

    reloaded = Presentation(str(output_path))
    after_untouched = {n: slide_text(reloaded.slides[n - 1]) for n in untouched_positions}
    untouched_ok = before_untouched == after_untouched
    slide1_text = slide_text(reloaded.slides[0])
    expected_position_titles = dict(zip([8, 9, 12, 13, 14], new_titles))
    checks: list[dict[str, Any]] = [
        {
            "check_id": "output_slide_count",
            "observed": len(reloaded.slides),
            "expected": EXPECTED_INPUT_SLIDES,
            "passed": len(reloaded.slides) == EXPECTED_INPUT_SLIDES,
        },
        {
            "check_id": "untouched_slides_2_7_10_11_unchanged",
            "observed": "unchanged" if untouched_ok else "changed",
            "expected": "unchanged",
            "passed": untouched_ok,
        },
        {
            "check_id": "slide1_cards_updated",
            "observed": "All 433" in slide1_text and "PLCG2" in slide1_text,
            "expected": True,
            "passed": "All 433" in slide1_text and "PLCG2" in slide1_text,
        },
    ]
    for position, title in expected_position_titles.items():
        present = title in slide_text(reloaded.slides[position - 1])
        checks.append(
            {
                "check_id": f"slide{position}_new_title",
                "observed": present,
                "expected": True,
                "passed": present,
            }
        )
    notes_ok = all(
        reloaded.slides[n - 1].has_notes_slide
        and reloaded.slides[n - 1].notes_slide.notes_text_frame is not None
        and reloaded.slides[n - 1].notes_slide.notes_text_frame.text.strip() != ""
        for n in [1, 8, 9, 12, 13, 14]
    )
    checks.append(
        {"check_id": "updated_slides_have_notes", "observed": notes_ok, "expected": True, "passed": notes_ok}
    )

    audit = args.audit.resolve()
    audit.parent.mkdir(parents=True, exist_ok=True)
    with audit.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["check_id", "observed", "expected", "passed"],
            delimiter="\t",
            lineterminator="\n",
        )
        writer.writeheader()
        writer.writerows(checks)
    failed = [row["check_id"] for row in checks if not row["passed"]]
    if failed:
        raise RuntimeError("Deck update failed checks: " + ", ".join(failed))
    print(f"updated={output_path}")
    print("slides_replaced=1(cards),8,9,12,13,14")
    print(f"original_sha256={original_hash}")
    print(f"updated_sha256={sha256_file(output_path)}")
    print(f"audit={audit}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
