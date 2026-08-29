#!/usr/bin/env python3
"""Repoint Part II of the Phase 20 deck at the SEA-AD simple aggregation.

The updater rewrites the Part II references on slides 1 and 2, keeps slides
3-10 byte-content identical, and appends a six-slide SEA-AD Part II: a
section divider, two setup slides, the top-five and driver-recurrence figure
slides, and a summary slide. Every displayed number is loaded from the
validated project-local SEA-AD tables and checked against the frozen
expected values before any slide is written.
"""

from __future__ import annotations

import argparse
import copy
import csv
import hashlib
import os
import tempfile
from collections import Counter
from pathlib import Path
from typing import Any, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
DEFAULT_DECK = ROOT / "docs" / "presentations" / "phase20_sex_apoe_kda_fine_broad.pptx"
RESULT_DIR = ROOT / "results" / "validation_human" / "11_sex_apoe_kda_simple_aggr"
MANIFEST_PATH = (
    ROOT
    / "results"
    / "validation_human"
    / "10_seaad_kda_rediscovery"
    / "10a_inputs"
    / "seaad_kda_run_manifest.tsv"
)
FIGURE_DIR = (
    ROOT / "results" / "figures" / "validation_human" / "phase_11_sex_apoe_simple_aggr"
)
RECURRENCE_PNG = (
    FIGURE_DIR / "driver_recurrence" / "phase11_seaad_simple_aggr_driver_recurrence.png"
)
TOP5_PNG = FIGURE_DIR / "top5_candidates" / "phase11_seaad_simple_aggr_top5_candidates.png"
RECURRENCE_DATA = (
    FIGURE_DIR
    / "driver_recurrence"
    / "phase11_seaad_simple_aggr_driver_recurrence_plot_data.tsv"
)
TOP5_DATA = (
    FIGURE_DIR
    / "top5_candidates"
    / "phase11_seaad_simple_aggr_top5_candidates_plot_data.tsv"
)
AUDIT_PATH = (
    ROOT
    / "results"
    / "presentations"
    / "phase20_sex_apoe_kda_fine_broad"
    / "phase11_seaad_part2_slide_update_checks.tsv"
)

TRUE_VALUES = {"TRUE", "T", "1", "YES"}
EXPECTED_INPUT_SLIDES = 10
UNCHANGED_SLIDES = range(3, 11)

NAVY = RGBColor(15, 35, 61)
NAVY_2 = RGBColor(30, 59, 91)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
TEAL = RGBColor(0, 158, 115)
GOLD = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(126, 76, 154)
TEAL_TEXT = RGBColor(0, 104, 77)
GOLD_TEXT = RGBColor(137, 86, 0)
VERMILION_TEXT = RGBColor(158, 58, 0)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)
LIGHT = RGBColor(224, 231, 239)
MID = RGBColor(108, 121, 136)
DARK = RGBColor(31, 39, 48)
GRAY = RGBColor(78, 90, 104)
PALE_SKY = RGBColor(235, 247, 252)
PALE_GREEN = RGBColor(228, 244, 238)
PALE_GOLD = RGBColor(253, 242, 217)
PALE_RED = RGBColor(253, 235, 228)
PALE_GRAY = RGBColor(241, 244, 247)

FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--audit", type=Path, default=AUDIT_PATH)
    return parser.parse_args()


def truth(value: Any) -> bool:
    return str(value).strip().upper() in TRUE_VALUES


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def read_tsv(path: Path) -> list[dict[str, str]]:
    if not path.is_file():
        raise FileNotFoundError(path)
    with path.open("r", encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle, delimiter="\t"))


def readable_accent(color: RGBColor) -> RGBColor:
    if color == TEAL:
        return TEAL_TEXT
    if color == GOLD:
        return GOLD_TEXT
    if color == VERMILION:
        return VERMILION_TEXT
    return color


def fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_run(
    run,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    font: str = FONT_BODY,
) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15,
    color: RGBColor = DARK,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.03,
    font: str = FONT_BODY,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic, font=font)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    color: RGBColor = WHITE,
    outline: RGBColor | None = LIGHT,
    radius: bool = True,
    width: float = 1.0,
):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color)
    if outline is None:
        shape.line.fill.background()
    else:
        line(shape, outline, width)
    return shape


def add_circle(
    slide,
    x: float,
    y: float,
    diameter: float,
    color: RGBColor,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(diameter),
        Inches(diameter),
    )
    fill(shape, color)
    shape.line.fill.background()
    return shape


def add_bullets(
    slide,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    *,
    size: float = 13.2,
    color: RGBColor = DARK,
    accent: RGBColor = BLUE,
    line_h: float = 0.55,
) -> None:
    current_y = y
    for item in items:
        add_circle(slide, x, current_y + 0.14, 0.085, accent)
        add_text(
            slide,
            item,
            x + 0.19,
            current_y,
            w - 0.19,
            line_h,
            size=size,
            color=color,
        )
        current_y += line_h


def add_metric(
    slide,
    value: str,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: RGBColor = BLUE,
    bg: RGBColor = WHITE,
) -> None:
    add_rect(slide, x, y, w, 1.12, color=bg, outline=LIGHT)
    value_size = 25 if len(value) <= 9 else 18
    add_text(
        slide,
        value,
        x + 0.16,
        y + 0.10,
        w - 0.32,
        0.47,
        size=value_size,
        color=readable_accent(accent),
        bold=True,
        font=FONT_HEAD,
    )
    add_text(
        slide,
        label,
        x + 0.16,
        y + 0.60,
        w - 0.32,
        0.43,
        size=8.8 if len(label) > 27 else 9.6,
        color=GRAY,
        bold=True,
    )


def add_panel_title(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: RGBColor,
) -> None:
    add_rect(slide, x, y + 0.03, 0.07, 0.31, color=accent, outline=None, radius=False)
    add_text(
        slide,
        text,
        x + 0.16,
        y,
        w - 0.16,
        0.36,
        size=15.2,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
    )


def add_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    widths: Sequence[float],
    *,
    row_h: float = 0.48,
    header_h: float = 0.48,
    font_size: float = 9.6,
    highlight_last: bool = False,
) -> None:
    current_x = x
    for header, width_value in zip(headers, widths, strict=True):
        add_rect(
            slide,
            current_x,
            y,
            width_value,
            header_h,
            color=NAVY,
            outline=WHITE,
            radius=False,
            width=0.6,
        )
        add_text(
            slide,
            header,
            current_x + 0.06,
            y + 0.06,
            width_value - 0.12,
            header_h - 0.10,
            size=font_size - 0.4,
            color=WHITE,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        current_x += width_value
    for row_index, row in enumerate(rows):
        current_x = x
        bg = PALE_GREEN if highlight_last and row_index == len(rows) - 1 else (
            WHITE if row_index % 2 == 0 else PALE_GRAY
        )
        for value, width_value in zip(row, widths, strict=True):
            add_rect(
                slide,
                current_x,
                y + header_h + row_index * row_h,
                width_value,
                row_h,
                color=bg,
                outline=LIGHT,
                radius=False,
                width=0.6,
            )
            add_text(
                slide,
                value,
                current_x + 0.06,
                y + header_h + row_index * row_h + 0.05,
                width_value - 0.12,
                row_h - 0.08,
                size=font_size,
                color=DARK,
                bold=highlight_last and row_index == len(rows) - 1,
                valign=MSO_ANCHOR.MIDDLE,
            )
            current_x += width_value


def set_alt_text(shape, description: str) -> None:
    c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("name", description)
    c_nv_pr.set("descr", description)


def add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    alt: str,
):
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    picture_w, picture_h = image_w * scale, image_h * scale
    picture_x = x + (w - picture_w) / 2
    picture_y = y + (h - picture_h) / 2
    picture = slide.shapes.add_picture(
        str(path),
        Inches(picture_x),
        Inches(picture_y),
        Inches(picture_w),
        Inches(picture_h),
    )
    set_alt_text(picture, alt)
    return picture


_NOTES_BODY_TEMPLATE = None


def set_notes_body_template(element) -> None:
    """Remember a notes body placeholder to clone into new notes slides.

    The deck's notes master carries no placeholders, so notes slides created
    for appended slides would otherwise have no body text frame.
    """
    global _NOTES_BODY_TEMPLATE
    _NOTES_BODY_TEMPLATE = element


def add_notes(
    slide,
    *,
    goal: str,
    walkthrough: str,
    boundary: str,
    transition: str,
) -> None:
    notes = slide.notes_slide
    if notes.notes_text_frame is None:
        if _NOTES_BODY_TEMPLATE is None:
            raise RuntimeError("No notes body template registered")
        notes.shapes._spTree.append(copy.deepcopy(_NOTES_BODY_TEMPLATE))
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        raise RuntimeError("Notes placeholder is unavailable")
    frame.text = (
        f"Teaching goal: {goal}\n\n"
        f"Walk through: {walkthrough}\n\n"
        f"Scientific boundary: {boundary}\n\n"
        f"Transition: {transition}"
    )


def blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    raise RuntimeError("The deck template has no Blank layout")


def new_slide(prs: Presentation, *, bg: RGBColor = OFF_WHITE):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_title_block(slide, title: str, subtitle: str | None = None) -> None:
    title_size = 21.0 if len(title) >= 64 else 24.5
    add_text(
        slide,
        title,
        0.58,
        0.45,
        11.65,
        0.52,
        size=title_size,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
    )
    if subtitle:
        add_text(slide, subtitle, 0.60, 0.97, 11.55, 0.30, size=11.2, color=GRAY)


def validate_inputs() -> dict[str, Any]:
    required = [
        RECURRENCE_PNG,
        TOP5_PNG,
        RECURRENCE_DATA,
        TOP5_DATA,
        MANIFEST_PATH,
        RESULT_DIR / "simple_status.tsv",
        RESULT_DIR / "simple_category_gene_aggregates.tsv",
        RESULT_DIR / "simple_category_summary.tsv",
    ]
    missing = [str(path) for path in required if not path.is_file()]
    if missing:
        raise FileNotFoundError("Missing SEA-AD source(s): " + ", ".join(missing))

    for figure_id in ("driver_recurrence", "top5_candidates"):
        status_path = (
            FIGURE_DIR / figure_id / f"phase11_seaad_simple_aggr_{figure_id}_status.tsv"
        )
        status = read_tsv(status_path)
        if (
            len(status) != 1
            or status[0]["validation_status"] != "validated_complete"
            or int(status[0]["failed_checks"]) != 0
            or status[0]["scope"] != "non_mt_driver"
        ):
            raise RuntimeError(f"Figure source is not validated: {status_path}")

    status_rows = read_tsv(RESULT_DIR / "simple_status.tsv")
    if len(status_rows) != 1:
        raise RuntimeError("simple_status.tsv must contain exactly one row")
    status = status_rows[0]
    if (
        status["analysis_id"] != "seaad_simple_returned_only_non_core_mt_acat_v1"
        or status["cohort"] != "SEAAD"
        or status["execution_status"] != "complete"
        or int(status["failed_check_count"]) != 0
    ):
        raise RuntimeError("SEA-AD simple aggregation status is not validated")

    manifest = read_tsv(MANIFEST_PATH)
    eligible = [row for row in manifest if row["eligibility_status"] == "eligible"]
    terminal_counts = Counter(row["terminal_status"] for row in manifest)
    eligible_groups = Counter(row["signature_group"] for row in eligible)
    supertypes = {row["supertype_id"] for row in manifest}

    categories = read_tsv(RESULT_DIR / "simple_category_gene_aggregates.tsv")
    non_mt = [
        row
        for row in categories
        if row["case_id"] == "non_mt_driver" and not truth(row["is_core_mito"])
    ]
    group_counts = Counter(row["signature_group"] for row in non_mt)
    network_counts = Counter(row["broad_network"] for row in non_mt)
    category_keys = {(row["signature_group"], row["broad_network"]) for row in non_mt}

    recurrence = read_tsv(RECURRENCE_DATA)
    top5 = read_tsv(TOP5_DATA)

    values: dict[str, Any] = {
        "planned_slots": len(manifest),
        "supertype_count": len(supertypes),
        "not_estimable_slots": terminal_counts["source_contrast_not_estimable"],
        "query_empty_slots": terminal_counts["query_empty"],
        "query_below_minimum_slots": terminal_counts["query_below_minimum"],
        "small_query_calls": terminal_counts["eligible_small_query"],
        "phase18_sized_calls": terminal_counts["eligible_phase18_sized"],
        "active_calls": len(eligible),
        "eligible_groups": eligible_groups,
        "significant_calls": int(status["completed_significant_call_count"]),
        "empty_calls": int(status["completed_empty_call_count"]),
        "stock_returned_rows": int(status["all_class_stock_returned_row_count"]),
        "mt_excluded_rows": int(status["mt_excluded_returned_row_count"]),
        "non_mt_returned_rows": int(status["non_mt_retained_returned_row_count"]),
        "non_mt_category_units": len(non_mt),
        "non_mt_unique_genes": len({row["current_symbol"] for row in non_mt}),
        "non_mt_categories": len(category_keys),
        "structural_categories": int(status["structural_category_count"]),
        "top5_rows": len(top5),
        "top5_unique_genes": len({row["current_symbol"] for row in top5}),
        "group_counts": group_counts,
        "network_counts": network_counts,
        "recurrence": recurrence,
        "top5": top5,
    }
    expected = {
        "planned_slots": 1548,
        "supertype_count": 129,
        "not_estimable_slots": 786,
        "query_empty_slots": 703,
        "query_below_minimum_slots": 17,
        "small_query_calls": 21,
        "phase18_sized_calls": 21,
        "active_calls": 42,
        "significant_calls": 27,
        "empty_calls": 15,
        "stock_returned_rows": 201,
        "mt_excluded_rows": 80,
        "non_mt_returned_rows": 121,
        "non_mt_category_units": 96,
        "non_mt_unique_genes": 91,
        "non_mt_categories": 4,
        "structural_categories": 42,
        "top5_rows": 18,
        "top5_unique_genes": 18,
    }
    for key, expected_value in expected.items():
        if values[key] != expected_value:
            raise RuntimeError(
                f"Source count drift for {key}: {values[key]} != {expected_value}"
            )
    if int(status["category_gene_unit_count"]) != len(non_mt):
        raise RuntimeError("Category unit count does not match the aggregate table")
    if eligible_groups != Counter({"M_e33": 40, "F_e33": 1, "F_e4": 1}):
        raise RuntimeError("Eligible-call group distribution drifted")
    if len(recurrence) != 20 or recurrence[0]["current_symbol"] != "HGSNAT":
        raise RuntimeError("Recurrence source does not match the validated top-20 contract")
    if int(recurrence[0]["category_count"]) != 2:
        raise RuntimeError("Top recurrent SEA-AD gene must appear in two categories")
    return values


def replace_run_text(slide, shape_name: str, old: str, new: str) -> None:
    matches = [shape for shape in slide.shapes if shape.name == shape_name]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one shape named {shape_name!r}, found {len(matches)}")
    runs = [
        run
        for paragraph in matches[0].text_frame.paragraphs
        for run in paragraph.runs
    ]
    if len(runs) != 1:
        raise RuntimeError(f"Expected one text run in {shape_name!r}, found {len(runs)}")
    if runs[0].text != old:
        raise RuntimeError(
            f"Pre-update contract failed for {shape_name!r}: found {runs[0].text!r}"
        )
    runs[0].text = new


def update_slide_1(slide, facts: dict[str, Any]) -> None:
    replace_run_text(
        slide,
        "TextBox 6",
        "Fine-cell reaggregation",
        "ROSMAP reaggregation",
    )
    replace_run_text(
        slide,
        "TextBox 7",
        "295 included runs → coverage + support + ACAT → 74 candidate units",
        "295 included runs → returned-only non-MT ACAT → 689 category units",
    )
    replace_run_text(
        slide,
        "TextBox 10",
        "Direct broad-cell KDA",
        "SEA-AD validation",
    )
    replace_run_text(
        slide,
        "TextBox 11",
        "3 completed runs → within-run BH → 12 direct candidate rows",
        f"{facts['active_calls']} active KDA calls → returned-only non-MT ACAT → "
        f"{facts['non_mt_category_units']} category units",
    )
    add_notes(
        slide,
        goal="Introduce the two cohorts and signal that both parts use the same returned-only aggregation rule.",
        walkthrough="Part I aggregates the 295 included ROSMAP fine-cell KDA runs into sex/APOE by broad-network categories with the returned-only non-MT ACAT rule, giving 689 category units. Part II applies the identical rule to the 42 active SEA-AD supertype calls, giving 96 category units.",
        boundary="The two unit counts are not interchangeable evidence tiers: both are exploratory post-selected aggregates, and the SEA-AD call distribution is heavily concentrated in M_e33.",
        transition="First compare the cohorts at a glance, then unpack the ROSMAP branch before turning to SEA-AD.",
    )


def update_slide_2(slide, facts: dict[str, Any]) -> None:
    replace_run_text(
        slide,
        "TextBox 2",
        "Two branches answer related—but not identical—questions",
        "Two cohorts, one returned-only aggregation rule",
    )
    replace_run_text(
        slide,
        "TextBox 3",
        "Both start with mitochondrial DEG queries; only the fine-cell branch combines evidence across runs.",
        "Both apply the same returned-only non-MT ACAT rule; the cohorts differ in run availability and balance.",
    )
    replace_run_text(
        slide,
        "TextBox 5",
        "FINE-CELL REAGGREGATION",
        "ROSMAP FINE-CELL REAGGREGATION",
    )
    replace_run_text(
        slide,
        "TextBox 14",
        "74",
        "689",
    )
    replace_run_text(
        slide,
        "TextBox 19",
        "Uses coverage, individual-run support, ACAT, then category BH.",
        "Aggregates stock-significant returns only; ACAT for ≥2 calls.",
    )
    replace_run_text(
        slide,
        "TextBox 21",
        "Final 74 units represent 37 distinct genes in 16 categories.",
        "Final 689 units represent 433 distinct genes in 32 categories.",
    )
    replace_run_text(
        slide,
        "TextBox 23",
        "DIRECT BROAD-CELL KDA",
        "SEA-AD VALIDATION",
    )
    replace_run_text(
        slide,
        "TextBox 24",
        "Which genes pass within one broad-cell direction?",
        "Do drivers return in an independent human cohort?",
    )
    replace_run_text(slide, "TextBox 26", "84", f"{facts['planned_slots']:,}")
    replace_run_text(
        slide,
        "TextBox 29",
        "3",
        str(facts["active_calls"]),
    )
    replace_run_text(
        slide,
        "TextBox 30",
        "completed KDA runs",
        "active KDA calls",
    )
    replace_run_text(slide, "TextBox 32", "12", str(facts["non_mt_category_units"]))
    replace_run_text(
        slide,
        "TextBox 33",
        "direct candidate rows",
        "gene × category units",
    )
    replace_run_text(
        slide,
        "TextBox 35",
        "Keeps broad-cell × group × direction results separate.",
        "Applies the same returned-only non-MT ACAT rule.",
    )
    replace_run_text(
        slide,
        "TextBox 37",
        "Uses complete explicit tests and within-run non-core-MT BH.",
        "Relaxed tier: donor ≥3 per arm, FDR-only query, n ≥ 3.",
    )
    replace_run_text(
        slide,
        "TextBox 39",
        "No coverage, ACAT, recurrence gate, or combined up/down q.",
        "40 of 42 calls are M_e33; 4 categories have returns.",
    )
    add_notes(
        slide,
        goal="Give the audience a mental map of the discovery cohort and the validation cohort before the details.",
        walkthrough="The left card summarizes the ROSMAP fine-cell reaggregation: 648 planned slots, 295 included runs, and 689 non-MT gene-by-category units representing 433 genes. The right card summarizes SEA-AD: 1,548 planned supertype slots, 42 active relaxed-tier calls, and 96 units representing 91 genes.",
        boundary="The SEA-AD tier is deliberately relaxed (donor at least 3 per arm, FDR-only query, no fold-change cutoff), and 40 of its 42 active calls are M_e33, so absence of a category is availability, not a tested null.",
        transition="Begin Part I with the ROSMAP counting units, then return to SEA-AD in Part II.",
    )


def append_section_slide(prs: Presentation, facts: dict[str, Any]):
    slide = new_slide(prs, bg=NAVY)
    add_text(slide, "PART 2", 0.78, 0.67, 2.4, 0.28, size=10.5, color=TEAL, bold=True)
    add_rect(slide, 0.78, 1.26, 0.10, 2.38, color=TEAL, outline=None, radius=False)
    add_text(
        slide,
        "SEA-AD sex/APOE KDA validation",
        1.18,
        1.36,
        8.65,
        1.34,
        size=32,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        f"{facts['supertype_count']} SEA-AD supertypes define the planned universe; "
        f"{facts['active_calls']} active relaxed-tier calls are aggregated with the same "
        "returned-only non-MT ACAT rule used in Part I.",
        1.20,
        3.00,
        8.30,
        0.84,
        size=15.0,
        color=RGBColor(204, 219, 234),
    )
    add_rect(slide, 9.78, 1.28, 2.70, 4.70, color=NAVY_2, outline=None)
    add_text(slide, "IN THIS PART", 10.08, 1.67, 2.10, 0.24, size=9.4, color=TEAL, bold=True)
    topics = ["Run universe", "Query + tier rules", "Returned-only rule", "Result figures"]
    current_y = 2.20
    for index, topic in enumerate(topics, start=1):
        add_circle(slide, 10.06, current_y + 0.03, 0.34, TEAL)
        add_text(
            slide,
            str(index),
            10.06,
            current_y + 0.08,
            0.34,
            0.17,
            size=8.2,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, topic, 10.53, current_y, 1.58, 0.52, size=11.1, color=WHITE, bold=True)
        current_y += 0.82
    add_notes(
        slide,
        goal="Introduce the SEA-AD validation branch and its sequence from setup to results.",
        walkthrough="The section first defines the SEA-AD run universe of 129 supertypes, six sex/APOE groups, and two query directions. It then explains the relaxed exploratory tier, the returned-only non-MT ACAT rule reused from Part I, and the two requested figures.",
        boundary="The branch reuses the validated 42 SEA-AD KDA calls without rerunning KDA; it validates the aggregation recipe, not a formally FDR-controlled candidate list.",
        transition="Start with how 1,548 planned directional slots reduce to 42 active calls.",
    )
    return slide


def append_run_universe_slide(prs: Presentation, facts: dict[str, Any]):
    slide = new_slide(prs)
    add_title_block(
        slide,
        f"{facts['supertype_count']} SEA-AD supertypes create "
        f"{facts['planned_slots']:,} planned directional slots",
        "The planned universe mirrors the ROSMAP design: supertype × sex/APOE contrast × AD-up/AD-down mito query.",
    )
    add_rect(slide, 0.70, 1.47, 3.70, 1.31, color=PALE_GREEN, outline=TEAL)
    add_text(
        slide,
        f"{facts['supertype_count']} SEA-AD supertypes",
        0.95, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "× 6 sex/APOE groups = 774 contrasts",
        0.95, 2.18, 3.20, 0.30, size=12.3, color=TEAL_TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, 4.82, 1.47, 3.70, 1.31, color=PALE_SKY, outline=SKY)
    add_text(
        slide,
        "774 signed contrasts",
        5.07, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"× up/down = {facts['planned_slots']:,} directional slots",
        5.07, 2.18, 3.20, 0.30, size=12.3, color=BLUE, bold=True, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, 8.94, 1.47, 3.70, 1.31, color=PALE_GOLD, outline=GOLD)
    add_text(
        slide,
        f"{facts['active_calls']} active KDA calls",
        9.19, 1.74, 3.20, 0.31, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "estimable source and query n ≥ 3",
        9.19, 2.18, 3.20, 0.30, size=12.3, color=GOLD_TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_table(
        slide,
        ["Mutually exclusive slot outcome", "Slots", "KDA called?"],
        [
            ["Source DEG contrast not estimable", f"{facts['not_estimable_slots']}", "No"],
            ["Estimable source; FDR-only query = 0 genes", f"{facts['query_empty_slots']}", "No"],
            ["Effective query 1–2", f"{facts['query_below_minimum_slots']}", "No"],
            ["Effective query 3–9", f"{facts['small_query_calls']}", "Yes (small query)"],
            ["Effective query ≥10", f"{facts['phase18_sized_calls']}", "Yes"],
            ["TOTAL", f"{facts['planned_slots']:,}", f"{facts['active_calls']} calls"],
        ],
        0.86, 3.10, [6.35, 1.22, 3.95], row_h=0.47, header_h=0.47,
        font_size=9.6, highlight_last=True,
    )
    add_text(
        slide,
        f"Group availability is unbalanced: {facts['eligible_groups']['M_e33']} of the "
        f"{facts['active_calls']} active calls are M_e33; F_e33 and F_e4 contribute one each.",
        0.88, 6.30, 11.60, 0.28, size=11.0, color=PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        goal="Show exactly how 1,548 planned SEA-AD directional slots become 42 active KDA calls.",
        walkthrough="The five outcomes are mutually exclusive. 786 slots inherit non-estimable source contrasts, 703 have no FDR-only mitochondrial query gene, and 17 keep only one or two genes. The remaining 42 slots meet the minimum of three effective query genes; 21 are small queries of three to nine genes and 21 are Phase 18-sized.",
        boundary="A skipped slot has no KDA result and must not be counted as a completed null run. The active calls concentrate in M_e33, so category-level breadth is bounded by availability.",
        transition="Next, the relaxed tier rules and the returned-only aggregation that reuses the Part I recipe.",
    )
    return slide


def append_rules_slide(prs: Presentation, facts: dict[str, Any]):
    slide = new_slide(prs)
    add_title_block(
        slide,
        "Relaxed-tier queries feed the same returned-only rule as Part I",
        "The SEA-AD calls come from a deliberately relaxed exploratory tier; the aggregation recipe is unchanged.",
    )
    add_metric(slide, str(facts["active_calls"]), "active KDA calls", 0.68, 1.42, 2.18, accent=TEAL)
    add_metric(slide, str(facts["significant_calls"]), "calls with significant returns", 3.03, 1.42, 2.18, accent=TEAL)
    add_metric(slide, str(facts["empty_calls"]), "completed-empty calls", 5.38, 1.42, 2.18, accent=TEAL)
    add_metric(slide, str(facts["stock_returned_rows"]), "stock returned rows", 7.73, 1.42, 2.18, accent=GOLD)
    add_metric(slide, str(facts["non_mt_returned_rows"]), "non-MT rows retained", 10.08, 1.42, 2.18, accent=GOLD)
    add_rect(slide, 0.72, 2.96, 5.86, 3.55, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Relaxed exploratory tier", 1.03, 3.28, 5.25, accent=TEAL)
    add_bullets(slide, [
        "Donor support ≥3 per disease arm in each supertype contrast.",
        "Mito DEG query: within-contrast BH FDR < 0.05 only.",
        "No fold-change cutoff (fdr_only_query_sensitivity).",
        "Effective KDA query needs ≥3 network-mapped genes.",
        "Frozen tier: posthoc_exploratory__fdr_only__donor3__query3__coverage80__q05.",
    ], 1.03, 3.81, 5.25, size=10.6, line_h=0.50, accent=TEAL)
    add_rect(slide, 6.82, 2.96, 5.80, 3.55, color=PALE_GOLD, outline=GOLD)
    add_panel_title(slide, "Returned-only rule (same as Part I)", 7.13, 3.28, 5.19, accent=GOLD)
    add_bullets(slide, [
        "Start from stock call_key_drivers() returns (within-call q ≤ 0.05).",
        f"Drop core-MitoCarta drivers: {facts['mt_excluded_rows']} of "
        f"{facts['stock_returned_rows']} rows excluded.",
        "One returned row → copy its within-call q unchanged.",
        "≥2 returned rows → equal-weight ACAT of returned q values.",
        "No P=1 backfill for unreturned calls; no across-gene BH.",
    ], 7.13, 3.81, 5.19, size=10.6, line_h=0.50, accent=GOLD)
    add_text(
        slide,
        "The resulting score is exploratory and post-selected—not a formally FDR-controlled cross-call q value.",
        0.88, 6.72, 11.60, 0.28, size=11.0, color=VERMILION_TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        goal="State the SEA-AD tier rules and confirm the aggregation recipe is identical to Part I.",
        walkthrough="The 42 active calls completed as 27 with significant returns and 15 valid empties, producing 201 stock returned rows. Removing 80 core-MitoCarta rows leaves 121 non-MT rows. Singleton genes keep their within-call BH q; genes with two or more returned rows get an equal-weight ACAT of only those returned q values.",
        boundary="Inputs are preselected within-call significant returns, ACAT is applied to adjusted values, and no final across-gene correction is made; use the score for exploratory ranking only.",
        transition="The two requested figures now display the top-five candidates and driver recurrence.",
    )
    return slide


def append_top5_slide(prs: Presentation, facts: dict[str, Any]):
    slide = new_slide(prs)
    add_title_block(
        slide,
        f"Top-five display: {facts['top5_rows']} non-MT entries across "
        f"{facts['non_mt_categories']} categories",
    )
    add_picture_contain(
        slide,
        TOP5_PNG,
        0.62,
        0.98,
        12.10,
        6.30,
        alt=(
            "SEA-AD simple returned-only female and male panels showing up to five "
            "non-MT genes per sex/APOE by broad-cell category"
        ),
    )
    add_notes(
        slide,
        goal="Interpret the SEA-AD non-MT top-five display and its sparse category structure.",
        walkthrough="The figure contains 18 tiles across 4 return-bearing categories: F_e33 excitatory neurons with three genes, and the M_e33 excitatory, inhibitory, and oligodendrocyte categories with five each. Blue means ACAT across at least two returned calls; orange means a one-call within-call BH q passthrough.",
        boundary="Top five is a display cap, not a significance threshold. Color identifies the calculation route, not an evidence tier, and 38 structural categories without a non-MT return are omitted.",
        transition="The next slide shows which genes recur across the four populated categories.",
    )
    return slide


def append_recurrence_slide(prs: Presentation, facts: dict[str, Any]):
    recurrence = facts["recurrence"]
    first = recurrence[0]
    runners = ", ".join(row["current_symbol"] for row in recurrence[1:5])
    slide = new_slide(prs)
    add_title_block(
        slide,
        f"{first['current_symbol']} recurs across {first['category_count']} "
        "SEA-AD returned-only categories",
        "Returned-only recurrence is descriptive across sex/APOE × broad-cell categories.",
    )
    add_rect(slide, 0.47, 1.27, 7.16, 5.40, color=WHITE, outline=LIGHT)
    add_picture_contain(
        slide,
        RECURRENCE_PNG,
        0.58,
        1.37,
        6.94,
        5.18,
        alt=(
            "SEA-AD simple returned-only bar chart of the twenty non-MT genes "
            "appearing in the most sex/APOE by broad-cell categories"
        ),
    )
    add_rect(slide, 7.91, 1.42, 4.75, 4.98, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "How to read it", 8.22, 1.74, 4.13, accent=TEAL)
    add_bullets(slide, [
        "Bar length = number of categories containing the gene.",
        "Fill = categories whose score combines ≥2 returned calls.",
        f"{first['current_symbol']}: {first['category_count']} categories, "
        f"{first['sex_apoe_group_count']} group, {first['broad_network_count']} networks; "
        f"ACAT in {first['acat_combined_category_count']}.",
        f"{runners} also recur in 2 categories.",
        "Top 20 by category count, best exploratory score, then symbol.",
    ], 8.23, 2.28, 4.02, size=11.3, line_h=0.61, accent=TEAL)
    add_rect(slide, 8.22, 5.55, 4.10, 0.56, color=WHITE, outline=LIGHT)
    add_text(
        slide,
        "Only 4 categories have returns—recurrence ceilings are low.",
        8.41, 5.66, 3.72, 0.36, size=10.1, color=VERMILION_TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        goal="Interpret SEA-AD returned-only recurrence without over-reading the small category universe.",
        walkthrough="Bar length counts each sex/APOE by broad-cell category once. HGSNAT appears in two categories (M_e33 excitatory and oligodendrocytes) with nine total returned calls; DYNLT1, ZNF706, METTL26, and GOLT1B also appear in two categories. The remaining displayed genes appear once and are ordered by best exploratory score.",
        boundary="With only four return-bearing categories, the maximum possible recurrence is small; fill shows the calculation route, not evidence strength, and recurrence is not independent replication.",
        transition="The summary slide translates the SEA-AD results into returned-row, unit, gene, and category counts.",
    )
    return slide


def append_summary_slide(prs: Presentation, facts: dict[str, Any]):
    slide = new_slide(prs)
    add_title_block(
        slide,
        f"SEA-AD output: {facts['non_mt_category_units']} non-MT category units "
        f"represent {facts['non_mt_unique_genes']} genes",
    )
    add_metric(slide, str(facts["non_mt_returned_rows"]), "non-MT returned call rows", 0.72, 1.52, 2.12, accent=TEAL)
    add_metric(slide, str(facts["non_mt_category_units"]), "non-MT gene × category units", 3.06, 1.52, 2.12, accent=TEAL)
    add_metric(slide, str(facts["non_mt_unique_genes"]), "distinct gene symbols", 5.40, 1.52, 2.12, accent=TEAL)
    add_metric(slide, str(facts["non_mt_categories"]), "categories with non-MT returns", 7.74, 1.52, 2.12, accent=BLUE)
    add_metric(slide, str(facts["top5_rows"]), "top-five displayed rows", 10.08, 1.52, 2.12, accent=GOLD)
    add_rect(slide, 0.72, 3.14, 5.76, 3.10, color=PALE_GREEN, outline=TEAL)
    add_panel_title(slide, "Where non-MT driver units occur", 1.02, 3.46, 5.16, accent=TEAL)
    add_text(slide, "By sex/APOE group", 1.03, 4.03, 2.05, 0.24, size=10.2, color=TEAL_TEXT, bold=True)
    group_counts: Counter[str] = facts["group_counts"]
    group_text = (
        f"M_e33 {group_counts['M_e33']}  •  F_e33 {group_counts['F_e33']}  •  "
        "all other groups 0"
    )
    add_text(slide, group_text, 1.03, 4.40, 4.97, 0.60, size=11.7, color=NAVY, bold=True)
    add_text(slide, "By broad network", 1.03, 5.16, 2.05, 0.24, size=10.2, color=TEAL_TEXT, bold=True)
    network_counts: Counter[str] = facts["network_counts"]
    network_text = (
        f"Excitatory {network_counts['Excitatory_neurons']} • "
        f"Inhibitory {network_counts['Inhibitory_neurons']} • "
        f"Oligo {network_counts['Oligodendrocytes']}\n"
        "Astrocytes, Microglia, OPCs, Vasculature 0"
    )
    add_text(slide, network_text, 1.03, 5.48, 4.97, 0.62, size=10.3, color=NAVY, bold=True)
    add_rect(slide, 6.76, 3.14, 5.48, 3.10, color=PALE_SKY, outline=SKY)
    add_panel_title(slide, "Most recurrent returned-only genes", 7.06, 3.46, 4.88, accent=BLUE)
    recurrence = facts["recurrence"][:6]
    maximum = max(int(row["category_count"]) for row in recurrence)
    for index, row in enumerate(recurrence):
        y = 4.02 + index * 0.32
        count = int(row["category_count"])
        add_text(slide, row["current_symbol"], 7.08, y, 1.42, 0.22, size=11.2, color=NAVY, bold=True)
        add_rect(slide, 8.63, y + 0.03, 2.62 * count / maximum, 0.16, color=BLUE, outline=None, radius=False)
        add_text(slide, str(count), 11.42, y, 0.35, 0.22, size=10.4, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    add_text(
        slide,
        "Counts are category presence—not independent replication or call counts.",
        7.08, 5.99, 4.74, 0.24, size=8.7, color=GRAY, italic=True,
    )
    add_text(
        slide,
        f"Validation breadth is bounded by availability: {facts['eligible_groups']['M_e33']} "
        f"of {facts['active_calls']} active SEA-AD calls are M_e33.",
        0.88, 6.60, 11.60, 0.28, size=11.0, color=PURPLE, bold=True, align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        goal="Summarize the scale and distribution of the SEA-AD returned-only non-MT result set.",
        walkthrough="The 121 non-MT returned call rows form 96 gene-by-category units representing 91 genes across 4 populated categories. The left panel shows the M_e33 concentration; the right panel shows the five genes recurring in two categories, led by HGSNAT.",
        boundary="These descriptive counts condition on stock-significant returns and on a call distribution dominated by M_e33. They do not use nonreturns, coverage, support gates, or a final FDR-controlled cross-gene family.",
        transition="This closes the SEA-AD validation part of the deck.",
    )
    return slide


def slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


def write_audit(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["check_id", "observed", "expected", "passed"],
            delimiter="\t",
            lineterminator="\n",
        )
        writer.writeheader()
        writer.writerows(rows)


def main() -> int:
    args = parse_args()
    input_path = args.input.resolve()
    output_path = args.output.resolve()
    if not input_path.is_file():
        raise FileNotFoundError(input_path)
    facts = validate_inputs()
    original_deck_hash = sha256_file(input_path)

    prs = Presentation(str(input_path))
    if len(prs.slides) != EXPECTED_INPUT_SLIDES:
        raise RuntimeError(
            f"Expected the {EXPECTED_INPUT_SLIDES}-slide deck, found {len(prs.slides)}"
        )
    before_texts = {
        number: slide_text(prs.slides[number - 1]) for number in UNCHANGED_SLIDES
    }
    before_shape_counts = {
        number: len(prs.slides[number - 1].shapes) for number in UNCHANGED_SLIDES
    }

    set_notes_body_template(prs.slides[0].notes_slide.notes_placeholder._element)
    update_slide_1(prs.slides[0], facts)
    update_slide_2(prs.slides[1], facts)
    append_section_slide(prs, facts)
    append_run_universe_slide(prs, facts)
    append_rules_slide(prs, facts)
    append_top5_slide(prs, facts)
    append_recurrence_slide(prs, facts)
    append_summary_slide(prs, facts)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    file_descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{output_path.name}.", suffix=".tmp", dir=output_path.parent
    )
    os.close(file_descriptor)
    temporary = Path(temporary_name)
    try:
        prs.save(str(temporary))
        temporary.replace(output_path)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise

    reloaded = Presentation(str(output_path))
    after_texts = {
        number: slide_text(reloaded.slides[number - 1]) for number in UNCHANGED_SLIDES
    }
    after_shape_counts = {
        number: len(reloaded.slides[number - 1].shapes) for number in UNCHANGED_SLIDES
    }
    unchanged_ok = before_texts == after_texts and (
        before_shape_counts == after_shape_counts
    )
    slide1_text = slide_text(reloaded.slides[0])
    slide2_text = slide_text(reloaded.slides[1])
    new_titles = {
        11: "SEA-AD sex/APOE KDA validation",
        12: f"{facts['supertype_count']} SEA-AD supertypes create "
        f"{facts['planned_slots']:,} planned directional slots",
        13: "Relaxed-tier queries feed the same returned-only rule as Part I",
        14: f"Top-five display: {facts['top5_rows']} non-MT entries across "
        f"{facts['non_mt_categories']} categories",
        15: "HGSNAT recurs across 2 SEA-AD returned-only categories",
        16: f"SEA-AD output: {facts['non_mt_category_units']} non-MT category units "
        f"represent {facts['non_mt_unique_genes']} genes",
    }
    title_results = {
        number: title in slide_text(reloaded.slides[number - 1])
        for number, title in new_titles.items()
    }
    picture_counts = {
        number: sum(
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            for shape in reloaded.slides[number - 1].shapes
        )
        for number in (14, 15)
    }

    checks = [
        {
            "check_id": "output_slide_count",
            "observed": len(reloaded.slides),
            "expected": 16,
            "passed": len(reloaded.slides) == 16,
        },
        {
            "check_id": "slides_3_to_10_text_and_shape_counts_unchanged",
            "observed": "unchanged" if unchanged_ok else "changed",
            "expected": "unchanged",
            "passed": unchanged_ok,
        },
        {
            "check_id": "slide1_part2_repointed",
            "observed": "SEA-AD validation" in slide1_text
            and "42 active KDA calls" in slide1_text,
            "expected": True,
            "passed": "SEA-AD validation" in slide1_text
            and "42 active KDA calls" in slide1_text,
        },
        {
            "check_id": "slide1_part1_updated",
            "observed": "689 category units" in slide1_text,
            "expected": True,
            "passed": "689 category units" in slide1_text,
        },
        {
            "check_id": "slide2_cohort_comparison",
            "observed": "SEA-AD VALIDATION" in slide2_text
            and "1,548" in slide2_text
            and "689" in slide2_text,
            "expected": True,
            "passed": "SEA-AD VALIDATION" in slide2_text
            and "1,548" in slide2_text
            and "689" in slide2_text,
        },
    ]
    for number, passed in title_results.items():
        checks.append(
            {
                "check_id": f"slide{number}_title",
                "observed": passed,
                "expected": True,
                "passed": passed,
            }
        )
    for number, count in picture_counts.items():
        checks.append(
            {
                "check_id": f"slide{number}_has_picture",
                "observed": count,
                "expected": 1,
                "passed": count == 1,
            }
        )
    write_audit(args.audit.resolve(), checks)
    failed = [row["check_id"] for row in checks if not row["passed"]]
    if failed:
        raise RuntimeError("Slide update failed checks: " + ", ".join(failed))
    print(f"updated={output_path}")
    print("slides_edited=1,2")
    print("slides_appended=11,12,13,14,15,16")
    print(f"original_sha256={original_deck_hash}")
    print(f"updated_sha256={sha256_file(output_path)}")
    print(f"audit={args.audit.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
