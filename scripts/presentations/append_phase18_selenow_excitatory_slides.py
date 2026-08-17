#!/usr/bin/env python3
"""Append the SELENOW excitatory-neuron pathway and STRING slides."""

from __future__ import annotations

import argparse
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN


REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO))

from scripts.presentations.append_phase18_cox7c_astrocyte_slides import (  # noqa: E402
    replace_text_preserve_style,
)
from scripts.presentations.build_phase18_key_driver_selection_deck import (  # noqa: E402
    GOLD,
    GRAY,
    LIGHT,
    NAVY,
    PALE_GREEN,
    TEAL,
    VERMILION,
    WHITE,
    add_bullets,
    add_header,
    add_metric,
    add_panel_title,
    add_picture_contain,
    add_rect,
    add_source,
    add_text,
    new_slide,
    trim_white,
)


DEFAULT_DECK = REPO / "docs/presentations/key_driver_selection_analysis 08182026.pptx"
FIG_ROOT = REPO / "results/figures/analysis/phase_18_key_driver_selection/SELENOW/excitatory"
PATHWAY_FIG = FIG_ROOT / "phase18_selenow_excitatory_consensus_network_pathways.png"
STRING_FIG = FIG_ROOT / "stringdb_full_medium_conf.png"


def append_slides(input_path: Path, output_path: Path) -> Path:
    for path in (input_path, PATHWAY_FIG, STRING_FIG):
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation(input_path)
    if len(prs.slides) != 29:
        raise RuntimeError(f"Expected a 29-slide source deck, found {len(prs.slides)}")

    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    if "SELENOW • EXCITATORY NEURONS" in all_text:
        raise RuntimeError("SELENOW excitatory-neuron slides already appear to be present")

    replace_text_preserve_style(prs.slides[1], "Slides 16–29", "Slides 16–31")
    replace_text_preserve_style(
        prs.slides[1],
        "Place APOE, RPL11, and COX7C on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
        "Place APOE, RPL11, COX7C, and SELENOW on cell-type network graphs, examine protein-level coherence, and define proteomics validation.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "Next: localize candidate evidence across sex/APOE strata, then examine APOE, RPL11, and COX7C in depth.",
        "Next: localize candidate evidence across sex/APOE strata, then examine APOE, RPL11, COX7C, and SELENOW in depth.",
    )
    replace_text_preserve_style(
        prs.slides[8],
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–29 show APOE, RPL11, and COX7C network/protein examples",
        "Slides 13–15 add independent human-genetic support; slide 17 introduces the validation panel; slides 18–31 show APOE, RPL11, COX7C, and SELENOW network/protein examples",
    )

    with tempfile.TemporaryDirectory(prefix="selenow_excitatory_deck_assets_") as temp_dir:
        assets = Path(temp_dir)
        pathway = trim_white(PATHWAY_FIG, assets)
        string = trim_white(STRING_FIG, assets)

        # 30 — SELENOW excitatory-neuron directed pathway figure
        slide = new_slide(prs)
        add_header(
            slide,
            "SELENOW • excitatory neurons",
            "SELENOW anchors a stable excitatory mitochondrial program",
            30,
            accent=TEAL,
            subtitle="Fourteen conservative-support runs across nine fine cell types yield complete leave-one-fine-cell-type candidate retention.",
        )
        add_picture_contain(
            slide,
            pathway,
            0.10,
            1.22,
            8.92,
            5.92,
            alt="SELENOW-centered excitatory-neuron consensus network with directed Bayesian-network edges and contextual pathway outlines",
        )
        add_rect(slide, 9.18, 1.42, 3.62, 5.48, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "Evidence scale", 9.45, 1.70, 3.07, accent=TEAL)
        add_metric(slide, "14", "conservative-support excitatory runs", 9.45, 2.18, 1.39, accent=TEAL)
        add_metric(slide, "5.75×10⁻⁶", "aggregate ACAT q", 10.99, 2.18, 1.54, accent=GOLD)
        add_panel_title(slide, "Network readout", 9.45, 3.48, 3.07, accent=GOLD)
        add_bullets(
            slide,
            [
                "25 nodes and 25 directed edges retain 13 of 14 observed mitochondrial query hits at the 4/14 display threshold.",
                "Upstream chain: LAMTOR5 → ATP5MPL → SELENOW; seven direct outputs include COA3, PRELID1, PSEN1, and LAGE3.",
                "CHCHD10 and PRELID1 recur in all 14 query overlaps; NDUFB11 recurs in 13 and COA3/SLC25A4 in 11.",
                "Respiration, mitochondrial translation, and selenium outlines are contextual; none has BH FDR < 0.05.",
            ],
            9.45,
            3.96,
            2.94,
            size=10.3,
            line_h=0.58,
            accent=GOLD,
        )
        add_text(
            slide,
            "Interpretation: strong, stable topology plus external AD biology; pathway labels organize the graph but do not establish enrichment or activity.",
            9.48,
            6.43,
            2.88,
            0.32,
            size=9.2,
            color=NAVY,
            bold=True,
        )
        add_source(
            slide,
            "Source: phase18_selenow_excitatory_consensus_network_pathways.png, caption, methods, ORA table, and gene-by-gene analysis",
        )

        # 31 — SELENOW excitatory-neuron STRING figure and mechanistic interpretation
        slide = new_slide(prs)
        add_header(
            slide,
            "SELENOW • excitatory neurons",
            "STRING separates SELENOW from the mitochondrial protein core",
            31,
            accent=TEAL,
            subtitle="Medium-confidence STRING associations are undirected and not excitatory-neuron- or AD-specific.",
        )
        add_rect(slide, 0.42, 1.39, 6.37, 5.45, color=WHITE, outline=LIGHT)
        add_picture_contain(
            slide,
            string,
            0.64,
            1.62,
            5.93,
            4.98,
            alt="STRING medium-confidence functional association network for SELENOW and excitatory-neuron consensus-neighborhood proteins",
        )
        add_rect(slide, 7.04, 1.39, 5.76, 5.45, color=WHITE, outline=LIGHT)
        add_panel_title(slide, "What the image supports", 7.34, 1.72, 5.15, accent=TEAL)
        add_bullets(
            slide,
            [
                "A coherent mitochondrial component spans respiratory proteins, mitochondrial ribosomal proteins, CLPP, and SLC25A4.",
                "SELENOW sits in a small association component with SELENOM and MIEN1 rather than inside the respiratory core.",
                "Several modeled outputs are weakly connected or isolated, so STRING supports target-module coherence more than a direct SELENOW→module chain.",
            ],
            7.34,
            2.20,
            5.05,
            size=11.0,
            line_h=0.60,
            accent=TEAL,
        )
        add_panel_title(slide, "Externally reinforced mechanism", 7.34, 4.34, 5.15, accent=GOLD)
        add_text(
            slide,
            "Direct AD-model evidence indicates that SELENOW binds tau, promotes ubiquitin–proteasome clearance, and improves tau pathology and memory when overexpressed. Independent work supports redox and respiration functions, providing a second route to the predicted mitochondrial phenotype.",
            7.35,
            4.80,
            5.02,
            0.76,
            size=10.5,
            color=GRAY,
        )
        add_rect(slide, 7.35, 5.70, 4.98, 0.78, color=PALE_GREEN, outline=TEAL)
        add_text(
            slide,
            "In APOE-isogenic excitatory neurons under Aβ or tau stress, use CRISPRi/CRISPRa and rescue with wild-type versus redox-site-mutant SELENOW; measure tau clearance, proteasome activity, ROS/glutathione, respiration, viability, and the frozen target module.",
            7.57,
            5.82,
            4.55,
            0.54,
            size=9.2,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "The rescue separates redox-dependent action from scaffolding/proteostasis effects.",
            7.43,
            6.61,
            4.82,
            0.17,
            size=8.7,
            color=VERMILION,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_source(
            slide,
            "Sources: stringdb_full_medium_conf.png • phase18_key_driver_gene_by_gene_initial_analysis.md • Ren 2024; Misra 2023; Jeong 2002",
        )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    validate_output(output_path)
    return output_path


def validate_output(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 31:
        raise RuntimeError(f"Expected 31 slides, found {len(prs.slides)}")
    image_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    if image_shapes != 20:
        raise RuntimeError(f"Expected 20 figure images, found {image_shapes}")

    endings = []
    for slide_index in range(len(prs.slides) - 2, len(prs.slides)):
        slide = prs.slides[slide_index]
        endings.append(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
        )
    if not all("SELENOW • EXCITATORY NEURONS" in text for text in endings):
        raise RuntimeError("The final two slides are not the expected SELENOW excitatory-neuron slides")

    with zipfile.ZipFile(path) as archive:
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels"))
        )
    if "filter_attrition" in package_text:
        raise RuntimeError("Deprecated filter_attrition content found in the PPTX package")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--output", type=Path, default=DEFAULT_DECK)
    args = parser.parse_args()
    result = append_slides(args.input, args.output)
    print(result)
